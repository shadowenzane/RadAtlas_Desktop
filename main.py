import os
import sys
import json
import math
import sqlite3
import warnings

# 抑制 Python 级别警告
# 注意：libpng iCCP 警告是C级别输出，Python无法捕获，不影响功能
warnings.filterwarnings("ignore")

# 版本信息
APP_VERSION = 'v1.0.0'
APP_AUTHOR = '楚雄州人民医院 医学影像中心 张兴文'
import shutil
import requests
import tempfile
from datetime import datetime

from PyQt5.QtWidgets import (
    QApplication, QMainWindow, QWidget, QVBoxLayout, QHBoxLayout,
    QLabel, QLineEdit, QPushButton, QTextBrowser, QTextEdit, QSplitter,
    QListWidget, QListWidgetItem, QTabWidget, QDialog,
    QMessageBox, QStatusBar, QFrame, QGroupBox, QSizePolicy,
    QScrollArea, QComboBox, QFileDialog, QMenu, QMenuBar,
    QAction, QCompleter, QToolBar, QSpinBox, QCheckBox,
    QInputDialog, QSlider, QGridLayout, QAbstractItemView,
    QStackedWidget, QSizeGrip, QRadioButton, QFormLayout
)
from PyQt5.QtCore import Qt, QSize, QStringListModel, pyqtSignal, QRect, QRectF, QPoint, QPointF, QTimer, QEvent, QPropertyAnimation, QEasingCurve, QUrl, QObject, QThread
from PyQt5.QtWidgets import QGraphicsOpacityEffect
from PyQt5.QtGui import QFont, QColor, QPalette, QIcon, QTextCursor, QTextListFormat, QTextCharFormat, QPixmap, QPainter, QPainterPath, QPen, QPolygon, QBrush, QTransform, QTextDocument, QImage

from models import (
    init_user_db, authenticate_user, init_db, load_data,
    get_all_users, create_user, delete_user, change_password,
    rename_user, admin_change_password,
    hash_password, DATABASE, USER_DB, APP_DIR, copy_public_to_user,
    add_disease, update_disease, delete_disease,
    add_image, delete_image, get_disease, search_diseases,
    generate_key, encrypt_image, decrypt_image, get_image_by_id, update_image_encryption,
    encrypt_data, decrypt_data, encrypt_text, decrypt_text,
    store_image_encrypted, load_image_decrypted,
    store_annotations_encrypted, load_annotations_decrypted,
    migrate_images_to_db
)
from ai_helper import (load_config, save_config, call_ai, PROVIDERS,
                       load_multi_config, call_diagnosis_multi,
                       KNOWLEDGE_PROVIDERS,
                       test_llm_connection, test_kb_connection)

# ── 主题定义 ──────────────────────────────────────────────
THEMES = {
    '深蓝暗夜': {
        'bg': '#141416', 'bg2': '#1a1a1e', 'bg3': '#1e1e22',
        'fg': '#d9d9dc', 'fg2': '#888890', 'fg3': '#666670',
        'accent': '#3f7bf7', 'accent_hover': '#5a91ff', 'accent_press': '#2a5fd6',
        'border': '#2a2a30', 'danger': '#e64a3a', 'muted': '#3a3a40',
        'select_bg': '#2a3a5a', 'select_fg': '#5a91ff',
        'title_bg': '#1a1a2e', 'title_fg': '#5a91ff',
    },
    '极光绿': {
        'bg': '#0d1b12', 'bg2': '#122018', 'bg3': '#162a1e',
        'fg': '#d0e8d4', 'fg2': '#7aaa82', 'fg3': '#4a7a52',
        'accent': '#2ecc71', 'accent_hover': '#4ade80', 'accent_press': '#22a85a',
        'border': '#1e3a24', 'danger': '#e74c3c', 'muted': '#1e3a24',
        'select_bg': '#1a4a2a', 'select_fg': '#4ade80',
        'title_bg': '#0d1b12', 'title_fg': '#4ade80',
    },
    '暮光紫': {
        'bg': '#161220', 'bg2': '#1c1830', 'bg3': '#221e38',
        'fg': '#d8d0e8', 'fg2': '#9088a8', 'fg3': '#686080',
        'accent': '#9b59b6', 'accent_hover': '#b370cf', 'accent_press': '#7d3f98',
        'border': '#2e2848', 'danger': '#e74c3c', 'muted': '#2e2848',
        'select_bg': '#3a2860', 'select_fg': '#b370cf',
        'title_bg': '#1c1830', 'title_fg': '#b370cf',
    },
    '暖沙': {
        'bg': '#1a1714', 'bg2': '#22201c', 'bg3': '#2a2824',
        'fg': '#e0d8cc', 'fg2': '#a09880', 'fg3': '#706850',
        'accent': '#e67e22', 'accent_hover': '#f0983a', 'accent_press': '#c06610',
        'border': '#3a3630', 'danger': '#e64a3a', 'muted': '#3a3630',
        'select_bg': '#4a3a20', 'select_fg': '#f0983a',
        'title_bg': '#22201c', 'title_fg': '#f0983a',
    },
    '浅色经典': {
        'bg': '#f5f5f8', 'bg2': '#ffffff', 'bg3': '#eeeef2',
        'fg': '#222230', 'fg2': '#666678', 'fg3': '#9999a8',
        'accent': '#3f7bf7', 'accent_hover': '#5a91ff', 'accent_press': '#2a5fd6',
        'border': '#d0d0d8', 'danger': '#e64a3a', 'muted': '#c0c0c8',
        'select_bg': '#d0e0ff', 'select_fg': '#3f7bf7',
        'title_bg': '#e8e8f0', 'title_fg': '#3f7bf7',
    },
    '赤焰红': {
        'bg': '#1a1012', 'bg2': '#221418', 'bg3': '#2a1a1e',
        'fg': '#e8d0d4', 'fg2': '#a88088', 'fg3': '#785060',
        'accent': '#e74c3c', 'accent_hover': '#f06050', 'accent_press': '#c0392b',
        'border': '#3a2028', 'danger': '#ff4444', 'muted': '#3a2028',
        'select_bg': '#4a2028', 'select_fg': '#f06050',
        'title_bg': '#221418', 'title_fg': '#f06050',
    },
    '冰川蓝': {
        'bg': '#0e1520', 'bg2': '#141e2c', 'bg3': '#1a2838',
        'fg': '#c8d8e8', 'fg2': '#7898b0', 'fg3': '#507088',
        'accent': '#3498db', 'accent_hover': '#5dade2', 'accent_press': '#2178b0',
        'border': '#1e3040', 'danger': '#e74c3c', 'muted': '#1e3040',
        'select_bg': '#1e3a5a', 'select_fg': '#5dade2',
        'title_bg': '#141e2c', 'title_fg': '#5dade2',
    },
    '樱花粉': {
        'bg': '#1a1218', 'bg2': '#221820', 'bg3': '#2a1e28',
        'fg': '#e8d0e0', 'fg2': '#a880a0', 'fg3': '#785070',
        'accent': '#e84393', 'accent_hover': '#f06aa8', 'accent_press': '#c03078',
        'border': '#3a2038', 'danger': '#e74c3c', 'muted': '#3a2038',
        'select_bg': '#4a2040', 'select_fg': '#f06aa8',
        'title_bg': '#221820', 'title_fg': '#f06aa8',
    },
    '琥珀金': {
        'bg': '#18160e', 'bg2': '#201e14', 'bg3': '#28261a',
        'fg': '#e8e0c8', 'fg2': '#a89868', 'fg3': '#787048',
        'accent': '#f1c40f', 'accent_hover': '#f4d03f', 'accent_press': '#d4ac0d',
        'border': '#3a3620', 'danger': '#e74c3c', 'muted': '#3a3620',
        'select_bg': '#4a4020', 'select_fg': '#f4d03f',
        'title_bg': '#201e14', 'title_fg': '#f4d03f',
    },
}


def build_stylesheet(t):
    """根据主题字典生成样式表"""
    return f"""
QMainWindow, QDialog {{
    background-color: {t['bg']};
}}
QWidget {{
    background-color: {t['bg']};
    color: {t['fg']};
    font-family: "Microsoft YaHei", "Segoe UI", sans-serif;
    font-size: 13px;
}}
QLabel {{
    background-color: transparent;
    color: {t['fg']};
}}
QLineEdit {{
    background-color: {t['bg2']};
    color: {t['fg']};
    border: 1px solid {t['border']};
    border-radius: 6px;
    padding: 8px 12px;
    font-size: 14px;
    selection-background-color: {t['accent']};
}}
QLineEdit:focus {{
    border-color: {t['accent']};
}}
QPushButton {{
    background-color: {t['accent']};
    color: #ffffff;
    border: none;
    border-radius: 6px;
    padding: 8px 16px;
    font-size: 13px;
    font-weight: bold;
}}
QPushButton:hover {{
    background-color: {t['accent_hover']};
}}
QPushButton:pressed {{
    background-color: {t['accent_press']};
}}
QPushButton#dangerBtn {{
    background-color: {t['danger']};
}}
QPushButton#dangerBtn:hover {{
    background-color: #f06050;
}}
QPushButton#mutedBtn {{
    background-color: {t['muted']};
}}
QPushButton#mutedBtn:hover {{
    background-color: {t['border']};
}}
QPushButton#tabBtn {{
    background-color: {t['bg2']};
    color: {t['fg2']};
    border: none;
    border-radius: 0px;
    padding: 10px 16px;
    font-weight: normal;
    border-bottom: 3px solid transparent;
}}
QPushButton#tabBtn:hover {{
    background-color: {t['bg3']};
    color: {t['fg']};
}}
QPushButton#tabBtn[active="true"] {{
    color: {t['accent']};
    border-bottom: 3px solid {t['accent']};
    font-weight: bold;
}}
QListWidget {{
    background-color: {t['bg2']};
    border: none;
    border-radius: 6px;
    outline: none;
    padding: 4px;
}}
QListWidget::item {{
    background-color: {t['bg3']};
    color: {t['fg']};
    border: none;
    border-radius: 4px;
    padding: 10px 12px;
    margin: 2px 0px;
}}
QListWidget::item:hover {{
    background-color: {t['border']};
}}
QListWidget::item:selected {{
    background-color: {t['select_bg']};
    color: {t['select_fg']};
}}
QListWidget::item:disabled {{
    color: {t['accent']};
    background-color: transparent;
    font-weight: bold;
    font-size: 12px;
    padding: 6px 12px;
}}
QTextBrowser {{
    background-color: {t['bg2']};
    color: {t['fg']};
    border: none;
    border-radius: 6px;
    padding: 12px;
    font-size: 13px;
}}
QStatusBar {{
    background-color: {t['bg2']};
    color: {t['fg2']};
    border-top: 1px solid {t['border']};
    font-size: 11px;
    padding: 4px 8px;
}}
QSplitter::handle {{
    background-color: {t['border']};
    width: 2px;
}}
QGroupBox {{
    border: 1px solid {t['border']};
    border-radius: 6px;
    margin-top: 12px;
    padding-top: 16px;
    font-weight: bold;
}}
QGroupBox::title {{
    subcontrol-origin: margin;
    left: 12px;
    padding: 0 6px;
    color: {t['accent']};
}}
QComboBox {{
    background-color: {t['bg2']};
    color: {t['fg']};
    border: 1px solid {t['border']};
    border-radius: 6px;
    padding: 6px 12px;
}}
QComboBox::drop-down {{
    border: none;
}}
QComboBox QAbstractItemView {{
    background-color: {t['bg2']};
    color: {t['fg']};
    selection-background-color: {t['select_bg']};
}}
QMenuBar {{
    background-color: {t['bg2']};
    color: {t['fg']};
    border-bottom: 1px solid {t['border']};
}}
QMenuBar::item:selected {{
    background-color: {t['select_bg']};
}}
QMenu {{
    background-color: {t['bg2']};
    color: {t['fg']};
    border: 1px solid {t['border']};
}}
QMenu::item:selected {{
    background-color: {t['select_bg']};
}}
QScrollArea {{
    border: none;
    background-color: {t['bg']};
}}
QSpinBox {{
    background-color: {t['bg2']};
    color: {t['fg']};
    border: 1px solid {t['border']};
    border-radius: 4px;
    padding: 4px;
}}
QSlider::groove:horizontal {{
    background-color: {t['border']};
    height: 6px;
    border-radius: 3px;
}}
QSlider::handle:horizontal {{
    background-color: {t['accent']};
    width: 14px;
    height: 14px;
    margin: -4px 0;
    border-radius: 7px;
}}
QSlider::sub-page:horizontal {{
    background-color: {t['accent']};
    border-radius: 3px;
}}
QToolTip {{
    background-color: {t['bg3']};
    color: {t['fg']};
    border: 1px solid {t['border']};
    border-radius: 4px;
    padding: 4px 8px;
    font-size: 12px;
}}
"""


# 默认主题
DARK_STYLE = build_stylesheet(THEMES['暖沙'])


class _TitleBtn(QPushButton):
    """标题栏自绘图标按钮"""
    def __init__(self, mode='min', parent=None):
        super().__init__(parent)
        self.mode = mode
        self.setFixedSize(36, 28)
        self.setCursor(Qt.PointingHandCursor)
        self._hover = False
        self._fg_color = QColor(180, 180, 185)
        self._is_restored = False  # 用于max按钮的状态

    def set_restored(self, restored):
        self._is_restored = restored
        self.update()

    def paintEvent(self, event):
        painter = QPainter(self)
        painter.setRenderHint(QPainter.Antialiasing)
        r = self.rect()

        # 背景
        if self._hover:
            if self.mode == 'close':
                painter.fillRect(r, QColor(232, 74, 58))
            else:
                painter.fillRect(r, QColor(255, 255, 255, 25))
        else:
            painter.fillRect(r, Qt.transparent)

        # 图标颜色
        icon_color = QColor(255, 255, 255) if (self._hover and self.mode == 'close') else self._fg_color
        painter.setPen(QPen(icon_color, 2))

        cx, cy = r.center().x(), r.center().y()

        if self.mode == 'min':
            # 水平线 —
            painter.drawLine(cx - 5, cy, cx + 5, cy)
        elif self.mode == 'max':
            if self._is_restored:
                # 还原图标：重叠矩形
                painter.setPen(QPen(icon_color, 1.5))
                painter.drawRect(cx - 5, cy - 1, 8, 8)  # 后面的矩形
                pts = [QPoint(cx - 3, cy - 1), QPoint(cx - 3, cy - 4),
                       QPoint(cx + 5, cy - 4), QPoint(cx + 5, cy + 3), QPoint(cx + 3, cy + 3)]
                painter.drawPolyline(QPolygon(pts))
            else:
                # 最大化图标：矩形轮廓
                painter.setPen(QPen(icon_color, 1.5))
                painter.drawRect(cx - 5, cy - 5, 10, 10)
        elif self.mode == 'close':
            # X 关闭
            painter.drawLine(cx - 4, cy - 4, cx + 4, cy + 4)
            painter.drawLine(cx + 4, cy - 4, cx - 4, cy + 4)
        elif self.mode == 'return':
            # 返回图标 ↩
            painter.setPen(QPen(icon_color, 1.5))
            # 弧线
            painter.drawArc(cx - 5, cy - 5, 10, 10, 30 * 16, 300 * 16)
            # 箭头
            pts = [QPoint(cx - 5, cy - 2), QPoint(cx - 5, cy + 3), QPoint(cx - 1, cy)]
            painter.setBrush(icon_color)
            painter.drawPolygon(QPolygon(pts))
            painter.setBrush(Qt.NoBrush)

        painter.end()

    def enterEvent(self, event):
        self._hover = True
        self.update()

    def leaveEvent(self, event):
        self._hover = False
        self.update()

    def set_fg_color(self, color):
        self._fg_color = color
        self.update()


class _PinBtn(QPushButton):
    """图钉按钮 - QPainter 绘制的图钉图标，用于标签页独立显示"""
    def __init__(self, parent=None):
        super().__init__(parent)
        self.setFixedSize(20, 28)
        self.setCursor(Qt.PointingHandCursor)
        self.setToolTip('独立显示此标签页')
        self._hover = False
        self._pinned = False
        self._fg_color = QColor(180, 180, 185)
        self.setStyleSheet("QPushButton { background: transparent; border: none; }")

    def set_pinned(self, pinned):
        """设置图钉状态：pinned=True 表示已脱离为独立窗口"""
        self._pinned = pinned
        self.setToolTip('已独立显示' if pinned else '独立显示此标签页')
        self.update()

    def paintEvent(self, event):
        painter = QPainter(self)
        painter.setRenderHint(QPainter.Antialiasing)
        r = self.rect()

        # 背景
        if self._hover:
            painter.fillRect(r, QColor(255, 255, 255, 20))
        else:
            painter.fillRect(r, Qt.transparent)

        # 图钉图标颜色
        icon_color = QColor(63, 123, 247) if self._pinned else self._fg_color
        painter.setPen(QPen(icon_color, 1.5))

        cx = r.center().x()
        # 图钉头部（圆形/菱形）
        if self._pinned:
            # 已固定：实心图钉
            painter.setBrush(icon_color)
            painter.drawEllipse(QPointF(cx, r.top() + 8), 4, 4)
            painter.setBrush(Qt.NoBrush)
            # 图钉针身
            painter.drawLine(cx, r.top() + 12, cx, r.bottom() - 4)
            # 图钉尖端
            painter.setPen(QPen(icon_color, 1))
            painter.drawLine(cx - 2, r.bottom() - 6, cx, r.bottom() - 4)
            painter.drawLine(cx + 2, r.bottom() - 6, cx, r.bottom() - 4)
        else:
            # 未固定：空心图钉轮廓
            painter.drawEllipse(QPointF(cx, r.top() + 8), 4, 4)
            # 图钉针身
            painter.drawLine(cx, r.top() + 12, cx, r.bottom() - 4)
            # 图钉尖端
            painter.setPen(QPen(icon_color, 1))
            painter.drawLine(cx - 2, r.bottom() - 6, cx, r.bottom() - 4)
            painter.drawLine(cx + 2, r.bottom() - 6, cx, r.bottom() - 4)

        painter.end()

    def enterEvent(self, event):
        self._hover = True
        self.update()

    def leaveEvent(self, event):
        self._hover = False
        self.update()


class _PasswordLineEdit(QLineEdit):
    """带悬浮可见性切换按钮的密码输入框

    - 眼睛图标悬浮在输入框右侧末端
    - 输入框为空时自动隐藏图标
    - 点击图标切换密码可见性
    """
    def __init__(self, parent=None):
        super().__init__(parent)
        self.setEchoMode(QLineEdit.Password)
        # 悬浮切换按钮
        self._toggle_btn = QPushButton(self)
        self._toggle_btn.setCursor(Qt.PointingHandCursor)
        self._toggle_btn.setText('\U0001F441')
        self._toggle_btn.setCheckable(True)
        self._toggle_btn.setFocusPolicy(Qt.NoFocus)
        self._toggle_btn.setStyleSheet("""
            QPushButton { border: none; background: transparent; color: #888890; font-size: 14px; padding: 0 4px; }
            QPushButton:hover { color: #c0c0c8; }
            QPushButton:checked { color: #3f7bf7; }
        """)
        self._toggle_btn.clicked.connect(self._on_toggle)
        self._toggle_btn.hide()
        # 文本变化时显示/隐藏按钮
        self.textChanged.connect(self._on_text_changed)
        # 留出右侧空间给按钮
        self._btn_width = 28

    def _on_text_changed(self, text):
        """有文本时显示按钮，无文本时隐藏"""
        if text:
            self._toggle_btn.show()
        else:
            self._toggle_btn.hide()
            # 恢复密码模式
            if self._toggle_btn.isChecked():
                self._toggle_btn.setChecked(False)
                self.setEchoMode(QLineEdit.Password)

    def _on_toggle(self, checked):
        """切换密码可见性"""
        self.setEchoMode(QLineEdit.Normal if checked else QLineEdit.Password)

    def resizeEvent(self, event):
        """重新定位悬浮按钮到右侧末端"""
        super().resizeEvent(event)
        h = self.height()
        self._toggle_btn.setFixedHeight(h)
        self._toggle_btn.setGeometry(self.width() - self._btn_width, 0, self._btn_width, h)

    def setEchoMode(self, mode):
        """重写以同步按钮状态"""
        super().setEchoMode(mode)


class _CustomTitleBar(QWidget):
    """自定义标题栏，与主题配色一致"""
    return_clicked = pyqtSignal()

    def __init__(self, parent_window, title='', icon_text='RadAtlas', show_return=False):
        super().__init__(parent_window)
        self.parent_window = parent_window
        self._drag_pos = None
        self.setFixedHeight(36)
        self.setObjectName('customTitleBar')

        layout = QHBoxLayout(self)
        layout.setContentsMargins(10, 0, 6, 0)
        layout.setSpacing(6)

        # 图标/名称 - 使用更美观的图标样式
        icon_label = QLabel('RadAtlas')
        icon_label.setFont(QFont("Microsoft YaHei", 9, QFont.Bold))
        icon_label.setStyleSheet("background: transparent; color: #3f7bf7;")
        icon_label.setFixedWidth(80)
        layout.addWidget(icon_label)

        # 标题
        self.title_label = QLabel(title)
        self.title_label.setFont(QFont("Microsoft YaHei", 9))
        self.title_label.setStyleSheet("background: transparent;")
        layout.addWidget(self.title_label, stretch=1)

        # 返回按钮（可选，仅脱离窗口显示）
        if show_return:
            self.btn_return = _TitleBtn('return', self)
            self.btn_return.setToolTip('返回主窗口')
            self.btn_return.clicked.connect(self.return_clicked.emit)
            layout.addWidget(self.btn_return)
        else:
            self.btn_return = None

        # 最小化按钮
        self.btn_min = _TitleBtn('min', self)
        self.btn_min.clicked.connect(parent_window.showMinimized)
        layout.addWidget(self.btn_min)

        # 最大化/还原按钮
        self.btn_max = _TitleBtn('max', self)
        self.btn_max.clicked.connect(self._toggle_maximize)
        layout.addWidget(self.btn_max)

        # 关闭按钮
        self.btn_close = _TitleBtn('close', self)
        self.btn_close.clicked.connect(parent_window.close)
        layout.addWidget(self.btn_close)

    def set_title(self, title):
        self.title_label.setText(title)

    def _toggle_maximize(self, checked=False):
        if self.parent_window.isMaximized():
            self.parent_window.showNormal()
            self.btn_max.set_restored(False)
        else:
            self.parent_window.showMaximized()
            self.btn_max.set_restored(True)

    def update_theme(self, t):
        """根据主题更新标题栏样式"""
        title_bg = t.get('title_bg', t['bg2'])
        title_fg = t.get('title_fg', t['accent'])
        fg2 = t['fg2']
        self.setStyleSheet(f"""
            QWidget#customTitleBar {{
                background-color: {title_bg};
                border-bottom: 1px solid {t['border']};
            }}
            QLabel {{
                color: {title_fg};
                background: transparent;
            }}
        """)
        # 更新标题按钮图标颜色
        self.btn_min.set_fg_color(QColor(fg2))
        self.btn_max.set_fg_color(QColor(fg2))
        self.btn_close.set_fg_color(QColor(fg2))
        if self.btn_return:
            self.btn_return.set_fg_color(QColor(fg2))

    def mousePressEvent(self, event):
        if event.button() == Qt.LeftButton:
            self._drag_pos = event.globalPos() - self.parent_window.frameGeometry().topLeft()

    def mouseMoveEvent(self, event):
        if self._drag_pos and event.buttons() & Qt.LeftButton:
            if self.parent_window.isMaximized():
                self.parent_window.showNormal()
                self.btn_max.set_restored(False)
                # 按比例还原窗口位置
                screen = QApplication.desktop().availableGeometry()
                ratio = self._drag_pos.x() / screen.width()
                self.parent_window.move(event.globalPos() - self._drag_pos)
                self._drag_pos = event.pos()  # 重新计算
            self.parent_window.move(event.globalPos() - self._drag_pos)

    def mouseReleaseEvent(self, event):
        self._drag_pos = None

    def mouseDoubleClickEvent(self, event):
        self._toggle_maximize()


class LoginDialog(QDialog):
    def __init__(self, parent=None):
        super().__init__(parent)
        self.user_info = None
        self.user_password = None
        self.setWindowTitle(f'RadAtlas {APP_VERSION} - 登录')
        self.setFixedSize(400, 420)
        self.setWindowFlags(Qt.Dialog | Qt.FramelessWindowHint)

        # 登录窗口淡入动画
        self._login_opacity = QGraphicsOpacityEffect(self)
        self.setGraphicsEffect(self._login_opacity)
        self._login_opacity.setOpacity(0)
        self._login_anim = QPropertyAnimation(self._login_opacity, b"opacity")
        self._login_anim.setDuration(400)
        self._login_anim.setStartValue(0.0)
        self._login_anim.setEndValue(1.0)
        self._login_anim.setEasingCurve(QEasingCurve.InOutCubic)
        QTimer.singleShot(100, self._login_anim.start)

        layout = QVBoxLayout(self)
        layout.setContentsMargins(0, 0, 0, 0)
        layout.setSpacing(0)

        # 自定义标题栏
        self.title_bar = _CustomTitleBar(self, f'RadAtlas {APP_VERSION} - 登录')
        self.title_bar.update_theme(THEMES['深蓝暗夜'])
        layout.addWidget(self.title_bar)

        # 内容区域
        content = QVBoxLayout()
        content.setContentsMargins(40, 30, 40, 20)
        content.setSpacing(12)

        title = QLabel('RadAtlas')
        title.setFont(QFont("Microsoft YaHei", 28, QFont.Bold))
        title.setStyleSheet("color: #3f7bf7; background: transparent;")
        title.setAlignment(Qt.AlignCenter)
        content.addWidget(title)

        subtitle = QLabel('影像图鉴助手')
        subtitle.setFont(QFont("Microsoft YaHei", 12))
        subtitle.setStyleSheet("color: #666670; background: transparent;")
        subtitle.setAlignment(Qt.AlignCenter)
        content.addWidget(subtitle)

        content.addSpacing(20)

        self.username_input = QLineEdit()
        self.username_input.setPlaceholderText('用户名')
        self.username_input.setFixedHeight(44)
        content.addWidget(self.username_input)

        self.password_input = _PasswordLineEdit()
        self.password_input.setPlaceholderText('密码')
        self.password_input.setFixedHeight(44)
        self.password_input.returnPressed.connect(self.do_login)
        content.addWidget(self.password_input)

        content.addSpacing(8)

        login_btn = QPushButton('登 录')
        login_btn.setFixedHeight(44)
        login_btn.setFont(QFont("Microsoft YaHei", 14, QFont.Bold))
        login_btn.clicked.connect(self.do_login)
        content.addWidget(login_btn)

        content.addStretch()

        # 版本信息和制作人
        version_label = QLabel(f'RadAtlas {APP_VERSION}')
        version_label.setStyleSheet("color: #555560; font-size: 10px; background: transparent;")
        version_label.setAlignment(Qt.AlignCenter)
        content.addWidget(version_label)

        author_label = QLabel(APP_AUTHOR)
        author_label.setStyleSheet("color: #555560; font-size: 10px; background: transparent;")
        author_label.setAlignment(Qt.AlignCenter)
        content.addWidget(author_label)

        layout.addLayout(content)

    def do_login(self, checked=False):
        username = self.username_input.text().strip()
        password = self.password_input.text().strip()
        if not username or not password:
            QMessageBox.warning(self, '提示', '请输入用户名和密码')
            return
        try:
            user = authenticate_user(username, password)
        except Exception as e:
            QMessageBox.critical(self, '错误', f'认证失败:\n{e}')
            return
        if user:
            self.user_info = user
            self.user_password = password
            self.accept()
        else:
            QMessageBox.warning(self, '错误', '用户名或密码错误')


# ── AI配置对话框 ──────────────────────────────────────────
class AIConfigDialog(QDialog):
    def __init__(self, parent=None):
        super().__init__(parent)
        self.setWindowTitle('AI 大模型配置')
        self.setMinimumSize(550, 500)
        self._provider_rows = []
        self._init_ui()

    def _init_ui(self):
        layout = QVBoxLayout(self)
        layout.setContentsMargins(20, 20, 20, 20)
        layout.setSpacing(10)

        layout.addWidget(QLabel('配置大模型API（最多3个，AI诊断查询时可同时使用）'))

        # 模型配置列表区域
        self.config_scroll = QScrollArea()
        self.config_scroll.setWidgetResizable(True)
        self.config_widget = QWidget()
        self.config_layout = QVBoxLayout(self.config_widget)
        self.config_layout.setSpacing(8)
        self.config_layout.setContentsMargins(4, 4, 4, 4)
        self.config_scroll.setWidget(self.config_widget)
        layout.addWidget(self.config_scroll, stretch=1)

        # 加载已有配置
        providers = load_multi_config()
        for pc in providers:
            self._add_provider_row(pc)
        if not self._provider_rows:
            self._add_provider_row()

        # 添加按钮
        add_btn_row = QHBoxLayout()
        add_btn = QPushButton('+ 添加模型')
        add_btn.setObjectName('mutedBtn')
        add_btn.clicked.connect(lambda: self._add_provider_row())
        add_btn_row.addWidget(add_btn)
        add_btn_row.addStretch()
        layout.addLayout(add_btn_row)

        layout.addSpacing(10)

        # 按钮
        btn_row = QHBoxLayout()
        save_btn = QPushButton('保存')
        save_btn.clicked.connect(self._save)
        cancel_btn = QPushButton('取消')
        cancel_btn.setObjectName('mutedBtn')
        cancel_btn.clicked.connect(self.reject)
        btn_row.addStretch()
        btn_row.addWidget(cancel_btn)
        btn_row.addWidget(save_btn)
        layout.addLayout(btn_row)

    def _add_provider_row(self, config=None):
        """添加一行模型配置"""
        if config is None:
            config = {}
        if len(self._provider_rows) >= 3:
            QMessageBox.information(self, '提示', '最多配置3个模型')
            return

        row_widget = QWidget()
        row_widget.setStyleSheet("background: rgba(255,255,255,4); border-radius: 6px; padding: 8px;")
        row_layout = QGridLayout(row_widget)
        row_layout.setSpacing(6)

        # 启用复选框
        enabled_cb = QCheckBox('启用')
        enabled_cb.setChecked(config.get('enabled', True))
        row_layout.addWidget(enabled_cb, 0, 0)

        # 提供商
        row_layout.addWidget(QLabel('提供商:'), 0, 1)
        provider_combo = QComboBox()
        for pid, pinfo in PROVIDERS.items():
            provider_combo.addItem(pinfo['name'], pid)
        cur_provider = config.get('provider', 'deepseek')
        idx = list(PROVIDERS.keys()).index(cur_provider) if cur_provider in PROVIDERS else 0
        provider_combo.setCurrentIndex(idx)
        row_layout.addWidget(provider_combo, 0, 2)

        # 模型（可编辑，支持自定义模型名）
        row_layout.addWidget(QLabel('模型:'), 1, 1)
        model_combo = QComboBox()
        model_combo.setEditable(True)
        row_layout.addWidget(model_combo, 1, 2)

        # API Key
        row_layout.addWidget(QLabel('API Key:'), 2, 1)
        api_key_input = _PasswordLineEdit()
        api_key_input.setText(config.get('api_key', ''))
        api_key_input.setPlaceholderText('输入API Key')
        row_layout.addWidget(api_key_input, 2, 2)

        # 自定义URL
        row_layout.addWidget(QLabel('自定义URL:'), 3, 1)
        custom_url_input = QLineEdit()
        custom_url_input.setText(config.get('custom_api_url', ''))
        custom_url_input.setPlaceholderText('留空使用默认')
        row_layout.addWidget(custom_url_input, 3, 2)

        # 删除按钮
        del_btn = QPushButton('删除')
        del_btn.setObjectName('dangerBtn')
        del_btn.setFixedWidth(50)
        del_btn.clicked.connect(lambda: self._remove_provider_row(row_widget))
        row_layout.addWidget(del_btn, 0, 3)

        # 更新模型列表
        def update_models(idx):
            model_combo.clear()
            pid = provider_combo.currentData()
            if pid and pid in PROVIDERS:
                for model in PROVIDERS[pid]['models']:
                    model_combo.addItem(model)

        provider_combo.currentIndexChanged.connect(update_models)
        update_models(0)

        # 设置当前模型
        cur_model = config.get('model', '')
        for i in range(model_combo.count()):
            if model_combo.itemText(i) == cur_model:
                model_combo.setCurrentIndex(i)
                break

        row = {
            'widget': row_widget,
            'enabled_cb': enabled_cb,
            'provider_combo': provider_combo,
            'model_combo': model_combo,
            'api_key_input': api_key_input,
            'custom_url_input': custom_url_input,
        }
        self._provider_rows.append(row)
        self.config_layout.addWidget(row_widget)

        # 测试按钮
        test_btn = QPushButton('测试')
        test_btn.setFixedWidth(50)
        test_btn.clicked.connect(lambda checked=False, r=row, b=test_btn: self._test_provider(r, b))
        row_layout.addWidget(test_btn, 1, 3)

    def _remove_provider_row(self, widget):
        """删除一行模型配置"""
        self.config_layout.removeWidget(widget)
        self._provider_rows = [r for r in self._provider_rows if r['widget'] != widget]
        widget.deleteLater()

    def _test_provider(self, row, btn):
        """测试单个LLM提供商联通性"""
        provider = row['provider_combo'].currentData() or 'deepseek'
        model = row['model_combo'].currentText() or 'deepseek-chat'
        api_key = row['api_key_input'].text().strip()
        custom_url = row['custom_url_input'].text().strip()

        btn.setEnabled(False)
        btn.setText('测试中...')

        self._test_thread = QThread()
        self._test_worker = _LLMTestWorker(provider, model, api_key, custom_url)
        self._test_worker.moveToThread(self._test_thread)
        self._test_thread.started.connect(self._test_worker.run)

        def _on_finished(result):
            btn.setEnabled(True)
            btn.setText('测试')
            if result['success']:
                QMessageBox.information(self, '测试成功', result['message'])
            else:
                QMessageBox.warning(self, '测试失败', result['message'])
            self._test_thread.quit()

        self._test_worker.finished.connect(_on_finished)
        self._test_thread.start()

    def _save(self, checked=False):
        """保存配置"""
        providers = []
        for row in self._provider_rows:
            pc = {
                'enabled': row['enabled_cb'].isChecked(),
                'provider': row['provider_combo'].currentData() or 'deepseek',
                'model': row['model_combo'].currentText() or 'deepseek-chat',
                'api_key': row['api_key_input'].text().strip(),
                'custom_api_url': row['custom_url_input'].text().strip(),
            }
            providers.append(pc)

        config = load_config()
        config['providers'] = providers
        # 兼容旧格式
        if providers:
            config['provider'] = providers[0].get('provider', 'deepseek')
            config['model'] = providers[0].get('model', 'deepseek-chat')
            config['api_key'] = providers[0].get('api_key', '')
            config['custom_api_url'] = providers[0].get('custom_api_url', '')
        save_config(config)
        QMessageBox.information(self, '成功', 'AI配置已保存')
        self.accept()


# ── AI影像诊断查询对话框 ──────────────────────────────────
class AIDiagnosisDialog(QDialog):
    """AI影像诊断查询对话框 - 支持多模型并行查询和知识库"""
    diagnosis_clicked = pyqtSignal(dict)  # 点击诊断条目时发出

    def __init__(self, parent=None):
        super().__init__(parent)
        self.setWindowTitle('AI 影像诊断查询')
        self.setMinimumSize(900, 700)
        self.resize(1100, 750)
        self.setWindowFlags(Qt.Window | Qt.FramelessWindowHint)
        self._results = {}
        self._init_ui()

    def _init_ui(self):
        layout = QVBoxLayout(self)
        layout.setContentsMargins(0, 0, 0, 0)
        layout.setSpacing(0)

        # 自定义标题栏
        self.title_bar = _CustomTitleBar(self, 'AI 影像诊断查询  |  楚雄州人民医院 医学影像中心 张兴文')
        layout.addWidget(self.title_bar)

        # 内容区域
        content = QWidget()
        content_layout = QVBoxLayout(content)
        content_layout.setContentsMargins(16, 16, 16, 16)
        content_layout.setSpacing(10)

        # ── 查询参数区 ──
        param_group = QGroupBox('查询参数')
        param_group.setStyleSheet("""
            QGroupBox { color: #b0b0b8; font-weight: bold; border: 1px solid #2a2a30;
                        border-radius: 6px; margin-top: 8px; padding-top: 16px; }
            QGroupBox::title { subcontrol-origin: margin; left: 12px; padding: 0 6px; }
        """)
        param_layout = QGridLayout(param_group)
        param_layout.setSpacing(8)

        # 检查类型
        param_layout.addWidget(QLabel('检查类型:'), 0, 0)
        self.exam_combo = QComboBox()
        self.exam_combo.addItems(['CT', 'X-Ray', 'MRI', 'PET-CT', '超声', 'DSA', 'ECG(心电图)'])
        self.exam_combo.setFixedWidth(120)
        param_layout.addWidget(self.exam_combo, 0, 1)

        # 关键字
        param_layout.addWidget(QLabel('关键征象:'), 0, 2)
        self.keywords_input = QLineEdit()
        self.keywords_input.setPlaceholderText('输入关键字，如：腹部 胃 包块 等密度...')
        self.keywords_input.returnPressed.connect(self._do_query)
        param_layout.addWidget(self.keywords_input, 0, 3, 1, 3)

        # 大模型选择
        param_layout.addWidget(QLabel('选择大模型:'), 1, 0)
        self.model_checks_layout = QHBoxLayout()
        self.model_checks = []
        providers = load_multi_config()
        for i, pc in enumerate(providers):
            if pc.get('api_key'):
                name = PROVIDERS.get(pc['provider'], {}).get('name', pc['provider'])
                cb = QCheckBox(f'{name} ({pc.get("model", "")})')
                # 默认只选中前2个模型
                cb.setChecked(i < 2)
                cb.setProperty('provider_config', pc)
                self.model_checks.append(cb)
                self.model_checks_layout.addWidget(cb)
        if not self.model_checks:
            no_model = QLabel('未配置AI模型，请先在"AI配置"中设置')
            no_model.setStyleSheet("color: #e64a3a;")
            self.model_checks_layout.addWidget(no_model)
        self.model_checks_layout.addStretch()
        param_layout.addLayout(self.model_checks_layout, 1, 1, 1, 5)

        # 知识库选择
        param_layout.addWidget(QLabel('知识库:'), 2, 0)
        kb_layout = QHBoxLayout()
        self.kb_none = QRadioButton('不使用')
        self.kb_tencent = QRadioButton('腾讯知识库')
        self.kb_volcengine = QRadioButton('火山方舟知识库')
        self.kb_notebooklm = QRadioButton('NotebookLM')
        # 默认选中火山方舟（若已配置），否则选第一个有配置的，否则选中"不使用"
        kb_configs = load_config().get('knowledge_bases', {})
        default_kb_selected = False
        for kb_type, kb_radio in [('volcengine', self.kb_volcengine), ('tencent', self.kb_tencent), ('notebooklm', self.kb_notebooklm)]:
            kb_cfg = kb_configs.get(kb_type, {})
            has_creds = kb_cfg.get('api_key') or (kb_cfg.get('access_key') and kb_cfg.get('secret_key'))
            if has_creds and not default_kb_selected:
                kb_radio.setChecked(True)
                default_kb_selected = True
        if not default_kb_selected:
            self.kb_none.setChecked(True)
        kb_layout.addWidget(self.kb_none)
        kb_layout.addWidget(self.kb_tencent)
        kb_layout.addWidget(self.kb_volcengine)
        kb_layout.addWidget(self.kb_notebooklm)

        self.kb_config_btn = QPushButton('知识库配置')
        self.kb_config_btn.setObjectName('mutedBtn')
        self.kb_config_btn.setFixedWidth(90)
        self.kb_config_btn.clicked.connect(self._show_kb_config)
        kb_layout.addWidget(self.kb_config_btn)
        self.help_btn = QPushButton('帮助')
        self.help_btn.setObjectName('mutedBtn')
        self.help_btn.setFixedWidth(60)
        self.help_btn.clicked.connect(self._show_help)
        kb_layout.addWidget(self.help_btn)
        kb_layout.addStretch()
        param_layout.addLayout(kb_layout, 2, 1, 1, 5)

        content_layout.addWidget(param_group)

        # ── 查询按钮 ──
        btn_row = QHBoxLayout()
        self.query_btn = QPushButton('开始查询')
        self.query_btn.setFixedSize(120, 36)
        self.query_btn.clicked.connect(self._do_query)
        btn_row.addStretch()
        btn_row.addWidget(self.query_btn)

        self.status_label = QLabel('')
        self.status_label.setStyleSheet("color: #888890; font-size: 12px;")
        btn_row.addWidget(self.status_label)
        content_layout.addLayout(btn_row)

        # ── 结果展示区 ──
        self.result_splitter = QSplitter(Qt.Horizontal)
        # 左侧：诊断条目列表（按大模型分列）— 包裹在 QScrollArea 中支持滚动
        self.result_list_panel = QWidget()
        outer_list_layout = QVBoxLayout(self.result_list_panel)
        outer_list_layout.setContentsMargins(0, 0, 0, 0)
        outer_list_layout.setSpacing(4)
        self.result_list_label = QLabel('诊断条目')
        self.result_list_label.setStyleSheet("color: #888890; font-size: 12px; font-weight: bold;")
        outer_list_layout.addWidget(self.result_list_label)
        self.result_list_scroll = QScrollArea()
        self.result_list_scroll.setWidgetResizable(True)
        self.result_list_scroll.setFrameShape(QFrame.NoFrame)
        self.result_list_scroll.setStyleSheet("QScrollArea { background: transparent; }")
        scroll_content = QWidget()
        scroll_content.setStyleSheet("background: transparent;")
        self.result_list_layout = QVBoxLayout(scroll_content)
        self.result_list_layout.setContentsMargins(0, 0, 0, 0)
        self.result_list_layout.setSpacing(4)
        self.result_list_layout.addWidget(QLabel('请输入关键字后点击"开始查询"'))
        self.result_list_scroll.setWidget(scroll_content)
        outer_list_layout.addWidget(self.result_list_scroll)
        self.result_splitter.addWidget(self.result_list_panel)

        # 右侧：诊断详情
        self.detail_panel = QWidget()
        detail_layout = QVBoxLayout(self.detail_panel)
        detail_layout.setContentsMargins(0, 0, 0, 0)
        detail_layout.setSpacing(4)
        # 详情工具栏（缩放 + 收藏）
        detail_toolbar = QHBoxLayout()
        detail_toolbar.setContentsMargins(0, 0, 0, 0)
        detail_label = QLabel('诊断详情')
        detail_label.setStyleSheet("color: #888890; font-size: 12px; font-weight: bold;")
        detail_toolbar.addWidget(detail_label)
        detail_toolbar.addStretch()
        self._detail_font_size = 13
        zoom_in_btn = QPushButton('＋')
        zoom_in_btn.setFixedSize(28, 24)
        zoom_in_btn.setToolTip('放大字体')
        zoom_in_btn.setStyleSheet("QPushButton { font-size:14px; font-weight:bold; padding:0; }")
        zoom_in_btn.clicked.connect(lambda checked=False: self._zoom_detail(1))
        detail_toolbar.addWidget(zoom_in_btn)
        zoom_out_btn = QPushButton('－')
        zoom_out_btn.setFixedSize(28, 24)
        zoom_out_btn.setToolTip('缩小字体')
        zoom_out_btn.setStyleSheet("QPushButton { font-size:14px; font-weight:bold; padding:0; }")
        zoom_out_btn.clicked.connect(lambda checked=False: self._zoom_detail(-1))
        detail_toolbar.addWidget(zoom_out_btn)
        self.fav_btn = QPushButton('☆ 收藏')
        self.fav_btn.setFixedHeight(24)
        self.fav_btn.setToolTip('收藏当前文档')
        self.fav_btn.setStyleSheet("QPushButton { font-size:11px; padding:2px 8px; }")
        self.fav_btn.clicked.connect(self._toggle_favorite)
        detail_toolbar.addWidget(self.fav_btn)
        detail_layout.addLayout(detail_toolbar)
        self.detail_browser = QTextBrowser()
        self.detail_browser.setOpenExternalLinks(False)
        self.detail_browser.anchorClicked.connect(self._on_detail_link_clicked)
        self.detail_browser.setStyleSheet("""
            QTextBrowser { background: #1a1a1e; border: 1px solid #2a2a30;
                          border-radius: 6px; color: #d0d0d4; padding: 12px; font-size: 13px; }
        """)
        self.detail_browser.setHtml('<p style="color:#666670;">点击左侧诊断条目查看详情</p>')
        detail_layout.addWidget(self.detail_browser)
        self.result_splitter.addWidget(self.detail_panel)
        self.result_splitter.setSizes([400, 600])

        # 初始化收藏状态
        self._current_kb_doc = None
        self._detail_font_size = 13
        self.fav_btn.setEnabled(False)

        content_layout.addWidget(self.result_splitter, stretch=1)

        layout.addWidget(content, stretch=1)

    def _do_query(self):
        """执行AI诊断查询"""
        keywords = self.keywords_input.text().strip()
        if not keywords:
            QMessageBox.warning(self, '提示', '请输入关键征象')
            return

        # 获取选中的大模型
        selected = [cb.property('provider_config') for cb in self.model_checks if cb.isChecked()]

        # 获取知识库配置
        kb_config = None
        if self.kb_tencent.isChecked():
            kb_config = self._load_kb_config('tencent')
        elif self.kb_volcengine.isChecked():
            kb_config = self._load_kb_config('volcengine')
        elif self.kb_notebooklm.isChecked():
            kb_config = self._load_kb_config('notebooklm')

        # 未选择大模型时，必须配置并选择知识库（支持直接知识库检索）
        if not selected:
            if not kb_config or not kb_config.get('api_key'):
                QMessageBox.warning(self, '提示', '请至少选择一个大模型，或选择并配置知识库后进行直接检索')
                return

        exam_type = self.exam_combo.currentText()
        self.query_btn.setEnabled(False)
        self.status_label.setText('正在查询中，请稍候...')
        self.status_label.setStyleSheet("color: #3f7bf7; font-size: 12px;")

        # 在后台线程中执行查询
        self._query_thread = QThread()
        self._query_worker = _DiagnosisWorker(exam_type, keywords, selected, kb_config)
        self._query_worker.moveToThread(self._query_thread)
        self._query_thread.started.connect(self._query_worker.run)
        self._query_worker.finished.connect(self._on_query_finished)
        self._query_worker.finished.connect(self._query_thread.quit)
        self._query_thread.start()

    def _on_query_finished(self, results):
        """查询完成回调"""
        self.query_btn.setEnabled(True)
        self._results = results

        # 提取按疾病分别检索的知识库结果
        self._kb_groups = results.pop('_kb_groups', [])
        kb_error = results.pop('_kb_error', '')

        # 收集所有知识库文档快照（用于在诊断详情中显示）
        self._kb_docs = []
        seen = set()
        for group in self._kb_groups:
            for doc in group.get('docs', []):
                key = (doc.get('doc_name', ''), doc.get('page', ''), doc.get('snippet', '')[:50])
                if key not in seen:
                    seen.add(key)
                    self._kb_docs.append(doc)
        kb_errors = [kb_error] if kb_error else []

        # 区分大模型结果与纯知识库检索模式
        provider_results = {k: v for k, v in results.items() if not k.startswith('_')}
        has_llm_results = len(provider_results) > 0

        # 检查是否全部失败
        if has_llm_results:
            all_failed = all(not r.get('success') for r in provider_results.values())
        else:
            all_failed = not self._kb_docs

        if all_failed:
            self.status_label.setText('查询失败')
            self.status_label.setStyleSheet("color: #e64a3a; font-size: 12px;")
            if has_llm_results:
                errors = '\n'.join(f'{k}: {v.get("error", "未知错误")}' for k, v in provider_results.items())
                msg = f'所有大模型查询均失败：\n{errors}'
            else:
                msg = f'知识库查询失败：\n{kb_error or "未知错误"}'
            QMessageBox.warning(self, '查询失败', msg)
            return

        success_count = sum(1 for r in provider_results.values() if r.get('success'))
        if self._kb_docs:
            kb_info = f'，引用{len(self._kb_docs)}篇文档'
        elif kb_errors:
            kb_info = f'，知识库查询失败'
        else:
            kb_info = ''
        if has_llm_results:
            self.status_label.setText(f'查询完成，{success_count}/{len(provider_results)}个模型成功{kb_info}')
        else:
            self.status_label.setText(f'知识库检索完成，引用{len(self._kb_docs)}篇文档')
        self.status_label.setStyleSheet("color: #4ade80; font-size: 12px;" if not kb_errors else "color: #f59e0b; font-size: 12px;")

        # 如果知识库查询失败，显示错误提示
        if kb_errors and not self._kb_docs:
            kb_type_name = ''
            if self.kb_tencent.isChecked():
                kb_type_name = '腾讯知识库'
            elif self.kb_volcengine.isChecked():
                kb_type_name = '火山方舟知识库'
            elif self.kb_notebooklm.isChecked():
                kb_type_name = 'Google NotebookLM'
            QMessageBox.warning(self, '知识库查询失败',
                f'{kb_type_name}查询失败：\n{kb_errors[0]}\n\n诊断结果仍已显示，但未包含知识库参考文档。')

        # 清空结果列表（result_list_label 已在外层布局，无需保留）
        while self.result_list_layout.count():
            item = self.result_list_layout.takeAt(0)
            w = item.widget()
            if w:
                w.deleteLater()

        # 按大模型分列显示结果
        for provider_name, result in provider_results.items():
            # 模型标题
            title = QLabel(f'  {provider_name}')
            title.setStyleSheet("""
                color: #3f7bf7; font-size: 13px; font-weight: bold;
                background: rgba(63,123,247,20); border-radius: 4px; padding: 4px 8px;
            """)
            self.result_list_layout.addWidget(title)

            if not result.get('success'):
                err_label = QLabel(f'    查询失败: {result.get("error", "未知错误")}')
                err_label.setStyleSheet("color: #e64a3a; font-size: 11px; padding: 2px 8px;")
                err_label.setWordWrap(True)
                self.result_list_layout.addWidget(err_label)
                continue

            for i, item_data in enumerate(result.get('data', [])):
                name = item_data.get('disease_name', f'诊断 #{i+1}')
                confidence = item_data.get('confidence', '')
                btn = QPushButton(f'    {name}  [{confidence}]')
                btn.setFlat(True)
                btn.setStyleSheet("""
                    QPushButton { text-align: left; color: #c0c0c8; padding: 5px 8px;
                                  border: none; border-radius: 4px; font-size: 12px; }
                    QPushButton:hover { background: rgba(63,123,247,30); color: #e0e0e4; }
                """)
                btn.setProperty('diagnosis_data', item_data)
                btn.setProperty('provider_name', provider_name)
                btn.clicked.connect(lambda checked, d=item_data, p=provider_name: self._show_detail(d, p))
                self.result_list_layout.addWidget(btn)

        # 知识库文档条目（按疾病分别检索，分组显示）
        if self._kb_groups:
            kb_title = QLabel('  知识库参考文档（按疾病分别检索）')
            kb_title.setStyleSheet("""
                color: #f59e0b; font-size: 13px; font-weight: bold;
                background: rgba(245,158,11,15); border-radius: 4px; padding: 4px 8px;
                margin-top: 8px;
            """)
            self.result_list_layout.addWidget(kb_title)

            for group in self._kb_groups:
                disease_name = group.get('disease', '')
                docs = group.get('docs', [])
                group_err = group.get('error', '')

                # 疾病分组标题
                group_label = QLabel(f'    ─ {disease_name}')
                group_label.setStyleSheet("color: #f59e0b; font-size: 12px; font-weight: bold; padding: 4px 8px;")
                self.result_list_layout.addWidget(group_label)

                if group_err and not docs:
                    err_label = QLabel(f'        检索失败: {group_err[:80]}')
                    err_label.setStyleSheet("color: #e64a3a; font-size: 11px; padding: 2px 8px;")
                    err_label.setWordWrap(True)
                    self.result_list_layout.addWidget(err_label)
                    continue

                for i, doc in enumerate(docs):
                    doc_name = doc.get('doc_name', f'文档 #{i+1}')
                    source = doc.get('source', '')
                    page = doc.get('page', '')
                    label_text = f'        {doc_name}'
                    if page:
                        label_text += f'  [第{page}页]'
                    elif source:
                        label_text += f'  [{source}]'
                    btn = QPushButton(label_text)
                    btn.setFlat(True)
                    btn.setStyleSheet("""
                        QPushButton { text-align: left; color: #f0c060; padding: 5px 8px;
                                      border: none; border-radius: 4px; font-size: 12px; }
                        QPushButton:hover { background: rgba(245,158,11,25); color: #ffd070; }
                    """)
                    btn.clicked.connect(lambda checked, d=doc: self._show_kb_detail(d))
                    self.result_list_layout.addWidget(btn)

        self.result_list_layout.addStretch()

    def _show_detail(self, data, provider_name):
        """显示诊断详情"""
        self._current_kb_doc = None
        self._update_fav_btn()
        fields = [
            ('disease_name', '疾病名称'),
            ('confidence', '置信度'),
            ('imaging_findings', '影像学表现'),
            ('report_template', '标准报告模板'),
            ('differential_diagnosis', '鉴别诊断'),
            ('clinical_manifestation', '临床表现'),
            ('pathophysiology', '病理生理及症状学特征'),
            ('treatment', '临床治疗方法'),
        ]
        html = f'<div style="font-family: Microsoft YaHei;">'
        html += f'<h3 style="color:#3f7bf7; border-bottom:1px solid #2a2a30; padding-bottom:8px;">'
        html += f'{data.get("disease_name", "未知")} <small style="color:#888890;">— {provider_name}</small></h3>'

        for key, label in fields:
            value = data.get(key, '')
            if not value or key == 'disease_name':
                continue
            html += f'<h4 style="color:#5a91ff; margin-top:16px;">{label}</h4>'
            html += f'<div style="color:#c0c0c8; line-height:1.8; padding:4px 8px; '
            html += f'background:rgba(255,255,255,3); border-radius:4px; white-space:pre-wrap;">'
            html += f'{value}</div>'

        # 知识库文档快照区域
        kb_docs = getattr(self, '_kb_docs', [])
        if kb_docs:
            html += f'<h4 style="color:#f59e0b; margin-top:20px; border-top:1px solid #2a2a30; padding-top:12px;">'
            html += f'&#128218; 知识库参考文档 ({len(kb_docs)}篇)</h4>'
            for i, doc in enumerate(kb_docs):
                doc_name = doc.get('doc_name', '未知文档')
                page = doc.get('page', '')
                snippet = doc.get('snippet', '')
                url = doc.get('url', '')
                source = doc.get('source', '')

                # 文档卡片
                html += f'<div style="margin:8px 0; padding:10px 12px; background:rgba(245,158,11,8); '
                html += f'border:1px solid rgba(245,158,11,25); border-radius:6px; '
                html += f'border-left:3px solid #f59e0b;">'

                # 文档名称和页码
                html += f'<div style="color:#f59e0b; font-weight:bold; font-size:13px;">'
                if url:
                    html += f'<a href="{url}" style="color:#f59e0b; text-decoration:none;">{doc_name}</a>'
                else:
                    html += f'{doc_name}'
                if page:
                    html += f' <span style="color:#888890; font-weight:normal; font-size:11px;">第{page}页</span>'
                html += f'</div>'

                # 来源标签
                if source:
                    html += f'<span style="display:inline-block; margin-top:4px; padding:1px 6px; '
                    html += f'background:rgba(63,123,247,15); color:#5a91ff; border-radius:3px; font-size:10px;">{source}</span>'

                # 内容片段
                if snippet:
                    html += f'<div style="color:#a0a0a8; font-size:12px; margin-top:6px; '
                    html += f'line-height:1.6; white-space:pre-wrap; max-height:80px; overflow:hidden;">{snippet}</div>'

                # 链接按钮
                if url:
                    html += f'<div style="margin-top:6px;">'
                    html += f'<a href="{url}" style="color:#3f7bf7; font-size:11px; text-decoration:none;">'
                    html += f'&#128279; 查看原文</a></div>'

                html += f'</div>'

        html += '</div>'
        self.detail_browser.setHtml(html)
        # 发出信号
        self.diagnosis_clicked.emit(data)

    def _show_kb_detail(self, doc):
        """显示知识库文档快照详情"""
        doc_name = doc.get('doc_name', '未知文档')
        page = doc.get('page', '')
        snippet = doc.get('snippet', '')
        url = doc.get('url', '')
        source = doc.get('source', '')

        html = '<div style="font-family: Microsoft YaHei;">'
        html += '<h3 style="color:#f59e0b; border-bottom:1px solid #2a2a30; padding-bottom:8px;">'
        html += f'&#128218; {doc_name}</h3>'

        # 元信息表格
        html += '<div style="margin:12px 0; padding:12px; background:rgba(245,158,11,8); '
        html += 'border:1px solid rgba(245,158,11,25); border-radius:6px; '
        html += 'border-left:3px solid #f59e0b;">'

        # 文档名称
        html += f'<div style="margin-bottom:8px;">'
        html += f'<span style="color:#888890; font-size:12px; display:inline-block; width:80px;">文档名称</span>'
        html += f'<span style="color:#f0c060; font-size:13px; font-weight:bold;">{doc_name}</span>'
        html += '</div>'

        # 来源
        if source:
            html += f'<div style="margin-bottom:8px;">'
            html += f'<span style="color:#888890; font-size:12px; display:inline-block; width:80px;">来源</span>'
            html += f'<span style="display:inline-block; padding:2px 8px; background:rgba(63,123,247,15); '
            html += f'color:#5a91ff; border-radius:3px; font-size:11px;">{source}</span>'
            html += '</div>'

        # 页码/章节
        if page:
            html += f'<div style="margin-bottom:8px;">'
            html += f'<span style="color:#888890; font-size:12px; display:inline-block; width:80px;">页码/位置</span>'
            html += f'<span style="color:#c0c0c8; font-size:12px;">第 {page} 页</span>'
            html += '</div>'

        # 链接
        if url:
            html += f'<div style="margin-bottom:8px;">'
            html += f'<span style="color:#888890; font-size:12px; display:inline-block; width:80px;">原文链接</span>'
            html += f'<a href="{url}" style="color:#3f7bf7; font-size:12px; text-decoration:none;">'
            html += f'&#128279; {url[:80]}{"..." if len(url) > 80 else ""}</a>'
            html += '</div>'

        html += '</div>'

        # 内容摘要 — 以链接形式显示，点击打开查看器
        if snippet:
            images = doc.get('images', []) or []
            img_count = len(images)
            html += '<h4 style="color:#5a91ff; margin-top:16px;">切片/快照内容</h4>'
            html += '<div style="padding:16px; background:rgba(255,255,255,3); border-radius:6px; '
            html += 'border-left:3px solid #5a91ff; text-align:center;">'
            # 简短预览（前80字）
            preview = snippet[:80].replace('\n', ' ').strip()
            if len(snippet) > 80:
                preview += '...'
            html += f'<p style="color:#888890; font-size:12px; margin:0 0 10px 0;">{preview}</p>'
            # 链接文本：有图片时提示含图片
            link_text = '&#128196; 查看完整切片内容（可缩放）'
            if img_count:
                link_text = f'&#128196; 查看完整切片内容（含{img_count}张原书图片，可缩放）'
            html += '<a href="snippet://view" style="display:inline-block; padding:8px 24px; '
            html += 'background:rgba(245,158,11,20); color:#f0c060; border-radius:6px; '
            html += f'text-decoration:none; font-size:13px; font-weight:bold;">{link_text}</a>'
            html += '</div>'

        # 如果有链接，添加查看原文按钮
        if url:
            html += '<div style="margin-top:16px; text-align:center;">'
            html += f'<a href="{url}" style="display:inline-block; padding:8px 24px; '
            html += 'background:rgba(63,123,247,20); color:#5a91ff; border-radius:6px; '
            html += 'text-decoration:none; font-size:13px;">&#128279; 查看原文</a>'
            html += '</div>'

        html += '</div>'
        self._current_kb_doc = doc
        self._update_fav_btn()
        self.detail_browser.setHtml(html)

    def _zoom_detail(self, delta):
        """放大/缩小详情字体"""
        self._detail_font_size = max(8, min(28, self._detail_font_size + delta))
        self.detail_browser.setStyleSheet(f"""
            QTextBrowser {{ background: #1a1a1e; border: 1px solid #2a2a30;
                          border-radius: 6px; color: #d0d0d4; padding: 12px;
                          font-size: {self._detail_font_size}px; }}
        """)

    def _on_detail_link_clicked(self, url):
        """处理详情面板中的链接点击"""
        from PyQt5.QtGui import QDesktopServices
        if url.scheme() == 'snippet':
            # 切片内容查看链接 → 打开查看器
            doc = getattr(self, '_current_kb_doc', None)
            if doc:
                viewer = _KBSnippetViewer(doc, self)
                viewer.exec_()
        else:
            # 外部链接 → 系统浏览器打开
            QDesktopServices.openUrl(url)

    def _toggle_favorite(self, checked=False):
        """收藏/取消收藏当前知识库文档"""
        doc = getattr(self, '_current_kb_doc', None)
        if not doc:
            return
        favorites = self._load_favorites()
        doc_key = f"{doc.get('doc_name', '')}|{doc.get('page', '')}|{doc.get('snippet', '')[:50]}"
        if doc_key in favorites:
            del favorites[doc_key]
        else:
            favorites[doc_key] = {
                'doc_name': doc.get('doc_name', ''),
                'page': doc.get('page', ''),
                'snippet': doc.get('snippet', ''),
                'url': doc.get('url', ''),
                'source': doc.get('source', ''),
            }
        self._save_favorites(favorites)
        self._update_fav_btn()

    def _update_fav_btn(self):
        """更新收藏按钮状态"""
        doc = getattr(self, '_current_kb_doc', None)
        if not doc:
            self.fav_btn.setText('☆ 收藏')
            self.fav_btn.setEnabled(False)
            return
        self.fav_btn.setEnabled(True)
        favorites = self._load_favorites()
        doc_key = f"{doc.get('doc_name', '')}|{doc.get('page', '')}|{doc.get('snippet', '')[:50]}"
        if doc_key in favorites:
            self.fav_btn.setText('★ 已收藏')
        else:
            self.fav_btn.setText('☆ 收藏')

    def _load_favorites(self):
        """加载收藏的知识库文档"""
        import json as _json
        fav_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'kb_favorites.json')
        if os.path.exists(fav_path):
            try:
                with open(fav_path, 'r', encoding='utf-8') as f:
                    return _json.load(f)
            except Exception:
                return {}
        return {}

    def _save_favorites(self, favorites):
        """保存收藏的知识库文档"""
        import json as _json
        fav_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'kb_favorites.json')
        try:
            with open(fav_path, 'w', encoding='utf-8') as f:
                _json.dump(favorites, f, ensure_ascii=False, indent=2)
        except Exception:
            pass

    def _show_kb_config(self, checked=False):
        """显示知识库配置"""
        dlg = _KBConfigDialog(self)
        if dlg.exec_() == QDialog.Accepted:
            pass

    def _show_help(self, checked=False):
        """显示帮助文档"""
        dlg = _HelpDialog(self)
        dlg.exec_()

    def _load_kb_config(self, kb_type):
        """加载知识库配置"""
        config = load_config()
        kb_configs = config.get('knowledge_bases', {})
        return kb_configs.get(kb_type, {'type': kb_type, 'api_key': ''})


class _LLMTestWorker(QObject):
    """LLM联通测试后台工作线程"""
    finished = pyqtSignal(dict)

    def __init__(self, provider, model, api_key, custom_url):
        super().__init__()
        self.provider = provider
        self.model = model
        self.api_key = api_key
        self.custom_url = custom_url

    def run(self):
        result = test_llm_connection(self.provider, self.model, self.api_key, self.custom_url)
        self.finished.emit(result)


class _KBTestWorker(QObject):
    """知识库联通测试后台工作线程"""
    finished = pyqtSignal(dict)

    def __init__(self, kb_config):
        super().__init__()
        self.kb_config = kb_config

    def run(self):
        result = test_kb_connection(self.kb_config)
        self.finished.emit(result)


class _DiagnosisWorker(QObject):
    """后台诊断查询工作线程"""
    finished = pyqtSignal(dict)

    def __init__(self, exam_type, keywords, selected_providers, kb_config):
        super().__init__()
        self.exam_type = exam_type
        self.keywords = keywords
        self.selected_providers = selected_providers
        self.kb_config = kb_config

    def run(self):
        try:
            results = call_diagnosis_multi(
                self.exam_type, self.keywords,
                selected_providers=self.selected_providers,
                use_knowledge_base=self.kb_config
            )
            self.finished.emit(results)
        except Exception as e:
            self.finished.emit({'错误': {'success': False, 'data': [], 'error': str(e)}})


class _HelpDialog(QDialog):
    """帮助文档对话框 - LLM与知识库配置说明"""
    def __init__(self, parent=None):
        super().__init__(parent)
        self.setWindowTitle('帮助文档 - 配置说明')
        self.setMinimumSize(720, 640)
        self._init_ui()

    def _init_ui(self):
        layout = QVBoxLayout(self)
        layout.setContentsMargins(0, 0, 0, 0)

        browser = QTextBrowser()
        browser.setOpenExternalLinks(True)
        browser.setStyleSheet("""
            QTextBrowser { background: #1a1a1e; border: none;
                          color: #d0d0d4; padding: 16px; font-size: 13px; }
        """)
        browser.setHtml(self._build_html())
        layout.addWidget(browser)

        btn_row = QHBoxLayout()
        btn_row.addStretch()
        close_btn = QPushButton('关闭')
        close_btn.setObjectName('mutedBtn')
        close_btn.setFixedWidth(80)
        close_btn.clicked.connect(self.accept)
        btn_row.addWidget(close_btn)
        layout.addLayout(btn_row)

    def _build_html(self):
        return """<div style="font-family: Microsoft YaHei; line-height: 1.8;">

<h2 style="color:#3f7bf7; border-bottom:2px solid #3f7bf7; padding-bottom:8px;">RadAtlas 配置帮助文档</h2>

<p style="color:#888890;">本文档详细介绍各家大模型(LLM)和知识库的配置流程及方法。软件默认主题为"暖沙"，默认知识库为火山方舟。</p>

<!-- ==================== LLM 配置 ==================== -->
<h2 style="color:#f59e0b; margin-top:24px;">一、大模型(LLM)配置</h2>
<p style="color:#c0c0c8;">在"AI配置"对话框中可添加多个大模型。支持同时选择多个模型并行查询，系统默认选中前2个已配置的模型。</p>

<h3 style="color:#5a91ff; margin-top:16px;">1. DeepSeek</h3>
<table style="color:#c0c0c8; font-size:12px;" cellpadding="4">
<tr><td style="color:#888890;">API地址</td><td>https://api.deepseek.com/v1/chat/completions</td></tr>
<tr><td style="color:#888890;">获取API Key</td><td><a href="https://platform.deepseek.com/api_keys" style="color:#3f7bf7;">https://platform.deepseek.com/api_keys</a></td></tr>
<tr><td style="color:#888890;">可用模型</td><td>deepseek-chat (通用对话), deepseek-reasoner (推理模型)</td></tr>
<tr><td style="color:#888890;">配置步骤</td><td>注册DeepSeek账号 → 创建API Key → 填入AI配置</td></tr>
<tr><td style="color:#888890;">费用</td><td>输入￥2/百万Token，输出￥8/百万Token（价格仅供参考）</td></tr>
</table>

<h3 style="color:#5a91ff; margin-top:16px;">2. 豆包(火山引擎)</h3>
<table style="color:#c0c0c8; font-size:12px;" cellpadding="4">
<tr><td style="color:#888890;">API地址</td><td>https://ark.cn-beijing.volces.com/api/v3/responses</td></tr>
<tr><td style="color:#888890;">获取API Key</td><td><a href="https://console.volcengine.com/ark" style="color:#3f7bf7;">https://console.volcengine.com/ark</a> → APIKey管理</td></tr>
<tr><td style="color:#888890;">可用模型</td><td>doubao-seed-2-0-pro, doubao-seed-1-6, doubao-1-5-pro-32k 等</td></tr>
<tr><td style="color:#888890;">备注</td><td>模型名称也可填入火山方舟的Endpoint ID（如 ep-xxxxxxxx）。doubao-seed系列为推理模型，已启用 reasoning effort=low 加速响应</td></tr>
<tr><td style="color:#888890;">配置步骤</td><td>注册火山引擎账号 → 开通方舟大模型服务 → 创建API Key → 填入AI配置</td></tr>
<tr><td style="color:#888890;">超时设置</td><td>测试连接超时60秒，诊断调用超时180秒（推理模型响应较慢）</td></tr>
</table>

<h3 style="color:#5a91ff; margin-top:16px;">3. OpenAI</h3>
<table style="color:#c0c0c8; font-size:12px;" cellpadding="4">
<tr><td style="color:#888890;">API地址</td><td>https://api.openai.com/v1/chat/completions</td></tr>
<tr><td style="color:#888890;">获取API Key</td><td><a href="https://platform.openai.com/api-keys" style="color:#3f7bf7;">https://platform.openai.com/api-keys</a></td></tr>
<tr><td style="color:#888890;">可用模型</td><td>gpt-4o, gpt-4o-mini, gpt-3.5-turbo</td></tr>
<tr><td style="color:#888890;">配置步骤</td><td>注册OpenAI账号 → 创建API Key → 填入AI配置（需支持海外网络访问）</td></tr>
</table>

<h3 style="color:#5a91ff; margin-top:16px;">4. 智谱AI</h3>
<table style="color:#c0c0c8; font-size:12px;" cellpadding="4">
<tr><td style="color:#888890;">API地址</td><td>https://open.bigmodel.cn/api/paas/v4/chat/completions</td></tr>
<tr><td style="color:#888890;">获取API Key</td><td><a href="https://open.bigmodel.cn/" style="color:#3f7bf7;">https://open.bigmodel.cn/</a></td></tr>
<tr><td style="color:#888890;">可用模型</td><td>glm-4-plus, glm-4-flash (免费), glm-4</td></tr>
<tr><td style="color:#888890;">配置步骤</td><td>注册智谱AI开放平台 → 创建API Key → 填入AI配置</td></tr>
</table>

<h3 style="color:#5a91ff; margin-top:16px;">5. 通义千问</h3>
<table style="color:#c0c0c8; font-size:12px;" cellpadding="4">
<tr><td style="color:#888890;">API地址</td><td>https://dashscope.aliyuncs.com/compatible-mode/v1/chat/completions</td></tr>
<tr><td style="color:#888890;">获取API Key</td><td><a href="https://dashscope.console.aliyun.com/" style="color:#3f7bf7;">https://dashscope.console.aliyun.com/</a></td></tr>
<tr><td style="color:#888890;">可用模型</td><td>qwen-max, qwen-plus, qwen-turbo</td></tr>
<tr><td style="color:#888890;">配置步骤</td><td>注册阿里云账号 → 开通DashScope → 创建API Key → 填入AI配置</td></tr>
</table>

<p style="color:#888890; font-size:12px; margin-top:8px;">提示：模型名称输入框可编辑，支持手动输入自定义模型名称。可使用"测试连接"按钮验证配置是否正确。</p>

<!-- ==================== 知识库配置 ==================== -->
<h2 style="color:#f59e0b; margin-top:24px;">二、知识库配置（重点）</h2>
<p style="color:#c0c0c8;">知识库用于检索医学影像诊断相关文档。支持三种知识库，每次查询选择一个。<b style="color:#f59e0b;">默认选中火山方舟知识库</b>（若已配置）。可单独使用知识库检索（不选大模型），也可与大模型配合使用。</p>

<h3 style="color:#5a91ff; margin-top:16px;">1. 火山方舟知识库（默认推荐）</h3>
<p style="color:#c0c0c8; font-size:12px;">火山方舟知识库支持两种检索模式，可同时配置。系统优先使用模式1（search_knowledge），失败自动降级到模式2（Responses API）。</p>

<p style="color:#c0c0c8; font-size:12px; margin-top:8px;"><b style="color:#f59e0b;">📋 配置字段对照表：</b></p>
<table style="color:#c0c0c8; font-size:11px;" cellpadding="4" border="1" bordercolor="#2a2a30">
<tr style="background:#1a1a1e;"><td style="color:#f59e0b;">配置对话框字段</td><td style="color:#f59e0b;">火山引擎控制台对应位置</td><td style="color:#f59e0b;">所属模式</td></tr>
<tr><td>Access Key</td><td>IAM密钥管理 → Access Key ID</td><td>模式1</td></tr>
<tr><td>Secret Key</td><td>IAM密钥管理 → Secret Access Key</td><td>模式1</td></tr>
<tr><td>Resource ID</td><td>火山方舟控制台 → 知识库详情页 → Resource ID</td><td>模式1</td></tr>
<tr><td>知识库名称</td><td>火山方舟控制台 → 知识库 → 集合名称（可选）</td><td>模式1</td></tr>
<tr><td>API Key</td><td>火山方舟控制台 → APIKey管理 → 创建API Key</td><td>模式2</td></tr>
<tr><td>旗舰版Endpoint ID</td><td>火山方舟控制台 → 旗舰版知识库 → Endpoint ID</td><td>模式2</td></tr>
</table>

<p style="color:#c0c0c8; font-size:12px; margin-top:12px;"><b style="color:#f59e0b;">🔑 步骤一：获取 Access Key / Secret Key（模式1必填）</b></p>
<ol style="color:#c0c0c8; font-size:12px;">
<li>访问 <a href="https://console.volcengine.com/iam/keymanage" style="color:#3f7bf7;">火山引擎IAM密钥管理</a>，登录火山引擎账号</li>
<li>点击"<b>创建密钥</b>"按钮，弹出密钥创建对话框</li>
<li>获取 <b>Access Key ID</b>（如 AKLTxxxxx）和 <b>Secret Access Key</b>（如 T1dKxxxxx）</li>
<li><b style="color:#f59e0b;">方法A（推荐）：CSV一键导入</b>
  <ul>
  <li>在密钥列表页点击"<b>导出CSV</b>"，下载CSV文件到本地</li>
  <li>打开RadAtlas知识库配置对话框</li>
  <li>点击顶部"<b>📂 导入CSV配置文件</b>"按钮，选择刚下载的CSV文件</li>
  <li>Access Key 和 Secret Key 自动填入，无需手动复制</li>
  </ul>
</li>
<li><b style="color:#f59e0b;">方法B：手动复制</b>
  <ul>
  <li>复制 Access Key ID → 填入配置对话框的"<b>Access Key</b>"输入框</li>
  <li>复制 Secret Access Key → 填入配置对话框的"<b>Secret Key</b>"输入框</li>
  <li>⚠️ Secret Access Key 仅在创建时显示一次，请务必保存</li>
  </ul>
</li>
</ol>

<p style="color:#c0c0c8; font-size:12px; margin-top:12px;"><b style="color:#f59e0b;">📚 步骤二：获取 Resource ID（模式1必填）</b></p>
<ol style="color:#c0c0c8; font-size:12px;">
<li>访问 <a href="https://console.volcengine.com/ark" style="color:#3f7bf7;">火山方舟控制台</a> → 左侧菜单"<b>知识库</b>"</li>
<li>点击"<b>创建知识库</b>"，上传医学影像诊断相关文档（PDF/EPUB/TXT等）</li>
<li>等待文档处理完成（分片、向量化）</li>
<li>进入知识库详情页，找到 <b>Resource ID</b>（格式如 <code>vcf-xxxxxxxx</code> 或 <code>kb-xxxxxxxx</code>）</li>
<li>复制 Resource ID → 填入配置对话框的"<b>Resource ID</b>"输入框</li>
</ol>

<p style="color:#c0c0c8; font-size:12px; margin-top:12px;"><b style="color:#f59e0b;">🔑 步骤三：获取 API Key（模式2可选）</b></p>
<ol style="color:#c0c0c8; font-size:12px;">
<li>访问 <a href="https://console.volcengine.com/ark" style="color:#3f7bf7;">火山方舟控制台</a> → 左侧菜单"<b>APIKey管理</b>"</li>
<li>点击"<b>创建API Key</b>" → 复制生成的API Key（如 ark-xxxxxxxx）</li>
<li>填入配置对话框的"<b>API Key</b>"输入框</li>
</ol>

<p style="color:#c0c0c8; font-size:12px; margin-top:12px;"><b style="color:#f59e0b;">📚 步骤四：获取旗舰版Endpoint ID（模式2可选）</b></p>
<ol style="color:#c0c0c8; font-size:12px;">
<li>在火山方舟控制台创建"<b>旗舰版知识库</b>"（需开通旗舰版服务）</li>
<li>进入旗舰版知识库详情页，获取 <b>Endpoint ID</b>（如 ep-xxxxxxxx）</li>
<li>填入配置对话框的"<b>旗舰版Endpoint ID</b>"输入框</li>
</ol>

<p style="color:#888890; font-size:12px; margin-top:12px;"><b style="color:#f59e0b;">⚠️ 常见问题排查：</b></p>
<ul style="color:#c0c0c8; font-size:12px;">
<li><b>HTTP 403 "check sign error"</b>：AK/SK凭证错误 → 使用CSV导入避免手动输入错误</li>
<li><b>HTTP 400 "collection not exist"</b>：Resource ID不存在 → 到火山方舟控制台确认正确的Resource ID</li>
<li><b>HTTP 429</b>：QPS限流 → 稍后重试，系统已自动添加重试机制</li>
<li>若AK/SK权限申请困难 → 建议仅使用模式2（API Key + 旗舰版Endpoint ID）</li>
<li>模式1返回的切片包含<b>原书图片</b>（预签名URL，4小时有效），可在切片查看器中直接查看</li>
</ul>

<h3 style="color:#5a91ff; margin-top:16px;">2. 腾讯IMA知识库</h3>
<p style="color:#c0c0c8; font-size:12px;">腾讯IMA知识库优先使用<b>知识库模块</b>（wiki/v1/search_knowledge），无结果时降级到<b>笔记模块</b>（note/v1/search_note_book）。</p>

<p style="color:#c0c0c8; font-size:12px; margin-top:8px;"><b style="color:#f59e0b;">📋 配置字段对照表：</b></p>
<table style="color:#c0c0c8; font-size:11px;" cellpadding="4" border="1" bordercolor="#2a2a30">
<tr style="background:#1a1a1e;"><td style="color:#f59e0b;">配置对话框字段</td><td style="color:#f59e0b;">IMA网站对应位置</td><td style="color:#f59e0b;">必填</td></tr>
<tr><td>Client ID</td><td>IMA开放平台 → API应用凭证 → Client ID</td><td>是</td></tr>
<tr><td>API Key</td><td>IMA开放平台 → API应用凭证 → API Key</td><td>是</td></tr>
<tr><td>知识库ID</td><td>IMA知识库管理页面 → 知识库ID</td><td>否（推荐填写）</td></tr>
</table>

<p style="color:#c0c0c8; font-size:12px; margin-top:12px;"><b style="color:#f59e0b;">🔑 步骤一：获取 Client ID 和 API Key</b></p>
<ol style="color:#c0c0c8; font-size:12px;">
<li>访问 <a href="https://ima.qq.com" style="color:#3f7bf7;">https://ima.qq.com</a> → 登录腾讯IMA账号</li>
<li>点击右上角头像 → 进入"<b>开放平台/Agent接口</b>"页面</li>
<li>或直接访问 <a href="https://ima.qq.com/agent-interface" style="color:#3f7bf7;">https://ima.qq.com/agent-interface</a></li>
<li>创建API应用凭证，获取：
  <ul>
  <li><b>Client ID</b> → 填入配置对话框的"<b>Client ID</b>"输入框</li>
  <li><b>API Key</b> → 填入配置对话框的"<b>API Key</b>"输入框</li>
  </ul>
</li>
</ol>

<p style="color:#c0c0c8; font-size:12px; margin-top:12px;"><b style="color:#f59e0b;">📚 步骤二：获取知识库ID（可选但推荐）</b></p>
<ol style="color:#c0c0c8; font-size:12px;">
<li>在IMA中创建知识库 → 上传医学影像诊断相关文档/PDF</li>
<li>进入知识库管理页面，找到目标知识库的 <b>知识库ID</b></li>
<li>复制知识库ID → 填入配置对话框的"<b>知识库ID</b>"输入框</li>
<li>⚠️ 若不填知识库ID，wiki模块不可用，系统将仅使用笔记模块搜索（可能无结果）</li>
</ol>

<p style="color:#888890; font-size:12px; margin-top:8px;">说明：检索策略为疾病名 → "检查类型+疾病名" → 核心词三级查询。建议在IMA中创建专门的医学影像诊断知识库并上传相关文档。</p>

<h3 style="color:#5a91ff; margin-top:16px;">3. Google NotebookLM (Gemini File Search)</h3>
<table style="color:#c0c0c8; font-size:12px;" cellpadding="4">
<tr><td style="color:#888890;">API地址</td><td>https://generativelanguage.googleapis.com/v1beta</td></tr>
<tr><td style="color:#888890;">认证方式</td><td>x-goog-api-key 请求头</td></tr>
<tr><td style="color:#888890;">获取API Key</td><td><a href="https://aistudio.google.com/apikey" style="color:#3f7bf7;">https://aistudio.google.com/apikey</a></td></tr>
<tr><td style="color:#888890;">需要填写</td><td>API Key（必填）+ File Search Store ID（可选）</td></tr>
<tr><td style="color:#888890;">使用模型</td><td>gemini-2.5-flash</td></tr>
</table>
<p style="color:#c0c0c8; font-size:12px; margin-top:8px;"><b style="color:#f59e0b;">两种检索模式：</b></p>
<ul style="color:#c0c0c8; font-size:12px;">
<li><b>File Search Store</b>：填写 Store ID，检索自有上传文档（需在Google AI Studio中创建File Search Store并上传文档）。</li>
<li><b>Google Search</b>：不填 Store ID，自动降级为Google网络搜索（检索公开网络信息）。</li>
</ul>
<p style="color:#888890; font-size:12px;">配置步骤：访问Google AI Studio → 创建API Key → （可选）创建File Search Store并上传文档 → 填入知识库配置（需支持海外网络访问）</p>

<!-- ==================== 使用说明 ==================== -->
<h2 style="color:#f59e0b; margin-top:24px;">三、使用说明</h2>

<h3 style="color:#5a91ff; margin-top:16px;">1. 检查类型</h3>
<p style="color:#c0c0c8; font-size:12px;">支持以下检查类型：CT、X-Ray、MRI、PET-CT、超声、DSA、<b style="color:#f59e0b;">ECG(心电图)</b>。</p>
<ul style="color:#c0c0c8; font-size:12px;">
<li><b>影像类</b>（CT/X-Ray/MRI等）：使用影像诊断专家角色，分析影像学表现</li>
<li><b>ECG(心电图)</b>：使用心电诊断专家角色，分析心率、心律、P波、PR间期、QRS波群、ST段、T波、QT间期等心电图特征</li>
</ul>

<h3 style="color:#5a91ff; margin-top:16px;">2. 大模型 + 知识库联合检索（推荐）</h3>
<ol style="color:#c0c0c8; font-size:12px;">
<li>选择1-2个大模型（默认选中前2个已配置的模型）</li>
<li>选择1个知识库（默认选中火山方舟，若已配置）</li>
<li>输入检查类型和关键征象，点击"开始查询"</li>
<li>系统先并行调用大模型获取诊断结果</li>
<li>从诊断结果中提取匹配度最高的1-3个疾病名</li>
<li>对每个疾病<b>分别</b>查询知识库，获取相关文档</li>
<li>知识库文档按疾病分组显示在诊断条目下方</li>
</ol>

<h3 style="color:#5a91ff; margin-top:16px;">3. 单独知识库检索</h3>
<ol style="color:#c0c0c8; font-size:12px;">
<li>不选择任何大模型</li>
<li>选择1个已配置的知识库</li>
<li>输入关键字，点击"开始查询"</li>
<li>系统直接用关键字检索知识库，返回匹配文档</li>
</ol>

<h3 style="color:#5a91ff; margin-top:16px;">4. 知识库文档查看</h3>
<ul style="color:#c0c0c8; font-size:12px;">
<li>点击左侧知识库文档条目，右侧显示文档快照详情（含元信息和切片预览）</li>
<li>详情包含：文档名称、来源、切片ID、原文链接</li>
<li><b>切片/快照查看</b>：点击"查看完整切片内容"链接，打开独立查看器窗口</li>
<li><b>原书图片</b>：火山方舟知识库返回的切片包含原书图片，在查看器中自动下载并显示（后台异步加载，带加载进度提示）</li>
<li><b>放大/缩小</b>：查看器有缩放按钮（＋/－/1:1重置），使用 QTextBrowser 内置缩放，<b>同时缩放文字和图片</b></li>
<li><b>收藏保存</b>：点击"☆ 收藏"按钮可保存当前文档，再次点击取消收藏</li>
<li>左侧结果列表过多时会自动显示滚动条，可滚动查看所有条目</li>
<li>点击"查看原文"链接可跳转到原始文档</li>
</ul>

<h3 style="color:#5a91ff; margin-top:16px;">5. 主题切换</h3>
<ul style="color:#c0c0c8; font-size:12px;">
<li>软件默认主题为"<b>暖沙</b>"</li>
<li>可通过菜单栏切换主题（深蓝暗夜、暖沙等）</li>
<li>切片查看器窗口自动跟随主程序主题配色</li>
</ul>

<h3 style="color:#5a91ff; margin-top:16px;">6. 联通测试</h3>
<ul style="color:#c0c0c8; font-size:12px;">
<li>在AI配置和知识库配置对话框中，每个配置项都有"测试连接"按钮</li>
<li>点击后会向对应服务发送测试请求，验证配置是否正确</li>
<li>测试期间请保持网络畅通</li>
</ul>

<p style="color:#666670; font-size:11px; margin-top:24px; border-top:1px solid #2a2a30; padding-top:12px;">
提示：所有API Key和密码输入框右侧都有眼睛图标，点击可查看输入内容，不输入时自动隐藏。
</p>

</div>"""


class _ImageLoadWorker(QThread):
    """后台下载原书图片的worker线程"""
    image_loaded = pyqtSignal(int, str)  # index, local_file_path
    all_done = pyqtSignal()

    def __init__(self, image_urls):
        super().__init__()
        self._urls = image_urls

    def run(self):
        import requests as _requests
        import tempfile as _tf
        import os as _os
        for i, url in enumerate(self._urls):
            try:
                resp = _requests.get(url, timeout=15)
                if resp.status_code == 200 and resp.content:
                    # 保存到临时文件
                    suffix = '.png'
                    ct = resp.headers.get('Content-Type', '')
                    if 'jpeg' in ct or 'jpg' in ct:
                        suffix = '.jpg'
                    elif 'webp' in ct:
                        suffix = '.webp'
                    fd, tmp_path = _tf.mkstemp(suffix=suffix, prefix='kb_img_')
                    with _os.fdopen(fd, 'wb') as f:
                        f.write(resp.content)
                    self.image_loaded.emit(i, tmp_path)
            except Exception:
                pass
        self.all_done.emit()


class _KBSnippetViewer(QDialog):
    """知识库文档切片/快照查看器 — 支持缩放查看图文内容，主题与主程序一致"""
    def __init__(self, doc, parent=None):
        super().__init__(parent)
        doc_name = doc.get('doc_name', '未知文档')
        self.setWindowTitle(f'切片查看 - {doc_name}')
        self.setMinimumSize(700, 600)
        self._doc = doc
        self._zoom_level = 0
        # 检测当前主题
        self._theme_name = self._detect_theme()
        self._t = THEMES[self._theme_name]
        self._init_ui()

    def _detect_theme(self):
        app = QApplication.instance()
        if app:
            ss = app.styleSheet()
            for name, t in THEMES.items():
                if t['bg'] in ss:
                    return name
        return '暖沙'

    def _init_ui(self):
        t = self._t
        layout = QVBoxLayout(self)
        layout.setContentsMargins(0, 0, 0, 0)
        layout.setSpacing(0)
        self.setStyleSheet(f"QDialog {{ background: {t['bg']}; }}")

        # 工具栏
        toolbar = QHBoxLayout()
        toolbar.setContentsMargins(12, 8, 12, 8)
        toolbar.setSpacing(6)
        title_label = QLabel(f'&#128218; {self._doc.get("doc_name", "未知文档")}')
        title_label.setStyleSheet(f"color: {t['accent_hover']}; font-size: 14px; font-weight: bold;")
        title_label.setTextFormat(Qt.RichText)
        toolbar.addWidget(title_label)
        toolbar.addStretch()

        # 页码信息
        page = self._doc.get('page', '')
        if page:
            page_label = QLabel(f'切片ID: {page}')
            page_label.setStyleSheet(f"color: {t['fg2']}; font-size: 12px;")
            toolbar.addWidget(page_label)

        # 缩放按钮
        btn_style = f"QPushButton {{ background: {t['bg3']}; color: {t['fg']}; border: 1px solid {t['border']}; border-radius: 4px; font-size:14px; font-weight:bold; padding:0; }} QPushButton:hover {{ background: {t['muted']}; }}"
        zoom_in_btn = QPushButton('＋')
        zoom_in_btn.setFixedSize(32, 28)
        zoom_in_btn.setToolTip('放大')
        zoom_in_btn.setStyleSheet(btn_style)
        zoom_in_btn.clicked.connect(lambda checked=False: self._zoom(1))
        toolbar.addWidget(zoom_in_btn)

        zoom_out_btn = QPushButton('－')
        zoom_out_btn.setFixedSize(32, 28)
        zoom_out_btn.setToolTip('缩小')
        zoom_out_btn.setStyleSheet(btn_style)
        zoom_out_btn.clicked.connect(lambda checked=False: self._zoom(-1))
        toolbar.addWidget(zoom_out_btn)

        reset_btn = QPushButton('1:1')
        reset_btn.setFixedSize(40, 28)
        reset_btn.setToolTip('重置缩放')
        reset_btn.setStyleSheet(btn_style.replace('font-size:14px', 'font-size:11px'))
        reset_btn.clicked.connect(self._reset_zoom)
        toolbar.addWidget(reset_btn)

        layout.addLayout(toolbar)

        # 分割线
        sep = QFrame()
        sep.setFrameShape(QFrame.HLine)
        sep.setStyleSheet(f"color: {t['border']};")
        layout.addWidget(sep)

        # 内容区
        self.browser = QTextBrowser()
        self.browser.setOpenExternalLinks(True)
        self.browser.setStyleSheet(f"""
            QTextBrowser {{ background: {t['bg2']}; border: none;
                          color: {t['fg']}; padding: 20px; font-size: 14px; line-height: 1.8; }}
        """)
        # 先显示文本内容，图片异步加载
        self.browser.setHtml(self._build_content_html())
        layout.addWidget(self.browser)

        # 异步下载并加载原书图片
        self._load_images_async()

        # 底部按钮栏
        btn_row = QHBoxLayout()
        btn_row.setContentsMargins(12, 8, 12, 8)
        btn_row.setSpacing(8)

        url = self._doc.get('url', '')
        if url:
            open_btn = QPushButton('🔗 查看原文')
            open_btn.setStyleSheet(f"""
                QPushButton {{ background: {t['bg3']}; color: {t['accent_hover']}; border: 1px solid {t['border']};
                              border-radius: 4px; padding: 6px 16px; font-size: 12px; }}
                QPushButton:hover {{ background: {t['muted']}; }}
            """)
            open_btn.clicked.connect(lambda: self.browser.setSource(QUrl(url)))
            btn_row.addWidget(open_btn)

        btn_row.addStretch()

        close_btn = QPushButton('关闭')
        close_btn.setStyleSheet(f"QPushButton {{ background: {t['bg3']}; color: {t['fg']}; border: 1px solid {t['border']}; border-radius: 4px; padding: 6px 16px; }} QPushButton:hover {{ background: {t['muted']}; }}")
        close_btn.setFixedWidth(80)
        close_btn.clicked.connect(self.accept)
        btn_row.addWidget(close_btn)
        layout.addLayout(btn_row)

    def _build_content_html(self):
        """构建切片内容HTML（含原书图片），使用主题配色"""
        t = self._t
        doc_name = self._doc.get('doc_name', '')
        page = self._doc.get('page', '')
        content = self._doc.get('full_content', '') or self._doc.get('snippet', '')
        source = self._doc.get('source', '')
        url = self._doc.get('url', '')
        images = self._doc.get('images', []) or []

        html = '<div style="font-family: Microsoft YaHei;">'

        # 元信息条
        meta_parts = []
        if source:
            meta_parts.append(f'<span style="color:{t["accent_hover"]};">{source}</span>')
        if page:
            meta_parts.append(f'<span style="color:{t["fg2"]};">切片ID: {page}</span>')
        if images:
            meta_parts.append(f'<span style="color:{t["accent_hover"]};">含 {len(images)} 张原书图片</span>')
        if meta_parts:
            html += f'<div style="margin-bottom:16px; padding:8px 12px; background:{t["bg3"]}; '
            html += f'border-radius:4px; border-left:3px solid {t["accent"]}; font-size:12px;">'
            html += ' &nbsp;|&nbsp; '.join(meta_parts)
            html += '</div>'

        # 切片文本内容（完整显示）
        if content:
            html += f'<div style="color:{t["fg"]}; line-height:2.0; white-space:pre-wrap; font-size:14px;">'
            import re as _re
            content_escaped = content.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
            content_escaped = _re.sub(
                r'(https?://[^\s<]+)',
                rf'<a href="\1" style="color:{t["accent_hover"]};">\1</a>',
                content_escaped
            )
            html += content_escaped
            html += '</div>'
        else:
            html += f'<p style="color:{t["fg3"]};">（无文本内容）</p>'

        # 原书图片
        if images:
            local_paths = getattr(self, '_image_local_paths', {})
            loaded_count = len(local_paths)
            html += f'<hr style="border:none; border-top:1px solid {t["border"]}; margin:20px 0;">'
            if loaded_count < len(images):
                html += f'<h3 style="color:{t["accent_hover"]}; margin:16px 0 12px 0;">原书图片（{len(images)}张，已加载{loaded_count}张...）</h3>'
            else:
                html += f'<h3 style="color:{t["accent_hover"]}; margin:16px 0 12px 0;">原书图片（{len(images)}张）</h3>'
            for i, img in enumerate(images):
                caption = img.get('caption', '')
                local_path = local_paths.get(i)
                if local_path:
                    img_src = QUrl.fromLocalFile(local_path).toString()
                    html += f'<div style="margin-bottom:24px; text-align:center; page-break-inside:avoid;">'
                    html += f'<img src="{img_src}" style="max-width:90%; max-height:500px; '
                    html += f'border:1px solid {t["border"]}; border-radius:6px; margin:8px 0; '
                    html += 'box-shadow:0 2px 8px rgba(0,0,0,0.3);" />'
                    if caption:
                        html += f'<p style="color:{t["fg2"]}; font-size:12px; margin-top:6px;">{caption}</p>'
                    else:
                        html += f'<p style="color:{t["fg3"]}; font-size:11px; margin-top:6px;">图 {i+1}</p>'
                    html += '</div>'
                else:
                    html += f'<div style="margin-bottom:20px; text-align:center; padding:40px; '
                    html += f'background:{t["bg3"]}; border:2px dashed {t["border"]}; border-radius:8px;">'
                    html += f'<p style="color:{t["fg3"]}; font-size:12px;">图 {i+1} 加载中...</p>'
                    html += '</div>'

        html += '</div>'
        return html

    def _load_images_async(self):
        """异步下载原书图片到临时文件，再通过 file:// URL 加载"""
        images = self._doc.get('images', []) or []
        if not images:
            return
        urls = [img.get('url', '') for img in images if img.get('url')]
        if not urls:
            return
        # 初始化 local_path 字段
        self._image_local_paths = {}
        # 启动后台下载线程
        self._image_worker = _ImageLoadWorker(urls)
        self._image_worker.image_loaded.connect(self._on_image_loaded)
        self._image_worker.all_done.connect(self._on_images_done)
        self._image_worker.start()

    def _on_image_loaded(self, idx, local_path):
        """单张图片下载完成 → 记录本地路径并刷新HTML"""
        self._image_local_paths[idx] = local_path
        self._refresh_content()

    def _on_images_done(self):
        """所有图片下载完成"""
        self._refresh_content()

    def _refresh_content(self):
        """刷新内容（保持缩放级别）"""
        self.browser.setHtml(self._build_content_html())
        if self._zoom_level > 0:
            self.browser.zoomIn(self._zoom_level)
        elif self._zoom_level < 0:
            self.browser.zoomOut(-self._zoom_level)

    def _zoom(self, delta):
        """缩放内容（QTextBrowser.zoomIn/zoomOut 同时缩放文字和图片）"""
        self._zoom_level += delta
        if delta > 0:
            self.browser.zoomIn(delta)
        else:
            self.browser.zoomOut(-delta)

    def _reset_zoom(self):
        """重置缩放"""
        if self._zoom_level > 0:
            self.browser.zoomOut(self._zoom_level)
        elif self._zoom_level < 0:
            self.browser.zoomIn(-self._zoom_level)
        self._zoom_level = 0

    def closeEvent(self, event):
        """关闭时清理临时图片文件"""
        import os as _os
        for path in getattr(self, '_image_local_paths', {}).values():
            try:
                if _os.path.exists(path):
                    _os.remove(path)
            except Exception:
                pass
        super().closeEvent(event)


class _KBConfigDialog(QDialog):
    """知识库配置对话框"""
    def __init__(self, parent=None):
        super().__init__(parent)
        self.setWindowTitle('知识库配置')
        self.setMinimumSize(520, 520)
        self._init_ui()

    def _init_ui(self):
        scroll = QScrollArea()
        scroll.setWidgetResizable(True)
        scroll.setFrameShape(QFrame.NoFrame)
        content = QWidget()
        layout = QVBoxLayout(content)
        layout.setContentsMargins(20, 20, 20, 20)
        layout.setSpacing(12)

        config = load_config()
        kb_configs = config.get('knowledge_bases', {})
        group_style = """
            QGroupBox { color: #b0b0b8; font-weight: bold; border: 1px solid #2a2a30;
                        border-radius: 6px; margin-top: 8px; padding-top: 16px; }
            QGroupBox::title { subcontrol-origin: margin; left: 12px; padding: 0 6px; }
        """
        # 输入框使用较小字号
        content.setStyleSheet("QLineEdit { font-size: 11px; } QLabel { font-size: 11px; }")

        # CSV一键导入按钮
        import_row = QHBoxLayout()
        import_hint = QLabel('一键导入火山方舟凭证：')
        import_hint.setStyleSheet("color: #888890; font-size: 12px;")
        import_row.addWidget(import_hint)
        import_row.addStretch()
        csv_btn = QPushButton('📂 导入CSV配置文件')
        csv_btn.setStyleSheet("""
            QPushButton { background: rgba(63,123,247,20); color: #5a91ff; border: 1px solid rgba(63,123,247,40);
                          border-radius: 4px; padding: 6px 14px; font-size: 12px; }
            QPushButton:hover { background: rgba(63,123,247,30); }
        """)
        csv_btn.clicked.connect(self._import_csv)
        import_row.addWidget(csv_btn)
        layout.addLayout(import_row)

        # 腾讯IMA知识库
        tencent_group = QGroupBox('腾讯IMA知识库')
        tencent_group.setStyleSheet(group_style)
        tencent_layout = QFormLayout(tencent_group)
        tencent_cfg = kb_configs.get('tencent', {})
        self.tencent_api_key = _PasswordLineEdit()
        self.tencent_api_key.setText(tencent_cfg.get('api_key', '') or tencent_cfg.get('client_id', ''))
        self.tencent_api_key.setPlaceholderText('IMA Client ID')
        tencent_layout.addRow('Client ID:', self.tencent_api_key)
        self.tencent_secret_key = _PasswordLineEdit()
        self.tencent_secret_key.setText(tencent_cfg.get('secret_key', '') or tencent_cfg.get('api_key_secret', ''))
        self.tencent_secret_key.setPlaceholderText('IMA API Key')
        tencent_layout.addRow('API Key:', self.tencent_secret_key)
        self.tencent_kb_id = QLineEdit(tencent_cfg.get('knowledge_base_id', ''))
        self.tencent_kb_id.setPlaceholderText('知识库ID (可选, 留空则搜索全部知识库)')
        tencent_layout.addRow('知识库ID:', self.tencent_kb_id)
        tencent_hint = QLabel('使用腾讯IMA个人知识库 OpenAPI。\n在 https://ima.qq.com/agent-interface 获取 Client ID 和 API Key。\n优先使用知识库模块(wiki/v1/search_knowledge)检索，无结果时降级到笔记模块(search_note_book)。\n知识库ID可在IMA知识库管理页面获取，留空则搜索全部知识库。')
        tencent_hint.setStyleSheet("color: #666670; font-size: 11px;")
        tencent_hint.setWordWrap(True)
        tencent_layout.addRow(tencent_hint)
        tencent_test_btn = QPushButton('测试连接')
        tencent_test_btn.clicked.connect(lambda checked=False: self._test_kb('tencent', tencent_test_btn))
        tencent_layout.addRow('', tencent_test_btn)
        layout.addWidget(tencent_group)

        # 火山方舟知识库
        volc_group = QGroupBox('火山方舟知识库')
        volc_group.setStyleSheet(group_style)
        volc_layout = QFormLayout(volc_group)
        volc_cfg = kb_configs.get('volcengine', {})
        self.volc_api_key = _PasswordLineEdit()
        self.volc_api_key.setText(volc_cfg.get('api_key', ''))
        self.volc_api_key.setPlaceholderText('API Key (用于Responses API模式)')
        volc_layout.addRow('API Key:', self.volc_api_key)
        self.volc_access_key = _PasswordLineEdit()
        self.volc_access_key.setText(volc_cfg.get('access_key', '') or volc_cfg.get('ak', ''))
        self.volc_access_key.setPlaceholderText('Access Key ID (用于search_knowledge模式)')
        volc_layout.addRow('Access Key:', self.volc_access_key)
        self.volc_secret_key = _PasswordLineEdit()
        self.volc_secret_key.setText(volc_cfg.get('secret_key', '') or volc_cfg.get('sk', ''))
        self.volc_secret_key.setPlaceholderText('Secret Access Key (用于search_knowledge模式)')
        volc_layout.addRow('Secret Key:', self.volc_secret_key)
        self.volc_resource_id = QLineEdit(volc_cfg.get('resource_id', ''))
        self.volc_resource_id.setPlaceholderText('知识库Resource ID (优先使用)')
        volc_layout.addRow('Resource ID:', self.volc_resource_id)
        self.volc_collection = QLineEdit(volc_cfg.get('collection_name', ''))
        self.volc_collection.setPlaceholderText('知识库名称 (与Resource ID二选一)')
        volc_layout.addRow('知识库名称:', self.volc_collection)
        self.volc_endpoint_id = QLineEdit(volc_cfg.get('endpoint_id', ''))
        self.volc_endpoint_id.setPlaceholderText('旗舰版知识库Endpoint ID (可选, 用于Responses API)')
        volc_layout.addRow('旗舰版Endpoint ID:', self.volc_endpoint_id)
        volc_hint = QLabel('模式1 search_knowledge（推荐）：需Access Key + Secret Key + Resource ID/集合名称\n'
                           '模式2 Responses API：需API Key + 旗舰版知识库ID\n'
                           '两种模式可同时配置，优先使用search_knowledge')
        volc_hint.setStyleSheet("color: #666670; font-size: 11px;")
        volc_hint.setWordWrap(True)
        volc_layout.addRow(volc_hint)
        volc_test_btn = QPushButton('测试连接')
        volc_test_btn.clicked.connect(lambda checked=False: self._test_kb('volcengine', volc_test_btn))
        volc_layout.addRow('', volc_test_btn)
        layout.addWidget(volc_group)

        # Google NotebookLM / Gemini File Search
        nblm_group = QGroupBox('Google NotebookLM (Gemini File Search)')
        nblm_group.setStyleSheet(group_style)
        nblm_layout = QFormLayout(nblm_group)
        nblm_cfg = kb_configs.get('notebooklm', {})
        self.nblm_api_key = _PasswordLineEdit()
        self.nblm_api_key.setText(nblm_cfg.get('api_key', ''))
        self.nblm_api_key.setPlaceholderText('Google AI Studio API Key')
        nblm_layout.addRow('API Key:', self.nblm_api_key)
        self.nblm_corpus_id = QLineEdit(nblm_cfg.get('file_search_store', '') or nblm_cfg.get('corpus_id', ''))
        self.nblm_corpus_id.setPlaceholderText('File Search Store ID (可选)')
        nblm_layout.addRow('Store ID:', self.nblm_corpus_id)
        nblm_hint = QLabel('使用 Gemini 2.5 Flash + File Search 工具检索。\n需在 Google AI Studio 获取 API Key。\n无 Store ID 时使用 Google Search 检索网络信息。')
        nblm_hint.setStyleSheet("color: #666670; font-size: 11px;")
        nblm_hint.setWordWrap(True)
        nblm_layout.addRow(nblm_hint)
        nblm_test_btn = QPushButton('测试连接')
        nblm_test_btn.clicked.connect(lambda checked=False: self._test_kb('notebooklm', nblm_test_btn))
        nblm_layout.addRow('', nblm_test_btn)
        layout.addWidget(nblm_group)

        layout.addSpacing(10)

        # 按钮
        btn_row = QHBoxLayout()
        save_btn = QPushButton('保存')
        save_btn.clicked.connect(self._save)
        cancel_btn = QPushButton('取消')
        cancel_btn.setObjectName('mutedBtn')
        cancel_btn.clicked.connect(self.reject)
        btn_row.addStretch()
        btn_row.addWidget(cancel_btn)
        btn_row.addWidget(save_btn)
        layout.addLayout(btn_row)

        scroll.setWidget(content)
        main_layout = QVBoxLayout(self)
        main_layout.setContentsMargins(0, 0, 0, 0)
        main_layout.addWidget(scroll)

    def _import_csv(self):
        """从火山引擎CSV配置文件一键导入AK/SK"""
        from PyQt5.QtWidgets import QFileDialog
        import csv as _csv

        file_path, _ = QFileDialog.getOpenFileName(
            self, '选择火山引擎CSV配置文件', '', 'CSV文件 (*.csv);;所有文件 (*.*)'
        )
        if not file_path:
            return

        try:
            with open(file_path, 'r', encoding='utf-8-sig') as f:
                reader = _csv.DictReader(f)
                ak = sk = tenant_id = ''
                for row in reader:
                    # 匹配多种可能的列名
                    ak = (row.get('Access Key ID') or row.get('access_key_id') or
                          row.get('AK') or row.get('ak') or row.get('AccessKeyId') or '').strip()
                    sk = (row.get('Secret Access Key') or row.get('secret_access_key') or
                          row.get('SK') or row.get('sk') or row.get('SecretAccessKey') or '').strip()
                    tenant_id = (row.get('所属主账号ID') or row.get('tenant_id') or
                                 row.get('主账号ID') or row.get('AccountId') or '').strip()
                    if ak and sk:
                        break

            if not ak or not sk:
                QMessageBox.warning(self, '导入失败',
                    '未能从CSV文件中找到 Access Key ID 和 Secret Access Key。\n'
                    '请确认CSV文件包含以下列（不区分大小写）：\n'
                    '  - Access Key ID / AK / AccessKeyId\n'
                    '  - Secret Access Key / SK / SecretAccessKey')
                return

            self.volc_access_key.setText(ak)
            self.volc_secret_key.setText(sk)
            QMessageBox.information(self, '导入成功',
                f'已成功导入火山方舟凭证：\n'
                f'  Access Key: {ak[:16]}...\n'
                f'  Secret Key: {sk[:8]}...{"（已填入配置）"}\n\n'
                f'请点击"保存"按钮保存配置，然后填写 Resource ID 后即可使用。')
        except Exception as e:
            QMessageBox.warning(self, '导入失败', f'读取CSV文件失败：{str(e)}')

    def _test_kb(self, kb_type, btn):
        """测试知识库联通性"""
        # 从UI收集当前配置
        if kb_type == 'tencent':
            kb_config = {
                'type': 'tencent',
                'api_key': self.tencent_api_key.text().strip(),
                'secret_key': self.tencent_secret_key.text().strip(),
                'knowledge_base_id': self.tencent_kb_id.text().strip(),
            }
        elif kb_type == 'volcengine':
            kb_config = {
                'type': 'volcengine',
                'api_key': self.volc_api_key.text().strip(),
                'access_key': self.volc_access_key.text().strip(),
                'secret_key': self.volc_secret_key.text().strip(),
                'resource_id': self.volc_resource_id.text().strip(),
                'collection_name': self.volc_collection.text().strip(),
                'endpoint_id': self.volc_endpoint_id.text().strip(),
            }
        elif kb_type == 'notebooklm':
            kb_config = {
                'type': 'notebooklm',
                'api_key': self.nblm_api_key.text().strip(),
                'file_search_store': self.nblm_corpus_id.text().strip(),
            }
        else:
            return

        btn.setEnabled(False)
        btn.setText('测试中...')

        self._kb_test_thread = QThread()
        self._kb_test_worker = _KBTestWorker(kb_config)
        self._kb_test_worker.moveToThread(self._kb_test_thread)
        self._kb_test_thread.started.connect(self._kb_test_worker.run)

        def _on_finished(result):
            btn.setEnabled(True)
            btn.setText('测试连接')
            if result['success']:
                QMessageBox.information(self, '测试成功', result['message'])
            else:
                QMessageBox.warning(self, '测试失败', result['message'])
            self._kb_test_thread.quit()

        self._kb_test_worker.finished.connect(_on_finished)
        self._kb_test_thread.start()

    def _save(self, checked=False):
        config = load_config()
        if 'knowledge_bases' not in config:
            config['knowledge_bases'] = {}
        config['knowledge_bases']['tencent'] = {
            'type': 'tencent',
            'api_key': self.tencent_api_key.text().strip(),
            'secret_key': self.tencent_secret_key.text().strip(),
            'knowledge_base_id': self.tencent_kb_id.text().strip(),
        }
        config['knowledge_bases']['volcengine'] = {
            'type': 'volcengine',
            'api_key': self.volc_api_key.text().strip(),
            'access_key': self.volc_access_key.text().strip(),
            'secret_key': self.volc_secret_key.text().strip(),
            'resource_id': self.volc_resource_id.text().strip(),
            'collection_name': self.volc_collection.text().strip(),
            'endpoint_id': self.volc_endpoint_id.text().strip(),
        }
        config['knowledge_bases']['notebooklm'] = {
            'type': 'notebooklm',
            'api_key': self.nblm_api_key.text().strip(),
            'file_search_store': self.nblm_corpus_id.text().strip(),
        }
        save_config(config)
        QMessageBox.information(self, '成功', '知识库配置已保存')
        self.accept()


# ── 添加/编辑疾病对话框 ──────────────────────────────────
class DiseaseDialog(QDialog):
    FIELDS = [
        ('name_cn', '中文名 *'),
        ('name_en', '英文名'),
        ('system', '系统'),
        ('category', '分类'),
        ('clinical', '临床表现'),
        ('diagnosis', '诊断要点'),
        ('primary_img', '主要影像方法'),
        ('secondary_img', '辅助影像方法'),
        ('xray_finding', 'X线所见'),
        ('ct_finding', 'CT所见'),
        ('mri_finding', 'MRI所见'),
        ('pet_finding', 'PET所见'),
        ('report_template', '报告模板'),
        ('differential_diagnosis', '鉴别诊断'),
        ('treatment', '治疗原则'),
    ]

    def __init__(self, parent=None, disease_data=None):
        super().__init__(parent)
        self.disease_data = disease_data
        self.setWindowTitle('编辑疾病' if disease_data else '添加疾病')
        self.setMinimumSize(600, 700)

        layout = QVBoxLayout(self)
        layout.setContentsMargins(20, 20, 20, 20)

        scroll = QScrollArea()
        scroll.setWidgetResizable(True)
        inner = QWidget()
        form = QVBoxLayout(inner)
        form.setSpacing(8)

        self.inputs = {}
        for key, label in self.FIELDS:
            lbl = QLabel(label)
            lbl.setStyleSheet("font-weight: bold; font-size: 12px;")
            form.addWidget(lbl)
            inp = QLineEdit()
            if disease_data:
                inp.setText(disease_data.get(key, '') or '')
            self.inputs[key] = inp
            form.addWidget(inp)

        scroll.setWidget(inner)
        layout.addWidget(scroll)

        # AI 填充按钮行
        ai_row = QHBoxLayout()
        self.ai_fill_btn = QPushButton('AI 一键填充')
        self.ai_fill_btn.clicked.connect(self._ai_fill)
        ai_row.addWidget(self.ai_fill_btn)

        self.ai_status = QLabel('')
        self.ai_status.setStyleSheet("color: #888890; font-size: 11px; background: transparent;")
        ai_row.addWidget(self.ai_status)
        ai_row.addStretch()
        layout.addLayout(ai_row)

        btn_row = QHBoxLayout()
        save_btn = QPushButton('保存')
        save_btn.clicked.connect(self.accept)
        cancel_btn = QPushButton('取消')
        cancel_btn.setObjectName('mutedBtn')
        cancel_btn.clicked.connect(self.reject)
        btn_row.addStretch()
        btn_row.addWidget(cancel_btn)
        btn_row.addWidget(save_btn)
        layout.addLayout(btn_row)

    def _ai_fill(self, checked=False):
        name_cn = self.inputs['name_cn'].text().strip()
        if not name_cn:
            QMessageBox.warning(self, '提示', '请先输入疾病中文名')
            return
        config = load_config()
        if not config.get('api_key'):
            ret = QMessageBox.question(self, '未配置AI',
                '尚未配置AI大模型API，是否现在配置？')
            if ret == QMessageBox.Yes:
                dlg = AIConfigDialog(self)
                if dlg.exec_() != QDialog.Accepted:
                    return
                config = load_config()
            else:
                return

        self.ai_fill_btn.setEnabled(False)
        self.ai_status.setText('正在查询AI，请稍候...')
        QApplication.processEvents()

        try:
            result = call_ai(name_cn, config)
            # 填充字段
            filled = 0
            for key, _ in self.FIELDS:
                val = result.get(key, '').strip() if isinstance(result.get(key), str) else ''
                if val and not self.inputs[key].text().strip():
                    self.inputs[key].setText(val)
                    filled += 1
                elif val:
                    # 如果已有内容，也覆盖（AI结果更完整）
                    self.inputs[key].setText(val)
                    filled += 1
            self.ai_status.setText(f'AI填充完成，已填入 {filled} 个字段')
        except ValueError as e:
            QMessageBox.warning(self, '配置错误', str(e))
            self.ai_status.setText('填充失败：未配置API Key')
        except requests.exceptions.Timeout:
            QMessageBox.warning(self, '超时', 'AI请求超时，请检查网络连接')
            self.ai_status.setText('填充失败：请求超时')
        except requests.exceptions.ConnectionError:
            QMessageBox.warning(self, '连接错误', '无法连接AI服务，请检查网络和API地址')
            self.ai_status.setText('填充失败：连接错误')
        except json.JSONDecodeError:
            QMessageBox.warning(self, '解析错误', 'AI返回的数据格式异常，请重试')
            self.ai_status.setText('填充失败：数据格式异常')
        except Exception as e:
            QMessageBox.critical(self, '错误', f'AI填充失败:\n{e}')
            self.ai_status.setText('填充失败')
        finally:
            self.ai_fill_btn.setEnabled(True)

    def get_data(self):
        return {key: inp.text().strip() for key, inp in self.inputs.items()}


# ── 用户管理对话框 ────────────────────────────────────────
class UserManageDialog(QDialog):
    def __init__(self, parent=None):
        super().__init__(parent)
        self.setWindowTitle('用户管理')
        self.setMinimumSize(520, 550)

        layout = QVBoxLayout(self)
        layout.setContentsMargins(20, 20, 20, 20)
        layout.setSpacing(12)

        # 用户列表
        self.user_list = QListWidget()
        self.user_list.currentItemChanged.connect(self._on_user_selected)
        self._load_users()
        layout.addWidget(QLabel('已有用户:'))
        layout.addWidget(self.user_list)

        # 操作按钮行
        op_row = QHBoxLayout()
        edit_btn = QPushButton('修改用户名/密码')
        edit_btn.clicked.connect(self._edit_user)
        op_row.addWidget(edit_btn)
        del_btn = QPushButton('删除选中用户')
        del_btn.setObjectName('dangerBtn')
        del_btn.clicked.connect(self._delete_user)
        op_row.addWidget(del_btn)
        layout.addLayout(op_row)

        layout.addSpacing(10)

        # 添加用户区
        layout.addWidget(QLabel('添加新用户:'))
        add_form = QHBoxLayout()
        self.new_username = QLineEdit()
        self.new_username.setPlaceholderText('用户名')
        add_form.addWidget(self.new_username)
        self.new_password = _PasswordLineEdit()
        self.new_password.setPlaceholderText('密码')
        add_form.addWidget(self.new_password)
        layout.addLayout(add_form)

        role_row = QHBoxLayout()
        role_row.addWidget(QLabel('角色:'))
        self.role_combo = QComboBox()
        self.role_combo.addItems(['user', 'admin'])
        role_row.addWidget(self.role_combo)
        role_row.addStretch()
        layout.addLayout(role_row)

        db_row = QHBoxLayout()
        db_row.addWidget(QLabel('数据库:'))
        self.db_combo = QComboBox()
        self.db_combo.addItems(['含基础数据库', '空数据库'])
        db_row.addWidget(self.db_combo)
        db_row.addStretch()
        layout.addLayout(db_row)

        add_btn = QPushButton('添加用户')
        add_btn.clicked.connect(self._add_user)
        layout.addWidget(add_btn)

    def _load_users(self):
        self.user_list.clear()
        for row in get_all_users():
            uid, username, role, database, created = row
            item = QListWidgetItem(f'[{role}] {username}  (ID:{uid})')
            item.setData(Qt.UserRole, uid)
            item.setData(Qt.UserRole + 1, username)
            item.setData(Qt.UserRole + 2, role)
            self.user_list.addItem(item)

    def _on_user_selected(self, current, previous):
        pass

    def _edit_user(self, checked=False):
        item = self.user_list.currentItem()
        if not item:
            QMessageBox.warning(self, '提示', '请先选择一个用户')
            return
        uid = item.data(Qt.UserRole)
        old_name = item.data(Qt.UserRole + 1)
        dlg = EditUserDialog(self, uid, old_name)
        if dlg.exec_() == QDialog.Accepted:
            self._load_users()

    def _add_user(self, checked=False):
        username = self.new_username.text().strip()
        password = self.new_password.text().strip()
        if not username or not password:
            QMessageBox.warning(self, '提示', '请输入用户名和密码')
            return
        role = self.role_combo.currentText()
        use_base = self.db_combo.currentIndex() == 0
        ok, db_path = create_user(username, password, role)
        if not ok:
            QMessageBox.warning(self, '错误', '用户名已存在')
            return
        init_db(db_path)
        if use_base:
            copy_public_to_user(db_path)
        self._load_users()
        self.new_username.clear()
        self.new_password.clear()
        QMessageBox.information(self, '成功', f'用户 {username} 已创建')

    def _delete_user(self, checked=False):
        item = self.user_list.currentItem()
        if not item:
            return
        uid = item.data(Qt.UserRole)
        if QMessageBox.question(self, '确认', '确定删除该用户？') == QMessageBox.Yes:
            delete_user(uid)
            self._load_users()


class EditUserDialog(QDialog):
    """修改用户名和密码"""
    def __init__(self, parent, user_id, current_username):
        super().__init__(parent)
        self.user_id = user_id
        self.current_username = current_username
        self.setWindowTitle(f'修改用户 - {current_username}')
        self.setFixedSize(400, 260)

        layout = QVBoxLayout(self)
        layout.setContentsMargins(20, 20, 20, 20)
        layout.setSpacing(12)

        # 修改用户名
        layout.addWidget(QLabel('修改用户名:'))
        name_row = QHBoxLayout()
        self.new_name_input = QLineEdit()
        self.new_name_input.setPlaceholderText('新用户名')
        self.new_name_input.setText(current_username)
        name_row.addWidget(self.new_name_input)
        rename_btn = QPushButton('修改')
        rename_btn.clicked.connect(self._rename)
        name_row.addWidget(rename_btn)
        layout.addLayout(name_row)

        layout.addSpacing(10)

        # 修改密码
        layout.addWidget(QLabel('修改密码:'))
        self.new_pwd_input = _PasswordLineEdit()
        self.new_pwd_input.setPlaceholderText('新密码')
        layout.addWidget(self.new_pwd_input)

        self.confirm_pwd_input = _PasswordLineEdit()
        self.confirm_pwd_input.setPlaceholderText('确认新密码')
        layout.addWidget(self.confirm_pwd_input)

        pwd_btn = QPushButton('修改密码')
        pwd_btn.clicked.connect(self._change_pwd)
        layout.addWidget(pwd_btn)

    def _rename(self, checked=False):
        new_name = self.new_name_input.text().strip()
        if not new_name:
            QMessageBox.warning(self, '提示', '用户名不能为空')
            return
        if new_name == self.current_username:
            return
        if rename_user(self.user_id, new_name):
            self.current_username = new_name
            QMessageBox.information(self, '成功', '用户名已修改')
            self.accept()
        else:
            QMessageBox.warning(self, '错误', '用户名已存在')

    def _change_pwd(self, checked=False):
        new_pwd = self.new_pwd_input.text().strip()
        confirm = self.confirm_pwd_input.text().strip()
        if not new_pwd:
            QMessageBox.warning(self, '提示', '请输入新密码')
            return
        if new_pwd != confirm:
            QMessageBox.warning(self, '提示', '两次密码输入不一致')
            return
        admin_change_password(self.user_id, new_pwd)
        QMessageBox.information(self, '成功', '密码已修改')
        self.accept()


# ── 图像浏览窗口 ──────────────────────────────────────────
class ImageViewerDialog(QDialog):
    """图像浏览窗口：缩放/平移/翻页，双击进入编辑模式"""

    def __init__(self, image_paths, current_index=0, parent=None, user_info=None, user_password=None, active_db=None):
        super().__init__(parent)
        self.setWindowTitle('图像浏览')
        self.setMinimumSize(900, 650)
        self.setWindowFlags(Qt.Window | Qt.FramelessWindowHint)
        self.image_paths = image_paths  # list of (img_id, img_path)
        self.current_index = current_index
        self.scale = 1.0
        self.current_theme = self._detect_theme()
        self.user_info = user_info
        self.user_password = user_password
        self.active_db = active_db
        # 浏览窗口淡入动画
        self._viewer_opacity = QGraphicsOpacityEffect(self)
        self.setGraphicsEffect(self._viewer_opacity)
        self._viewer_opacity.setOpacity(0)
        self._viewer_anim = QPropertyAnimation(self._viewer_opacity, b"opacity")
        self._viewer_anim.setDuration(300)
        self._viewer_anim.setStartValue(0.0)
        self._viewer_anim.setEndValue(1.0)
        self._viewer_anim.setEasingCurve(QEasingCurve.InOutCubic)
        QTimer.singleShot(50, self._viewer_anim.start)

        layout = QVBoxLayout(self)
        layout.setContentsMargins(0, 0, 0, 0)
        layout.setSpacing(0)

        # 自定义标题栏
        self.title_bar = _CustomTitleBar(self, '图像浏览')
        layout.addWidget(self.title_bar)

        # 顶部工具栏
        top_bar = QHBoxLayout()
        top_bar.setContentsMargins(8, 6, 8, 6)
        top_bar.setSpacing(6)

        btn_prev = QPushButton('< 上一张')
        btn_prev.setFixedHeight(30)
        btn_prev.setFont(QFont("SimSun", 9))
        btn_prev.clicked.connect(self._prev_image)
        top_bar.addWidget(btn_prev)

        self.page_label = QLabel()
        self.page_label.setFont(QFont("SimSun", 9))
        self.page_label.setStyleSheet("background: transparent; font-size: 9px;")
        top_bar.addWidget(self.page_label)

        btn_next = QPushButton('下一张 >')
        btn_next.setFixedHeight(30)
        btn_next.setFont(QFont("SimSun", 9))
        btn_next.clicked.connect(self._next_image)
        top_bar.addWidget(btn_next)

        top_bar.addSpacing(12)

        btn_zoom_out = QPushButton('−')
        btn_zoom_out.setFixedSize(28, 28)
        btn_zoom_out.setFont(QFont("SimSun", 12))
        btn_zoom_out.clicked.connect(self._zoom_out)
        top_bar.addWidget(btn_zoom_out)

        self.zoom_label = QLabel('100%')
        self.zoom_label.setFixedWidth(44)
        self.zoom_label.setFont(QFont("SimSun", 9))
        self.zoom_label.setStyleSheet("background: transparent; font-size: 9px;")
        top_bar.addWidget(self.zoom_label)

        btn_zoom_in = QPushButton('+')
        btn_zoom_in.setFixedSize(28, 28)
        btn_zoom_in.setFont(QFont("SimSun", 12))
        btn_zoom_in.clicked.connect(self._zoom_in)
        top_bar.addWidget(btn_zoom_in)

        btn_fit = QPushButton('适应')
        btn_fit.setObjectName('mutedBtn')
        btn_fit.setFixedSize(44, 28)
        btn_fit.setFont(QFont("SimSun", 9))
        btn_fit.clicked.connect(self._zoom_fit)
        top_bar.addWidget(btn_fit)

        btn_orig = QPushButton('1:1')
        btn_orig.setObjectName('mutedBtn')
        btn_orig.setFixedSize(36, 28)
        btn_orig.setFont(QFont("SimSun", 9))
        btn_orig.clicked.connect(self._zoom_original)
        top_bar.addWidget(btn_orig)

        top_bar.addStretch()

        btn_edit = QPushButton('编辑图片')
        btn_edit.setObjectName('accentBtn')
        btn_edit.setFixedHeight(30)
        btn_edit.setFont(QFont("SimSun", 10, QFont.Bold))
        btn_edit.clicked.connect(self._enter_edit)
        top_bar.addWidget(btn_edit)

        btn_close = QPushButton('关闭')
        btn_close.setObjectName('mutedBtn')
        btn_close.setFixedSize(52, 28)
        btn_close.setFont(QFont("SimSun", 9))
        btn_close.clicked.connect(self.reject)
        top_bar.addWidget(btn_close)

        layout.addLayout(top_bar)

        # 画布
        self.viewer_canvas = _ViewerCanvas(self, user_info=self.user_info, user_password=self.user_password, active_db=self.active_db)
        layout.addWidget(self.viewer_canvas, stretch=1)

        # 底部提示
        hint = QLabel('  双击图片进入编辑模式  |  滚轮缩放  |  拖拽平移')
        hint.setFixedHeight(24)
        hint.setFont(QFont("SimSun", 9))
        hint.setObjectName('editorInfoBar')
        layout.addWidget(hint)

        self._apply_theme()
        self._load_current()

    def _detect_theme(self):
        app = QApplication.instance()
        if app:
            ss = app.styleSheet()
            for name, t in THEMES.items():
                if t['bg'] in ss:
                    return name
        return '暖沙'

    def _apply_theme(self):
        t = THEMES[self.current_theme]
        is_light = self.current_theme == '浅色经典'
        bg = t['bg']
        bg2 = t['bg2']
        fg = t['fg']
        fg2 = t['fg2']
        border = t['border']
        accent = t['accent']
        muted = t['muted']
        self.setStyleSheet(f"""
        QDialog {{ background-color: {bg}; }}
        QWidget#editorInfoBar {{
            background-color: {bg2}; color: {fg2};
            border-top: 1px solid {border}; font-size: 9px;
            font-family: "SimSun";
        }}
        QPushButton {{
            background-color: {muted}; color: {fg};
            border: 1px solid {border}; border-radius: 4px; font-size: 10px;
            font-family: "SimSun";
        }}
        QPushButton:hover {{ background-color: {border}; }}
        QPushButton#mutedBtn {{ background-color: {muted}; }}
        QPushButton#mutedBtn:hover {{ background-color: {border}; }}
        QPushButton#accentBtn {{ background-color: {accent}; color: #fff; border: none; }}
        QPushButton#accentBtn:hover {{ background-color: {accent}; opacity: 0.85; }}
        QLabel {{ background-color: transparent; color: {fg}; font-family: "SimSun"; }}
        """)
        self.viewer_canvas.bg_color = QColor(245, 245, 248) if is_light else QColor(20, 20, 22)
        self.viewer_canvas.update()
        # 更新标题栏主题
        if hasattr(self, 'title_bar'):
            self.title_bar.update_theme(t)

    def _load_current(self):
        if 0 <= self.current_index < len(self.image_paths):
            img_id, img_path = self.image_paths[self.current_index]
            self.viewer_canvas.load_image(img_path, img_id)
            self.page_label.setText(f' {self.current_index + 1} / {len(self.image_paths)} ')
            self.zoom_label.setText('100%')
            self.scale = 1.0
            self.viewer_canvas.scale = 1.0
            self.viewer_canvas.offset = QPoint(0, 0)
            self.viewer_canvas.update()

    def _prev_image(self, checked=False):
        if self.current_index > 0:
            self.current_index -= 1
            self._load_current()

    def _next_image(self, checked=False):
        if self.current_index < len(self.image_paths) - 1:
            self.current_index += 1
            self._load_current()

    def _zoom_in(self, checked=False):
        self.scale = min(self.scale * 1.25, 5.0)
        self.viewer_canvas.scale = self.scale
        self.zoom_label.setText(f'{int(self.scale * 100)}%')
        self.viewer_canvas.update()

    def _zoom_out(self, checked=False):
        self.scale = max(self.scale / 1.25, 0.1)
        self.viewer_canvas.scale = self.scale
        self.zoom_label.setText(f'{int(self.scale * 100)}%')
        self.viewer_canvas.update()

    def _zoom_fit(self, checked=False):
        if self.viewer_canvas.pixmap:
            cw = self.viewer_canvas.width()
            ch = self.viewer_canvas.height()
            pw = self.viewer_canvas.pixmap.width()
            ph = self.viewer_canvas.pixmap.height()
            if pw > 0 and ph > 0:
                self.scale = min(cw / pw, ch / ph) * 0.9
                self.viewer_canvas.scale = self.scale
                self.viewer_canvas.offset = QPoint(0, 0)
                self.zoom_label.setText(f'{int(self.scale * 100)}%')
                self.viewer_canvas.update()

    def _zoom_original(self, checked=False):
        self.scale = 1.0
        self.viewer_canvas.scale = 1.0
        self.viewer_canvas.offset = QPoint(0, 0)
        self.zoom_label.setText('100%')
        self.viewer_canvas.update()

    def _enter_edit(self, checked=False):
        self._open_editor()

    def _open_editor(self):
        if 0 <= self.current_index < len(self.image_paths):
            img_id, img_path = self.image_paths[self.current_index]
            if os.path.exists(img_path) or (img_id and self.active_db and self.user_password):
                editor = ImageEditorDialog(img_path, self,
                                           active_db=self.active_db,
                                           user_password=self.user_password,
                                           image_id=img_id)
                editor.image_updated.connect(self._on_image_updated)
                editor.exec_()
                # 编辑后重新加载
                self.viewer_canvas.load_image(img_path, img_id=img_id)
    
    def _on_image_updated(self, img_path):
        """图片已更新，重新加载
        """
        # 更新当前显示的图片
        self.viewer_canvas.load_image(img_path)
        self.viewer_canvas.update()


class _ViewerCanvas(QWidget):
    """浏览画布：支持缩放/平移/双击"""
    def __init__(self, viewer, parent=None, user_info=None, user_password=None, active_db=None):
        super().__init__(parent)
        self.viewer = viewer
        self.pixmap = None
        self.image_path = None
        self.image_id = None
        self.scale = 1.0
        self.offset = QPoint(0, 0)
        self.dragging = False
        self.last_pos = None
        self.bg_color = QColor(20, 20, 22)
        self.user_info = user_info
        self.user_password = user_password
        self.active_db = active_db
        self.setMouseTracking(True)

    def load_image(self, path, img_id=None):
        self.image_path = path
        self.image_id = img_id
        pixmap = None
        # 优先从数据库加载加密图片
        if self.active_db and img_id and self.user_password:
            try:
                key = generate_key(self.user_password)
                decrypted_data = load_image_decrypted(self.active_db, img_id, key)
                if decrypted_data:
                    pixmap = QPixmap()
                    pixmap.loadFromData(decrypted_data)
            except Exception as e:
                print(f"从数据库加载图片失败: {e}")

        # 如果数据库加载失败，回退到文件系统
        if (pixmap is None or pixmap.isNull()) and path and os.path.exists(path):
            # 检查图片是否加密（文件级别）
            if self.active_db and img_id:
                try:
                    img_info = get_image_by_id(self.active_db, img_id)
                    if img_info and img_info.get('encrypted'):
                        current_user_id = self.user_info.get('id') if self.user_info else None
                        if img_info.get('owner_id') is not None and img_info.get('owner_id') != current_user_id:
                            self.pixmap = None
                            self.update()
                            return
                        if self.user_password:
                            try:
                                key = generate_key(self.user_password)
                                decrypted_data = decrypt_image(path, key)
                                if decrypted_data:
                                    pixmap = QPixmap()
                                    pixmap.loadFromData(decrypted_data)
                            except Exception as e:
                                print(f"解密失败: {e}")
                except Exception as e:
                    print(f"加载图片信息失败: {e}")

            if pixmap is None or pixmap.isNull():
                preview_path = os.path.splitext(path)[0] + '.preview.png'
                if os.path.exists(preview_path):
                    pixmap = QPixmap(preview_path)
                if pixmap is None or pixmap.isNull():
                    pixmap = QPixmap(path)

        self.pixmap = pixmap
        if pixmap and not pixmap.isNull():
            QTimer.singleShot(50, self.viewer._zoom_fit)
        self.update()

    def paintEvent(self, event):
        painter = QPainter(self)
        painter.setRenderHint(QPainter.Antialiasing)
        painter.fillRect(self.rect(), self.bg_color)
        if not self.pixmap:
            painter.setPen(QColor(100, 100, 100))
            painter.setFont(QFont("Microsoft YaHei", 14))
            painter.drawText(self.rect(), Qt.AlignCenter, '无图片')
            painter.end()
            return
        scaled = self.pixmap.scaled(
            self.pixmap.size() * self.scale,
            Qt.KeepAspectRatio, Qt.SmoothTransformation
        )
        x = (self.width() - scaled.width()) // 2 + self.offset.x()
        y = (self.height() - scaled.height()) // 2 + self.offset.y()
        painter.drawPixmap(x, y, scaled)
        painter.end()

    def wheelEvent(self, event):
        delta = event.angleDelta().y()
        if delta > 0:
            self.viewer._zoom_in()
        else:
            self.viewer._zoom_out()

    def mousePressEvent(self, event):
        if event.button() == Qt.LeftButton:
            self.dragging = True
            self.last_pos = event.pos()
            self.setCursor(Qt.ClosedHandCursor)

    def mouseMoveEvent(self, event):
        if self.dragging and self.last_pos:
            delta = event.pos() - self.last_pos
            self.offset += delta
            self.last_pos = event.pos()
            self.update()

    def mouseReleaseEvent(self, event):
        if event.button() == Qt.LeftButton:
            self.dragging = False
            self.last_pos = None
            self.setCursor(Qt.ArrowCursor)

    def mouseDoubleClickEvent(self, event):
        self.viewer._open_editor()


# ── 图片编辑器 ────────────────────────────────────────────

# 预定义颜色（在 _StyleBar 之前定义）


class _ToolButton(QWidget):
    """浮动工具栏中的图标按钮"""
    clicked = pyqtSignal(bool)

    def __init__(self, tool_id, parent=None):
        super().__init__(parent)
        self.tool_id = tool_id
        self.setFixedSize(36, 36)
        self._hover = False
        self._selected = False
        self.setCursor(Qt.PointingHandCursor)

    def set_selected(self, sel):
        self._selected = sel
        self.update()

    def paintEvent(self, event):
        painter = QPainter(self)
        painter.setRenderHint(QPainter.Antialiasing)
        r = self.rect()
        # 背景
        if self._selected:
            painter.setPen(Qt.NoPen)
            painter.setBrush(QColor(63, 123, 247, 50))
            painter.drawRoundedRect(r.adjusted(2, 2, -2, -2), 6, 6)
            painter.setBrush(Qt.NoBrush)
        elif self._hover:
            painter.fillRect(r, QColor(255, 255, 255, 30))
        # 绘制图标
        fg = QColor(190, 195, 200)
        painter.setPen(QPen(fg, 1.5))
        cx, cy = r.center().x(), r.center().y()
        tid = self.tool_id
        if tid == 'select':
            # 光标箭头 - 更精致
            painter.setPen(QPen(fg, 1.3))
            pts = [QPoint(cx - 4, cy - 7), QPoint(cx - 4, cy + 4),
                   QPoint(cx - 1, cy + 1), QPoint(cx + 1, cy + 6),
                   QPoint(cx + 3, cy + 4), QPoint(cx + 1, cy),
                   QPoint(cx + 5, cy)]
            painter.drawPolyline(QPolygon(pts))
        elif tid == 'pan':
            # 四向移动图标
            painter.setPen(QPen(fg, 1.3, Qt.SolidLine, Qt.RoundCap, Qt.RoundJoin))
            # 十字线
            painter.drawLine(cx, cy - 7, cx, cy + 7)
            painter.drawLine(cx - 7, cy, cx + 7, cy)
            # 四个箭头
            for dx, dy in [(0, -7), (0, 7), (-7, 0), (7, 0)]:
                ax, ay = cx + dx, cy + dy
                if dx == 0:  # 上下箭头
                    s = 1 if dy < 0 else -1
                    painter.drawLine(ax, ay, ax - 2, ay + s * 3)
                    painter.drawLine(ax, ay, ax + 2, ay + s * 3)
                else:  # 左右箭头
                    s = 1 if dx < 0 else -1
                    painter.drawLine(ax, ay, ax + s * 3, ay - 2)
                    painter.drawLine(ax, ay, ax + s * 3, ay + 2)
        elif tid == 'text':
            painter.setFont(QFont("SimSun", 12, QFont.Bold))
            painter.setPen(QPen(fg, 1))
            painter.drawText(r.adjusted(0, 0, 0, -4), Qt.AlignCenter, 'T')
            painter.setPen(QPen(fg, 1))
            painter.drawLine(cx, cy + 3, cx, cy + 8)
        elif tid == 'arrow':
            painter.setPen(QPen(fg, 1.3))
            painter.drawLine(cx - 6, cy + 6, cx + 5, cy - 5)
            pts = [QPoint(cx + 5, cy - 5), QPoint(cx + 1, cy - 2), QPoint(cx + 3, cy - 0)]
            painter.setBrush(fg)
            painter.drawPolygon(QPolygon(pts))
            painter.setBrush(Qt.NoBrush)
        elif tid == 'line':
            painter.setPen(QPen(fg, 1.3))
            painter.drawLine(cx - 7, cy + 7, cx + 7, cy - 7)
        elif tid == 'rect':
            painter.setPen(QPen(fg, 1.3))
            painter.drawRect(cx - 7, cy - 5, 14, 10)
        elif tid == 'circle':
            painter.setPen(QPen(fg, 1.3))
            painter.drawEllipse(cx - 6, cy - 6, 12, 12)
        elif tid == 'measure':
            painter.setPen(QPen(fg, 1.3))
            # 水平线+两端竖线
            painter.drawLine(cx - 7, cy, cx + 7, cy)
            painter.drawLine(cx - 7, cy - 3, cx - 7, cy + 3)
            painter.drawLine(cx + 7, cy - 3, cx + 7, cy + 3)
            # 中间刻度
            for i in range(3):
                tx = cx - 3 + i * 3
                h = 2 if i == 1 else 1
                painter.drawLine(tx, cy - h, tx, cy + h)
        elif tid == 'eraser':
            painter.setPen(QPen(fg, 1.3))
            # 简洁橡皮擦 - 圆角矩形
            painter.drawRoundedRect(cx - 6, cy - 4, 12, 8, 2, 2)
            # 分隔线
            painter.drawLine(cx - 6, cy - 1, cx + 6, cy - 1)
        painter.end()

    def enterEvent(self, event):
        self._hover = True
        self.update()

    def leaveEvent(self, event):
        self._hover = False
        self.update()

    def mousePressEvent(self, event):
        if event.button() == Qt.LeftButton:
            self.clicked.emit(True)


class _ActionBtn(QWidget):
    """浮动工具栏中的操作按钮（撤销、清空、保存、返回）"""
    clicked = pyqtSignal(bool)

    def __init__(self, action_id, parent=None):
        super().__init__(parent)
        self.action_id = action_id
        self.setFixedSize(38, 26)
        self._hover = False
        self.setCursor(Qt.PointingHandCursor)

    def paintEvent(self, event):
        painter = QPainter(self)
        painter.setRenderHint(QPainter.Antialiasing)
        r = self.rect()
        # 背景
        if self._hover:
            painter.setPen(QPen(QColor(255, 255, 255, 40), 1))
            painter.setBrush(QColor(255, 255, 255, 20))
        else:
            painter.setPen(QPen(QColor(80, 80, 85), 1))
            painter.setBrush(QColor(60, 60, 65))
        painter.drawRoundedRect(r.adjusted(1, 1, -1, -1), 4, 4)
        painter.setBrush(Qt.NoBrush)
        # 绘制图标
        fg = QColor(190, 195, 200)
        painter.setPen(QPen(fg, 1.3))
        cx, cy = r.center().x(), r.center().y()
        aid = self.action_id
        if aid == 'undo':
            painter.setPen(QPen(fg, 1.3))
            painter.drawArc(cx - 5, cy - 4, 10, 10, 30 * 16, 300 * 16)
            pts = [QPoint(cx - 4, cy - 3), QPoint(cx - 4, cy - 7), QPoint(cx - 1, cy - 3)]
            painter.setBrush(fg)
            painter.drawPolygon(QPolygon(pts))
            painter.setBrush(Qt.NoBrush)
        elif aid == 'clear':
            painter.setPen(QPen(fg, 1.3))
            painter.drawRect(cx - 4, cy - 1, 8, 7)
            painter.drawLine(cx - 5, cy - 1, cx + 5, cy - 1)
            painter.drawLine(cx - 2, cy - 3, cx + 2, cy - 3)
            painter.drawLine(cx - 2, cy + 1, cx - 2, cy + 5)
            painter.drawLine(cx + 2, cy + 1, cx + 2, cy + 5)
        elif aid == 'save':
            painter.setPen(QPen(fg, 1.3))
            painter.drawRect(cx - 5, cy - 5, 10, 10)
            painter.setBrush(QColor(60, 60, 65))
            painter.drawRect(cx - 3, cy - 5, 6, 4)
            painter.setBrush(Qt.NoBrush)
            painter.drawRect(cx - 3, cy + 1, 6, 4)
        elif aid == 'back':
            painter.setPen(QPen(fg, 1.3))
            pts = [QPoint(cx + 3, cy - 4), QPoint(cx - 3, cy), QPoint(cx + 3, cy + 4)]
            painter.setBrush(fg)
            painter.drawPolygon(QPolygon(pts))
            painter.setBrush(Qt.NoBrush)
            painter.drawLine(cx + 5, cy - 5, cx + 5, cy + 5)
        painter.end()

    def enterEvent(self, event):
        self._hover = True
        self.update()

    def leaveEvent(self, event):
        self._hover = False
        self.update()

    def mousePressEvent(self, event):
        if event.button() == Qt.LeftButton:
            self.clicked.emit(True)


class _FloatingToolbar(QWidget):
    """浮动可拖拽工具栏 - 左侧工具按钮"""
    tool_changed = pyqtSignal(str)
    color_changed = pyqtSignal(QColor)
    line_width_changed = pyqtSignal(int)
    line_style_changed = pyqtSignal(str)
    action_triggered = pyqtSignal(str)

    # 工具分类
    SELECT_TOOLS = [('select', 'S', '选择'), ('pan', 'V', '移动')]
    DRAW_TOOLS = [('text', 'T', '文字'), ('arrow', 'A', '箭头'), ('line', 'L', '直线'),
                  ('rect', 'R', '矩形'), ('circle', 'C', '圆形'), ('measure', 'M', '测量')]
    EDIT_TOOLS = [('eraser', 'E', '橡皮擦')]

    def __init__(self, parent=None):
        super().__init__(parent)
        self.current_tool = 'pan'
        self.current_color_idx = 0
        self.current_width_idx = 1
        self.current_style_idx = 0
        self._drag_pos = None
        self.zoom_pct = 100

        self.setFixedWidth(48)
        self.setAttribute(Qt.WA_StyledBackground, False)
        self.setAutoFillBackground(False)
        self.setCursor(Qt.ArrowCursor)

        layout = QVBoxLayout(self)
        layout.setContentsMargins(6, 6, 6, 6)
        layout.setSpacing(2)

        # === 选择工具 ===
        self.tool_btns = {}
        for tid, shortcut, name in self.SELECT_TOOLS:
            btn = _ToolButton(tid, self)
            btn.setToolTip(f'{name} ({shortcut})')
            btn.clicked.connect(lambda checked, t=tid: self._on_tool(t))
            self.tool_btns[tid] = btn
            layout.addWidget(btn, alignment=Qt.AlignCenter)

        self._add_separator(layout)

        # === 绘图工具 ===
        for tid, shortcut, name in self.DRAW_TOOLS:
            btn = _ToolButton(tid, self)
            btn.setToolTip(f'{name} ({shortcut})')
            btn.clicked.connect(lambda checked, t=tid: self._on_tool(t))
            self.tool_btns[tid] = btn
            layout.addWidget(btn, alignment=Qt.AlignCenter)

        self._add_separator(layout)

        # === 编辑工具 ===
        for tid, shortcut, name in self.EDIT_TOOLS:
            btn = _ToolButton(tid, self)
            btn.setToolTip(f'{name} ({shortcut})')
            btn.clicked.connect(lambda checked, t=tid: self._on_tool(t))
            self.tool_btns[tid] = btn
            layout.addWidget(btn, alignment=Qt.AlignCenter)

        self._add_separator(layout)

        # === 缩放 ===
        zoom_row = QHBoxLayout()
        zoom_row.setSpacing(2)
        btn_zo = self._create_action_btn('-', '缩小', 'zoom_out')
        zoom_row.addWidget(btn_zo)
        self.zoom_label = QLabel('100%')
        self.zoom_label.setAlignment(Qt.AlignCenter)
        self.zoom_label.setFont(QFont("SimSun", 7))
        self.zoom_label.setStyleSheet("color: #aaa; background: transparent;")
        self.zoom_label.setFixedWidth(28)
        zoom_row.addWidget(self.zoom_label)
        btn_zi = self._create_action_btn('+', '放大', 'zoom_in')
        zoom_row.addWidget(btn_zi)
        layout.addLayout(zoom_row)

        zoom_row2 = QHBoxLayout()
        zoom_row2.setSpacing(2)
        btn_fit = self._create_action_btn('⊡', '适应窗口', 'zoom_fit')
        zoom_row2.addWidget(btn_fit)
        btn_orig = self._create_action_btn('1:1', '原始大小', 'zoom_orig')
        zoom_row2.addWidget(btn_orig)
        layout.addLayout(zoom_row2)

        self._add_separator(layout)

        # === 操作 ===
        action_icons = [
            ('undo', '撤销 (Ctrl+Z)'), ('clear', '清空'),
            ('save', '保存 (Ctrl+Y)'), ('back', '返回 (Esc)'),
        ]
        self.action_btns = {}
        for aid, tip in action_icons:
            btn = _ActionBtn(aid)
            btn.setToolTip(tip)
            btn.clicked.connect(lambda checked, a=aid: self.action_triggered.emit(a))
            self.action_btns[aid] = btn
            layout.addWidget(btn, alignment=Qt.AlignCenter)

        layout.addStretch()
        self.tool_btns['pan'].set_selected(True)

    def _add_separator(self, layout):
        sep = QFrame()
        sep.setFrameShape(QFrame.HLine)
        sep.setFixedHeight(1)
        sep.setStyleSheet("background-color: rgba(255,255,255,20);")
        layout.addWidget(sep)

    def _create_action_btn(self, icon, tooltip, action):
        btn = QPushButton(icon)
        btn.setFixedSize(20, 20)
        btn.setFont(QFont("SimSun", 9))
        btn.setCursor(Qt.PointingHandCursor)
        btn.setToolTip(tooltip)
        btn.setStyleSheet("background-color: #555; color: #ddd; border: 1px solid #666; border-radius: 3px;")
        btn.clicked.connect(lambda checked, a=action: self.action_triggered.emit(a))
        return btn

    def _on_tool(self, tid):
        self.current_tool = tid
        for k, btn in self.tool_btns.items():
            btn.set_selected(k == tid)
        self.tool_changed.emit(tid)

    def set_zoom_pct(self, pct):
        self.zoom_pct = pct
        self.zoom_label.setText(f'{pct}%')

    def get_color(self):
        return _ANNOTATION_COLORS[self.current_color_idx]

    def get_line_width(self):
        return _ANNOTATION_WIDTHS[self.current_width_idx]

    def get_line_style(self):
        return _ANNOTATION_STYLES[self.current_style_idx][0]

    def paintEvent(self, event):
        painter = QPainter(self)
        painter.setRenderHint(QPainter.Antialiasing)
        painter.setBrush(QColor(30, 30, 35, 220))
        painter.setPen(QPen(QColor(255, 255, 255, 30), 1))
        painter.drawRoundedRect(self.rect().adjusted(1, 1, -1, -1), 8, 8)
        painter.end()

    def mousePressEvent(self, event):
        if event.button() == Qt.LeftButton:
            self._drag_pos = event.pos()

    def mouseMoveEvent(self, event):
        if self._drag_pos and event.buttons() & Qt.LeftButton:
            delta = event.pos() - self._drag_pos
            self.move(self.pos() + delta)

    def mouseReleaseEvent(self, event):
        self._drag_pos = None


# 默认注释颜色、线宽、线型
_ANNOTATION_COLORS = [
    QColor(255, 80, 80), QColor(80, 200, 80), QColor(80, 140, 255),
    QColor(255, 200, 50), QColor(200, 80, 255), QColor(255, 150, 50),
]
_ANNOTATION_WIDTHS = [1, 2, 3, 5, 8, 10]
_ANNOTATION_STYLES = [('solid', '实线'), ('dashed', '虚线'), ('dotted', '点线'), ('wavy', '波浪')]


class _WidthBtn(QWidget):
    """线宽选择按钮 - 用不同粗细的水平线表示"""
    clicked = pyqtSignal(bool)

    def __init__(self, width_val, parent=None):
        super().__init__(parent)
        self.width_val = width_val
        self.setFixedSize(22, 22)
        self._hover = False
        self._selected = False
        self.setCursor(Qt.PointingHandCursor)

    def set_selected(self, sel):
        self._selected = sel
        self.update()

    def paintEvent(self, event):
        painter = QPainter(self)
        painter.setRenderHint(QPainter.Antialiasing)
        r = self.rect()
        # 背景
        if self._selected:
            painter.setPen(Qt.NoPen)
            painter.setBrush(QColor(63, 123, 247, 60))
            painter.drawRoundedRect(r.adjusted(1, 1, -1, -1), 3, 3)
        elif self._hover:
            painter.setPen(Qt.NoPen)
            painter.setBrush(QColor(255, 255, 255, 20))
            painter.drawRoundedRect(r.adjusted(1, 1, -1, -1), 3, 3)
        # 绘制不同粗细的水平线
        pen_w = min(self.width_val, 6)  # 限制最大显示宽度
        color = QColor(220, 220, 225) if not self._selected else QColor(160, 200, 255)
        painter.setPen(QPen(color, pen_w, Qt.SolidLine, Qt.RoundCap))
        cx, cy = r.center().x(), r.center().y()
        painter.drawLine(cx - 7, cy, cx + 7, cy)
        painter.end()

    def enterEvent(self, event):
        self._hover = True
        self.update()

    def leaveEvent(self, event):
        self._hover = False
        self.update()

    def mousePressEvent(self, event):
        if event.button() == Qt.LeftButton:
            self.clicked.emit(True)


class _LineStyleBtn(QWidget):
    """线型选择按钮 - 用不同样式的线段表示"""
    clicked = pyqtSignal(bool)

    def __init__(self, style_id, parent=None):
        super().__init__(parent)
        self.style_id = style_id
        self.setFixedSize(36, 22)
        self._hover = False
        self._selected = False
        self.setCursor(Qt.PointingHandCursor)

    def set_selected(self, sel):
        self._selected = sel
        self.update()

    def paintEvent(self, event):
        painter = QPainter(self)
        painter.setRenderHint(QPainter.Antialiasing)
        r = self.rect()
        # 背景
        if self._selected:
            painter.setPen(Qt.NoPen)
            painter.setBrush(QColor(63, 123, 247, 60))
            painter.drawRoundedRect(r.adjusted(1, 1, -1, -1), 3, 3)
        elif self._hover:
            painter.setPen(Qt.NoPen)
            painter.setBrush(QColor(255, 255, 255, 20))
            painter.drawRoundedRect(r.adjusted(1, 1, -1, -1), 3, 3)

        color = QColor(220, 220, 225) if not self._selected else QColor(160, 200, 255)
        cx, cy = r.center().x(), r.center().y()
        pen = QPen(color, 2.5, Qt.SolidLine, Qt.RoundCap)

        if self.style_id == 'solid':
            # 实线
            pen.setStyle(Qt.SolidLine)
            painter.setPen(pen)
            painter.drawLine(cx - 13, cy, cx + 13, cy)
        elif self.style_id == 'dashed':
            # 虚线
            pen.setStyle(Qt.DashLine)
            painter.setPen(pen)
            painter.drawLine(cx - 13, cy, cx + 13, cy)
        elif self.style_id == 'dotted':
            # 点线
            pen.setStyle(Qt.DotLine)
            painter.setPen(pen)
            painter.drawLine(cx - 13, cy, cx + 13, cy)
        elif self.style_id == 'wavy':
            # 波浪线 - 使用QPainterPath绘制平滑曲线
            pen.setStyle(Qt.SolidLine)
            painter.setPen(pen)
            from math import sin, pi
            path = QPainterPath()
            first = True
            for x in range(-13, 14):
                y = sin(x * pi / 6) * 5  # 增大振幅和频率
                pt = QPointF(cx + x, cy + y)
                if first:
                    path.moveTo(pt)
                    first = False
                else:
                    path.lineTo(pt)
            painter.drawPath(path)
        painter.end()

    def enterEvent(self, event):
        self._hover = True
        self.update()

    def leaveEvent(self, event):
        self._hover = False
        self.update()

    def mousePressEvent(self, event):
        if event.button() == Qt.LeftButton:
            self.clicked.emit(True)


class _StyleBar(QWidget):
    """顶部样式栏 - 颜色、线宽、线型选择"""
    color_changed = pyqtSignal(QColor)
    line_width_changed = pyqtSignal(int)
    line_style_changed = pyqtSignal(str)

    def __init__(self, parent=None):
        super().__init__(parent)
        self.current_color_idx = 0
        self.current_width_idx = 1
        self.current_style_idx = 0
        self._drag_pos = None

        self.setFixedHeight(36)
        self.setAttribute(Qt.WA_StyledBackground, False)
        self.setAutoFillBackground(False)
        self.setCursor(Qt.ArrowCursor)

        layout = QHBoxLayout(self)
        layout.setContentsMargins(8, 4, 8, 4)
        layout.setSpacing(6)

        # 颜色
        cl = QLabel('颜色')
        cl.setFont(QFont("SimSun", 8))
        cl.setStyleSheet("color: #aaa; background: transparent;")
        layout.addWidget(cl)

        self.color_btns = []
        for i, c in enumerate(_ANNOTATION_COLORS):
            btn = QPushButton()
            btn.setFixedSize(14, 14)
            btn.setCursor(Qt.PointingHandCursor)
            btn.setToolTip(c.name())
            border = '#3f7bf7' if i == 0 else '#555'
            btn.setStyleSheet(
                f"background-color: {c.name()}; border: 2px solid {border}; border-radius: 7px;"
            )
            btn.clicked.connect(lambda checked, idx=i: self._on_color(idx))
            self.color_btns.append(btn)
            layout.addWidget(btn)

        # 分隔
        sep1 = QFrame()
        sep1.setFrameShape(QFrame.VLine)
        sep1.setFixedWidth(1)
        sep1.setStyleSheet("background-color: rgba(255,255,255,20);")
        layout.addWidget(sep1)

        # 线宽
        self.width_btns = []
        for i, w in enumerate(_ANNOTATION_WIDTHS):
            btn = _WidthBtn(w)
            btn.setToolTip(f'{w}px')
            btn.set_selected(i == self.current_width_idx)
            btn.clicked.connect(lambda checked, idx=i: self._on_width(idx))
            self.width_btns.append(btn)
            layout.addWidget(btn)

        # 分隔
        sep2 = QFrame()
        sep2.setFrameShape(QFrame.VLine)
        sep2.setFixedWidth(1)
        sep2.setStyleSheet("background-color: rgba(255,255,255,20);")
        layout.addWidget(sep2)

        # 线型
        self.style_btns = []
        for i, (style_id, style_text) in enumerate(_ANNOTATION_STYLES):
            btn = _LineStyleBtn(style_id)
            btn.setToolTip(style_text)
            btn.set_selected(i == self.current_style_idx)
            btn.clicked.connect(lambda checked, idx=i: self._on_style(idx))
            self.style_btns.append(btn)
            layout.addWidget(btn)

        layout.addStretch()

    def _on_color(self, idx):
        self.current_color_idx = idx
        for i, btn in enumerate(self.color_btns):
            border = '#3f7bf7' if i == idx else '#555'
            c = _ANNOTATION_COLORS[i]
            btn.setStyleSheet(
                f"background-color: {c.name()}; border: 2px solid {border}; border-radius: 7px;"
            )
        self.color_changed.emit(_ANNOTATION_COLORS[idx])

    def _on_width(self, idx):
        self.current_width_idx = idx
        for i, btn in enumerate(self.width_btns):
            btn.set_selected(i == idx)
        self.line_width_changed.emit(_ANNOTATION_WIDTHS[idx])

    def _on_style(self, idx):
        self.current_style_idx = idx
        for i, btn in enumerate(self.style_btns):
            btn.set_selected(i == idx)
        self.line_style_changed.emit(_ANNOTATION_STYLES[idx][0])

    def get_color(self):
        return _ANNOTATION_COLORS[self.current_color_idx]

    def get_line_width(self):
        return _ANNOTATION_WIDTHS[self.current_width_idx]

    def get_line_style(self):
        return _ANNOTATION_STYLES[self.current_style_idx][0]

    def paintEvent(self, event):
        painter = QPainter(self)
        painter.setRenderHint(QPainter.Antialiasing)
        painter.setBrush(QColor(30, 30, 35, 210))
        painter.setPen(Qt.NoPen)
        painter.drawRoundedRect(self.rect(), 6, 6)
        painter.end()

    def mousePressEvent(self, event):
        if event.button() == Qt.LeftButton:
            self._drag_pos = event.pos()

    def mouseMoveEvent(self, event):
        if self._drag_pos and event.buttons() & Qt.LeftButton:
            delta = event.pos() - self._drag_pos
            self.move(self.pos() + delta)

    def mouseReleaseEvent(self, event):
        self._drag_pos = None


class _PropertyPanel(QWidget):
    """选中注释时的属性面板"""
    color_changed = pyqtSignal(QColor)
    bg_color_changed = pyqtSignal(object)  # QColor or None
    opacity_changed = pyqtSignal(int)
    rotation_changed = pyqtSignal(float)
    delete_requested = pyqtSignal()
    line_style_changed = pyqtSignal(str)
    width_changed = pyqtSignal(int)
    height_changed = pyqtSignal(int)
    size_changed = pyqtSignal(int, int)  # width, height

    def __init__(self, parent=None):
        super().__init__(parent)
        self.setFixedWidth(150)
        self.setAttribute(Qt.WA_StyledBackground, False)
        self.setAutoFillBackground(False)
        self.hide()
        self.current_color = None
        self.current_bg_color = None

        layout = QVBoxLayout(self)
        layout.setContentsMargins(8, 8, 8, 8)
        layout.setSpacing(6)

        title = QLabel('属性面板')
        title.setFont(QFont("SimSun", 9, QFont.Bold))
        title.setStyleSheet("color: #ccc; background: transparent;")
        layout.addWidget(title)

        # 分割线
        sep0 = QFrame()
        sep0.setFrameShape(QFrame.HLine)
        sep0.setFixedHeight(1)
        sep0.setStyleSheet("background-color: rgba(255,255,255,20);")
        layout.addWidget(sep0)

        # 透明度
        ol = QLabel('透明度')
        ol.setFont(QFont("SimSun", 8))
        ol.setStyleSheet("color: #aaa; background: transparent;")
        layout.addWidget(ol)
        self.opacity_slider = QSlider(Qt.Horizontal)
        self.opacity_slider.setRange(0, 100)
        self.opacity_slider.setValue(55)
        self.opacity_slider.setToolTip('调整透明度')
        self.opacity_slider.setFixedWidth(130)
        self.opacity_slider.valueChanged.connect(lambda v: self.opacity_changed.emit(v))
        layout.addWidget(self.opacity_slider)
        self.opacity_label = QLabel('55%')
        self.opacity_label.setFont(QFont("SimSun", 7))
        self.opacity_label.setStyleSheet("color: #888; background: transparent;")
        self.opacity_label.setAlignment(Qt.AlignCenter)
        layout.addWidget(self.opacity_label)

        # 分割线
        sep1 = QFrame()
        sep1.setFrameShape(QFrame.HLine)
        sep1.setFixedHeight(1)
        sep1.setStyleSheet("background-color: rgba(255,255,255,20);")
        layout.addWidget(sep1)

        # 旋转
        rl = QLabel('旋转角度')
        rl.setFont(QFont("SimSun", 8))
        rl.setStyleSheet("color: #aaa; background: transparent;")
        layout.addWidget(rl)
        self.rotation_input = QSpinBox()
        self.rotation_input.setRange(0, 360)
        self.rotation_input.setValue(0)
        self.rotation_input.setFont(QFont("SimSun", 9))
        self.rotation_input.setSuffix('°')
        self.rotation_input.setToolTip('旋转角度')
        self.rotation_input.setFixedWidth(130)
        self.rotation_input.valueChanged.connect(lambda v: self.rotation_changed.emit(float(v)))
        layout.addWidget(self.rotation_input)

        # 线型
        ll = QLabel('线条样式')
        ll.setFont(QFont("SimSun", 8))
        ll.setStyleSheet("color: #aaa; background: transparent;")
        layout.addWidget(ll)
        self.line_style_combo = QComboBox()
        self.line_style_combo.setFont(QFont("SimSun", 9))
        self.line_style_combo.addItems(['实线', '虚线', '点线', '波浪线'])
        self.line_style_combo.setToolTip('选择线型')
        self.line_style_combo.setFixedWidth(130)
        self.line_style_combo.currentIndexChanged.connect(self._on_line_style)
        layout.addWidget(self.line_style_combo)

        # 线宽
        wl = QLabel('线宽(px)')
        wl.setFont(QFont("SimSun", 8))
        wl.setStyleSheet("color: #aaa; background: transparent;")
        layout.addWidget(wl)
        self.width_spin = QSpinBox()
        self.width_spin.setRange(1, 20)
        self.width_spin.setValue(2)
        self.width_spin.setFont(QFont("SimSun", 9))
        self.width_spin.setSuffix('px')
        self.width_spin.setToolTip('线宽像素')
        self.width_spin.setFixedWidth(130)
        self.width_spin.valueChanged.connect(lambda v: self.width_changed.emit(v))
        layout.addWidget(self.width_spin)

        # 分割线
        sep2 = QFrame()
        sep2.setFrameShape(QFrame.HLine)
        sep2.setFixedHeight(1)
        sep2.setStyleSheet("background-color: rgba(255,255,255,20);")
        layout.addWidget(sep2)

        # 像素尺寸（仅对矩形和圆形有效）
        sl = QLabel('尺寸(px)')
        sl.setFont(QFont("SimSun", 8))
        sl.setStyleSheet("color: #aaa; background: transparent;")
        layout.addWidget(sl)
        size_row = QHBoxLayout()
        size_row.setSpacing(4)

        self.size_w_spin = QSpinBox()
        self.size_w_spin.setRange(5, 500)
        self.size_w_spin.setValue(50)
        self.size_w_spin.setFont(QFont("SimSun", 8))
        self.size_w_spin.setSuffix('w')
        self.size_w_spin.setToolTip('宽度')
        self.size_w_spin.setFixedWidth(60)
        size_row.addWidget(self.size_w_spin)

        size_row.addWidget(QLabel('×'))

        self.size_h_spin = QSpinBox()
        self.size_h_spin.setRange(5, 500)
        self.size_h_spin.setValue(50)
        self.size_h_spin.setFont(QFont("SimSun", 8))
        self.size_h_spin.setSuffix('h')
        self.size_h_spin.setToolTip('高度')
        self.size_h_spin.setFixedWidth(60)
        size_row.addWidget(self.size_h_spin)

        self.size_w_spin.valueChanged.connect(self._on_size_change)
        self.size_h_spin.valueChanged.connect(self._on_size_change)
        layout.addLayout(size_row)

        # 分割线
        sep3 = QFrame()
        sep3.setFrameShape(QFrame.HLine)
        sep3.setFixedHeight(1)
        sep3.setStyleSheet("background-color: rgba(255,255,255,20);")
        layout.addWidget(sep3)

        # 删除
        del_btn = QPushButton('删除元件')
        del_btn.setFont(QFont("SimSun", 9))
        del_btn.setCursor(Qt.PointingHandCursor)
        del_btn.setStyleSheet("background-color: #d9534f; color: white; border: none; border-radius: 4px; padding: 4px;")
        del_btn.clicked.connect(self.delete_requested.emit)
        layout.addWidget(del_btn)

        # 更新透明度标签
        self.opacity_slider.valueChanged.connect(lambda v: self.opacity_label.setText(f'{v}%'))

    def _on_color_click(self, idx):
        self.current_color = _ANNOTATION_COLORS[idx]
        self.color_changed.emit(_ANNOTATION_COLORS[idx])

    def _on_bg_color_click(self, idx):
        self.current_bg_color = _ANNOTATION_COLORS[idx]
        self.bg_color_changed.emit(_ANNOTATION_COLORS[idx])

    def _on_bg_color_none(self):
        self.current_bg_color = None
        self.bg_color_changed.emit(None)

    def _update_color_btn_states(self):
        pass  # 颜色选择已移至顶部样式栏

    def _update_bg_color_btn_states(self):
        pass  # 颜色选择已移至顶部样式栏

    def _on_line_style(self, idx):
        styles = ['solid', 'dashed', 'dotted', 'wavy']
        if 0 <= idx < len(styles):
            self.line_style_changed.emit(styles[idx])

    def _on_size_change(self):
        w = self.size_w_spin.value()
        h = self.size_h_spin.value()
        self.size_changed.emit(w, h)

    def set_annotation(self, data):
        """根据注释数据更新面板"""
        self.rotation_input.blockSignals(True)
        self.opacity_slider.blockSignals(True)
        self.width_spin.blockSignals(True)
        self.line_style_combo.blockSignals(True)
        self.size_w_spin.blockSignals(True)
        self.size_h_spin.blockSignals(True)

        # 更新数值控件
        self.rotation_input.setValue(int(data.get('rotation', 0)))
        bg_opacity = data.get('bg_opacity', 140)
        opacity_val = int(bg_opacity * 100 / 255)
        self.opacity_slider.setValue(opacity_val)
        self.opacity_label.setText(f'{opacity_val}%')
        self.width_spin.setValue(int(data.get('width', 2)))

        line_style = data.get('line_style', 'solid')
        style_idx = {'solid': 0, 'dashed': 1, 'dotted': 2, 'wavy': 3}.get(line_style, 0)
        self.line_style_combo.setCurrentIndex(style_idx)

        # 更新尺寸（仅对矩形和圆形）
        if 'w' in data and 'h' in data:
            self.size_w_spin.setValue(max(5, int(data['w'])))
            self.size_h_spin.setValue(max(5, int(data['h'])))
            self.size_w_spin.setEnabled(True)
            self.size_h_spin.setEnabled(True)
        else:
            self.size_w_spin.setEnabled(False)
            self.size_h_spin.setEnabled(False)

        self.rotation_input.blockSignals(False)
        self.opacity_slider.blockSignals(False)
        self.width_spin.blockSignals(False)
        self.line_style_combo.blockSignals(False)
        self.size_w_spin.blockSignals(False)
        self.size_h_spin.blockSignals(False)

    def paintEvent(self, event):
        painter = QPainter(self)
        painter.setRenderHint(QPainter.Antialiasing)
        painter.setBrush(QColor(30, 30, 35, 220))
        painter.setPen(Qt.NoPen)
        painter.drawRoundedRect(self.rect(), 8, 8)
        painter.end()


class ImageEditorDialog(QDialog):
    """图片编辑器：浮动工具栏 + 选择/变换系统"""

    TOOL_SHORTCUTS = {
        Qt.Key_V: 'pan', Qt.Key_T: 'text', Qt.Key_A: 'arrow',
        Qt.Key_L: 'line', Qt.Key_R: 'rect', Qt.Key_C: 'circle',
        Qt.Key_M: 'measure', Qt.Key_E: 'eraser', Qt.Key_S: 'select',
    }

    image_updated = pyqtSignal(str)  # 信号：图片已更新，传递图片路径

    def __init__(self, image_path, parent=None, active_db=None, user_password=None, image_id=None):
        super().__init__(parent)
        self.setWindowTitle('图片编辑器')
        self.setMinimumSize(1050, 720)
        self.setWindowFlags(Qt.Window | Qt.FramelessWindowHint)
        self.image_path = image_path
        self.active_db = active_db
        self.user_password = user_password
        self.image_id = image_id
        self.scale = 1.0
        self.tool = 'pan'
        self.current_color = _ANNOTATION_COLORS[0]
        self.current_line_width = 2
        # 编辑器淡入动画
        self._editor_opacity = QGraphicsOpacityEffect(self)
        self.setGraphicsEffect(self._editor_opacity)
        self._editor_opacity.setOpacity(0)
        self._editor_anim = QPropertyAnimation(self._editor_opacity, b"opacity")
        self._editor_anim.setDuration(300)
        self._editor_anim.setStartValue(0.0)
        self._editor_anim.setEndValue(1.0)
        self._editor_anim.setEasingCurve(QEasingCurve.InOutCubic)
        QTimer.singleShot(50, self._editor_anim.start)
        self.current_line_style = 'solid'
        self.current_theme = self._detect_theme()
        # 实时保存定时器
        self._auto_save_timer = QTimer(self)
        self._auto_save_timer.setSingleShot(True)
        self._auto_save_timer.timeout.connect(self._auto_save)

        layout = QVBoxLayout(self)
        layout.setContentsMargins(0, 0, 0, 0)
        layout.setSpacing(0)

        # 自定义标题栏
        self.title_bar = _CustomTitleBar(self, '图片编辑器  |  楚雄州人民医院 医学影像中心 张兴文')
        layout.addWidget(self.title_bar)

        # 信息栏
        info_bar = QLabel(f'  {os.path.basename(image_path)}  |  当前工具: 移动  |  Esc返回浏览')
        info_bar.setFixedHeight(26)
        info_bar.setObjectName('editorInfoBar')
        info_bar.setFont(QFont("SimSun", 9))
        self.info_bar = info_bar
        layout.addWidget(info_bar)

        # 画布容器（用于叠放浮动工具栏和属性面板）
        self.canvas_container = QWidget()
        self.canvas_container.setObjectName('canvasContainer')

        self.canvas = _EditorCanvas(self, active_db=active_db, user_password=user_password, image_id=image_id)
        self.canvas.image_path = image_path
        self.canvas.load_image()

        # 浮动工具栏（左侧，只有工具按钮）
        self.toolbar = _FloatingToolbar(self.canvas_container)
        self.toolbar.move(10, 10)
        self.toolbar.show()
        self.toolbar.tool_changed.connect(self._set_tool)
        self.toolbar.action_triggered.connect(self._on_action)

        # 样式栏（顶部，颜色/线宽/线型）
        self.style_bar = _StyleBar(self.canvas_container)
        self.style_bar.move(70, 10)
        self.style_bar.show()
        self.style_bar.color_changed.connect(self._on_color_changed)
        self.style_bar.line_width_changed.connect(self._on_width_changed)
        self.style_bar.line_style_changed.connect(self._on_style_changed)

        # 属性面板
        self.prop_panel = _PropertyPanel(self.canvas_container)
        self.prop_panel.move(70, 52)
        self.prop_panel.color_changed.connect(self._on_prop_color)
        self.prop_panel.bg_color_changed.connect(self._on_prop_bg_color)
        self.prop_panel.opacity_changed.connect(self._on_prop_opacity)
        self.prop_panel.rotation_changed.connect(self._on_prop_rotation)
        self.prop_panel.delete_requested.connect(self._on_prop_delete)
        self.prop_panel.line_style_changed.connect(self._on_prop_line_style)
        self.prop_panel.width_changed.connect(self._on_prop_width)
        self.prop_panel.size_changed.connect(self._on_prop_size)

        # 使用布局管理 canvas，浮动工具栏和属性面板通过 move 定位叠在上面
        container_layout = QVBoxLayout(self.canvas_container)
        container_layout.setContentsMargins(0, 0, 0, 0)
        container_layout.addWidget(self.canvas)

        layout.addWidget(self.canvas_container, stretch=1)

        self._apply_theme()

    def keyPressEvent(self, event):
        key = event.key()
        if key in self.TOOL_SHORTCUTS:
            self._set_tool(self.TOOL_SHORTCUTS[key])
        elif key == Qt.Key_Escape:
            self.accept()
        elif key == Qt.Key_Z and event.modifiers() & Qt.ControlModifier:
            self._undo()
        elif key == Qt.Key_Y and event.modifiers() & Qt.ControlModifier:
            self._save()
        elif key == Qt.Key_Plus or key == Qt.Key_Equal:
            self._zoom_in()
        elif key == Qt.Key_Minus:
            self._zoom_out()
        elif key == Qt.Key_0 and event.modifiers() & Qt.ControlModifier:
            self._zoom_fit()
        elif key == Qt.Key_Delete:
            self._delete_selected()
        else:
            super().keyPressEvent(event)

    def _detect_theme(self):
        app = QApplication.instance()
        if app:
            ss = app.styleSheet()
            for name, t in THEMES.items():
                if t['bg'] in ss:
                    return name
        return '暖沙'

    def _apply_theme(self):
        t = THEMES[self.current_theme]
        is_light = self.current_theme == '浅色经典'
        bg = t['bg']
        bg2 = t['bg2']
        fg = t['fg']
        fg2 = t['fg2']
        border = t['border']
        accent = t['accent']
        muted = t['muted']

        self.setStyleSheet(f"""
        QDialog {{ background-color: {bg}; }}
        QWidget#editorInfoBar {{
            background-color: {bg2}; color: {fg2};
            border-bottom: 1px solid {border}; font-size: 9px;
            font-family: "SimSun";
        }}
        QPushButton {{
            background-color: {muted}; color: {fg};
            border: 1px solid {border}; border-radius: 4px; font-size: 10px;
            font-family: "SimSun";
        }}
        QPushButton:hover {{ background-color: {border}; }}
        QLabel {{ background-color: transparent; color: {fg}; }}
        QSpinBox {{
            background-color: {muted}; color: {fg};
            border: 1px solid {border}; border-radius: 4px;
            padding: 2px; font-size: 9px;
        }}
        QSlider::groove:horizontal {{
            background: {muted}; height: 6px; border-radius: 3px;
        }}
        QSlider::handle:horizontal {{
            background: {accent}; width: 12px; margin: -3px 0;
            border-radius: 6px;
        }}
        """)
        self.canvas.bg_color = QColor(245, 245, 248) if is_light else QColor(20, 20, 22)
        self.canvas.update()
        # 更新标题栏主题
        if hasattr(self, 'title_bar'):
            self.title_bar.update_theme(t)
        # 确保浮动工具栏和属性面板在最上层
        self.toolbar.raise_()
        self.style_bar.raise_()
        self.prop_panel.raise_()

    def _get_color(self):
        return self.style_bar.get_color()

    def _get_line_width(self):
        return self.style_bar.get_line_width()

    def _get_line_style(self):
        return self.style_bar.get_line_style()

    def _set_tool(self, tool):
        self.tool = tool
        self.canvas.tool = tool
        tool_names = {
            'select': '选择', 'pan': '移动', 'text': '文字',
            'arrow': '箭头', 'line': '直线', 'rect': '矩形',
            'circle': '圆形', 'measure': '测量', 'eraser': '橡皮擦',
        }
        self.info_bar.setText(
            f'  {os.path.basename(self.image_path)}  |  '
            f'当前工具: {tool_names.get(tool, tool)}  |  Esc返回浏览')

    def _on_color_changed(self, color):
        self.current_color = color
        # 同步更新选中元件的颜色
        idx = self.canvas.selected_index
        if 0 <= idx < len(self.canvas.annotations):
            self.canvas.annotations[idx][1]['color'] = color
            self.prop_panel.current_color = color
            self.prop_panel._update_color_btn_states()
            self.canvas.update()
            self._trigger_auto_save()

    def _on_width_changed(self, width):
        self.current_line_width = width
        # 同步更新选中元件的线宽
        idx = self.canvas.selected_index
        if 0 <= idx < len(self.canvas.annotations):
            self.canvas.annotations[idx][1]['width'] = width
            self.prop_panel.width_spin.blockSignals(True)
            self.prop_panel.width_spin.setValue(width)
            self.prop_panel.width_spin.blockSignals(False)
            self.canvas.update()
            self._trigger_auto_save()

    def _on_style_changed(self, style):
        self.current_line_style = style
        # 同步更新选中元件的线型
        idx = self.canvas.selected_index
        if 0 <= idx < len(self.canvas.annotations):
            self.canvas.annotations[idx][1]['line_style'] = style
            style_idx = {'solid': 0, 'dashed': 1, 'dotted': 2, 'wavy': 3}.get(style, 0)
            self.prop_panel.line_style_combo.blockSignals(True)
            self.prop_panel.line_style_combo.setCurrentIndex(style_idx)
            self.prop_panel.line_style_combo.blockSignals(False)
            self.canvas.update()
            self._trigger_auto_save()

    def _on_action(self, action):
        if action == 'undo':
            self._undo()
        elif action == 'clear':
            self._clear()
        elif action == 'save':
            self._save()
        elif action == 'back':
            self.accept()
        elif action == 'zoom_in':
            self._zoom_in()
        elif action == 'zoom_out':
            self._zoom_out()
        elif action == 'zoom_fit':
            self._zoom_fit()
        elif action == 'zoom_orig':
            self._zoom_original()

    def _zoom_in(self):
        self.scale = min(self.scale * 1.25, 5.0)
        self.canvas.scale = self.scale
        self.toolbar.set_zoom_pct(int(self.scale * 100))
        self.canvas.update()

    def _zoom_out(self):
        self.scale = max(self.scale / 1.25, 0.1)
        self.canvas.scale = self.scale
        self.toolbar.set_zoom_pct(int(self.scale * 100))
        self.canvas.update()

    def _zoom_fit(self):
        if self.canvas.pixmap:
            cw = self.canvas.width()
            ch = self.canvas.height()
            pw = self.canvas.pixmap.width()
            ph = self.canvas.pixmap.height()
            if pw > 0 and ph > 0:
                self.scale = min(cw / pw, ch / ph) * 0.9
                self.canvas.scale = self.scale
                self.canvas.offset = QPoint(0, 0)
                self.toolbar.set_zoom_pct(int(self.scale * 100))
                self.canvas.update()

    def _zoom_original(self, checked=False):
        self.scale = 1.0
        self.canvas.scale = 1.0
        self.canvas.offset = QPoint(0, 0)
        self.toolbar.set_zoom_pct(100)
        self.canvas.update()

    def _undo(self):
        if self.canvas.annotations:
            self.canvas.annotations.pop()
            self.canvas.selected_index = -1
            self.prop_panel.hide()
            self.canvas.update()

    def _clear(self):
        if self.canvas.annotations:
            if QMessageBox.question(self, '确认', '清空所有注释？') == QMessageBox.Yes:
                self.canvas.annotations.clear()
                self.canvas.selected_index = -1
                self.prop_panel.hide()
                self.canvas.update()

    def _auto_save(self):
        """自动保存当前编辑"""
        self.canvas.save_annotated(self.image_path)
        self.image_updated.emit(self.image_path)

    def _trigger_auto_save(self):
        """触发自动保存（延迟1秒）"""
        self._auto_save_timer.start(1000)

    def closeEvent(self, event):
        """关闭前保存注释"""
        self._auto_save_timer.stop()
        self.canvas.save_annotated(self.image_path)
        super().closeEvent(event)

    def _save(self):
        self.canvas.save_annotated(self.image_path)
        self.image_updated.emit(self.image_path)
        self.accept()

    def _delete_selected(self):
        if 0 <= self.canvas.selected_index < len(self.canvas.annotations):
            self.canvas.annotations.pop(self.canvas.selected_index)
            self.canvas.selected_index = -1
            self.prop_panel.hide()
            self.canvas.update()

    # 属性面板回调
    def _on_prop_color(self, color):
        idx = self.canvas.selected_index
        if 0 <= idx < len(self.canvas.annotations):
            self.canvas.annotations[idx][1]['color'] = color
            # 同步更新属性面板的颜色状态
            self.prop_panel.current_color = color
            self.prop_panel._update_color_btn_states()
            self.canvas.update()
            self._trigger_auto_save()

    def _on_prop_bg_color(self, color):
        idx = self.canvas.selected_index
        if 0 <= idx < len(self.canvas.annotations):
            self.canvas.annotations[idx][1]['bg_color'] = color
            self.prop_panel.current_bg_color = color
            self.prop_panel._update_bg_color_btn_states()
            self.canvas.update()
            self._trigger_auto_save()

    def _on_prop_opacity(self, val):
        idx = self.canvas.selected_index
        if 0 <= idx < len(self.canvas.annotations):
            self.canvas.annotations[idx][1]['bg_opacity'] = int(val * 255 / 100)
            self.canvas.update()
            self._trigger_auto_save()

    def _on_prop_rotation(self, deg):
        idx = self.canvas.selected_index
        if 0 <= idx < len(self.canvas.annotations):
            self.canvas.annotations[idx][1]['rotation'] = deg
            self._update_anchor(idx)
            self.canvas.update()
            self._trigger_auto_save()

    def _on_prop_delete(self):
        self._delete_selected()
        self._trigger_auto_save()

    def _on_prop_line_style(self, style):
        idx = self.canvas.selected_index
        if 0 <= idx < len(self.canvas.annotations):
            self.canvas.annotations[idx][1]['line_style'] = style
            self.canvas.update()
            self._trigger_auto_save()

    def _on_prop_width(self, w):
        idx = self.canvas.selected_index
        if 0 <= idx < len(self.canvas.annotations):
            self.canvas.annotations[idx][1]['width'] = w
            self.canvas.update()
            self._trigger_auto_save()

    def _on_prop_size(self, w, h):
        """修改元件尺寸（仅对矩形和圆形有效）"""
        idx = self.canvas.selected_index
        if 0 <= idx < len(self.canvas.annotations):
            atype, data = self.canvas.annotations[idx]
            if 'w' in data and 'h' in data:
                data['w'] = w
                data['h'] = h
                # 更新锚点
                data['anchor_x'] = data['x'] + data['w'] / 2
                data['anchor_y'] = data['y'] + data['h'] / 2
                self.canvas.update()
                self._trigger_auto_save()

    def _update_anchor(self, idx):
        """更新注释的旋转/缩放锚点为中心"""
        atype, data = self.canvas.annotations[idx]
        cx, cy = self.canvas._get_annotation_center(data)
        data['anchor_x'] = cx
        data['anchor_y'] = cy

    def show_property_panel(self, data):
        # 不再显示属性面板，颜色和样式选择在顶部样式栏中
        pass

    def hide_property_panel(self):
        pass


class _EditorCanvas(QWidget):
    """编辑画布：支持选择、旋转、缩放等多种标注变换"""
    HANDLE_SIZE = 8
    ROTATE_HANDLE_DIST = 20

    def __init__(self, editor, parent=None, active_db=None, user_password=None, image_id=None):
        super().__init__(parent)
        self.editor = editor
        self.image_path = None
        self.pixmap = None
        self.scale = 1.0
        self.tool = 'pan'
        self.annotations = []
        self.selected_index = -1
        self.drawing = False
        self.draw_start = None
        self.draw_current = None
        self.offset = QPoint(0, 0)
        self.dragging = False
        self.last_pos = None
        self.bg_color = QColor(20, 20, 22)
        self._drag_handle = -1  # -1=none, 0-7=resize, 8=rotate
        self._drag_start_pos = None
        self._drag_start_data = None
        self._dragging_annotation = False  # 是否在拖动整个注释
        self._drag_annot_start_img_pos = None  # 注释拖动起始图片坐标
        self.active_db = active_db
        self.user_password = user_password
        self.image_id = image_id
        self.setMouseTracking(True)
        self.setFocusPolicy(Qt.StrongFocus)
    
    def _notify_modified(self):
        """通知编辑器注释已修改，触发自动保存"""
        if hasattr(self.editor, '_trigger_auto_save'):
            self.editor._trigger_auto_save()

    def _annot_json_path(self):
        """获取注释数据文件路径"""
        if self.image_path:
            return os.path.splitext(self.image_path)[0] + '.annotations.json'
        return None

    def _serialize_annotation(self, ann):
        """将单个注释序列化为可JSON化的字典"""
        atype, data = ann
        d = {'type': atype}
        for k, v in data.items():
            if isinstance(v, QColor):
                d[k] = {'__color__': v.name()}
            else:
                d[k] = v
        return d

    def _deserialize_annotation(self, d):
        """从字典反序列化单个注释"""
        atype = d.pop('type')
        data = {}
        for k, v in d.items():
            if isinstance(v, dict) and '__color__' in v:
                data[k] = QColor(v['__color__'])
            else:
                data[k] = v
        return (atype, data)

    def save_annotations_json(self):
        """保存注释数据（优先加密存储到数据库，回退到JSON文件）"""
        serialized = [self._serialize_annotation(ann) for ann in self.annotations] if self.annotations else []
        json_str = json.dumps(serialized, ensure_ascii=False, indent=2)

        # 优先加密存储到数据库
        if self.active_db and self.user_password and self.image_id:
            try:
                key = generate_key(self.user_password)
                if serialized:
                    store_annotations_encrypted(self.active_db, self.image_id, json_str, key)
                else:
                    # 没有注释时清空数据库中的加密注释
                    conn = sqlite3.connect(self.active_db)
                    c = conn.cursor()
                    c.execute("UPDATE images SET annotations_encrypted=NULL WHERE id=?", (self.image_id,))
                    conn.commit()
                    conn.close()
                return
            except Exception as e:
                print(f"加密保存注释到数据库失败: {e}")

        # 回退到JSON文件
        json_path = self._annot_json_path()
        if not json_path or not self.annotations:
            if json_path and os.path.exists(json_path):
                os.remove(json_path)
            return
        try:
            with open(json_path, 'w', encoding='utf-8') as f:
                json.dump(serialized, f, ensure_ascii=False, indent=2)
        except Exception as e:
            print(f"保存注释数据失败: {e}")

    def load_annotations_json(self):
        """加载注释数据（优先从数据库解密，回退到JSON文件）"""
        # 优先从数据库加载加密注释
        if self.active_db and self.user_password and self.image_id:
            try:
                key = generate_key(self.user_password)
                decrypted = load_annotations_decrypted(self.active_db, self.image_id, key)
                if decrypted:
                    serialized = json.loads(decrypted)
                    self.annotations = [self._deserialize_annotation(d) for d in serialized]
                    return
            except Exception as e:
                print(f"从数据库加载注释失败: {e}")

        # 回退到JSON文件
        json_path = self._annot_json_path()
        if not json_path or not os.path.exists(json_path):
            return
        try:
            with open(json_path, 'r', encoding='utf-8') as f:
                serialized = json.load(f)
            self.annotations = [self._deserialize_annotation(d) for d in serialized]
        except Exception as e:
            print(f"加载注释数据失败: {e}")
            self.annotations = []

    def load_image(self):
        # 优先从数据库加载加密图片
        if self.active_db and self.image_id and self.user_password:
            try:
                key = generate_key(self.user_password)
                decrypted_data = load_image_decrypted(self.active_db, self.image_id, key)
                if decrypted_data:
                    pixmap = QPixmap()
                    pixmap.loadFromData(decrypted_data)
                    if not pixmap.isNull():
                        self.pixmap = pixmap
                        self.load_annotations_json()
                        self.update()
                        return
            except Exception as e:
                print(f"从数据库加载图片失败: {e}")

        # 回退到文件系统
        if self.image_path and os.path.exists(self.image_path):
            self.pixmap = QPixmap(self.image_path)
            self.load_annotations_json()
            self.update()

    def _img_offset(self):
        """返回图片在画布上的偏移量 (x, y)"""
        if not self.pixmap:
            return 0, 0
        scaled = self.pixmap.scaled(
            self.pixmap.size() * self.scale,
            Qt.KeepAspectRatio, Qt.SmoothTransformation
        )
        x = (self.width() - scaled.width()) // 2 + self.offset.x()
        y = (self.height() - scaled.height()) // 2 + self.offset.y()
        return x, y

    def _img_coords(self, pos):
        """将画布坐标转换为图片坐标"""
        x, y = self._img_offset()
        return pos.x() - x, pos.y() - y

    def _to_screen(self, img_x, img_y):
        """将图片坐标转换为画布坐标"""
        x, y = self._img_offset()
        return img_x + x, img_y + y

    @staticmethod
    def _get_annotation_center(data):
        """获取注释的中心点（图片坐标）"""
        if 'x' in data and 'w' in data:
            return data['x'] + data['w'] / 2, data['y'] + data['h'] / 2
        elif 'x1' in data:
            return (data['x1'] + data['x2']) / 2, (data['y1'] + data['y2']) / 2
        elif 'x' in data and 'y' in data:
            return data['x'], data['y']
        return 0, 0

    @staticmethod
    def _get_annotation_bbox(data):
        """获取注释的边界框 (x, y, w, h) 图片坐标"""
        if 'x' in data and 'w' in data:
            return data['x'], data['y'], data['w'], data['h']
        elif 'x1' in data:
            x = min(data['x1'], data['x2'])
            y = min(data['y1'], data['y2'])
            w = abs(data['x2'] - data['x1'])
            h = abs(data['y2'] - data['y1'])
            return x, y, w, h
        elif 'x' in data and 'y' in data:
            return data['x'] - 20, data['y'] - 20, 40, 40
        return 0, 0, 0, 0

    def _reverse_rotate_point(self, px, py, data):
        """将画布坐标点逆旋转到注释的局部坐标系"""
        rotation = data.get('rotation', 0)
        if rotation == 0:
            return px, py
        cx, cy = data.get('anchor_x', 0), data.get('anchor_y', 0)
        sx, sy = self._to_screen(cx, cy)
        dx = px - sx
        dy = py - sy
        angle = -rotation * math.pi / 180
        rx = dx * math.cos(angle) - dy * math.sin(angle) + sx
        ry = dx * math.sin(angle) + dy * math.cos(angle) + sy
        return rx, ry

    def _get_handles(self, data):
        """获取选中注释的8个缩放手柄和1个旋转手柄位置（画布坐标）"""
        bx, by, bw, bh = self._get_annotation_bbox(data)
        sx, sy = self._to_screen(bx, by)
        sw, sh = bw * self.scale, bh * self.scale
        hs = self.HANDLE_SIZE
        # 8个缩放手柄: TL, T, TR, R, BR, B, BL, L
        positions = [
            (sx, sy), (sx + sw / 2, sy), (sx + sw, sy),
            (sx + sw, sy + sh / 2), (sx + sw, sy + sh),
            (sx + sw / 2, sy + sh), (sx, sy + sh), (sx, sy + sh / 2),
        ]
        # 旋转手柄
        rot_pos = (sx + sw / 2, sy - self.ROTATE_HANDLE_DIST)
        return positions, rot_pos

    def _hit_handle(self, pos, data):
        """检测点击了哪个手柄，返回 -1=无, 0-7=缩放, 8=旋转"""
        positions, rot_pos = self._get_handles(data)
        hs = self.HANDLE_SIZE
        for i, (hx, hy) in enumerate(positions):
            if abs(pos.x() - hx) <= hs and abs(pos.y() - hy) <= hs:
                return i
        rx, ry = rot_pos
        if abs(pos.x() - rx) <= hs and abs(pos.y() - ry) <= hs:
            return 8
        return -1

    def _find_annotation_at(self, img_x, img_y, threshold=15):
        """查找点击位置对应的注释索引（考虑旋转）"""
        for i in range(len(self.annotations) - 1, -1, -1):
            atype, data = self.annotations[i]
            rotation = data.get('rotation', 0)
            # 逆旋转点击点
            if rotation != 0:
                cx, cy = data.get('anchor_x', 0), data.get('anchor_y', 0)
                dx = img_x - cx
                dy = img_y - cy
                angle = -rotation * math.pi / 180
                rx = dx * math.cos(angle) - dy * math.sin(angle) + cx
                ry = dx * math.sin(angle) + dy * math.cos(angle) + cy
            else:
                rx, ry = img_x, img_y

            if atype == 'text':
                if abs(data['x'] - rx) < 50 and abs(data['y'] - ry) < 30:
                    return i
            elif atype in ('arrow', 'line', 'measure'):
                mx = (data['x1'] + data['x2']) / 2
                my = (data['y1'] + data['y2']) / 2
                if abs(mx - rx) < threshold * 3 and abs(my - ry) < threshold * 3:
                    return i
            elif atype in ('rect', 'circle'):
                cx_a = data['x'] + data['w'] / 2
                cy_a = data['y'] + data['h'] / 2
                if abs(cx_a - rx) < max(data['w'] / 2, threshold) and \
                   abs(cy_a - ry) < max(data['h'] / 2, threshold):
                    return i
        return -1

    def paintEvent(self, event):
        painter = QPainter(self)
        painter.setRenderHint(QPainter.Antialiasing)
        painter.fillRect(self.rect(), self.bg_color)

        if not self.pixmap:
            painter.end()
            return

        scaled = self.pixmap.scaled(
            self.pixmap.size() * self.scale,
            Qt.KeepAspectRatio, Qt.SmoothTransformation
        )
        x, y = self._img_offset()
        painter.drawPixmap(x, y, scaled)

        color = self.editor._get_color()
        line_w = self.editor._get_line_width()

        # 绘制已有注释
        for ann_idx, ann in enumerate(self.annotations):
            atype, data = ann
            ann_color = data.get('color', color)
            ann_width = data.get('width', line_w)
            rotation = data.get('rotation', 0)
            anchor_x = data.get('anchor_x', 0)
            anchor_y = data.get('anchor_y', 0)

            if rotation != 0:
                sx, sy = self._to_screen(anchor_x, anchor_y)
                painter.save()
                painter.translate(sx, sy)
                painter.rotate(rotation)
                painter.translate(-sx, -sy)

            # 应用线型
            line_style = data.get('line_style', 'solid')
            style_map = {'solid': Qt.SolidLine, 'dashed': Qt.DashLine,
                         'dotted': Qt.DotLine, 'wavy': Qt.DashDotDotLine}
            pen = QPen(ann_color, ann_width)
            pen.setStyle(style_map.get(line_style, Qt.SolidLine))
            painter.setPen(pen)

            bg_color_ann = data.get('bg_color')
            bg_opacity_ann = data.get('bg_opacity', 140)

            if atype == 'text':
                painter.setPen(QPen(ann_color, 1))
                font_size = max(10, int(14 * self.scale))
                painter.setFont(QFont("SimSun", font_size, QFont.Bold))
                fm = painter.fontMetrics()
                text_rect = fm.boundingRect(data['text'])
                scr_x, scr_y = self._to_screen(data['x'], data['y'])
                bg_rect = QRect(scr_x - 2, scr_y - text_rect.height() - 2,
                                text_rect.width() + 4, text_rect.height() + 4)
                if bg_color_ann:
                    painter.fillRect(bg_rect, QColor(bg_color_ann.red(), bg_color_ann.green(), bg_color_ann.blue(), bg_opacity_ann))
                else:
                    painter.fillRect(bg_rect, QColor(0, 0, 0, bg_opacity_ann))
                painter.setPen(QPen(ann_color, 1))
                painter.drawText(scr_x, scr_y, data['text'])
            elif atype == 'arrow':
                sx1, sy1 = self._to_screen(data['x1'], data['y1'])
                sx2, sy2 = self._to_screen(data['x2'], data['y2'])
                # 使用带线型的箭头绘制
                painter.setPen(pen)
                painter.drawLine(int(sx1), int(sy1), int(sx2), int(sy2))
                angle = math.atan2(sy2 - sy1, sx2 - sx1)
                arrow_size = max(10, ann_width * 5)
                p1 = QPoint(int(sx2 - arrow_size * math.cos(angle - 0.4)),
                             int(sy2 - arrow_size * math.sin(angle - 0.4)))
                p2 = QPoint(int(sx2 - arrow_size * math.cos(angle + 0.4)),
                             int(sy2 - arrow_size * math.sin(angle + 0.4)))
                painter.setBrush(ann_color)
                painter.drawPolygon(QPolygon([QPoint(int(sx2), int(sy2)), p1, p2]))
                painter.setBrush(Qt.NoBrush)
            elif atype == 'line':
                sx1, sy1 = self._to_screen(data['x1'], data['y1'])
                sx2, sy2 = self._to_screen(data['x2'], data['y2'])
                painter.drawLine(int(sx1), int(sy1), int(sx2), int(sy2))
            elif atype == 'rect':
                sx, sy = self._to_screen(data['x'], data['y'])
                sw, sh = data['w'] * self.scale, data['h'] * self.scale
                fill_color = bg_color_ann if bg_color_ann else ann_color
                painter.setBrush(QColor(fill_color.red(), fill_color.green(), fill_color.blue(), bg_opacity_ann))
                painter.drawRect(QRectF(sx, sy, sw, sh))
                painter.setBrush(Qt.NoBrush)
            elif atype == 'circle':
                sx, sy = self._to_screen(data['x'], data['y'])
                sw, sh = data['w'] * self.scale, data['h'] * self.scale
                fill_color = bg_color_ann if bg_color_ann else ann_color
                painter.setBrush(QColor(fill_color.red(), fill_color.green(), fill_color.blue(), bg_opacity_ann))
                painter.drawEllipse(QRectF(sx, sy, sw, sh))
                painter.setBrush(Qt.NoBrush)
            elif atype == 'measure':
                sx1, sy1 = self._to_screen(data['x1'], data['y1'])
                sx2, sy2 = self._to_screen(data['x2'], data['y2'])
                painter.drawLine(int(sx1), int(sy1), int(sx2), int(sy2))
                painter.setBrush(ann_color)
                painter.drawEllipse(QPoint(int(sx1), int(sy1)), 3, 3)
                painter.drawEllipse(QPoint(int(sx2), int(sy2)), 3, 3)
                painter.setBrush(Qt.NoBrush)
                dist = data.get('distance', '')
                if dist:
                    mx = (sx1 + sx2) / 2
                    my = (sy1 + sy2) / 2 - 8
                    painter.setFont(QFont("SimSun", max(9, int(11 * self.scale)), QFont.Bold))
                    painter.setPen(QPen(ann_color, 1))
                    fm = painter.fontMetrics()
                    tr = fm.boundingRect(dist)
                    painter.fillRect(QRect(int(mx - 2), int(my - tr.height() - 2), tr.width() + 4, tr.height() + 4),
                                     QColor(0, 0, 0, 160))
                    painter.drawText(int(mx), int(my), dist)

            if rotation != 0:
                painter.restore()

        # 绘制选中注释的选择框和手柄
        if 0 <= self.selected_index < len(self.annotations):
            atype, data = self.annotations[self.selected_index]
            rotation = data.get('rotation', 0)
            anchor_x = data.get('anchor_x', 0)
            anchor_y = data.get('anchor_y', 0)

            if rotation != 0:
                sx, sy = self._to_screen(anchor_x, anchor_y)
                painter.save()
                painter.translate(sx, sy)
                painter.rotate(rotation)
                painter.translate(-sx, -sy)

            bx, by, bw, bh = self._get_annotation_bbox(data)
            sbx, sby = self._to_screen(bx, by)
            sbw, sbh = bw * self.scale, bh * self.scale

            # 虚线边框
            pen = QPen(QColor(63, 123, 247), 1, Qt.DashLine)
            painter.setPen(pen)
            painter.setBrush(Qt.NoBrush)
            painter.drawRect(QRectF(sbx - 2, sby - 2, sbw + 4, sbh + 4))

            # 缩放手柄
            positions, rot_pos = self._get_handles(data)
            hs = self.HANDLE_SIZE
            painter.setPen(QPen(QColor(63, 123, 247), 1))
            painter.setBrush(QColor(240, 240, 245))
            for hx, hy in positions:
                painter.drawRect(QRectF(hx - hs / 2, hy - hs / 2, hs, hs))

            # 旋转手柄
            rx, ry = rot_pos
            painter.setBrush(QColor(63, 123, 247))
            painter.drawEllipse(QPointF(rx, ry), hs / 2, hs / 2)
            # 连接线
            painter.setPen(QPen(QColor(63, 123, 247), 1, Qt.DashLine))
            top_center_x = sbx + sbw / 2
            top_center_y = sby
            painter.drawLine(int(top_center_x), int(top_center_y), int(rx), int(ry))

            if rotation != 0:
                painter.restore()

        # 绘制当前正在绘制的图形
        if self.drawing and self.draw_start and self.draw_current:
            pen = QPen(color, line_w)
            painter.setPen(pen)
            if self.tool == 'arrow':
                self._draw_arrow(painter, self.draw_start.x(), self.draw_start.y(),
                                 self.draw_current.x(), self.draw_current.y(), color, line_w)
            elif self.tool == 'line':
                painter.drawLine(self.draw_start, self.draw_current)
            elif self.tool == 'rect':
                rx = min(self.draw_start.x(), self.draw_current.x())
                ry = min(self.draw_start.y(), self.draw_current.y())
                rw = abs(self.draw_current.x() - self.draw_start.x())
                rh = abs(self.draw_current.y() - self.draw_start.y())
                painter.setBrush(QColor(color.red(), color.green(), color.blue(), 30))
                painter.drawRect(rx, ry, rw, rh)
                painter.setBrush(Qt.NoBrush)
            elif self.tool == 'circle':
                rx = min(self.draw_start.x(), self.draw_current.x())
                ry = min(self.draw_start.y(), self.draw_current.y())
                rw = abs(self.draw_current.x() - self.draw_start.x())
                rh = abs(self.draw_current.y() - self.draw_start.y())
                painter.setBrush(QColor(color.red(), color.green(), color.blue(), 30))
                painter.drawEllipse(rx, ry, rw, rh)
                painter.setBrush(Qt.NoBrush)
            elif self.tool == 'measure':
                painter.drawLine(self.draw_start, self.draw_current)
                painter.setBrush(color)
                painter.drawEllipse(self.draw_start, 3, 3)
                painter.drawEllipse(self.draw_current, 3, 3)
                painter.setBrush(Qt.NoBrush)
                dx = abs(self.draw_current.x() - self.draw_start.x()) / self.scale
                dy = abs(self.draw_current.y() - self.draw_start.y()) / self.scale
                dist_px = (dx ** 2 + dy ** 2) ** 0.5
                painter.setFont(QFont("SimSun", 11, QFont.Bold))
                painter.setPen(QPen(color, 1))
                mx = (self.draw_start.x() + self.draw_current.x()) / 2
                my = (self.draw_start.y() + self.draw_current.y()) / 2 - 8
                painter.drawText(int(mx), int(my), f'{dist_px:.1f}px')

        painter.end()

    def _draw_arrow(self, painter, x1, y1, x2, y2, color=None, width=2):
        pen = QPen(color or QColor(255, 80, 80), width)
        painter.setPen(pen)
        painter.drawLine(int(x1), int(y1), int(x2), int(y2))
        angle = math.atan2(y2 - y1, x2 - x1)
        arrow_size = max(10, width * 5)
        p1 = QPoint(int(x2 - arrow_size * math.cos(angle - 0.4)),
                     int(y2 - arrow_size * math.sin(angle - 0.4)))
        p2 = QPoint(int(x2 - arrow_size * math.cos(angle + 0.4)),
                     int(y2 - arrow_size * math.sin(angle + 0.4)))
        painter.setBrush(color or QColor(255, 80, 80))
        painter.drawPolygon(QPolygon([QPoint(int(x2), int(y2)), p1, p2]))
        painter.setBrush(Qt.NoBrush)

    def wheelEvent(self, event):
        delta = event.angleDelta().y()
        if delta > 0:
            self.editor._zoom_in()
        else:
            self.editor._zoom_out()

    def mousePressEvent(self, event):
        if event.button() != Qt.LeftButton:
            return

        # 选择工具
        if self.tool == 'select':
            # 先检查是否点击了选中注释的手柄
            if 0 <= self.selected_index < len(self.annotations):
                _, sel_data = self.annotations[self.selected_index]
                handle = self._hit_handle(event.pos(), sel_data)
                if handle >= 0:
                    self._drag_handle = handle
                    self._drag_start_pos = event.pos()
                    self._drag_start_data = dict(sel_data)
                    return

            # 检查是否点击了某个注释
            img_x, img_y = self._img_coords(event.pos())
            idx = self._find_annotation_at(img_x, img_y)
            if idx >= 0:
                self.selected_index = idx
                self.editor.show_property_panel(self.annotations[idx][1])
                # 开始拖动
                self._dragging_annotation = True
                self._drag_start_pos = event.pos()
                self._drag_start_data = dict(self.annotations[idx][1])
                self._drag_annot_start_img_pos = (img_x, img_y)
                self.setCursor(Qt.ClosedHandCursor)
            else:
                self.selected_index = -1
                self.editor.hide_property_panel()
            self.update()
            return

        if self.tool == 'pan':
            self.dragging = True
            self.last_pos = event.pos()
            self.setCursor(Qt.ClosedHandCursor)
        elif self.tool == 'text':
            ix, iy = self._img_coords(event.pos())
            text, ok = QInputDialog.getText(self, '添加文字', '注释文字:')
            if ok and text:
                color = self.editor._get_color()
                width = self.editor._get_line_width()
                self.annotations.append(('text', {
                    'x': ix, 'y': iy, 'text': text,
                    'color': color, 'width': width,
                    'rotation': 0, 'scale_x': 1.0, 'scale_y': 1.0,
                    'bg_color': None, 'bg_opacity': 140,
                    'anchor_x': ix, 'anchor_y': iy,
                    'line_style': 'solid',
                }))
                self.selected_index = len(self.annotations) - 1
                self.editor.show_property_panel(self.annotations[-1][1])
                self.editor._set_tool('select')
                self._notify_modified()
                self.update()
        elif self.tool == 'eraser':
            ix, iy = self._img_coords(event.pos())
            idx = self._find_annotation_at(ix, iy)
            if idx >= 0:
                self.annotations.pop(idx)
                if self.selected_index == idx:
                    self.selected_index = -1
                    self.editor.hide_property_panel()
                elif self.selected_index > idx:
                    self.selected_index -= 1
                self._notify_modified()
                self.update()
        else:
            self.drawing = True
            self.draw_start = event.pos()
            self.draw_current = event.pos()

    def mouseMoveEvent(self, event):
        # 拖拽手柄
        if self._drag_handle >= 0 and self._drag_start_pos:
            pos = event.pos()
            if self._drag_handle == 8:
                # 旋转手柄
                _, data = self.annotations[self.selected_index]
                anchor_x = data.get('anchor_x', 0)
                anchor_y = data.get('anchor_y', 0)
                sx, sy = self._to_screen(anchor_x, anchor_y)
                angle = math.atan2(pos.y() - sy, pos.x() - sx)
                deg = math.degrees(angle) + 90  # 0度朝上
                if deg < 0:
                    deg += 360
                data['rotation'] = deg % 360
                # 更新属性面板
                self.editor.prop_panel.rotation_input.blockSignals(True)
                self.editor.prop_panel.rotation_input.setValue(int(deg % 360))
                self.editor.prop_panel.rotation_input.blockSignals(False)
            else:
                # 缩放手柄
                self._resize_annotation(pos)
            self.update()
            return

        # 拖拽整个注释
        if self._dragging_annotation and self._drag_start_pos and self._drag_annot_start_img_pos:
            if 0 <= self.selected_index < len(self.annotations):
                current_img_x, current_img_y = self._img_coords(event.pos())
                start_img_x, start_img_y = self._drag_annot_start_img_pos
                dx = current_img_x - start_img_x
                dy = current_img_y - start_img_y
                atype, data = self.annotations[self.selected_index]
                # 根据注释类型移动坐标
                if atype == 'text':
                    data['x'] = self._drag_start_data['x'] + dx
                    data['y'] = self._drag_start_data['y'] + dy
                    data['anchor_x'] = data['x']
                    data['anchor_y'] = data['y']
                elif atype in ('arrow', 'line', 'measure'):
                    data['x1'] = self._drag_start_data['x1'] + dx
                    data['y1'] = self._drag_start_data['y1'] + dy
                    data['x2'] = self._drag_start_data['x2'] + dx
                    data['y2'] = self._drag_start_data['y2'] + dy
                    data['anchor_x'] = (data['x1'] + data['x2']) / 2
                    data['anchor_y'] = (data['y1'] + data['y2']) / 2
                elif atype in ('rect', 'circle'):
                    data['x'] = self._drag_start_data['x'] + dx
                    data['y'] = self._drag_start_data['y'] + dy
                    data['anchor_x'] = data['x'] + data['w'] / 2
                    data['anchor_y'] = data['y'] + data['h'] / 2
                self.update()
                self._notify_modified()
            return

        if self.dragging and self.last_pos:
            delta = event.pos() - self.last_pos
            self.offset += delta
            self.last_pos = event.pos()
            self.update()
        elif self.drawing:
            self.draw_current = event.pos()
            self.update()
        elif self.tool == 'select':
            # 选择工具时根据鼠标位置切换光标
            if 0 <= self.selected_index < len(self.annotations):
                _, sel_data = self.annotations[self.selected_index]
                handle = self._hit_handle(event.pos(), sel_data)
                if handle >= 0:
                    # 手柄光标
                    self.setCursor(self._cursor_for_handle(handle))
                else:
                    # 检查是否在注释内部
                    img_x, img_y = self._img_coords(event.pos())
                    idx = self._find_annotation_at(img_x, img_y)
                    if idx == self.selected_index:
                        self.setCursor(Qt.OpenHandCursor)
                    else:
                        self.setCursor(Qt.ArrowCursor)
            else:
                self.setCursor(Qt.ArrowCursor)

    def _cursor_for_handle(self, handle):
        """根据手柄索引返回对应的光标样式"""
        # 手柄顺序: 0=TL, 1=T, 2=TR, 3=R, 4=BR, 5=B, 6=BL, 7=L, 8=旋转
        cursors = {
            0: Qt.SizeFDiagCursor,   # TL
            1: Qt.SizeVerCursor,     # T
            2: Qt.SizeBDiagCursor,   # TR
            3: Qt.SizeHorCursor,     # R
            4: Qt.SizeFDiagCursor,   # BR
            5: Qt.SizeVerCursor,     # B
            6: Qt.SizeBDiagCursor,   # BL
            7: Qt.SizeHorCursor,     # L
            8: Qt.CrossCursor,       # 旋转
        }
        return cursors.get(handle, Qt.ArrowCursor)

    def mouseReleaseEvent(self, event):
        if event.button() != Qt.LeftButton:
            return

        # 结束手柄拖拽
        if self._drag_handle >= 0:
            self._drag_handle = -1
            self._drag_start_pos = None
            self._drag_start_data = None
            return

        # 结束注释拖拽
        if self._dragging_annotation:
            self._dragging_annotation = False
            self._drag_start_pos = None
            self._drag_start_data = None
            self._drag_annot_start_img_pos = None
            self.setCursor(Qt.OpenHandCursor)
            return

        if self.dragging:
            self.dragging = False
            self.last_pos = None
            self.setCursor(Qt.ArrowCursor)
        elif self.drawing and self.draw_start and self.draw_current:
            ix1, iy1 = self._img_coords(self.draw_start)
            ix2, iy2 = self._img_coords(self.draw_current)
            color = self.editor._get_color()
            width = self.editor._get_line_width()
            line_style = self.editor._get_line_style()
            base_props = {
                'color': color, 'width': width,
                'rotation': 0, 'scale_x': 1.0, 'scale_y': 1.0,
                'bg_color': None, 'bg_opacity': 140,
                'line_style': line_style,
            }
            if self.tool == 'arrow':
                base_props.update({'x1': ix1, 'y1': iy1, 'x2': ix2, 'y2': iy2})
                base_props['anchor_x'] = (ix1 + ix2) / 2
                base_props['anchor_y'] = (iy1 + iy2) / 2
                self.annotations.append(('arrow', base_props))
            elif self.tool == 'line':
                base_props.update({'x1': ix1, 'y1': iy1, 'x2': ix2, 'y2': iy2})
                base_props['anchor_x'] = (ix1 + ix2) / 2
                base_props['anchor_y'] = (iy1 + iy2) / 2
                self.annotations.append(('line', base_props))
            elif self.tool == 'rect':
                rx = min(ix1, ix2)
                ry = min(iy1, iy2)
                base_props.update({'x': rx, 'y': ry, 'w': abs(ix2 - ix1), 'h': abs(iy2 - iy1)})
                base_props['anchor_x'] = rx + abs(ix2 - ix1) / 2
                base_props['anchor_y'] = ry + abs(iy2 - iy1) / 2
                self.annotations.append(('rect', base_props))
            elif self.tool == 'circle':
                rx = min(ix1, ix2)
                ry = min(iy1, iy2)
                base_props.update({'x': rx, 'y': ry, 'w': abs(ix2 - ix1), 'h': abs(iy2 - iy1)})
                base_props['anchor_x'] = rx + abs(ix2 - ix1) / 2
                base_props['anchor_y'] = ry + abs(iy2 - iy1) / 2
                self.annotations.append(('circle', base_props))
            elif self.tool == 'measure':
                dx = abs(ix2 - ix1)
                dy = abs(iy2 - iy1)
                dist = (dx ** 2 + dy ** 2) ** 0.5
                base_props.update({
                    'x1': ix1, 'y1': iy1, 'x2': ix2, 'y2': iy2,
                    'distance': f'{dist:.1f}px',
                })
                base_props['anchor_x'] = (ix1 + ix2) / 2
                base_props['anchor_y'] = (iy1 + iy2) / 2
                self.annotations.append(('measure', base_props))

            # 自动选中新创建的注释并切换到选择工具
            self.selected_index = len(self.annotations) - 1
            self.editor.show_property_panel(self.annotations[-1][1])
            self.editor._set_tool('select')
            self._notify_modified()

            self.drawing = False
            self.draw_start = None
            self.draw_current = None
            self.update()

    def _resize_annotation(self, pos):
        """根据拖拽的缩放手柄调整注释大小"""
        if self.selected_index < 0 or not self._drag_start_data:
            return
        atype, data = self.annotations[self.selected_index]
        sd = self._drag_start_data
        handle = self._drag_handle

        # 计算鼠标移动的图片坐标偏移
        img_x, img_y = self._img_coords(pos)
        start_img_x = self._drag_start_pos.x() - self._img_offset()[0]
        start_img_y = self._drag_start_pos.y() - self._img_offset()[1]
        dx_img = img_x - start_img_x
        dy_img = img_y - start_img_y

        if 'w' in sd:
            # rect/circle 类型
            x, y, w, h = sd['x'], sd['y'], sd['w'], sd['h']
            if handle == 0:  # TL
                data['x'] = x + dx_img
                data['y'] = y + dy_img
                data['w'] = w - dx_img
                data['h'] = h - dy_img
            elif handle == 1:  # T
                data['y'] = y + dy_img
                data['h'] = h - dy_img
            elif handle == 2:  # TR
                data['y'] = y + dy_img
                data['w'] = w + dx_img
                data['h'] = h - dy_img
            elif handle == 3:  # R
                data['w'] = w + dx_img
            elif handle == 4:  # BR
                data['w'] = w + dx_img
                data['h'] = h + dy_img
            elif handle == 5:  # B
                data['h'] = h + dy_img
            elif handle == 6:  # BL
                data['x'] = x + dx_img
                data['w'] = w - dx_img
                data['h'] = h + dy_img
            elif handle == 7:  # L
                data['x'] = x + dx_img
                data['w'] = w - dx_img
            # 确保宽高不为负
            if data['w'] < 5:
                data['w'] = 5
            if data['h'] < 5:
                data['h'] = 5
            # 更新锚点
            data['anchor_x'] = data['x'] + data['w'] / 2
            data['anchor_y'] = data['y'] + data['h'] / 2
        elif 'x1' in sd:
            # arrow/line/measure 类型
            x1, y1, x2, y2 = sd['x1'], sd['y1'], sd['x2'], sd['y2']
            if handle in (0, 6, 7):  # 左侧手柄影响起点
                data['x1'] = x1 + dx_img
                data['y1'] = y1 + dy_img
            if handle in (2, 3, 4):  # 右侧手柄影响终点
                data['x2'] = x2 + dx_img
                data['y2'] = y2 + dy_img
            if handle == 1:  # 上中
                data['y1'] = y1 + dy_img
            if handle == 5:  # 下中
                data['y2'] = y2 + dy_img
            data['anchor_x'] = (data['x1'] + data['x2']) / 2
            data['anchor_y'] = (data['y1'] + data['y2']) / 2
        
        self._notify_modified()

    def save_annotated(self, path):
        if not self.pixmap:
            return
        # 保存注释数据（优先加密存储到数据库）
        self.save_annotations_json()
        # 生成带注释的预览图
        if not self.annotations:
            # 没有注释时删除旧的预览图
            preview_path = os.path.splitext(path)[0] + '.preview.png'
            if os.path.exists(preview_path):
                os.remove(preview_path)
            return
        result = QPixmap(self.pixmap.size())
        result.fill(QColor(0, 0, 0, 0))
        painter = QPainter(result)
        painter.drawPixmap(0, 0, self.pixmap)

        for ann in self.annotations:
            atype, data = ann
            color = data.get('color', QColor(255, 80, 80))
            width = data.get('width', 2)
            rotation = data.get('rotation', 0)
            anchor_x = data.get('anchor_x', 0)
            anchor_y = data.get('anchor_y', 0)

            if rotation != 0:
                painter.save()
                painter.translate(anchor_x, anchor_y)
                painter.rotate(rotation)
                painter.translate(-anchor_x, -anchor_y)

            pen = QPen(color, width)
            line_style = data.get('line_style', 'solid')
            style_map = {'solid': Qt.SolidLine, 'dashed': Qt.DashLine,
                         'dotted': Qt.DotLine, 'wavy': Qt.DashDotDotLine}
            pen.setStyle(style_map.get(line_style, Qt.SolidLine))
            painter.setPen(pen)
            bg_color_ann = data.get('bg_color')
            bg_opacity_ann = data.get('bg_opacity', 140)
            if atype == 'text':
                painter.setPen(QPen(color, 1))
                painter.setFont(QFont("SimSun", 14, QFont.Bold))
                fm = painter.fontMetrics()
                text_rect = fm.boundingRect(data['text'])
                bg_rect = QRect(data['x'] - 2, data['y'] - text_rect.height() - 2,
                                text_rect.width() + 4, text_rect.height() + 4)
                if bg_color_ann:
                    painter.fillRect(bg_rect, QColor(bg_color_ann.red(), bg_color_ann.green(), bg_color_ann.blue(), bg_opacity_ann))
                else:
                    painter.fillRect(bg_rect, QColor(0, 0, 0, bg_opacity_ann))
                painter.setPen(QPen(color, 1))
                painter.drawText(data['x'], data['y'], data['text'])
                painter.setPen(pen)
            elif atype == 'arrow':
                self._draw_arrow(painter, data['x1'], data['y1'], data['x2'], data['y2'], color, width)
            elif atype == 'line':
                painter.drawLine(int(data['x1']), int(data['y1']), int(data['x2']), int(data['y2']))
            elif atype == 'rect':
                fill_color = bg_color_ann if bg_color_ann else color
                painter.setBrush(QColor(fill_color.red(), fill_color.green(), fill_color.blue(), bg_opacity_ann))
                painter.drawRect(data['x'], data['y'], data['w'], data['h'])
                painter.setBrush(Qt.NoBrush)
            elif atype == 'circle':
                fill_color = bg_color_ann if bg_color_ann else color
                painter.setBrush(QColor(fill_color.red(), fill_color.green(), fill_color.blue(), bg_opacity_ann))
                painter.drawEllipse(data['x'], data['y'], data['w'], data['h'])
                painter.setBrush(Qt.NoBrush)
            elif atype == 'measure':
                painter.drawLine(int(data['x1']), int(data['y1']), int(data['x2']), int(data['y2']))
                painter.setBrush(color)
                painter.drawEllipse(QPoint(int(data['x1']), int(data['y1'])), 3, 3)
                painter.drawEllipse(QPoint(int(data['x2']), int(data['y2'])), 3, 3)
                painter.setBrush(Qt.NoBrush)
                dist = data.get('distance', '')
                if dist:
                    painter.setFont(QFont("SimSun", 11, QFont.Bold))
                    mx = (data['x1'] + data['x2']) / 2
                    my = (data['y1'] + data['y2']) / 2 - 8
                    painter.drawText(int(mx), int(my), dist)

            if rotation != 0:
                painter.restore()

        painter.end()
        # 保存带注释的预览图（不覆盖原图）
        preview_path = os.path.splitext(path)[0] + '.preview.png'
        result.save(preview_path)


# ── 内联图片控件（支持悬停缩放与双击浏览） ─────────────────
class _ImageLabel(QWidget):
    """内联图片控件：悬停时滚轮缩放，双击打开图片编辑器"""
    double_clicked = pyqtSignal(int, str)  # img_id, img_path

    def __init__(self, pixmap, img_id=0, img_path='', max_width=400, parent=None):
        super().__init__(parent)
        self._pixmap = pixmap
        self.img_id = img_id
        self.img_path = img_path
        self._max_width = max_width
        self._scale = 1.0
        self._hover = False
        self._anim_timer = QTimer(self)
        self._anim_timer.setSingleShot(True)
        self._anim_timer.timeout.connect(self._animate_zoom)
        self._target_scale = 1.0
        self.setMouseTracking(True)
        self.setCursor(Qt.PointingHandCursor)
        self._update_size()

    def _update_size(self):
        if self._pixmap.isNull():
            self.setFixedSize(0, 0)
            return
        w = min(self._pixmap.width(), self._max_width)
        ratio = w / self._pixmap.width()
        h = int(self._pixmap.height() * ratio * self._scale)
        w = int(w * self._scale)
        self.setFixedSize(w, h)
        self.update()

    def paintEvent(self, event):
        if self._pixmap.isNull():
            return
        painter = QPainter(self)
        painter.setRenderHint(QPainter.SmoothPixmapTransform)
        w = min(self._pixmap.width(), self._max_width)
        scaled = self._pixmap.scaled(w, int(w * self._pixmap.height() / self._pixmap.width()),
                                      Qt.KeepAspectRatio, Qt.SmoothTransformation)
        if self._scale != 1.0:
            scaled = scaled.scaled(int(scaled.width() * self._scale),
                                   int(scaled.height() * self._scale),
                                   Qt.KeepAspectRatio, Qt.SmoothTransformation)
        x = (self.width() - scaled.width()) // 2
        y = (self.height() - scaled.height()) // 2
        painter.drawPixmap(x, y, scaled)
        # 悬停时绘制蓝色边框
        if self._hover:
            pen = QPen(QColor(63, 123, 247, 160), 2)
            painter.setPen(pen)
            painter.setBrush(Qt.NoBrush)
            painter.drawRect(1, 1, self.width() - 2, self.height() - 2)
        painter.end()

    def wheelEvent(self, event):
        if not self._hover:
            event.ignore()
            return
        delta = event.angleDelta().y()
        if delta > 0:
            self._target_scale = min(self._scale * 1.1, 3.0)
        else:
            self._target_scale = max(self._scale / 1.1, 1.0)
        self._animate_zoom()
        event.accept()

    def _animate_zoom(self):
        step = 0.05 if self._target_scale > self._scale else -0.05
        new_scale = self._scale + step
        if (step > 0 and new_scale >= self._target_scale) or \
           (step < 0 and new_scale <= self._target_scale):
            self._scale = self._target_scale
        else:
            self._scale = new_scale
            self._anim_timer.start(16)
        self._update_size()

    def enterEvent(self, event):
        self._hover = True
        self.update()
        super().enterEvent(event)

    def leaveEvent(self, event):
        self._hover = False
        if self._scale > 1.0:
            self._target_scale = 1.0
            self._animate_zoom()
        else:
            self.update()
        super().leaveEvent(event)

    def mouseDoubleClickEvent(self, event):
        self.double_clicked.emit(self.img_id, self.img_path)


class _DraggableTextBox(QWidget):
    """可拖拽的文本框组件（Office 风格，8方向手柄）"""
    removed = pyqtSignal()
    HANDLE_SIZE = 8

    def __init__(self, parent=None):
        super().__init__(parent)
        self.setMinimumSize(150, 80)
        self._dragging = False
        self._resize_mode = None  # 当前拖动的手柄
        self._drag_start = None
        self._resize_start_pos = None
        self._resize_start_rect = None
        self.setMouseTracking(True)

        layout = QVBoxLayout(self)
        layout.setContentsMargins(6, 6, 6, 6)

        # 文本编辑器
        self.text_edit = QTextEdit()
        self.text_edit.setPlaceholderText('输入文本...')
        self.text_edit.setContextMenuPolicy(Qt.NoContextMenu)
        self.text_edit.setStyleSheet("""
            QTextEdit {
                background: transparent; border: none;
                color: #e0e0e4; font-size: 12px;
            }
        """)
        layout.addWidget(self.text_edit)
        # 将text_edit的右键事件转发到父组件
        self.text_edit.installEventFilter(self)

        # 右键菜单
        self.setContextMenuPolicy(Qt.CustomContextMenu)
        self.customContextMenuRequested.connect(self._show_context_menu)

    def _get_handles(self):
        """获取8个手柄的位置和类型"""
        r = self.rect()
        hs = self.HANDLE_SIZE
        cx = r.left() + r.width() // 2
        cy = r.top() + r.height() // 2
        return {
            'top-left': QRect(r.left(), r.top(), hs, hs),
            'top': QRect(cx - hs // 2, r.top(), hs, hs),
            'top-right': QRect(r.right() - hs, r.top(), hs, hs),
            'right': QRect(r.right() - hs, cy - hs // 2, hs, hs),
            'bottom-right': QRect(r.right() - hs, r.bottom() - hs, hs, hs),
            'bottom': QRect(cx - hs // 2, r.bottom() - hs, hs, hs),
            'bottom-left': QRect(r.left(), r.bottom() - hs, hs, hs),
            'left': QRect(r.left(), cy - hs // 2, hs, hs),
        }

    def _hit_handle(self, pos):
        """检查鼠标位置是否在手柄上"""
        for name, hrect in self._get_handles().items():
            if hrect.contains(pos):
                return name
        return None

    def _cursor_for_handle(self, handle):
        """根据手柄位置返回光标样式"""
        return {
            'top-left': Qt.SizeFDiagCursor,
            'bottom-right': Qt.SizeFDiagCursor,
            'top-right': Qt.SizeBDiagCursor,
            'bottom-left': Qt.SizeBDiagCursor,
            'top': Qt.SizeVerCursor,
            'bottom': Qt.SizeVerCursor,
            'left': Qt.SizeHorCursor,
            'right': Qt.SizeHorCursor,
        }.get(handle, Qt.ArrowCursor)

    def eventFilter(self, obj, event):
        """将text_edit的右键事件转发到_DraggableTextBox"""
        if obj == self.text_edit:
            if event.type() == event.ContextMenu:
                # 转发右键菜单事件
                self._show_context_menu(event.pos())
                return True
        return super().eventFilter(obj, event)

    def paintEvent(self, event):
        painter = QPainter(self)
        painter.setRenderHint(QPainter.Antialiasing)
        # 半透明背景
        painter.setPen(Qt.NoPen)
        painter.setBrush(QColor(40, 40, 46, 200))
        painter.drawRoundedRect(self.rect(), 4, 4)
        # 边框
        painter.setPen(QPen(QColor(63, 123, 247), 2))
        painter.setBrush(Qt.NoBrush)
        painter.drawRoundedRect(self.rect(), 4, 4)
        # 8个手柄
        hs = self.HANDLE_SIZE
        painter.setPen(QPen(QColor(63, 123, 247), 1))
        painter.setBrush(QColor(255, 255, 255))
        for name, hrect in self._get_handles().items():
            painter.drawRect(hrect)

    def mousePressEvent(self, event):
        if event.button() == Qt.LeftButton:
            # 检查是否点击手柄
            handle = self._hit_handle(event.pos())
            if handle:
                self._resize_mode = handle
                self._resize_start_pos = event.globalPos()
                self._resize_start_rect = QRect(self.pos(), self.size())
                self.setCursor(self._cursor_for_handle(handle))
                event.accept()
                return
            # 否则开始拖动
            self._dragging = True
            self._drag_start = event.globalPos() - self.pos()
            self.setCursor(Qt.ClosedHandCursor)
            event.accept()
        super().mousePressEvent(event)

    def mouseMoveEvent(self, event):
        if self._resize_mode and self._resize_start_pos and self._resize_start_rect:
            delta = event.globalPos() - self._resize_start_pos
            sr = self._resize_start_rect
            mode = self._resize_mode
            min_w, min_h = 150, 80

            new_x, new_y = sr.x(), sr.y()
            new_w, new_h = sr.width(), sr.height()

            # 横向缩放
            if mode == 'right':
                new_w = max(min_w, sr.width() + delta.x())
            elif mode == 'left':
                new_w = max(min_w, sr.width() - delta.x())
                new_x = sr.right() - new_w
            # 纵向缩放
            elif mode == 'bottom':
                new_h = max(min_h, sr.height() + delta.y())
            elif mode == 'top':
                new_h = max(min_h, sr.height() - delta.y())
                new_y = sr.bottom() - new_h
            # 角点等比缩放
            elif mode in ('top-left', 'top-right', 'bottom-left', 'bottom-right'):
                if mode in ('bottom-right', 'top-right'):
                    new_w = max(min_w, sr.width() + delta.x())
                if mode in ('bottom-left', 'top-left'):
                    new_w = max(min_w, sr.width() - delta.x())
                    new_x = sr.right() - new_w
                if mode in ('bottom', 'bottom-left', 'bottom-right'):
                    new_h = max(min_h, sr.height() + delta.y())
                if mode in ('top', 'top-left', 'top-right'):
                    new_h = max(min_h, sr.height() - delta.y())
                    new_y = sr.bottom() - new_h
                # 等比缩放
                ratio = sr.width() / max(1, sr.height())
                new_h = max(min_h, int(new_w / ratio))
                if mode in ('top', 'top-left', 'top-right'):
                    new_y = sr.bottom() - new_h

            self.move(new_x, new_y)
            self.resize(new_w, new_h)
            event.accept()
        elif self._dragging:
            self.move(event.globalPos() - self._drag_start)
            event.accept()
        else:
            # 悬停时切换光标
            handle = self._hit_handle(event.pos())
            if handle:
                self.setCursor(self._cursor_for_handle(handle))
            elif self.rect().contains(event.pos()):
                self.setCursor(Qt.OpenHandCursor)
            else:
                self.setCursor(Qt.ArrowCursor)
        super().mouseMoveEvent(event)

    def mouseReleaseEvent(self, event):
        if self._resize_mode:
            self._resize_mode = None
            self._resize_start_pos = None
            self._resize_start_rect = None
        if self._dragging:
            self._dragging = False
            self._drag_start = None
            self.setCursor(Qt.OpenHandCursor)
        super().mouseReleaseEvent(event)
        
    def _show_context_menu(self, pos):
        menu = QMenu(self)
        menu.setStyleSheet("""
            QMenu {
                background-color: #2a2a30; color: #c8c8cc;
                border: 1px solid #444; padding: 4px;
            }
            QMenu::item { padding: 6px 20px; border-radius: 3px; }
            QMenu::item:selected { background-color: #3f7bf7; }
        """)
        
        act_font = menu.addAction('字体设置')
        act_remove = menu.addAction('删除文本框')
        
        action = menu.exec_(self.mapToGlobal(pos))
        if action == act_font:
            self._show_font_dialog()
        elif action == act_remove:
            self.removed.emit()
            self.deleteLater()
            
    def _show_font_dialog(self):
        """简单字体设置对话框"""
        from PyQt5.QtWidgets import QFontDialog
        font, ok = QFontDialog.getFont(self.text_edit.font(), self)
        if ok:
            self.text_edit.setFont(font)


class _OfficeImageToolbar(QWidget):
    """Office 风格浮动图片工具栏"""
    toolbar_action = pyqtSignal(str)
    
    def __init__(self, parent=None):
        super().__init__(parent)
        self.setWindowFlags(Qt.Tool | Qt.FramelessWindowHint | Qt.WindowStaysOnTopHint)
        self.setAttribute(Qt.WA_TranslucentBackground)
        self.setFixedSize(360, 48)
        
        layout = QHBoxLayout(self)
        layout.setContentsMargins(8, 6, 8, 6)
        layout.setSpacing(2)
        
        # 尺寸调整
        layout.addWidget(QLabel('尺寸:'))
        self.width_spin = QSpinBox()
        self.width_spin.setRange(10, 5000)
        self.width_spin.setSuffix(' px')
        self.width_spin.setFixedWidth(80)
        self.width_spin.valueChanged.connect(lambda: self.toolbar_action.emit('resize'))
        layout.addWidget(self.width_spin)
        
        self.height_spin = QSpinBox()
        self.height_spin.setRange(10, 5000)
        self.height_spin.setSuffix(' px')
        self.height_spin.setFixedWidth(80)
        self.height_spin.valueChanged.connect(lambda: self.toolbar_action.emit('resize'))
        layout.addWidget(self.height_spin)
        
        # 环绕方式
        layout.addSpacing(8)
        layout.addWidget(QLabel('环绕:'))
        self.wrap_combo = QComboBox()
        self.wrap_combo.addItems(['嵌入型', '四周型', '紧密型', '浮于上方', '衬于下方'])
        self.wrap_combo.setFixedWidth(90)
        self.wrap_combo.currentIndexChanged.connect(lambda: self.toolbar_action.emit('wrap'))
        layout.addWidget(self.wrap_combo)
        
        # 应用通用样式
        for w in [self.width_spin, self.height_spin, self.wrap_combo]:
            w.setStyleSheet("""
                QSpinBox, QComboBox {
                    background: rgba(50,50,56,250); color: #e0e0e4;
                    border: 1px solid #555; border-radius: 3px;
                    padding: 2px 4px; font-size: 11px;
                }
            """)
        
    def paintEvent(self, event):
        painter = QPainter(self)
        painter.setRenderHint(QPainter.Antialiasing)
        painter.setPen(Qt.NoPen)
        painter.setBrush(QColor(40, 40, 46, 245))
        painter.drawRoundedRect(self.rect(), 6, 6)
        painter.setPen(QPen(QColor(255, 255, 255, 40), 1))
        painter.setBrush(Qt.NoBrush)
        painter.drawRoundedRect(self.rect().adjusted(1, 1, -1, -1), 6, 6)
        
    def set_image_size(self, width, height):
        self.width_spin.blockSignals(True)
        self.height_spin.blockSignals(True)
        self.width_spin.setValue(int(width))
        self.height_spin.setValue(int(height))
        self.width_spin.blockSignals(False)
        self.height_spin.blockSignals(False)
        
    def get_image_size(self):
        return self.width_spin.value(), self.height_spin.value()
        
    def get_wrap_mode(self):
        modes = ['inline', 'square', 'tight', 'front', 'behind']
        return modes[self.wrap_combo.currentIndex()]


class _ImageEditableTextEdit(QTextEdit):
    """支持图片选中、调整尺寸和右键菜单的富文本编辑器（Office 风格）"""
    image_selected = pyqtSignal(str, int, int)  # 图片名, 宽度, 高度

    def __init__(self, parent=None):
        super().__init__(parent)
        # 选中图片状态
        self._selected_image_name = None  # 选中图片的资源名
        self._selected_image_pos = -1     # 选中图片在文档中的位置
        self._selected_image_rect = None  # 图片在视口中的矩形
        # 调整尺寸状态
        self._resize_handle_size = 8
        self._resize_mode = None
        self._resize_start_pos = None
        self._resize_start_rect = None
        # 拖动状态
        self._dragging_image = False
        self._drag_start_pos = None
        self._drag_cursor_pos = None
        # 当前工具（保留接口，但默认就是选中模式）
        self._current_tool = 'image_select'
        self._wrap_mode = 'inline'
        # Office 风格工具栏
        self._office_toolbar = None

        # 设置滚动条策略
        self.setVerticalScrollBarPolicy(Qt.ScrollBarAsNeeded)
        self.setHorizontalScrollBarPolicy(Qt.ScrollBarAsNeeded)
        self.setSizePolicy(QSizePolicy.Expanding, QSizePolicy.Expanding)
        self.setLineWrapMode(QTextEdit.WidgetWidth)
        self.setMouseTracking(True)
        self.viewport().setMouseTracking(True)
        # 安装事件过滤器以捕获viewport上的鼠标移动事件（用于光标切换）
        self.viewport().installEventFilter(self)

    def set_image_tool(self, tool_id):
        """设置当前图像编辑工具（兼容接口）"""
        self._current_tool = tool_id
        if tool_id == 'image_delete' and self._selected_image_name:
            self._delete_selected_image()
            return
        # 其它工具默认行为：保持选中模式
        self.setCursor(Qt.IBeamCursor)

    def eventFilter(self, obj, event):
        """捕获viewport上的鼠标移动事件，用于手柄光标切换"""
        if obj == self.viewport():
            if event.type() == event.MouseMove and not self._resize_mode and not self._dragging_image:
                if self._selected_image_rect:
                    handle = self._hit_resize_handle(event.pos())
                    if handle:
                        self.viewport().setCursor(self._cursor_for_handle(handle))
                    elif self._selected_image_rect.contains(event.pos()):
                        self.viewport().setCursor(Qt.OpenHandCursor)
                    else:
                        self.viewport().setCursor(Qt.IBeamCursor)
                else:
                    self.viewport().setCursor(Qt.IBeamCursor)
        return super().eventFilter(obj, event)

    def set_image_wrap_mode(self, mode):
        """设置图片环绕方式 - 通过修改文档HTML实现"""
        self._wrap_mode = mode
        if self._selected_image_pos < 0 or not self._selected_image_name:
            return
        # 保存当前光标位置
        saved_cursor = self.textCursor()
        # 获取当前图片格式
        cur = QTextCursor(self.document())
        cur.setPosition(self._selected_image_pos)
        cur.movePosition(QTextCursor.Right, QTextCursor.KeepAnchor)
        fmt = cur.charFormat()
        if not fmt.isImageFormat():
            return
        img_fmt = fmt.toImageFormat()
        img_name = img_fmt.name()
        img_w = int(img_fmt.width()) if img_fmt.width() > 0 else 0
        img_h = int(img_fmt.height()) if img_fmt.height() > 0 else 0

        # 构建新的HTML片段
        size_attr = ''
        if img_w > 0 and img_h > 0:
            size_attr = f' width="{img_w}" height="{img_h}"'

        html_map = {
            'inline': f'<img src="{img_name}"{size_attr} />',
            'square': f'<img src="{img_name}"{size_attr} style="float:left; margin:0 8px 8px 0;" />',
            'tight': f'<img src="{img_name}"{size_attr} style="float:left; margin:0 4px 4px 0;" />',
            'through': f'<img src="{img_name}"{size_attr} style="float:left; margin:0;" />',
            'topAndBottom': f'<img src="{img_name}"{size_attr} style="display:block; margin:8px auto;" />',
            'behind': f'<img src="{img_name}"{size_attr} style="float:left; position:relative; z-index:-1;" />',
            'front': f'<img src="{img_name}"{size_attr} style="float:left; position:relative; z-index:1;" />',
        }
        html = html_map.get(mode, f'<img src="{img_name}"{size_attr} />')

        # 替换图片
        cur.insertHtml(html)
        # 重新查找图片位置
        self._selected_image_pos = -1
        self._selected_image_rect = None
        # 重新定位图片
        block = self.document().begin()
        while block.isValid():
            it = block.begin()
            while not it.atEnd():
                frag = it.fragment()
                if frag.isValid() and frag.charFormat().isImageFormat():
                    if frag.charFormat().toImageFormat().name() == img_name:
                        self._selected_image_pos = frag.position()
                        self._selected_image_rect = self._image_view_rect(frag.position())
                        break
                it += 1
            if self._selected_image_pos >= 0:
                break
            block = block.next()
        # 恢复光标
        self.setTextCursor(saved_cursor)
        self.viewport().update()

    def _find_image_at(self, pos):
        """在视口的指定位置查找图片，返回 (image_name, doc_position, view_rect) 或 (None, -1, None)"""
        cursor = self.cursorForPosition(pos)
        doc_pos = cursor.position()
        # 检查光标右侧字符
        cur_right = QTextCursor(self.document())
        cur_right.setPosition(doc_pos)
        cur_right.movePosition(QTextCursor.Right, QTextCursor.KeepAnchor)
        fmt = cur_right.charFormat()
        if fmt.isImageFormat():
            img_name = fmt.toImageFormat().name()
            rect = self._image_view_rect(doc_pos)
            return img_name, doc_pos, rect
        # 检查光标左侧字符
        if doc_pos > 0:
            cur_left = QTextCursor(self.document())
            cur_left.setPosition(doc_pos - 1)
            cur_left.movePosition(QTextCursor.Right, QTextCursor.KeepAnchor)
            fmt = cur_left.charFormat()
            if fmt.isImageFormat():
                img_name = fmt.toImageFormat().name()
                rect = self._image_view_rect(doc_pos - 1)
                return img_name, doc_pos - 1, rect
        # 遍历文档中所有图片，检查点击位置是否在其矩形内
        block = self.document().begin()
        while block.isValid():
            it = block.begin()
            while not it.atEnd():
                frag = it.fragment()
                if frag.isValid() and frag.charFormat().isImageFormat():
                    img_fmt = frag.charFormat().toImageFormat()
                    frag_pos = frag.position()
                    rect = self._image_view_rect(frag_pos)
                    if rect and rect.contains(pos):
                        return img_fmt.name(), frag_pos, rect
                it += 1
            block = block.next()
        return None, -1, None

    def _image_view_rect(self, doc_pos):
        """根据文档中的图片位置计算其在视口中的实际矩形"""
        if doc_pos < 0:
            return None
        # 获取图片格式
        cur2 = QTextCursor(self.document())
        cur2.setPosition(doc_pos)
        cur2.movePosition(QTextCursor.Right, QTextCursor.KeepAnchor)
        fmt = cur2.charFormat()
        if not fmt.isImageFormat():
            return None
        img_fmt = fmt.toImageFormat()
        # 图片宽高
        w = int(img_fmt.width()) if img_fmt.width() > 0 else 0
        h = int(img_fmt.height()) if img_fmt.height() > 0 else 0
        if w <= 0 or h <= 0:
            res = self.document().resource(QTextDocument.ImageResource, QUrl(img_fmt.name()))
            if res is not None:
                try:
                    pix = QPixmap()
                    if isinstance(res, QPixmap):
                        pix = res
                    elif isinstance(res, QImage):
                        pix = QPixmap.fromImage(res)
                    elif isinstance(res, (bytes, bytearray)):
                        pix.loadFromData(bytes(res))
                    if not pix.isNull():
                        if w <= 0:
                            w = pix.width()
                        if h <= 0:
                            h = pix.height()
                except Exception:
                    pass
        if w <= 0 or h <= 0:
            return None
        # 使用 cursorRect 获取图片左上角位置
        cur = QTextCursor(self.document())
        cur.setPosition(doc_pos)
        rect_start = self.cursorRect(cur)
        return QRect(rect_start.topLeft(), QSize(w, h))

    def _select_image_at_pos(self, doc_pos):
        """根据文档位置选中图片"""
        if doc_pos < 0:
            self._clear_image_selection()
            return False
        cur = QTextCursor(self.document())
        cur.setPosition(doc_pos)
        cur.movePosition(QTextCursor.Right, QTextCursor.KeepAnchor)
        fmt = cur.charFormat()
        if not fmt.isImageFormat():
            self._clear_image_selection()
            return False
        self._selected_image_name = fmt.toImageFormat().name()
        self._selected_image_pos = doc_pos
        self._selected_image_rect = self._image_view_rect(doc_pos)
        self.viewport().update()
        return True

    def _clear_image_selection(self):
        self._selected_image_name = None
        self._selected_image_pos = -1
        self._selected_image_rect = None
        self.viewport().update()

    def _hit_resize_handle(self, pos):
        """检查鼠标位置是否在调整大小手柄上，返回手柄类型或 None"""
        if not self._selected_image_rect:
            return None
        hs = self._resize_handle_size
        rect = self._selected_image_rect
        # 8 个手柄（4角 + 4边中点）
        handles = {
            'top-left': QRect(rect.left() - hs, rect.top() - hs, hs * 2, hs * 2),
            'top': QRect(rect.left() + rect.width() // 2 - hs, rect.top() - hs, hs * 2, hs * 2),
            'top-right': QRect(rect.right() - hs, rect.top() - hs, hs * 2, hs * 2),
            'right': QRect(rect.right() - hs, rect.top() + rect.height() // 2 - hs, hs * 2, hs * 2),
            'bottom-right': QRect(rect.right() - hs, rect.bottom() - hs, hs * 2, hs * 2),
            'bottom': QRect(rect.left() + rect.width() // 2 - hs, rect.bottom() - hs, hs * 2, hs * 2),
            'bottom-left': QRect(rect.left() - hs, rect.bottom() - hs, hs * 2, hs * 2),
            'left': QRect(rect.left() - hs, rect.top() + rect.height() // 2 - hs, hs * 2, hs * 2),
        }
        for name, hrect in handles.items():
            if hrect.contains(pos):
                return name
        return None

    def _cursor_for_handle(self, handle):
        """根据手柄位置返回对应的光标样式"""
        return {
            'top-left': Qt.SizeFDiagCursor,
            'bottom-right': Qt.SizeFDiagCursor,
            'top-right': Qt.SizeBDiagCursor,
            'bottom-left': Qt.SizeBDiagCursor,
            'top': Qt.SizeVerCursor,
            'bottom': Qt.SizeVerCursor,
            'left': Qt.SizeHorCursor,
            'right': Qt.SizeHorCursor,
        }.get(handle, Qt.ArrowCursor)

    def mousePressEvent(self, event):
        if event.button() == Qt.LeftButton:
            # 优先检查是否点击在已选中图片的调整手柄上
            if self._selected_image_rect:
                handle = self._hit_resize_handle(event.pos())
                if handle:
                    self._resize_mode = handle
                    self._resize_start_pos = event.pos()
                    self._resize_start_rect = QRect(self._selected_image_rect)
                    self.viewport().setCursor(self._cursor_for_handle(handle))
                    event.accept()
                    return
                # 点击在选中图片内部 - 开始拖动
                if self._selected_image_rect.contains(event.pos()):
                    self._dragging_image = True
                    self._drag_start_pos = event.pos()
                    self.viewport().setCursor(Qt.ClosedHandCursor)
                    event.accept()
                    return
            # 检查是否点击在某张图片上
            img_name, doc_pos, rect = self._find_image_at(event.pos())
            if img_name:
                self._selected_image_name = img_name
                self._selected_image_pos = doc_pos
                self._selected_image_rect = rect
                self.viewport().update()
                # 显示 Office 风格工具栏
                self._show_office_toolbar()
                event.accept()
                return
            # 点击空白处清除选中
            self._clear_image_selection()
            self._hide_office_toolbar()
        super().mousePressEvent(event)

    def mouseMoveEvent(self, event):
        # 调整大小过程中
        if self._resize_mode and self._resize_start_pos and self._resize_start_rect:
            self._apply_resize(event.pos())
            event.accept()
            return
        # 拖动图片过程中
        if self._dragging_image and self._drag_start_pos:
            delta = event.pos() - self._drag_start_pos
            if delta.manhattanLength() > 4:
                self._move_image_by_drag(event.pos())
                self._drag_start_pos = event.pos()
            event.accept()
            return
        # 鼠标悬停在选中图片的手柄上时切换光标
        if self._selected_image_rect:
            handle = self._hit_resize_handle(event.pos())
            if handle:
                self.viewport().setCursor(self._cursor_for_handle(handle))
            else:
                if self._selected_image_rect.contains(event.pos()):
                    self.viewport().setCursor(Qt.OpenHandCursor)
                else:
                    self.viewport().setCursor(Qt.IBeamCursor)
        super().mouseMoveEvent(event)

    def mouseReleaseEvent(self, event):
        if self._resize_mode:
            # 提交尺寸修改到文档
            if self._selected_image_pos >= 0 and self._selected_image_rect:
                cur = QTextCursor(self.document())
                cur.setPosition(self._selected_image_pos)
                cur.movePosition(QTextCursor.Right, QTextCursor.KeepAnchor)
                fmt = cur.charFormat()
                if fmt.isImageFormat():
                    img_fmt = fmt.toImageFormat()
                    img_fmt.setWidth(self._selected_image_rect.width())
                    img_fmt.setHeight(self._selected_image_rect.height())
                    cur.setCharFormat(img_fmt)
            self._resize_mode = None
            self._resize_start_pos = None
            self._resize_start_rect = None
            # 重新计算选中矩形
            if self._selected_image_pos >= 0:
                self._selected_image_rect = self._image_view_rect(self._selected_image_pos)
            # 同步工具栏尺寸
            if self._office_toolbar and self._selected_image_rect:
                self._office_toolbar.set_image_size(
                    self._selected_image_rect.width(),
                    self._selected_image_rect.height()
                )
            self.viewport().update()
            event.accept()
            return
        if self._dragging_image:
            self._dragging_image = False
            self._drag_start_pos = None
            self.viewport().setCursor(Qt.OpenHandCursor)
            event.accept()
            return
        super().mouseReleaseEvent(event)
    
    def _move_image_by_drag(self, current_pos):
        """通过拖动移动图片到光标位置（剪切+插入到新位置）"""
        if self._selected_image_pos < 0 or not self._selected_image_name:
            return
        # 获取目标光标位置
        target_cursor = self.cursorForPosition(current_pos)
        target_pos = target_cursor.position()
        # 不能拖到自己位置
        if target_pos == self._selected_image_pos:
            return
        # 获取当前图片格式
        cur = QTextCursor(self.document())
        cur.setPosition(self._selected_image_pos)
        cur.movePosition(QTextCursor.Right, QTextCursor.KeepAnchor)
        fmt = cur.charFormat()
        if not fmt.isImageFormat():
            return
        img_fmt = fmt.toImageFormat()
        # 删除原图片
        cur.removeSelectedText()
        # 调整目标位置（如果在原位置之后，需要减1）
        if target_pos > self._selected_image_pos:
            target_pos -= 1
        # 插入图片到新位置
        new_cur = QTextCursor(self.document())
        new_cur.setPosition(target_pos)
        new_cur.insertImage(img_fmt)
        # 更新选中状态
        self._selected_image_pos = target_pos
        self._selected_image_rect = self._image_view_rect(target_pos)
        self.viewport().update()
    
    def _show_office_toolbar(self):
        """显示 Office 风格工具栏在图片上方"""
        if not self._selected_image_rect:
            return
        if self._office_toolbar is None:
            self._office_toolbar = _OfficeImageToolbar(self)
            self._office_toolbar.toolbar_action.connect(self._on_toolbar_action)
        # 同步当前图片尺寸
        self._office_toolbar.set_image_size(
            self._selected_image_rect.width(),
            self._selected_image_rect.height()
        )
        # 计算工具栏位置（图片上方）
        global_pos = self.viewport().mapToGlobal(
            QPoint(self._selected_image_rect.left(), self._selected_image_rect.top())
        )
        toolbar_y = global_pos.y() - self._office_toolbar.height() - 4
        if toolbar_y < 0:
            toolbar_y = global_pos.y() + self._selected_image_rect.height() + 4
        self._office_toolbar.move(global_pos.x(), toolbar_y)
        self._office_toolbar.show()
        self._office_toolbar.raise_()
    
    def _hide_office_toolbar(self):
        """隐藏 Office 风格工具栏"""
        if self._office_toolbar:
            self._office_toolbar.hide()
    
    def _on_toolbar_action(self, action):
        """处理工具栏操作"""
        if self._selected_image_pos < 0:
            return
        if action == 'resize':
            w, h = self._office_toolbar.get_image_size()
            cur = QTextCursor(self.document())
            cur.setPosition(self._selected_image_pos)
            cur.movePosition(QTextCursor.Right, QTextCursor.KeepAnchor)
            fmt = cur.charFormat()
            if fmt.isImageFormat():
                img_fmt = fmt.toImageFormat()
                img_fmt.setWidth(w)
                img_fmt.setHeight(h)
                cur.setCharFormat(img_fmt)
                self._selected_image_rect = self._image_view_rect(self._selected_image_pos)
                self.viewport().update()
        elif action == 'wrap':
            mode = self._office_toolbar.get_wrap_mode()
            self.set_image_wrap_mode(mode)

    def _apply_resize(self, current_pos):
        """根据鼠标拖动距离调整图片视觉矩形（不修改文档，释放时才提交）
        
        手柄行为：
        - 角点手柄：等比缩放
        - 左/右中间手柄：仅横向缩放
        - 上/下中间手柄：仅纵向缩放
        """
        if not self._resize_start_rect:
            return
        delta = current_pos - self._resize_start_pos
        sr = self._resize_start_rect
        new_w = sr.width()
        new_h = sr.height()
        ratio = sr.width() / max(1, sr.height())
        min_size = 20

        mode = self._resize_mode

        # 横向缩放
        if mode == 'right':
            new_w = max(min_size, sr.width() + delta.x())
        elif mode == 'left':
            new_w = max(min_size, sr.width() - delta.x())
        # 纵向缩放
        elif mode == 'bottom':
            new_h = max(min_size, sr.height() + delta.y())
        elif mode == 'top':
            new_h = max(min_size, sr.height() - delta.y())
        # 角点等比缩放
        elif mode in ('top-left', 'top-right', 'bottom-left', 'bottom-right'):
            if mode in ('bottom-right', 'top-right'):
                new_w = max(min_size, sr.width() + delta.x())
            if mode in ('bottom-left', 'top-left'):
                new_w = max(min_size, sr.width() - delta.x())
            if mode in ('bottom', 'bottom-left', 'bottom-right'):
                new_h = max(min_size, sr.height() + delta.y())
            if mode in ('top', 'top-left', 'top-right'):
                new_h = max(min_size, sr.height() - delta.y())
            # 等比缩放：以较大的变化量为准
            new_h = max(min_size, int(new_w / ratio))

        # 只更新视觉矩形，不修改文档
        self._selected_image_rect = QRect(sr.left(), sr.top(), new_w, new_h)
        self.viewport().update()

    def _delete_selected_image(self):
        """删除选中的图片"""
        if self._selected_image_pos < 0:
            return
        cur = QTextCursor(self.document())
        cur.setPosition(self._selected_image_pos)
        cur.movePosition(QTextCursor.Right, QTextCursor.KeepAnchor)
        cur.removeSelectedText()
        self._clear_image_selection()

    def keyPressEvent(self, event):
        # 选中图片后按 Delete/Backspace 删除
        if self._selected_image_name and event.key() in (Qt.Key_Delete, Qt.Key_Backspace):
            self._delete_selected_image()
            event.accept()
            return
        super().keyPressEvent(event)

    def wheelEvent(self, event):
        """滚轮事件 - 不处理图片缩放，让滚动条正常工作"""
        super().wheelEvent(event)

    def _compress_image_if_needed(self, image_path):
        """如果图像像素太大，自动压缩图像"""
        max_pixels = 4096 * 4096
        max_dimension = 4096
        try:
            from PIL import Image
            img = Image.open(image_path)
            width, height = img.size
            total_pixels = width * height
            if total_pixels > max_pixels or width > max_dimension or height > max_dimension:
                scale = min(max_dimension / width, max_dimension / height, 1.0)
                if scale < 1.0:
                    new_width = int(width * scale)
                    new_height = int(height * scale)
                    img = img.resize((new_width, new_height), Image.LANCZOS)
                    img.save(image_path, quality=90)
            return True
        except Exception:
            return False

    def paintEvent(self, event):
        super().paintEvent(event)
        # 调整大小过程中不刷新选中矩形（使用_apply_resize设置的视觉矩形）
        if not self._resize_mode:
            if self._selected_image_pos >= 0:
                new_rect = self._image_view_rect(self._selected_image_pos)
                if new_rect is not None:
                    self._selected_image_rect = new_rect
        if not self._selected_image_rect:
            return
        rect = self._selected_image_rect
        painter = QPainter(self.viewport())
        painter.setRenderHint(QPainter.Antialiasing)
        # 蓝色虚线边框
        pen = QPen(QColor(63, 123, 247), 1, Qt.DashLine)
        painter.setPen(pen)
        painter.setBrush(Qt.NoBrush)
        painter.drawRect(rect)
        # 8 个手柄（4角 + 4边中点）
        hs = self._resize_handle_size
        painter.setPen(QPen(QColor(63, 123, 247), 1))
        painter.setBrush(QColor(255, 255, 255))
        cx = rect.left() + rect.width() // 2
        cy = rect.top() + rect.height() // 2
        for x, y in [
            (rect.left(), rect.top()),       # top-left
            (cx, rect.top()),                 # top
            (rect.right(), rect.top()),       # top-right
            (rect.right(), cy),               # right
            (rect.right(), rect.bottom()),    # bottom-right
            (cx, rect.bottom()),              # bottom
            (rect.left(), rect.bottom()),     # bottom-left
            (rect.left(), cy),                # left
        ]:
            painter.drawRect(x - hs // 2, y - hs // 2, hs, hs)
        painter.end()
    
class _FormatBtn(QWidget):
    """富文本编辑器工具栏的图形化格式按钮"""
    clicked = pyqtSignal(bool)

    def __init__(self, fmt_id, fg_color='#c8c8cc', parent=None):
        super().__init__(parent)
        self.fmt_id = fmt_id
        self.fg_color = fg_color
        self.setFixedSize(34, 30)
        self._hover = False
        self.setCursor(Qt.PointingHandCursor)

    def paintEvent(self, event):
        painter = QPainter(self)
        painter.setRenderHint(QPainter.Antialiasing)
        r = self.rect()
        # 背景
        if self._hover:
            painter.setPen(Qt.NoPen)
            painter.setBrush(QColor(255, 255, 255, 25))
            painter.drawRoundedRect(r.adjusted(1, 1, -1, -1), 4, 4)
            painter.setBrush(Qt.NoBrush)
        fg = QColor(self.fg_color)
        painter.setPen(QPen(fg, 1.8))
        cx, cy = r.center().x(), r.center().y()
        fid = self.fmt_id
        if fid == 'bold':
            painter.setFont(QFont("SimSun", 14, QFont.Bold))
            painter.drawText(r, Qt.AlignCenter, 'B')
        elif fid == 'italic':
            painter.setFont(QFont("SimSun", 13, QFont.Bold, True))
            painter.drawText(r, Qt.AlignCenter, 'I')
        elif fid == 'underline':
            painter.setFont(QFont("SimSun", 13, QFont.Bold))
            painter.drawText(r.adjusted(0, 0, 0, -4), Qt.AlignCenter, 'U')
            painter.drawLine(cx - 6, cy + 6, cx + 6, cy + 6)
        elif fid == 'strikethrough':
            painter.setFont(QFont("SimSun", 13, QFont.Bold))
            painter.drawText(r.adjusted(0, 0, 0, -4), Qt.AlignCenter, 'S')
            painter.drawLine(cx - 6, cy + 2, cx + 6, cy + 2)
        elif fid == 'bullet_list':
            # 3行带圆点
            for i in range(3):
                y = cy - 6 + i * 6
                painter.setBrush(fg)
                painter.drawEllipse(cx - 7, y - 1, 3, 3)
                painter.setBrush(Qt.NoBrush)
                painter.drawLine(cx - 3, y, cx + 8, y)
        elif fid == 'ordered_list':
            # 3行带数字
            painter.setFont(QFont("SimSun", 6))
            for i in range(3):
                y = cy - 6 + i * 6
                painter.drawText(cx - 9, y - 3, 10, 10, Qt.AlignCenter, str(i + 1))
                painter.setPen(QPen(fg, 1.8))
                painter.drawLine(cx - 1, y, cx + 8, y)
        elif fid == 'todo_list':
            # 复选框图标
            painter.drawRect(cx - 7, cy - 6, 8, 8)
            # 对勾
            painter.drawLine(cx - 5, cy - 2, cx - 3, cy)
            painter.drawLine(cx - 3, cy, cx + 1, cy - 5)
            # 文字线
            painter.drawLine(cx + 3, cy - 4, cx + 8, cy - 4)
            painter.drawLine(cx + 3, cy + 1, cx + 8, cy + 1)
        elif fid == 'align_left':
            for i in range(3):
                y = cy - 6 + i * 6
                painter.drawLine(cx - 7, y, cx + 7, y)
        elif fid == 'align_center':
            for i in range(3):
                y = cy - 6 + i * 6
                painter.drawLine(cx - 4, y, cx + 4, y)
        elif fid == 'align_right':
            for i in range(3):
                y = cy - 6 + i * 6
                painter.drawLine(cx - 7, y, cx + 7, y)
        elif fid == 'text_color':
            painter.setFont(QFont("SimSun", 12, QFont.Bold))
            painter.drawText(r.adjusted(0, 0, 0, -6), Qt.AlignCenter, 'A')
            # 颜色条
            painter.setPen(Qt.NoPen)
            painter.setBrush(QColor(255, 80, 80))
            painter.drawRect(cx - 6, cy + 5, 12, 3)
            painter.setBrush(Qt.NoBrush)
        elif fid == 'bg_color':
            # 背景色块
            painter.setPen(Qt.NoPen)
            painter.setBrush(QColor(255, 220, 60, 150))
            painter.drawRect(cx - 7, cy - 5, 14, 12)
            painter.setBrush(Qt.NoBrush)
            painter.setPen(QPen(fg, 1.8))
            painter.setFont(QFont("SimSun", 10, QFont.Bold))
            painter.drawText(r.adjusted(0, 0, 0, -4), Qt.AlignCenter, 'A')
        elif fid == 'insert_image':
            # 图片图标 - 山和太阳
            painter.drawRect(cx - 8, cy - 6, 16, 12)
            # 太阳
            painter.setBrush(fg)
            painter.drawEllipse(cx + 2, cy - 4, 4, 4)
            painter.setBrush(Qt.NoBrush)
            # 山
            pts = [QPoint(cx - 6, cy + 4), QPoint(cx - 1, cy - 1), QPoint(cx + 4, cy + 4)]
            painter.drawPolygon(QPolygon(pts))
        elif fid == 'insert_link':
            # 链接图标
            painter.drawArc(cx - 8, cy - 4, 10, 8, 90 * 16, 180 * 16)
            painter.drawArc(cx - 2, cy - 4, 10, 8, 270 * 16, 180 * 16)
        elif fid == 'insert_table':
            # 表格网格
            painter.drawRect(cx - 7, cy - 6, 14, 12)
            painter.drawLine(cx - 7, cy - 1, cx + 7, cy - 1)
            painter.drawLine(cx - 7, cy + 4, cx + 7, cy + 4)
            painter.drawLine(cx, cy - 6, cx, cy + 6)
        elif fid == 'insert_hr':
            # 分隔线
            painter.drawLine(cx - 8, cy, cx + 8, cy)
        elif fid == 'insert_code':
            # 代码括号 </> 
            painter.setFont(QFont("Consolas", 9, QFont.Bold))
            painter.drawText(r, Qt.AlignCenter, '</>')
        elif fid == 'save':
            # 软盘图标
            painter.drawRect(cx - 6, cy - 7, 12, 14)
            painter.setBrush(QColor(40, 40, 45))
            painter.drawRect(cx - 3, cy - 7, 6, 5)
            painter.setBrush(Qt.NoBrush)
            painter.drawRect(cx - 4, cy + 1, 8, 6)
        elif fid == 'superscript':
            # 上标图标 x²
            painter.setFont(QFont("SimSun", 10))
            painter.drawText(cx - 4, cy - 5, 8, 8, Qt.AlignCenter, 'x')
            painter.setFont(QFont("SimSun", 6))
            painter.drawText(cx + 2, cy - 7, 6, 6, Qt.AlignCenter, '2')
        elif fid == 'subscript':
            # 下标图标 x₂
            painter.setFont(QFont("SimSun", 10))
            painter.drawText(cx - 4, cy - 2, 8, 8, Qt.AlignCenter, 'x')
            painter.setFont(QFont("SimSun", 6))
            painter.drawText(cx + 2, cy + 2, 6, 6, Qt.AlignCenter, '2')
        elif fid == 'code':
            # 行内代码图标 `<>`
            painter.setFont(QFont("Consolas", 8, QFont.Bold))
            painter.drawText(r, Qt.AlignCenter, '<>')
        elif fid == 'undo':
            # 撤销图标 ↩
            painter.drawArc(cx - 6, cy - 6, 12, 12, 30 * 16, 240 * 16)
            pts = [QPoint(cx - 6, cy + 2), QPoint(cx - 6, cy + 6), QPoint(cx - 1, cy + 3)]
            painter.drawPolyline(QPolygon(pts))
        elif fid == 'redo':
            # 重做图标 ↪
            painter.drawArc(cx - 6, cy - 6, 12, 12, 210 * 16, 240 * 16)
            pts = [QPoint(cx + 6, cy + 2), QPoint(cx + 6, cy + 6), QPoint(cx + 1, cy + 3)]
            painter.drawPolyline(QPolygon(pts))
        elif fid == 'image_select':
            # 选择工具 - 十字准星
            painter.drawLine(cx - 6, cy, cx + 6, cy)
            painter.drawLine(cx, cy - 6, cx, cy + 6)
        elif fid == 'image_move':
            # 移动工具 - 十字箭头
            painter.drawLine(cx - 6, cy, cx + 6, cy)
            painter.drawLine(cx, cy - 6, cx, cy + 6)
            painter.drawLine(cx - 4, cy - 4, cx - 1, cy - 1)
            painter.drawLine(cx + 4, cy - 4, cx + 1, cy - 1)
            painter.drawLine(cx - 4, cy + 4, cx - 1, cy + 1)
            painter.drawLine(cx + 4, cy + 4, cx + 1, cy + 1)
        elif fid == 'image_resize':
            # 缩放工具 - 对角箭头
            painter.drawLine(cx - 6, cy - 6, cx + 6, cy + 6)
            painter.drawLine(cx + 6, cy - 6, cx - 6, cy + 6)
            painter.drawLine(cx + 4, cy - 6, cx + 6, cy - 4)
            painter.drawLine(cx - 6, cy + 4, cx - 4, cy + 6)
        elif fid == 'image_rotate':
            # 旋转工具 - 弧形箭头
            painter.drawArc(cx - 6, cy - 6, 12, 12, 0, 270 * 16)
            painter.drawLine(cx - 6, cy, cx - 9, cy - 3)
            painter.drawLine(cx, cy - 6, cx + 3, cy - 9)
        elif fid == 'image_crop':
            # 裁剪工具 - 方形裁剪框
            painter.drawRect(cx - 6, cy - 6, 12, 12)
            painter.drawLine(cx - 2, cy - 6, cx - 2, cy - 2)
            painter.drawLine(cx + 2, cy - 6, cx + 2, cy - 2)
            painter.drawLine(cx - 6, cy - 2, cx - 2, cy - 2)
            painter.drawLine(cx + 2, cy - 2, cx + 6, cy - 2)
        elif fid == 'image_delete':
            # 删除工具 - 垃圾桶
            painter.drawRect(cx - 5, cy - 3, 10, 8)
            painter.drawRect(cx - 3, cy + 5, 6, 4)
            painter.drawLine(cx - 4, cy - 3, cx - 4, cy + 3)
            painter.drawLine(cx + 4, cy - 3, cx + 4, cy + 3)
        elif fid == 'image_caption':
            # 图片标题 - 图片下方加文字
            # 图片框
            painter.drawRect(cx - 6, cy - 5, 12, 8)
            # 下方文字线
            painter.drawLine(cx - 5, cy + 4, cx + 5, cy + 4)
            painter.drawLine(cx - 4, cy + 7, cx + 4, cy + 7)
        elif fid == 'arrow':
            # 箭头工具
            painter.setPen(QPen(fg, 1.5))
            painter.drawLine(cx - 6, cy + 6, cx + 4, cy - 4)
            pts = [QPoint(cx + 4, cy - 4), QPoint(cx + 1, cy - 1), QPoint(cx + 2, cy + 1)]
            painter.setBrush(fg)
            painter.drawPolygon(QPolygon(pts))
            painter.setBrush(Qt.NoBrush)
        elif fid == 'rect':
            # 矩形标注工具
            painter.setPen(QPen(fg, 1.5))
            painter.drawRect(cx - 7, cy - 5, 14, 10)
        elif fid == 'circle':
            # 圆形标注工具
            painter.setPen(QPen(fg, 1.5))
            painter.drawEllipse(cx - 6, cy - 6, 12, 12)
        elif fid == 'line':
            # 直线工具
            painter.setPen(QPen(fg, 1.5))
            painter.drawLine(cx - 7, cy + 5, cx + 7, cy - 5)
        elif fid == 'text':
            # 文字标注工具
            painter.setFont(QFont("SimSun", 12, QFont.Bold))
            painter.setPen(QPen(fg, 1))
            painter.drawText(r.adjusted(0, 0, 0, -4), Qt.AlignCenter, 'T')
            painter.setPen(QPen(fg, 1))
            painter.drawLine(cx, cy + 3, cx, cy + 7)
        elif fid == 'insert_textbox':
            # 文本框工具 - 方框内带T字
            painter.setPen(QPen(fg, 1.5))
            painter.drawRect(cx - 7, cy - 6, 14, 12)
            painter.setFont(QFont("SimSun", 8, QFont.Bold))
            painter.setPen(QPen(fg, 1))
            painter.drawText(r, Qt.AlignCenter, 'T')
        painter.end()

    def enterEvent(self, event):
        self._hover = True
        self.update()

    def leaveEvent(self, event):
        self._hover = False
        self.update()

    def mousePressEvent(self, event):
        if event.button() == Qt.LeftButton:
            self.clicked.emit(True)


class _InlineFormatBar(QWidget):
    """浮动格式工具栏 - 选中文本时出现在选区附近（Trilium Notes 风格）"""
    format_action = pyqtSignal(str)  # 格式动作信号

    def __init__(self, editor, parent=None):
        super().__init__(parent)
        self.editor = editor
        self.setWindowFlags(Qt.ToolTip | Qt.FramelessWindowHint)
        self.setAttribute(Qt.WA_TranslucentBackground, False)
        self.setAttribute(Qt.WA_ShowWithoutActivating)
        self.setFocusPolicy(Qt.NoFocus)
        self._hover = False

        # 淡入动画
        self._opacity_effect = QGraphicsOpacityEffect(self)
        self.setGraphicsEffect(self._opacity_effect)
        self._opacity_effect.setOpacity(0)
        self._show_anim = QPropertyAnimation(self._opacity_effect, b"opacity")
        self._show_anim.setDuration(150)
        self._show_anim.setStartValue(0.0)
        self._show_anim.setEndValue(1.0)
        self._show_anim.setEasingCurve(QEasingCurve.OutCubic)

        layout = QHBoxLayout(self)
        layout.setContentsMargins(6, 4, 6, 4)
        layout.setSpacing(1)

        # 标题级别下拉
        self._heading_combo = QComboBox()
        self._heading_combo.addItems(['正文', 'H1', 'H2', 'H3', 'H4'])
        self._heading_combo.setFixedWidth(58)
        self._heading_combo.setFixedHeight(24)
        self._heading_combo.setFocusPolicy(Qt.NoFocus)
        self._heading_combo.setStyleSheet("""
            QComboBox { background: rgba(255,255,255,15); color: #c8c8cc; border: 1px solid #555;
                        border-radius: 3px; font-size: 11px; padding: 0 4px; }
            QComboBox::drop-down { border: none; width: 14px; }
            QComboBox QAbstractItemView { background: #2a2a30; color: #c8c8cc; selection-background-color: #3f7bf7; }
        """)
        self._heading_combo.currentIndexChanged.connect(lambda idx: self.format_action.emit(f'heading_{idx}'))
        layout.addWidget(self._heading_combo)

        self._add_sep(layout)

        # 文本样式
        for fmt_id, tip in [('bold', '加粗'), ('italic', '斜体'), ('underline', '下划线'), ('strikethrough', '删除线')]:
            btn = _FormatBtn(fmt_id, fg_color='#d0d0d4')
            btn.setToolTip(tip)
            btn.clicked.connect(lambda checked, fid=fmt_id: self.format_action.emit(fid))
            layout.addWidget(btn)

        self._add_sep(layout)

        # 列表
        for fmt_id, tip in [('bullet_list', '无序列表'), ('ordered_list', '有序列表')]:
            btn = _FormatBtn(fmt_id, fg_color='#d0d0d4')
            btn.setToolTip(tip)
            btn.clicked.connect(lambda checked, fid=fmt_id: self.format_action.emit(fid))
            layout.addWidget(btn)

        self._add_sep(layout)

        # 对齐
        for fmt_id, tip in [('align_left', '左对齐'), ('align_center', '居中'), ('align_right', '右对齐')]:
            btn = _FormatBtn(fmt_id, fg_color='#d0d0d4')
            btn.setToolTip(tip)
            btn.clicked.connect(lambda checked, fid=fmt_id: self.format_action.emit(fid))
            layout.addWidget(btn)

        self._add_sep(layout)

        # 颜色
        btn_tc = _FormatBtn('text_color', fg_color='#d0d0d4')
        btn_tc.setToolTip('文字颜色')
        btn_tc.clicked.connect(lambda checked: self.format_action.emit('text_color'))
        layout.addWidget(btn_tc)

        btn_bg = _FormatBtn('bg_color', fg_color='#d0d0d4')
        btn_bg.setToolTip('背景颜色')
        btn_bg.clicked.connect(lambda checked: self.format_action.emit('bg_color'))
        layout.addWidget(btn_bg)

        self._add_sep(layout)

        # 上下标
        for fmt_id, tip in [('superscript', '上标'), ('subscript', '下标')]:
            btn = _FormatBtn(fmt_id, fg_color='#d0d0d4')
            btn.setToolTip(tip)
            btn.clicked.connect(lambda checked, fid=fmt_id: self.format_action.emit(fid))
            layout.addWidget(btn)

        self._add_sep(layout)

        # 代码格式化
        btn_code = _FormatBtn('code', fg_color='#d0d0d4')
        btn_code.setToolTip('行内代码')
        btn_code.clicked.connect(lambda checked: self.format_action.emit('code'))
        layout.addWidget(btn_code)

        self._add_sep(layout)

        # 链接
        btn_link = _FormatBtn('insert_link', fg_color='#d0d0d4')
        btn_link.setToolTip('插入链接')
        btn_link.clicked.connect(lambda checked: self.format_action.emit('insert_link'))
        layout.addWidget(btn_link)

        self._add_sep(layout)

        # 撤销/重做
        for fmt_id, tip in [('undo', '撤销'), ('redo', '重做')]:
            btn = _FormatBtn(fmt_id, fg_color='#d0d0d4')
            btn.setToolTip(tip)
            btn.clicked.connect(lambda checked, fid=fmt_id: self.format_action.emit(fid))
            layout.addWidget(btn)

        self.setFixedHeight(34)

    def _add_sep(self, layout):
        sep = QFrame()
        sep.setFrameShape(QFrame.VLine)
        sep.setFixedWidth(1)
        sep.setStyleSheet("background-color: rgba(255,255,255,30);")
        layout.addWidget(sep)

    def paintEvent(self, event):
        painter = QPainter(self)
        painter.setRenderHint(QPainter.Antialiasing)
        painter.setPen(Qt.NoPen)
        painter.setBrush(QColor(35, 35, 40, 230))
        painter.drawRoundedRect(self.rect().adjusted(1, 1, -1, -1), 8, 8)
        # 边框
        painter.setPen(QPen(QColor(255, 255, 255, 30), 1))
        painter.setBrush(Qt.NoBrush)
        painter.drawRoundedRect(self.rect().adjusted(1, 1, -1, -1), 8, 8)
        painter.end()

    def show_at_selection(self):
        """在选区上方显示工具栏（带淡入动画）"""
        cursor = self.editor.textCursor()
        if not cursor.hasSelection():
            self.hide_bar()
            return
        # 获取选区的屏幕坐标
        start_rect = self.editor.cursorRect(cursor.selectionStart())
        end_rect = self.editor.cursorRect(cursor.selectionEnd())
        mid_x = (start_rect.topLeft().x() + end_rect.topRight().x()) // 2
        top_y = min(start_rect.topLeft().y(), end_rect.topLeft().y())
        # 转换为全局坐标
        global_pos = self.editor.viewport().mapToGlobal(QPoint(mid_x, top_y))
        bar_w = self.width()
        bar_h = self.height()
        # 显示在选区上方
        pos_x = global_pos.x() - bar_w // 2
        pos_y = global_pos.y() - bar_h - 8
        # 确保不超出屏幕
        screen = QApplication.desktop().availableGeometry(self)
        if pos_x < screen.left():
            pos_x = screen.left()
        if pos_x + bar_w > screen.right():
            pos_x = screen.right() - bar_w
        if pos_y < screen.top():
            pos_y = global_pos.y() + 20  # 改为显示在选区下方
        self.move(pos_x, pos_y)
        self.show()
        self.raise_()
        self._show_anim.start()

    def hide_bar(self):
        """隐藏工具栏"""
        self._show_anim.stop()
        self._opacity_effect.setOpacity(0)
        self.hide()


class _BlockToolbar(QWidget):
    """左侧块工具栏 - "+"按钮打开插入菜单（Trilium Notes 风格）"""
    action_triggered = pyqtSignal(str)

    def __init__(self, editor, parent=None):
        super().__init__(parent)
        self.editor = editor
        self.setWindowFlags(Qt.ToolTip | Qt.FramelessWindowHint)
        self.setAttribute(Qt.WA_TranslucentBackground, False)
        self.setAttribute(Qt.WA_ShowWithoutActivating)
        self.setFocusPolicy(Qt.NoFocus)
        self._menu_open = False

        layout = QVBoxLayout(self)
        layout.setContentsMargins(2, 2, 2, 2)
        layout.setSpacing(0)

        self._plus_btn = QWidget()
        self._plus_btn.setFixedSize(28, 28)
        self._plus_btn.setCursor(Qt.PointingHandCursor)
        self._plus_btn.setToolTip('插入内容块')
        self._plus_btn._hover = False
        layout.addWidget(self._plus_btn)

        self.setFixedSize(32, 32)

    def paintEvent(self, event):
        painter = QPainter(self)
        painter.setRenderHint(QPainter.Antialiasing)
        r = self.rect()
        # 背景
        painter.setPen(Qt.NoPen)
        painter.setBrush(QColor(35, 35, 40, 200))
        painter.drawRoundedRect(r.adjusted(1, 1, -1, -1), 6, 6)
        # "+"号
        painter.setPen(QPen(QColor(180, 185, 190), 2))
        cx, cy = r.center().x(), r.center().y()
        painter.drawLine(cx - 6, cy, cx + 6, cy)
        painter.drawLine(cx, cy - 6, cx, cy + 6)
        painter.end()

    def mousePressEvent(self, event):
        if event.button() == Qt.LeftButton:
            self._show_insert_menu()

    def _show_insert_menu(self):
        menu = QMenu(self)
        menu.setWindowFlags(menu.windowFlags() | Qt.NoFocus)
        menu.setStyleSheet("""
            QMenu { background-color: #2a2a30; color: #c8c8cc; border: 1px solid #444; padding: 4px; }
            QMenu::item { padding: 6px 20px; border-radius: 4px; }
            QMenu::item:selected { background-color: #3f7bf7; }
        """)
        # 图片工具
        image_menu = menu.addMenu('🖼️ 图片')
        image_actions = [
            ('insert_image', '插入图片'),
            ('insert_image_from_db', '从数据库选择'),
            ('image_left', '图片左对齐'),
            ('image_center', '图片居中'),
            ('image_right', '图片右对齐'),
            ('image_full', '图片全屏宽度'),
            ('image_resize_small', '缩小图片'),
            ('image_resize_large', '放大图片'),
        ]
        for action_id, text in image_actions:
            act = image_menu.addAction(text)
            act.setData(action_id)
        # 文字工具
        text_menu = menu.addMenu('📝 文字')
        text_actions = [
            ('insert_table', '插入表格'),
            ('insert_code', '插入代码块'),
            ('insert_hr', '插入分隔线'),
            ('insert_link', '插入链接'),
            ('insert_todo', '待办事项'),
            ('insert_note', '嵌入笔记'),
        ]
        for action_id, text in text_actions:
            act = text_menu.addAction(text)
            act.setData(action_id)
        # 在按钮下方显示菜单
        global_pos = self.mapToGlobal(QPoint(0, self.height()))
        action = menu.exec_(global_pos)
        if action and action.data():
            self.action_triggered.emit(action.data())

    def update_position(self):
        """更新位置到编辑器当前光标行的左侧边距（Trilium Notes 风格）"""
        if not self.editor.isVisible():
            self.hide()
            return
        cursor = self.editor.textCursor()
        # 获取当前光标行的矩形
        cursor_rect = self.editor.cursorRect(cursor)
        # 行的顶部y坐标，x为编辑器视口左侧
        line_top = cursor_rect.top()
        # 转换为全局坐标：x 在编辑器左侧边距外，y 在当前行顶部居中
        global_pos = self.editor.viewport().mapToGlobal(QPoint(0, line_top))
        # 放置在编辑器左侧边距位置
        self.move(global_pos.x() - self.width() - 4, global_pos.y() + (cursor_rect.height() - self.height()) // 2)
        self.show()
        self.raise_()


class _ImageEditPopup(QWidget):
    """图片编辑弹出工具栏 - 点击编辑器中的图片时出现"""
    edit_image = pyqtSignal(str)  # 打开完整编辑器
    annotate_image = pyqtSignal(str, str)  # 快速标注 (图片路径, 工具)

    def __init__(self, parent=None):
        super().__init__(parent)
        self.setWindowFlags(Qt.ToolTip | Qt.FramelessWindowHint)
        self.setAttribute(Qt.WA_TranslucentBackground, False)
        self.setAttribute(Qt.WA_ShowWithoutActivating)
        self.setFocusPolicy(Qt.NoFocus)
        self._image_path = None

        main_layout = QVBoxLayout(self)
        main_layout.setContentsMargins(6, 4, 6, 4)
        main_layout.setSpacing(2)

        # 工具行
        tool_layout = QHBoxLayout()
        tool_layout.setSpacing(1)
        for tid in ['select', 'arrow', 'line', 'rect', 'circle', 'text', 'measure', 'eraser']:
            btn = _ToolButton(tid)
            btn.setFixedSize(28, 28)
            btn.clicked.connect(lambda checked, t=tid: self._on_tool(t))
            tool_layout.addWidget(btn)
        main_layout.addLayout(tool_layout)

        # 样式行
        self._style_bar = _StyleBar(self)
        self._style_bar.setFixedHeight(30)
        main_layout.addWidget(self._style_bar)

        # 打开编辑器按钮
        btn_open = QPushButton('打开编辑器')
        btn_open.setFixedHeight(22)
        btn_open.setFocusPolicy(Qt.NoFocus)
        btn_open.setStyleSheet("""
            QPushButton { background: rgba(63,123,247,180); color: white; border: none;
                          border-radius: 3px; font-size: 11px; padding: 0 8px; }
            QPushButton:hover { background: rgba(63,123,247,220); }
        """)
        btn_open.clicked.connect(self._on_open_editor)
        main_layout.addWidget(btn_open)

        self.setFixedWidth(280)

    def paintEvent(self, event):
        painter = QPainter(self)
        painter.setRenderHint(QPainter.Antialiasing)
        painter.setPen(Qt.NoPen)
        painter.setBrush(QColor(35, 35, 40, 230))
        painter.drawRoundedRect(self.rect().adjusted(1, 1, -1, -1), 8, 8)
        painter.setPen(QPen(QColor(255, 255, 255, 30), 1))
        painter.setBrush(Qt.NoBrush)
        painter.drawRoundedRect(self.rect().adjusted(1, 1, -1, -1), 8, 8)
        painter.end()

    def show_at_image(self, image_path, global_pos):
        """在图片位置显示"""
        self._image_path = image_path
        pos_x = global_pos.x() + 10
        pos_y = global_pos.y() + 10
        # 确保不超出屏幕
        screen = QApplication.desktop().availableGeometry(self)
        if pos_x + self.width() > screen.right():
            pos_x = global_pos.x() - self.width() - 10
        if pos_y + self.height() > screen.bottom():
            pos_y = global_pos.y() - self.height() - 10
        self.move(pos_x, pos_y)
        self.show()
        self.raise_()

    def _on_tool(self, tool_id):
        if self._image_path:
            self.annotate_image.emit(self._image_path, tool_id)

    def _on_open_editor(self, checked=False):
        if self._image_path:
            self.edit_image.emit(self._image_path)

    def get_style(self):
        return {
            'color': self._style_bar.get_color(),
            'width': self._style_bar.get_line_width(),
            'style': self._style_bar.get_line_style(),
        }


class _RichTextEditorDialog(QDialog):
    """图文编辑对话框：参考Trilium Notes设计，即时保存+完整工具栏+笔记列表"""
    def __init__(self, parent=None, record_type='medical', disease_id=None, active_db=None, user_password=None):
        super().__init__(parent)
        self.setWindowTitle('图文编辑器')
        self.setMinimumSize(1000, 650)
        self.setWindowFlags(Qt.Window | Qt.FramelessWindowHint)
        self.record_type = record_type
        self.disease_id = disease_id
        self.active_db = active_db or DATABASE
        self.user_password = user_password
        self._current_record_id = None
        self._save_timer = QTimer(self)
        self._save_timer.setSingleShot(True)
        self._save_timer.timeout.connect(self._auto_save)

        layout = QVBoxLayout(self)
        layout.setContentsMargins(0, 0, 0, 0)
        layout.setSpacing(0)

        # 自定义标题栏
        self.title_bar = _CustomTitleBar(self, '图文编辑器')
        layout.addWidget(self.title_bar)

        # 主内容区：左右分栏
        main_splitter = QSplitter(Qt.Horizontal)

        # ── 左侧：笔记列表 ──
        left_panel = QWidget()
        left_layout = QVBoxLayout(left_panel)
        left_layout.setContentsMargins(0, 0, 0, 0)
        left_layout.setSpacing(0)

        # 笔记列表头部
        list_header = QWidget()
        list_header.setFixedHeight(36)
        list_header_layout = QHBoxLayout(list_header)
        list_header_layout.setContentsMargins(8, 4, 8, 4)
        list_label = QLabel('笔记列表')
        list_label.setFont(QFont("Microsoft YaHei", 10, QFont.Bold))
        list_header_layout.addWidget(list_label)
        btn_new = QPushButton('+ 新建')
        btn_new.setFixedHeight(26)
        btn_new.clicked.connect(self._new_record)
        list_header_layout.addWidget(btn_new)
        left_layout.addWidget(list_header)

        # 笔记列表
        self.note_list = QListWidget()
        self.note_list.setContextMenuPolicy(Qt.CustomContextMenu)
        self.note_list.customContextMenuRequested.connect(self._show_note_context_menu)
        self.note_list.currentRowChanged.connect(self._on_note_selected)
        # 设置样式支持层级缩进
        self.note_list.setStyleSheet("""
            QListWidget { border: none; background: #1e1e22; outline: none; }
            QListWidget::item { padding: 4px 8px; border-bottom: 1px solid rgba(255,255,255,8);
                                color: #a0a0a8; font-size: 12px; }
            QListWidget::item:selected { background-color: rgba(63,123,247,30); color: #e0e0e4; }
            QListWidget::item:hover { background-color: rgba(255,255,255,8); color: #c8c8cc; }
        """)
        left_layout.addWidget(self.note_list)

        # 删除按钮
        btn_del = QPushButton('删除选中笔记')
        btn_del.setObjectName('dangerBtn')
        btn_del.setFixedHeight(30)
        btn_del.clicked.connect(self._delete_current_record)
        left_layout.addWidget(btn_del)

        left_panel.setFixedWidth(200)
        main_splitter.addWidget(left_panel)

        # ── 右侧：编辑区 ──
        right_panel = QWidget()
        right_layout = QVBoxLayout(right_panel)
        right_layout.setContentsMargins(0, 0, 0, 0)
        right_layout.setSpacing(0)

        # 标题输入区
        title_bar = QWidget()
        title_bar.setFixedHeight(50)
        title_layout = QHBoxLayout(title_bar)
        title_layout.setContentsMargins(12, 8, 12, 8)
        self.title_input = QLineEdit()
        self.title_input.setPlaceholderText('输入标题...')
        self.title_input.setFont(QFont("Microsoft YaHei", 14, QFont.Bold))
        self.title_input.textChanged.connect(self._on_content_changed)
        title_layout.addWidget(self.title_input)
        # 保存按钮
        btn_save = QPushButton('保存')
        btn_save.setFixedHeight(28)
        btn_save.setFixedWidth(60)
        btn_save.setObjectName('accentBtn')
        btn_save.setToolTip('保存 (Ctrl+S)')
        btn_save.clicked.connect(self._manual_save)
        title_layout.addWidget(btn_save)
        # 保存状态指示
        self.save_indicator = QLabel('已保存')
        self.save_indicator.setStyleSheet("color: #4ade80; font-size: 11px; background: transparent;")
        title_layout.addWidget(self.save_indicator)
        right_layout.addWidget(title_bar)

        # 格式化工具栏
        self._build_toolbar(right_layout)

        # 富文本编辑器
        self.text_edit = QTextEdit()
        self.text_edit.setAcceptRichText(True)
        self.text_edit.textChanged.connect(self._on_content_changed)
        # 启用右键菜单（包含粘贴选项）
        self.text_edit.setContextMenuPolicy(Qt.CustomContextMenu)
        self.text_edit.customContextMenuRequested.connect(self._show_edit_context_menu)
        right_layout.addWidget(self.text_edit, stretch=1)

        # 标签输入区
        tag_bar = QWidget()
        tag_bar.setFixedHeight(36)
        tag_layout = QHBoxLayout(tag_bar)
        tag_layout.setContentsMargins(12, 4, 12, 4)
        tag_label = QLabel('标签：')
        tag_layout.addWidget(tag_label)
        self.tag_input = QLineEdit()
        self.tag_input.setPlaceholderText('输入标签，用逗号分隔...')
        self.tag_input.textChanged.connect(self._on_content_changed)
        tag_layout.addWidget(self.tag_input)
        right_layout.addWidget(tag_bar)

        main_splitter.addWidget(right_panel)
        main_splitter.setStretchFactor(0, 0)
        main_splitter.setStretchFactor(1, 1)

        layout.addWidget(main_splitter)

        # 加载笔记列表
        self._load_note_list()

    def _build_toolbar(self, parent_layout):
        """构建格式化工具栏，参考Trilium Notes的工具栏设计"""
        toolbar_widget = QWidget()
        toolbar_widget.setFixedHeight(36)
        tb = QHBoxLayout(toolbar_widget)
        tb.setContentsMargins(8, 2, 8, 2)
        tb.setSpacing(2)

        # 标题级别
        heading_combo = QComboBox()
        heading_combo.addItems(['正文', '标题1', '标题2', '标题3', '标题4'])
        heading_combo.setFixedWidth(90)
        heading_combo.setFixedHeight(28)
        heading_combo.currentIndexChanged.connect(self._set_heading)
        tb.addWidget(heading_combo)
        self._heading_combo = heading_combo

        tb.addSpacing(6)
        # 分隔线
        sep1 = QFrame()
        sep1.setFrameShape(QFrame.VLine)
        sep1.setStyleSheet("color: #444;")
        tb.addWidget(sep1)
        tb.addSpacing(6)

        # 文本样式按钮
        style_buttons = [
            ('bold', '加粗', self._toggle_bold),
            ('italic', '斜体', self._toggle_italic),
            ('underline', '下划线', self._toggle_underline),
            ('strikethrough', '删除线', self._toggle_strikethrough),
        ]
        for fmt_id, tip, callback in style_buttons:
            btn = _FormatBtn(fmt_id)
            btn.setToolTip(tip)
            btn.clicked.connect(callback)
            tb.addWidget(btn)

        tb.addSpacing(6)
        sep2 = QFrame()
        sep2.setFrameShape(QFrame.VLine)
        sep2.setStyleSheet("color: #444;")
        tb.addWidget(sep2)
        tb.addSpacing(6)

        # 列表按钮
        list_buttons = [
            ('bullet_list', '无序列表', self._toggle_bullet_list),
            ('ordered_list', '有序列表', self._toggle_ordered_list),
            ('todo_list', '待办列表', self._toggle_todo_list),
        ]
        for fmt_id, tip, callback in list_buttons:
            btn = _FormatBtn(fmt_id)
            btn.setToolTip(tip)
            btn.clicked.connect(callback)
            tb.addWidget(btn)

        tb.addSpacing(6)
        sep3 = QFrame()
        sep3.setFrameShape(QFrame.VLine)
        sep3.setStyleSheet("color: #444;")
        tb.addWidget(sep3)
        tb.addSpacing(6)

        # 对齐按钮
        align_buttons = [
            ('align_left', '左对齐', self._align_left),
            ('align_center', '居中', self._align_center),
            ('align_right', '右对齐', self._align_right),
        ]
        for fmt_id, tip, callback in align_buttons:
            btn = _FormatBtn(fmt_id)
            btn.setToolTip(tip)
            btn.clicked.connect(callback)
            tb.addWidget(btn)

        tb.addSpacing(6)
        sep4 = QFrame()
        sep4.setFrameShape(QFrame.VLine)
        sep4.setStyleSheet("color: #444;")
        tb.addWidget(sep4)
        tb.addSpacing(6)

        # 颜色按钮
        self._color_btn = _FormatBtn('text_color')
        self._color_btn.setToolTip('文字颜色')
        self._color_btn.clicked.connect(self._set_text_color)
        tb.addWidget(self._color_btn)

        self._bg_color_btn = _FormatBtn('bg_color')
        self._bg_color_btn.setToolTip('背景颜色')
        self._bg_color_btn.clicked.connect(self._set_bg_color)
        tb.addWidget(self._bg_color_btn)

        tb.addSpacing(6)
        sep5 = QFrame()
        sep5.setFrameShape(QFrame.VLine)
        sep5.setStyleSheet("color: #444;")
        tb.addWidget(sep5)
        tb.addSpacing(6)

        # 插入按钮
        insert_buttons = [
            ('insert_image', '插入图片', self._insert_image),
            ('image_caption', '图片标题', self._add_image_caption),
            ('insert_link', '插入链接', self._insert_link),
            ('insert_table', '插入表格', self._insert_table),
            ('insert_hr', '插入分隔线', self._insert_hr),
        ]
        for fmt_id, tip, callback in insert_buttons:
            btn = _FormatBtn(fmt_id)
            btn.setToolTip(tip)
            btn.clicked.connect(callback)
            tb.addWidget(btn)

        tb.addStretch()

        # 代码块
        btn_code = _FormatBtn('insert_code')
        btn_code.setToolTip('代码块')
        btn_code.clicked.connect(self._insert_code_block)
        tb.addWidget(btn_code)

        # 保存按钮
        btn_save = _FormatBtn('save')
        btn_save.setToolTip('保存 (Ctrl+S)')
        btn_save.clicked.connect(self._manual_save)
        tb.addWidget(btn_save)

        parent_layout.addWidget(toolbar_widget)

    # ── 格式化操作 ──
    def _set_heading(self, idx):
        sizes = [10, 20, 16, 14, 12]
        weights = [QFont.Normal, QFont.Bold, QFont.Bold, QFont.Bold, QFont.Bold]
        fmt = self.text_edit.currentCharFormat()
        fmt.setFontPointSize(sizes[idx])
        fmt.setFontWeight(weights[idx])
        self.text_edit.setCurrentCharFormat(fmt)

    def _toggle_bold(self, checked=False):
        self.text_edit.setFontWeight(
            QFont.Normal if self.text_edit.fontWeight() == QFont.Bold else QFont.Bold
        )

    def _toggle_italic(self, checked=False):
        self.text_edit.setFontItalic(not self.text_edit.fontItalic())

    def _toggle_underline(self, checked=False):
        self.text_edit.setFontUnderline(not self.text_edit.fontUnderline())

    def _toggle_strikethrough(self, checked=False):
        fmt = self.text_edit.currentCharFormat()
        fmt.setFontStrikeOut(not fmt.fontStrikeOut())
        self.text_edit.setCurrentCharFormat(fmt)

    def _toggle_bullet_list(self, checked=False):
        self.text_edit.insertList(QTextListFormat.ListDisc)

    def _toggle_ordered_list(self, checked=False):
        self.text_edit.insertList(QTextListFormat.ListDecimal)

    def _toggle_todo_list(self, checked=False):
        cursor = self.text_edit.textCursor()
        cursor.insertList(QTextListFormat.ListDisc)
        cursor.insertText('☐ ')

    def _align_left(self, checked=False):
        self.text_edit.setAlignment(Qt.AlignLeft)

    def _align_center(self, checked=False):
        self.text_edit.setAlignment(Qt.AlignCenter)

    def _align_right(self, checked=False):
        self.text_edit.setAlignment(Qt.AlignRight)

    def _show_edit_context_menu(self, pos):
        """显示编辑器右键菜单（支持粘贴图片和表格）"""
        from PyQt5.QtWidgets import QMenu
        from PyQt5.QtGui import QClipboard
        
        menu = QMenu(self)
        menu.setStyleSheet("QMenu { background-color: #2a2a30; color: #c8c8cc; border: 1px solid #444; }"
                          "QMenu::item:selected { background-color: #3f7bf7; }")
        
        act_undo = menu.addAction('撤销')
        act_redo = menu.addAction('重做')
        menu.addSeparator()
        act_cut = menu.addAction('剪切')
        act_copy = menu.addAction('复制')
        act_paste = menu.addAction('粘贴')
        act_paste_image = menu.addAction('粘贴图片')
        act_paste_table = menu.addAction('粘贴表格')
        menu.addSeparator()
        act_select_all = menu.addAction('全选')

        action = menu.exec_(self.text_edit.viewport().mapToGlobal(pos))
        
        clipboard = QApplication.clipboard()
        
        if action == act_undo:
            self.text_edit.undo()
        elif action == act_redo:
            self.text_edit.redo()
        elif action == act_cut:
            self.text_edit.cut()
        elif action == act_copy:
            self.text_edit.copy()
        elif action == act_paste:
            self.text_edit.paste()
        elif action == act_paste_image:
            self._paste_image_from_clipboard()
        elif action == act_paste_table:
            self._paste_table_from_clipboard()
        elif action == act_select_all:
            self.text_edit.selectAll()

    def _paste_image_from_clipboard(self):
        """从剪贴板粘贴图片"""
        clipboard = QApplication.clipboard()
        mime_data = clipboard.mimeData()
        
        if mime_data.hasImage():
            pixmap = clipboard.pixmap()
            if not pixmap.isNull():
                # 保存图片到临时文件
                import uuid
                img_dir = os.path.join(APP_DIR, 'images')
                os.makedirs(img_dir, exist_ok=True)
                filename = str(uuid.uuid4())[:8] + '.png'
                dest = os.path.join(img_dir, filename)
                pixmap.save(dest)
                
                # 压缩图片
                self._compress_image_if_needed(dest)
                
                # 插入到编辑器
                self.text_edit.insertHtml(f'<img src="{dest}" style="max-width:100%; margin:8px 0;">')
                
    def _paste_table_from_clipboard(self):
        """从剪贴板粘贴表格（从Excel等复制）"""
        clipboard = QApplication.clipboard()
        mime_data = clipboard.mimeData()
        
        if mime_data.hasText():
            text = mime_data.text()
            # 尝试解析表格数据（Tab分隔行，换行分隔列）
            rows = text.strip().split('\n')
            if len(rows) > 0:
                # 构建HTML表格
                html = '<table border="1" style="border-collapse:collapse; margin:8px 0;">'
                for row in rows:
                    cells = row.split('\t')
                    html += '<tr>'
                    for cell in cells:
                        html += f'<td style="padding:4px 8px; border:1px solid #444;">{cell}</td>'
                    html += '</tr>'
                html += '</table>'
                self.text_edit.insertHtml(html)
    
    def _compress_image_if_needed(self, image_path):
        """如果图像像素太大，自动压缩图像"""
        max_pixels = 4096 * 4096
        max_dimension = 4096
        
        try:
            from PIL import Image
            img = Image.open(image_path)
            width, height = img.size
            
            total_pixels = width * height
            if total_pixels > max_pixels or width > max_dimension or height > max_dimension:
                scale = min(max_dimension / width, max_dimension / height, 1.0)
                if scale < 1.0:
                    new_width = int(width * scale)
                    new_height = int(height * scale)
                    img = img.resize((new_width, new_height), Image.LANCZOS)
                    img.save(image_path, quality=90)
            return True
        except Exception as e:
            return False

    def _set_text_color(self, checked=False):
        from PyQt5.QtWidgets import QColorDialog
        color = QColorDialog.getColor(Qt.black, self, '选择文字颜色')
        if color.isValid():
            self.text_edit.setTextColor(color)

    def _set_bg_color(self, checked=False):
        from PyQt5.QtWidgets import QColorDialog
        color = QColorDialog.getColor(Qt.white, self, '选择背景颜色')
        if color.isValid():
            self.text_edit.setTextBackgroundColor(color)

    def _insert_image(self, checked=False):
        file_paths, _ = QFileDialog.getOpenFileNames(
            self, '选择图片', '', '图片文件 (*.png *.jpg *.jpeg *.bmp *.gif)'
        )
        for fp in file_paths:
            img_dir = os.path.join(APP_DIR, 'images')
            os.makedirs(img_dir, exist_ok=True)
            import uuid
            ext = os.path.splitext(fp)[1]
            filename = str(uuid.uuid4())[:8] + ext
            dest = os.path.join(img_dir, filename)
            shutil.copy2(fp, dest)

            # 尝试加密存储到数据库
            user_password = getattr(self, 'user_password', None)
            active_db = getattr(self, 'active_db', None)
            if active_db and user_password:
                try:
                    key = generate_key(user_password)
                    with open(dest, 'rb') as f:
                        raw_data = f.read()
                    image_data_blob = encrypt_data(raw_data, key)
                    # 存储到数据库（不关联疾病，仅用于富文本引用）
                    add_image(
                        active_db, None, filename, '', '', '', 'image',
                        encrypted=True, image_data=image_data_blob
                    )
                except Exception as e:
                    print(f"加密存储图片失败: {e}")

            # 使用本地文件路径显示图片（保留文件用于显示）
            self.text_edit.insertHtml(
                f'<img src="{dest}" style="max-width:100%; margin:8px 0;">'
            )

    def _insert_link(self, checked=False):
        from PyQt5.QtWidgets import QInputDialog
        url, ok = QInputDialog.getText(self, '插入链接', 'URL:')
        if ok and url:
            self.text_edit.insertHtml(f'<a href="{url}" style="color:#3f7bf7;">{url}</a>')

    def _insert_table(self, checked=False):
        self.text_edit.insertHtml(
            '<table border="1" style="border-collapse:collapse; width:100%; margin:8px 0;">'
            '<tr><td style="padding:6px; border:1px solid #444;">&nbsp;</td>'
            '<td style="padding:6px; border:1px solid #444;">&nbsp;</td></tr>'
            '<tr><td style="padding:6px; border:1px solid #444;">&nbsp;</td>'
            '<td style="padding:6px; border:1px solid #444;">&nbsp;</td></tr>'
            '</table>'
        )

    def _insert_hr(self, checked=False):
        self.text_edit.insertHtml('<hr style="border:1px solid #444; margin:8px 0;">')

    def _add_image_caption(self, checked=False):
        """为图片添加标题和注释（类似Trilium Notes）"""
        from PyQt5.QtWidgets import QDialog, QVBoxLayout, QLineEdit, QTextEdit, QLabel, QPushButton, QHBoxLayout
        
        # 创建对话框
        dialog = QDialog(self)
        dialog.setWindowTitle('图片标题')
        dialog.setFixedSize(400, 250)
        
        layout = QVBoxLayout(dialog)
        layout.setContentsMargins(20, 20, 20, 20)
        layout.setSpacing(12)
        
        # 标题输入
        layout.addWidget(QLabel('图片标题：'))
        title_input = QLineEdit()
        title_input.setPlaceholderText('输入图片标题...')
        layout.addWidget(title_input)
        
        # 注释输入
        layout.addWidget(QLabel('注释说明：'))
        caption_input = QTextEdit()
        caption_input.setPlaceholderText('输入图片注释（可选）...')
        caption_input.setFixedHeight(80)
        layout.addWidget(caption_input)
        
        # 按钮
        btn_layout = QHBoxLayout()
        btn_ok = QPushButton('确定')
        btn_ok.setObjectName('accentBtn')
        btn_ok.clicked.connect(dialog.accept)
        btn_cancel = QPushButton('取消')
        btn_cancel.clicked.connect(dialog.reject)
        btn_layout.addWidget(btn_ok)
        btn_layout.addWidget(btn_cancel)
        layout.addLayout(btn_layout)
        
        if dialog.exec_() == QDialog.Accepted:
            title = title_input.text().strip()
            caption = caption_input.toPlainText().strip()
            
            if title or caption:
                # 创建图片标题HTML结构
                html = '<div style="margin:8px 0; text-align:center;">'
                if title:
                    html += f'<div style="font-size:12px; color:#888; font-style:italic; margin-top:4px;">{title}</div>'
                if caption:
                    html += f'<div style="font-size:11px; color:#666; margin-top:2px; text-align:left; padding:4px 8px; background:rgba(0,0,0,20); border-radius:3px;">{caption}</div>'
                html += '</div>'
                
                self.text_edit.insertHtml(html)

    def _insert_code_block(self, checked=False):
        self.text_edit.insertHtml(
            '<pre style="background:#1a1a1e; color:#c8c8cc; padding:12px; '
            'border-radius:4px; font-family:Consolas,monospace; font-size:13px; '
            'overflow-x:auto; margin:8px 0;">代码</pre>'
        )

    # ── 笔记管理 ──
    def _load_note_list(self):
        """加载笔记列表（支持层级显示）"""
        self.note_list.clear()
        conn = sqlite3.connect(self.active_db)
        c = conn.cursor()
        
        if self.record_type == 'medical':
            c.execute("SELECT id, parent_id, title FROM medical_records WHERE disease_id=? ORDER BY parent_id, id",
                      (self.disease_id,))
        else:
            c.execute("SELECT id, parent_id, title FROM anatomy_records ORDER BY parent_id, id")
        
        rows = c.fetchall()
        conn.close()
        
        # 构建层级结构
        def build_tree(items, parent_id, level=0):
            for rid, pid, title in items:
                if pid == parent_id:
                    indent = '    ' * level
                    display_title = f'{indent}▶ {title}' if level > 0 else title
                    item = QListWidgetItem(display_title or '无标题')
                    item.setData(Qt.UserRole, rid)
                    item.setData(Qt.UserRole + 1, pid)  # 存储父ID
                    item.setData(Qt.UserRole + 2, level)  # 存储层级
                    self.note_list.addItem(item)
                    # 递归添加子笔记
                    build_tree(items, rid, level + 1)
        
        build_tree(rows, 0, 0)
        
        # 自动选中第一个
        if self.note_list.count() > 0:
            self.note_list.setCurrentRow(0)

    def _on_note_selected(self, row):
        """选中笔记时加载内容"""
        if row < 0:
            return
        item = self.note_list.item(row)
        if not item:
            return
        rid = item.data(Qt.UserRole)
        self._current_record_id = rid
        conn = sqlite3.connect(self.active_db)
        c = conn.cursor()
        if self.record_type == 'medical':
            c.execute("SELECT title, content FROM medical_records WHERE id=?", (rid,))
        else:
            c.execute("SELECT title, content FROM anatomy_records WHERE id=?", (rid,))
        row_data = c.fetchone()
        conn.close()
        if row_data:
            # 阻止自动保存触发
            self._save_timer.stop()
            self.title_input.blockSignals(True)
            self.text_edit.blockSignals(True)
            self.tag_input.blockSignals(True)
            self.title_input.setText(row_data[0] or '')
            self.text_edit.setHtml(row_data[1] or '')
            self.tag_input.setText('')
            self.title_input.blockSignals(False)
            self.text_edit.blockSignals(False)
            self.tag_input.blockSignals(False)
            self.save_indicator.setText('已保存')
            self.save_indicator.setStyleSheet("color: #4ade80; font-size: 11px; background: transparent;")

    def _new_record(self, checked=False):
        """新建笔记（顶级）"""
        self._create_child_record(0)

    def _create_child_record(self, parent_id):
        """创建子笔记"""
        conn = sqlite3.connect(self.active_db)
        c = conn.cursor()
        if self.record_type == 'medical':
            c.execute("INSERT INTO medical_records (disease_id, parent_id, title, content) VALUES (?, ?, ?, ?)",
                      (self.disease_id, parent_id, '新建笔记', ''))
        else:
            c.execute("INSERT INTO anatomy_records (parent_id, title, content) VALUES (?, ?, ?)",
                      (parent_id, '新建笔记', ''))
        conn.commit()
        new_id = c.lastrowid
        conn.close()
        self._current_record_id = new_id
        self._load_note_list()
        # 选中新建的笔记
        for i in range(self.note_list.count()):
            if self.note_list.item(i).data(Qt.UserRole) == new_id:
                self.note_list.setCurrentRow(i)
                break
        # 加载新笔记内容
        self.title_input.setText('新建笔记')
        self.text_edit.setHtml('')
        self.tag_input.setText('')
        self.title_input.setFocus()
        self.title_input.selectAll()

    def _show_note_context_menu(self, pos):
        """显示笔记列表右键菜单"""
        menu = QMenu(self)
        menu.setStyleSheet("""
            QMenu { background-color: #2a2a30; color: #c8c8cc; border: 1px solid #444; padding: 4px; }
            QMenu::item { padding: 6px 24px; border-radius: 4px; }
            QMenu::item:selected { background-color: #3f7bf7; }
        """)
        
        item = self.note_list.itemAt(pos)
        if item:
            # 添加子笔记
            act_add_child = menu.addAction('添加子笔记')
            menu.addSeparator()
            # 删除笔记
            act_delete = menu.addAction('删除笔记')
            menu.addSeparator()
            # 重命名
            act_rename = menu.addAction('重命名')
            
            action = menu.exec_(self.note_list.mapToGlobal(pos))
            if action == act_add_child:
                parent_id = item.data(Qt.UserRole)
                self._create_child_record(parent_id)
            elif action == act_delete:
                self._delete_record(item.data(Qt.UserRole))
            elif action == act_rename:
                self._rename_record(item)
    
    def _delete_record(self, record_id):
        """删除笔记及其所有子笔记"""
        reply = QMessageBox.question(self, '确认删除', '确定要删除这条笔记及其所有子笔记吗？',
                                     QMessageBox.Yes | QMessageBox.No)
        if reply == QMessageBox.No:
            return
            
        conn = sqlite3.connect(self.active_db)
        c = conn.cursor()
        
        # 删除子笔记（递归删除）
        if self.record_type == 'medical':
            c.execute("DELETE FROM medical_records WHERE id=? OR parent_id=?", (record_id, record_id))
        else:
            c.execute("DELETE FROM anatomy_records WHERE id=? OR parent_id=?", (record_id, record_id))
        
        conn.commit()
        conn.close()
        self._load_note_list()
        self.title_input.setText('')
        self.text_edit.setHtml('')
        self._current_record_id = None

    def _rename_record(self, item):
        """重命名笔记"""
        from PyQt5.QtWidgets import QInputDialog
        current_title = item.text()
        new_title, ok = QInputDialog.getText(self, '重命名', '输入新标题:', text=current_title)
        if ok and new_title.strip():
            record_id = item.data(Qt.UserRole)
            conn = sqlite3.connect(self.active_db)
            c = conn.cursor()
            if self.record_type == 'medical':
                c.execute("UPDATE medical_records SET title=? WHERE id=?", (new_title.strip(), record_id))
            else:
                c.execute("UPDATE anatomy_records SET title=? WHERE id=?", (new_title.strip(), record_id))
            conn.commit()
            conn.close()
            item.setText(new_title.strip())

    def _delete_current_record(self, checked=False):
        """删除当前笔记"""
        if self._current_record_id is None:
            return
        reply = QMessageBox.question(
            self, '确认删除', '确定要删除此笔记吗？',
            QMessageBox.Yes | QMessageBox.No, QMessageBox.No
        )
        if reply == QMessageBox.Yes:
            conn = sqlite3.connect(self.active_db)
            c = conn.cursor()
            if self.record_type == 'medical':
                c.execute("DELETE FROM medical_records WHERE id=?", (self._current_record_id,))
            else:
                c.execute("DELETE FROM anatomy_records WHERE id=?", (self._current_record_id,))
            conn.commit()
            conn.close()
            self._current_record_id = None
            self._load_note_list()

    def _on_content_changed(self):
        """内容变更时触发自动保存（防抖500ms）"""
        self.save_indicator.setText('编辑中...')
        self.save_indicator.setStyleSheet("color: #f0a030; font-size: 11px; background: transparent;")
        # 更新列表项标题
        if self._current_record_id is not None:
            current_item = self.note_list.currentItem()
            if current_item:
                current_item.setText(self.title_input.text() or '无标题')
        self._save_timer.start(500)

    def _auto_save(self):
        """自动保存到数据库"""
        if self._current_record_id is None:
            return
        title = self.title_input.text().strip() or '无标题'
        content = self.text_edit.toHtml()
        try:
            conn = sqlite3.connect(self.active_db)
            c = conn.cursor()
            if self.record_type == 'medical':
                c.execute("UPDATE medical_records SET title=?, content=? WHERE id=?",
                          (title, content, self._current_record_id))
            else:
                c.execute("UPDATE anatomy_records SET title=?, content=? WHERE id=?",
                          (title, content, self._current_record_id))
            conn.commit()
            conn.close()
            self.save_indicator.setText('已保存')
            self.save_indicator.setStyleSheet("color: #4ade80; font-size: 11px; background: transparent;")
        except Exception as e:
            self.save_indicator.setText(f'保存失败: {e}')
            self.save_indicator.setStyleSheet("color: #e64a3a; font-size: 11px; background: transparent;")

    def _manual_save(self, checked=False):
        """手动保存"""
        self._save_timer.stop()
        self._auto_save()

    def keyPressEvent(self, event):
        if event.key() == Qt.Key_S and event.modifiers() & Qt.ControlModifier:
            self._manual_save()
            event.accept()
        else:
            super().keyPressEvent(event)

    def closeEvent(self, event):
        """关闭前保存"""
        self._save_timer.stop()
        self._auto_save()
        super().closeEvent(event)


# ── 详情面板 ──────────────────────────────────────────────
class DetailPanel(QWidget):
    def __init__(self, parent=None, user_info=None, user_password=None):
        super().__init__(parent)
        self.current_disease_id = None
        self.active_db = DATABASE
        self.disease_data = None
        self.current_tab = 0
        self.user_info = user_info
        self.user_password = user_password
        self._current_record_type = 'medical'
        self._editing_record_id = None
        self._is_edit_mode = False

        layout = QVBoxLayout(self)
        layout.setContentsMargins(0, 0, 0, 0)
        layout.setSpacing(0)

        # 标签栏
        tab_bar = QHBoxLayout()
        tab_bar.setContentsMargins(0, 0, 0, 0)
        tab_bar.setSpacing(0)

        self.tab_buttons = []
        self.tab_pin_btns = []
        self._detached_windows = {}
        self._tab_drag_idx = None
        self._tab_drag_start = None
        tab_names = ['临床与诊断', '影像所见及标准报告参考', '影像相关病例', '医学资料', '影像解剖图谱与资料']
        for i, name in enumerate(tab_names):
            # 标签按钮容器
            tab_item = QWidget()
            tab_item_layout = QHBoxLayout(tab_item)
            tab_item_layout.setContentsMargins(0, 0, 0, 0)
            tab_item_layout.setSpacing(0)
            # 图钉按钮（单击独立显示）
            pin_btn = _PinBtn()
            pin_btn.clicked.connect(lambda checked, idx=i: self._detach_tab(idx))
            self.tab_pin_btns.append(pin_btn)
            tab_item_layout.addWidget(pin_btn)
            # 标签按钮（支持拖动脱离）
            btn = QPushButton(name)
            btn.setObjectName('tabBtn')
            btn.setProperty('active', i == 0)
            btn.setCursor(Qt.PointingHandCursor)
            btn.clicked.connect(lambda checked, idx=i: self.switch_tab(idx))
            self.tab_buttons.append(btn)
            tab_item_layout.addWidget(btn)
            tab_bar.addWidget(tab_item)

        layout.addLayout(tab_bar)

        # 分隔线
        line = QFrame()
        line.setFrameShape(QFrame.HLine)
        line.setStyleSheet("background-color: #2a2a30; max-height: 1px;")
        layout.addWidget(line)

        # 文本浏览器（临床与诊断、影像所见及标准报告参考标签页使用）
        self.content_browser = QTextBrowser()
        self.content_browser.setOpenExternalLinks(True)
        self.content_browser.setStyleSheet("border: none; padding: 16px;")
        layout.addWidget(self.content_browser, stretch=1)

        # 记录标签页（影像所见、医学资料、影像解剖）：左侧列表+右侧内容
        self.record_splitter = QSplitter(Qt.Horizontal)
        self.record_splitter.setStyleSheet("QSplitter::handle { background: #2a2a30; width: 1px; }")

        # 左侧笔记列表（紧凑优雅风格）
        self.record_list_panel = QWidget()
        self.record_list_panel.setStyleSheet("background: transparent;")
        record_list_layout = QVBoxLayout(self.record_list_panel)
        record_list_layout.setContentsMargins(0, 0, 0, 0)
        record_list_layout.setSpacing(0)
        # 笔记列表标题栏：标题 + 菜单按钮
        record_list_header = QWidget()
        record_list_header.setFixedHeight(28)
        record_list_header.setStyleSheet("background: transparent; border-bottom: 1px solid rgba(255,255,255,15);")
        header_hl = QHBoxLayout(record_list_header)
        header_hl.setContentsMargins(8, 0, 4, 0)
        header_hl.setSpacing(0)
        header_label = QLabel('笔记')
        header_label.setFont(QFont("Microsoft YaHei", 9, QFont.Bold))
        header_label.setStyleSheet("background: transparent; color: #888890;")
        header_hl.addWidget(header_label)
        header_hl.addStretch()
        # 菜单按钮（⋮）
        self.record_menu_btn = QPushButton('⋮')
        self.record_menu_btn.setFixedSize(24, 24)
        self.record_menu_btn.setFont(QFont("Microsoft YaHei", 12))
        self.record_menu_btn.setCursor(Qt.PointingHandCursor)
        self.record_menu_btn.setStyleSheet("""
            QPushButton { background: transparent; border: none; color: #666670; border-radius: 3px; }
            QPushButton:hover { background: rgba(255,255,255,15); color: #c8c8cc; }
        """)
        self.record_menu_btn.clicked.connect(lambda checked: self._show_record_menu())
        header_hl.addWidget(self.record_menu_btn)
        record_list_layout.addWidget(record_list_header)
        self.record_list = QListWidget()
        self.record_list.currentRowChanged.connect(self._on_record_list_clicked)
        self.record_list.itemDoubleClicked.connect(self._on_record_list_double_clicked)
        self.record_list.setContextMenuPolicy(Qt.CustomContextMenu)
        self.record_list.customContextMenuRequested.connect(self._show_record_item_menu)
        self.record_list.setStyleSheet("""
            QListWidget { border: none; background: transparent; outline: none; }
            QListWidget::item { padding: 4px 8px; border-bottom: 1px solid rgba(255,255,255,8);
                                color: #a0a0a8; font-size: 12px; }
            QListWidget::item:selected { background-color: rgba(63,123,247,30); color: #e0e0e4; }
            QListWidget::item:hover { background-color: rgba(255,255,255,8); color: #c8c8cc; }
        """)
        record_list_layout.addWidget(self.record_list)
        self.record_list_panel.setFixedWidth(140)
        self.record_splitter.addWidget(self.record_list_panel)

        # 右侧内容：堆叠组件（查看模式 / 编辑模式）
        self.record_stack = QStackedWidget()

        # Page 0: 查看模式 - QTextBrowser（优化排版）
        self.record_content_browser = QTextBrowser()
        self.record_content_browser.setOpenExternalLinks(True)
        self.record_content_browser.setStyleSheet("""
            QTextBrowser { border: none; background: transparent; padding: 20px 28px; color: #d0d0d4; }
            QTextBrowser h1 { color: #3f7bf7; margin-bottom: 4px; }
            QTextBrowser h2 { color: #5a91ff; margin-top: 20px; margin-bottom: 8px; }
            QTextBrowser a { color: #3f7bf7; }
        """)
        self.record_stack.addWidget(self.record_content_browser)

        # Page 1: 编辑模式 - 标题输入 + 工具面板 + 编辑器（工具固定在上方）
        edit_widget = QWidget()
        edit_layout = QVBoxLayout(edit_widget)
        edit_layout.setContentsMargins(0, 0, 0, 0)
        edit_layout.setSpacing(0)

        # 顶部工具面板（常驻显示）
        self._tool_panel = QWidget()
        self._tool_panel.setStyleSheet("background: #222226;")
        tool_panel_layout = QHBoxLayout(self._tool_panel)
        tool_panel_layout.setContentsMargins(8, 4, 8, 4)
        tool_panel_layout.setSpacing(2)

        # 格式工具
        format_tools = [
            ('bold', '加粗'), ('italic', '斜体'), ('underline', '下划线'), ('strikethrough', '删除线'),
            ('superscript', '上标'), ('subscript', '下标'), ('code', '行内代码'),
        ]
        for fmt_id, tip in format_tools:
            btn = _FormatBtn(fmt_id, fg_color='#c8c8cc')
            btn.setToolTip(tip)
            btn.clicked.connect(lambda checked, fid=fmt_id: self._on_format_action(fid))
            tool_panel_layout.addWidget(btn)

        self._add_tool_sep(tool_panel_layout)

        # 对齐工具
        align_tools = [
            ('align_left', '左对齐'), ('align_center', '居中'), ('align_right', '右对齐'),
        ]
        for fmt_id, tip in align_tools:
            btn = _FormatBtn(fmt_id, fg_color='#c8c8cc')
            btn.setToolTip(tip)
            btn.clicked.connect(lambda checked, fid=fmt_id: self._on_format_action(fid))
            tool_panel_layout.addWidget(btn)

        self._add_tool_sep(tool_panel_layout)

        # 列表工具
        list_tools = [
            ('bullet_list', '无序列表'), ('ordered_list', '有序列表'),
        ]
        for fmt_id, tip in list_tools:
            btn = _FormatBtn(fmt_id, fg_color='#c8c8cc')
            btn.setToolTip(tip)
            btn.clicked.connect(lambda checked, fid=fmt_id: self._on_format_action(fid))
            tool_panel_layout.addWidget(btn)

        self._add_tool_sep(tool_panel_layout)

        # 插入工具
        insert_tools = [
            ('insert_image', '插入图片'), ('insert_table', '插入表格'),
            ('insert_hr', '分隔线'), ('insert_link', '插入链接'),
            ('insert_sticker', '插入贴图'),
        ]
        for fmt_id, tip in insert_tools:
            btn = _FormatBtn(fmt_id, fg_color='#c8c8cc')
            btn.setToolTip(tip)
            btn.clicked.connect(lambda checked, fid=fmt_id: self._on_block_action(fid))
            tool_panel_layout.addWidget(btn)

        self._add_tool_sep(tool_panel_layout)

        # 撤销/重做
        undo_redo_tools = [('undo', '撤销'), ('redo', '重做')]
        for fmt_id, tip in undo_redo_tools:
            btn = _FormatBtn(fmt_id, fg_color='#c8c8cc')
            btn.setToolTip(tip)
            btn.clicked.connect(lambda checked, fid=fmt_id: self._on_format_action(fid))
            tool_panel_layout.addWidget(btn)

        self._add_tool_sep(tool_panel_layout)

        # 图像编辑工具
        tool_panel_layout.addWidget(QLabel('图像:'))
        image_tools = [
            ('image_select', '选择'), ('image_delete', '删除'),
        ]
        for fmt_id, tip in image_tools:
            btn = _FormatBtn(fmt_id, fg_color='#c8c8cc')
            btn.setToolTip(tip)
            btn.clicked.connect(lambda checked, fid=fmt_id: self._on_image_action(fid))
            tool_panel_layout.addWidget(btn)

        self._add_tool_sep(tool_panel_layout)
        
        # 文本框和标注工具
        tool_panel_layout.addWidget(QLabel('标注:'))
        annotation_tools = [
            ('insert_textbox', '文本框'), ('arrow', '箭头'), ('rect', '矩形'), 
            ('circle', '圆形'), ('line', '直线'), ('text', '文字'),
        ]
        for fmt_id, tip in annotation_tools:
            btn = _FormatBtn(fmt_id, fg_color='#c8c8cc')
            btn.setToolTip(tip)
            btn.clicked.connect(lambda checked, fid=fmt_id: self._on_annotation_action(fid))
            tool_panel_layout.addWidget(btn)

        self._add_tool_sep(tool_panel_layout)

        # 图片环绕方式
        tool_panel_layout.addWidget(QLabel('环绕:'))
        self._wrap_combo = QComboBox()
        self._wrap_combo.addItems(['嵌入型', '四周环绕', '紧密环绕', '穿越环绕', '上下型', '衬于下方', '浮于上方'])
        self._wrap_combo.setFixedWidth(100)
        self._wrap_combo.setStyleSheet("""
            QComboBox { background: rgba(255,255,255,10); color: #c8c8cc; border: 1px solid #444; border-radius: 3px; font-size: 12px; }
            QComboBox::drop-down { border: none; }
            QComboBox QAbstractItemView { background: #2a2a30; color: #c8c8cc; }
        """)
        self._wrap_combo.currentIndexChanged.connect(self._on_wrap_changed)
        tool_panel_layout.addWidget(self._wrap_combo)

        tool_panel_layout.addStretch()

        # 将工具面板放入水平滚动区域
        tool_scroll = QScrollArea()
        tool_scroll.setWidget(self._tool_panel)
        tool_scroll.setWidgetResizable(False)
        tool_scroll.setFixedHeight(56)
        tool_scroll.setHorizontalScrollBarPolicy(Qt.ScrollBarAsNeeded)
        tool_scroll.setVerticalScrollBarPolicy(Qt.ScrollBarAlwaysOff)
        tool_scroll.setStyleSheet("""
            QScrollArea { border: none; background: #222226; }
            QScrollBar:horizontal { height: 6px; background: #222226; }
            QScrollBar::handle:horizontal { background: #555; border-radius: 3px; min-width: 30px; }
            QScrollBar::add-line:horizontal, QScrollBar::sub-line:horizontal { width: 0; }
        """)
        edit_layout.addWidget(tool_scroll)

        # 工具面板分隔线
        tool_sep = QFrame()
        tool_sep.setFrameShape(QFrame.HLine)
        tool_sep.setStyleSheet("background-color: #2a2a30;")
        edit_layout.addWidget(tool_sep)

        # 编辑区域内容
        edit_content = QWidget()
        edit_content_layout = QVBoxLayout(edit_content)
        edit_content_layout.setContentsMargins(0, 0, 0, 0)
        edit_content_layout.setSpacing(0)

        # 编辑模式标题栏（紧凑）
        edit_title_bar = QWidget()
        edit_title_bar.setFixedHeight(42)
        edit_title_layout = QHBoxLayout(edit_title_bar)
        edit_title_layout.setContentsMargins(12, 4, 12, 4)
        edit_title_layout.setSpacing(8)
        self.inline_title_input = QLineEdit()
        self.inline_title_input.setPlaceholderText('输入标题...')
        self.inline_title_input.setFont(QFont("Microsoft YaHei", 13, QFont.Bold))
        self.inline_title_input.setStyleSheet("""
            QLineEdit { background: transparent; border: none; border-bottom: 1px solid #3a3a40;
                        color: #e0e0e4; padding: 2px 0; }
            QLineEdit:focus { border-bottom-color: #3f7bf7; }
        """)
        edit_title_layout.addWidget(self.inline_title_input, stretch=1)
        # 保存状态指示
        self.inline_save_indicator = QLabel('已保存')
        self.inline_save_indicator.setStyleSheet("color: #4ade80; font-size: 11px; background: transparent;")
        edit_title_layout.addWidget(self.inline_save_indicator)
        # 保存按钮
        self.btn_inline_save = QPushButton('保存')
        self.btn_inline_save.setFixedHeight(26)
        self.btn_inline_save.setFixedWidth(52)
        self.btn_inline_save.setObjectName('accentBtn')
        self.btn_inline_save.setToolTip('保存 (Ctrl+S)')
        self.btn_inline_save.clicked.connect(lambda checked: self._exit_edit_mode(save=True))
        edit_title_layout.addWidget(self.btn_inline_save)
        # 取消按钮
        self.btn_inline_cancel = QPushButton('取消')
        self.btn_inline_cancel.setFixedHeight(26)
        self.btn_inline_cancel.setFixedWidth(52)
        self.btn_inline_cancel.clicked.connect(lambda checked: self._exit_edit_mode(save=False))
        edit_title_layout.addWidget(self.btn_inline_cancel)
        edit_content_layout.addWidget(edit_title_bar)

        # 分隔线
        title_sep = QFrame()
        title_sep.setFrameShape(QFrame.HLine)
        title_sep.setStyleSheet("background-color: #2a2a30; max-height: 1px;")
        edit_content_layout.addWidget(title_sep)

        # 内嵌富文本编辑器（支持图片选中、缩放、拖动、调整尺寸）
        self.inline_text_edit = _ImageEditableTextEdit(self)
        self.inline_text_edit.setAcceptRichText(True)
        self.inline_text_edit.setContextMenuPolicy(Qt.CustomContextMenu)
        self.inline_text_edit.customContextMenuRequested.connect(self._show_inline_context_menu)
        self.inline_text_edit.setStyleSheet("""
            QTextEdit { border: none; background: transparent; padding: 16px 24px; color: #d0d0d4; }
        """)
        edit_content_layout.addWidget(self.inline_text_edit, stretch=1)

        edit_layout.addWidget(edit_content, stretch=1)

        self.record_stack.addWidget(edit_widget)

        self.record_splitter.addWidget(self.record_stack)
        self.record_splitter.setStretchFactor(0, 0)
        self.record_splitter.setStretchFactor(1, 1)
        self.record_splitter.setSizes([160, 500])
        self.record_splitter.hide()
        layout.addWidget(self.record_splitter, stretch=1)

        # 存储当前影像页面的图片列表（用于双击浏览）
        self._current_image_list = []

        # 自动保存定时器
        self._inline_save_timer = QTimer(self)
        self._inline_save_timer.setSingleShot(True)
        self._inline_save_timer.timeout.connect(self._inline_auto_save)

        # 图片编辑弹窗（用于图片标注）
        self._image_edit_popup = _ImageEditPopup()
        self._image_edit_popup.edit_image.connect(self._open_image_editor)
        self._image_edit_popup.annotate_image.connect(self._quick_annotate_image)
        self._image_edit_popup.hide()

        # 连接编辑器信号
        self.inline_text_edit.textChanged.connect(self._on_edit_text_changed)

        # 安装事件过滤器以检测图片点击 + 标签拖动脱离
        self.inline_text_edit.viewport().installEventFilter(self)
        for btn in self.tab_buttons:
            btn.installEventFilter(self)

    def _add_tool_sep(self, layout):
        """添加工具面板分隔线"""
        sep = QFrame()
        sep.setFrameShape(QFrame.VLine)
        sep.setFixedWidth(1)
        sep.setStyleSheet("background-color: rgba(255,255,255,20);")
        layout.addWidget(sep)
        layout.addSpacing(4)

    def _on_image_action(self, action_id):
        """处理图像编辑工具的动作"""
        self.inline_text_edit.set_image_tool(action_id)

    def _on_annotation_action(self, action_id):
        """处理标注工具的动作（文本框、箭头、矩形等）"""
        if action_id == 'insert_textbox':
            self._insert_textbox()
        elif action_id in ('arrow', 'rect', 'circle', 'line', 'text'):
            # 在当前选中的图片上启动标注模式
            self._start_image_annotation(action_id)
    
    def _insert_textbox(self):
        """在编辑器中插入可拖拽文本框"""
        if not hasattr(self, '_textboxes'):
            self._textboxes = []
        textbox = _DraggableTextBox(self.inline_text_edit.viewport())
        # 默认位置：编辑器中心
        editor_rect = self.inline_text_edit.viewport().rect()
        textbox.resize(200, 100)
        textbox.move(
            editor_rect.width() // 2 - 100,
            editor_rect.height() // 2 - 50
        )
        textbox.removed.connect(lambda: self._textboxes.remove(textbox) if textbox in self._textboxes else None)
        textbox.show()
        textbox.raise_()
        self._textboxes.append(textbox)
    
    def _start_image_annotation(self, tool_id):
        """在选中图片上启动标注模式：打开完整图片编辑器"""
        edit = self.inline_text_edit
        if not edit._selected_image_name:
            QMessageBox.information(self, '提示', '请先选中要标注的图片')
            return
        # 通过完整图片编辑器进行标注
        img_path = edit._selected_image_name
        if hasattr(self, '_image_edit_popup'):
            self._image_edit_popup.annotate_image.emit(img_path, tool_id)

    def _on_wrap_changed(self, idx):
        """处理图片环绕方式变化"""
        wrap_modes = ['inline', 'square', 'tight', 'through', 'topAndBottom', 'behind', 'front']
        if idx < len(wrap_modes):
            self.inline_text_edit.set_image_wrap_mode(wrap_modes[idx])

    def _show_record_menu(self, checked=False):
        """显示笔记列表操作菜单"""
        menu = QMenu(self)
        menu.setStyleSheet("""
            QMenu { background-color: #2a2a30; color: #c8c8cc; border: 1px solid #444; padding: 4px; }
            QMenu::item { padding: 6px 24px; border-radius: 4px; }
            QMenu::item:selected { background-color: #3f7bf7; }
        """)
        act_add = menu.addAction('+ 添加笔记')
        act_add_child = menu.addAction('+ 添加子笔记')
        act_edit = menu.addAction('编辑选中')
        menu.addSeparator()
        act_delete = menu.addAction('删除选中')
        action = menu.exec_(self.record_menu_btn.mapToGlobal(
            QPoint(self.record_menu_btn.width(), 0)))
        if action == act_add:
            self._add_record()
        elif action == act_add_child:
            current = self.record_list.currentItem()
            if current:
                parent_id = current.data(Qt.UserRole)
                self._add_record(parent_id=parent_id)
            else:
                self._add_record()
        elif action == act_edit:
            self._edit_record()
        elif action == act_delete:
            self._delete_record()

    def _on_edit_text_changed(self):
        """编辑器文本变化时触发自动保存"""
        if self._is_edit_mode:
            self.inline_save_indicator.setText('编辑中...')
            self.inline_save_indicator.setStyleSheet("color: #f0a030; font-size: 11px; background: transparent;")
            self._inline_save_timer.start(2000)

    def _on_format_action(self, action_id):
        """处理格式栏的动作"""
        if action_id.startswith('heading_'):
            idx = int(action_id.split('_')[1])
            self._inline_set_heading(idx)
        elif action_id == 'bold':
            self._inline_toggle_bold()
        elif action_id == 'italic':
            self._inline_toggle_italic()
        elif action_id == 'underline':
            self._inline_toggle_underline()
        elif action_id == 'strikethrough':
            self._inline_toggle_strikethrough()
        elif action_id == 'bullet_list':
            self._inline_toggle_bullet_list()
        elif action_id == 'ordered_list':
            self._inline_toggle_ordered_list()
        elif action_id == 'align_left':
            self._inline_align_left()
        elif action_id == 'align_center':
            self._inline_align_center()
        elif action_id == 'align_right':
            self._inline_align_right()
        elif action_id == 'text_color':
            self._inline_set_text_color()
        elif action_id == 'bg_color':
            self._inline_set_bg_color()
        elif action_id == 'superscript':
            self._inline_toggle_superscript()
        elif action_id == 'subscript':
            self._inline_toggle_subscript()
        elif action_id == 'code':
            self._inline_insert_code()
        elif action_id == 'insert_link':
            self._inline_insert_link()
        elif action_id == 'undo':
            self.inline_text_edit.undo()
        elif action_id == 'redo':
            self.inline_text_edit.redo()

    def _on_block_action(self, action_id):
        """处理块工具栏的动作"""
        if action_id == 'insert_image':
            self._inline_insert_image()
        elif action_id == 'insert_image_from_db':
            self._inline_insert_image_from_db()
        elif action_id == 'image_left':
            self._inline_image_align('left')
        elif action_id == 'image_center':
            self._inline_image_align('center')
        elif action_id == 'image_right':
            self._inline_image_align('right')
        elif action_id == 'image_full':
            self._inline_image_align('full')
        elif action_id == 'image_resize_small':
            self._inline_image_resize(0.8)
        elif action_id == 'image_resize_large':
            self._inline_image_resize(1.25)
        elif action_id == 'insert_table':
            self._inline_insert_table()
        elif action_id == 'insert_code':
            self._inline_insert_code_block()
        elif action_id == 'insert_hr':
            self._inline_insert_hr()
        elif action_id == 'insert_link':
            self._inline_insert_link()
        elif action_id == 'insert_todo':
            self._inline_insert_todo()
        elif action_id == 'insert_note':
            self._inline_insert_note()
        elif action_id == 'insert_sticker':
            self._insert_sticker()
    
    def _insert_sticker(self):
        """插入贴图：从素材库选择并作为可拖动文本框形式插入"""
        from PyQt5.QtWidgets import QDialog, QVBoxLayout, QGridLayout, QPushButton, QLabel
        # 简易素材库：常用医学符号
        stickers = [
            ('★', '星标'), ('✓', '正确'), ('✗', '错误'), ('!', '警告'),
            ('?', '疑问'), ('→', '右箭头'), ('←', '左箭头'), ('↑', '上箭头'),
            ('↓', '下箭头'), ('●', '实心圆'), ('○', '空心圆'), ('■', '实心方'),
        ]
        dlg = QDialog(self)
        dlg.setWindowTitle('选择贴图')
        dlg.setFixedSize(320, 240)
        dlg.setStyleSheet("""
            QDialog { background: #2a2a30; }
            QPushButton { background: #3a3a40; color: #e0e0e4; 
                         border: 1px solid #555; border-radius: 4px;
                         font-size: 24px; min-height: 40px; }
            QPushButton:hover { background: #3f7bf7; border-color: #3f7bf7; }
            QLabel { color: #c8c8cc; }
        """)
        v = QVBoxLayout(dlg)
        v.addWidget(QLabel('选择贴图素材:'))
        grid = QGridLayout()
        selected = {'val': None}
        for i, (sym, name) in enumerate(stickers):
            btn = QPushButton(sym)
            btn.setToolTip(name)
            btn.clicked.connect(lambda checked, s=sym: (selected.update({'val': s}), dlg.accept()))
            grid.addWidget(btn, i // 4, i % 4)
        v.addLayout(grid)
        if dlg.exec_() == QDialog.Accepted and selected['val']:
            # 作为可拖动文本框插入
            if not hasattr(self, '_textboxes'):
                self._textboxes = []
            tb = _DraggableTextBox(self.inline_text_edit.viewport())
            tb.text_edit.setPlainText(selected['val'])
            tb.text_edit.setFont(QFont('Microsoft YaHei', 24, QFont.Bold))
            tb.text_edit.setAlignment(Qt.AlignCenter)
            tb.resize(80, 80)
            tb.move(50, 50)
            tb.removed.connect(lambda: self._textboxes.remove(tb) if tb in self._textboxes else None)
            tb.show()
            tb.raise_()
            self._textboxes.append(tb)

    def eventFilter(self, obj, event):
        """事件过滤器 - 检测编辑器中图片的点击 + 标签拖动脱离"""
        # 标签按钮拖动脱离
        if obj in self.tab_buttons and obj.isEnabled():
            tab_idx = self.tab_buttons.index(obj)
            if event.type() == event.MouseButtonPress and event.button() == Qt.LeftButton:
                self._tab_drag_idx = tab_idx
                self._tab_drag_start = event.globalPos()
            elif event.type() == event.MouseMove and self._tab_drag_idx is not None and self._tab_drag_start is not None:
                delta = event.globalPos() - self._tab_drag_start
                if delta.manhattanLength() > 15:
                    # 拖动距离足够，脱离为独立窗口
                    drag_idx = self._tab_drag_idx
                    self._tab_drag_idx = None
                    self._tab_drag_start = None
                    self._detach_tab(drag_idx)
                    return True
            elif event.type() == event.MouseButtonRelease:
                self._tab_drag_idx = None
                self._tab_drag_start = None
        # 编辑器图片点击检测 - 不拦截，让 _ImageEditableTextEdit 自行处理选中逻辑
        # 仅在双击图片时打开完整编辑器
        if obj == self.inline_text_edit.viewport() and self._is_edit_mode:
            if event.type() == event.MouseButtonDblClick and event.button() == Qt.LeftButton:
                cursor = self.inline_text_edit.cursorForPosition(event.pos())
                cursor.movePosition(QTextCursor.Right)
                char_fmt = cursor.charFormat()
                if char_fmt.isImageFormat():
                    img_path = char_fmt.toImageFormat().name()
                    if img_path:
                        self._open_image_editor(img_path)
                        return True
        return super().eventFilter(obj, event)

    def _open_image_editor(self, image_path):
        """打开完整图片编辑器"""
        self._image_edit_popup.hide()
        if not image_path or not os.path.exists(image_path):
            QMessageBox.warning(self, '错误', f'图片文件不存在: {image_path}')
            return
        try:
            dlg = ImageEditorDialog(image_path, self, active_db=self.active_db,
                                    user_password=self.user_password)
            dlg.image_updated.connect(self._on_image_updated)
            dlg.exec_()
        except Exception as e:
            QMessageBox.critical(self, '错误', f'打开图片编辑器失败: {e}')

    def _quick_annotate_image(self, image_path, tool_id):
        """快速标注图片 - 打开编辑器并选中对应工具"""
        self._image_edit_popup.hide()
        # 检查图片路径是否存在
        if not image_path or not os.path.exists(image_path):
            QMessageBox.warning(self, '错误', f'图片文件不存在: {image_path}')
            return
        try:
            dlg = ImageEditorDialog(image_path, self, active_db=self.active_db,
                                    user_password=self.user_password)
            dlg.image_updated.connect(self._on_image_updated)
            # 设置工具（会同步canvas.tool和工具栏选中状态）
            dlg._set_tool(tool_id)
            # 同步样式
            style = self._image_edit_popup.get_style()
            dlg.current_color = style['color']
            dlg.current_line_width = style['width']
            dlg.current_line_style = style['style']
            dlg.exec_()
        except Exception as e:
            QMessageBox.critical(self, '错误', f'打开图片编辑器失败: {e}')

    def _on_image_updated(self, image_path):
        """图片编辑器更新后刷新编辑器中的图片"""
        if not image_path:
            return
        # 重新加载文档中的图片资源
        doc = self.inline_text_edit.document()
        if os.path.exists(image_path):
            pixmap = QPixmap(image_path)
            if not pixmap.isNull():
                # 用文件路径作为资源名重新加载
                url = QUrl.fromLocalFile(image_path)
                doc.addResource(QTextDocument.ImageResource, url, pixmap)
                # 同时用原始名称（可能不含file:///前缀）重新加载
                doc.addResource(QTextDocument.ImageResource, QUrl(image_path), pixmap)
        # 刷新编辑器显示
        doc.markContentsDirty(0, doc.characterCount())
        # 自动保存笔记内容
        self._inline_save_timer.stop()
        self._inline_save_current()

    # ── 内嵌编辑器格式化操作 ──
    def _inline_set_heading(self, idx):
        sizes = [10, 20, 16, 14, 12]
        weights = [QFont.Normal, QFont.Bold, QFont.Bold, QFont.Bold, QFont.Bold]
        fmt = self.inline_text_edit.currentCharFormat()
        fmt.setFontPointSize(sizes[idx])
        fmt.setFontWeight(weights[idx])
        self.inline_text_edit.setCurrentCharFormat(fmt)

    def _inline_toggle_bold(self):
        self.inline_text_edit.setFontWeight(
            QFont.Normal if self.inline_text_edit.fontWeight() == QFont.Bold else QFont.Bold
        )

    def _inline_toggle_italic(self):
        self.inline_text_edit.setFontItalic(not self.inline_text_edit.fontItalic())

    def _inline_toggle_underline(self):
        self.inline_text_edit.setFontUnderline(not self.inline_text_edit.fontUnderline())

    def _inline_toggle_strikethrough(self):
        fmt = self.inline_text_edit.currentCharFormat()
        fmt.setFontStrikeOut(not fmt.fontStrikeOut())
        self.inline_text_edit.setCurrentCharFormat(fmt)

    def _inline_toggle_bullet_list(self):
        cursor = self.inline_text_edit.textCursor()
        cursor.insertList(QTextListFormat.ListDisc)

    def _inline_toggle_ordered_list(self):
        cursor = self.inline_text_edit.textCursor()
        cursor.insertList(QTextListFormat.ListDecimal)

    def _inline_align_left(self):
        self.inline_text_edit.setAlignment(Qt.AlignLeft)

    def _inline_align_center(self):
        self.inline_text_edit.setAlignment(Qt.AlignCenter)

    def _inline_align_right(self):
        self.inline_text_edit.setAlignment(Qt.AlignRight)

    def _inline_set_text_color(self):
        from PyQt5.QtWidgets import QColorDialog
        color = QColorDialog.getColor(Qt.black, self, '选择文字颜色')
        if color.isValid():
            self.inline_text_edit.setTextColor(color)

    def _inline_set_bg_color(self):
        from PyQt5.QtWidgets import QColorDialog
        color = QColorDialog.getColor(Qt.white, self, '选择背景颜色')
        if color.isValid():
            self.inline_text_edit.setTextBackgroundColor(color)

    def _inline_insert_image(self):
        file_paths, _ = QFileDialog.getOpenFileNames(
            self, '选择图片', '', '图片文件 (*.png *.jpg *.jpeg *.bmp *.gif)'
        )
        for fp in file_paths:
            img_dir = os.path.join(APP_DIR, 'images')
            os.makedirs(img_dir, exist_ok=True)
            import uuid
            ext = os.path.splitext(fp)[1]
            filename = str(uuid.uuid4())[:8] + ext
            dest = os.path.join(img_dir, filename)
            
            # 先复制文件，然后检查是否需要压缩
            shutil.copy2(fp, dest)
            
            # 如果图像像素太大，自动压缩
            self.inline_text_edit._compress_image_if_needed(dest)
            
            # 尝试加密存储到数据库
            if self.active_db and self.user_password:
                try:
                    key = generate_key(self.user_password)
                    with open(dest, 'rb') as f:
                        raw_data = f.read()
                    image_data_blob = encrypt_data(raw_data, key)
                    add_image(
                        self.active_db, None, filename, '', '', '', 'image',
                        encrypted=True, image_data=image_data_blob
                    )
                except Exception as e:
                    print(f"加密存储图片失败: {e}")
            self.inline_text_edit.insertHtml(
                f'<img src="{dest}" style="max-width:100%; margin:8px 0;">'
            )

    def _inline_insert_link(self):
        from PyQt5.QtWidgets import QInputDialog
        url, ok = QInputDialog.getText(self, '插入链接', 'URL:')
        if ok and url:
            self.inline_text_edit.insertHtml(f'<a href="{url}" style="color:#3f7bf7;">{url}</a>')

    def _inline_insert_table(self):
        self.inline_text_edit.insertHtml(
            '<table border="1" style="border-collapse:collapse; width:100%; margin:8px 0;">'
            '<tr><td style="padding:6px; border:1px solid #444;">&nbsp;</td>'
            '<td style="padding:6px; border:1px solid #444;">&nbsp;</td></tr>'
            '<tr><td style="padding:6px; border:1px solid #444;">&nbsp;</td>'
            '<td style="padding:6px; border:1px solid #444;">&nbsp;</td></tr>'
            '</table>'
        )

    def _inline_insert_hr(self):
        self.inline_text_edit.insertHtml('<hr style="border:1px solid #444; margin:8px 0;">')

    def _inline_insert_code_block(self):
        self.inline_text_edit.insertHtml(
            '<pre style="background:#1a1a1e; color:#c8c8cc; padding:12px; '
            'border-radius:4px; font-family:Consolas,monospace; font-size:13px; '
            'overflow-x:auto; margin:8px 0;">代码</pre>'
        )

    def _inline_toggle_superscript(self):
        """切换上标"""
        fmt = self.inline_text_edit.currentCharFormat()
        fmt.setVerticalAlignment(QTextCharFormat.AlignSuperScript if fmt.verticalAlignment() != QTextCharFormat.AlignSuperScript else QTextCharFormat.AlignNormal)
        self.inline_text_edit.setCurrentCharFormat(fmt)

    def _inline_toggle_subscript(self):
        """切换下标"""
        fmt = self.inline_text_edit.currentCharFormat()
        fmt.setVerticalAlignment(QTextCharFormat.AlignSubScript if fmt.verticalAlignment() != QTextCharFormat.AlignSubScript else QTextCharFormat.AlignNormal)
        self.inline_text_edit.setCurrentCharFormat(fmt)

    def _inline_insert_code(self):
        """插入行内代码"""
        cursor = self.inline_text_edit.textCursor()
        selected_text = cursor.selectedText()
        if selected_text:
            cursor.insertHtml(f'<code style="background:#333; color:#e0e0e4; padding:2px 6px; border-radius:3px; font-family:Consolas,monospace;">{selected_text}</code>')
        else:
            cursor.insertHtml('<code style="background:#333; color:#e0e0e4; padding:2px 6px; border-radius:3px; font-family:Consolas,monospace;">代码</code>')

    def _inline_insert_image_from_db(self):
        """从数据库选择图片插入"""
        if not self.active_db:
            QMessageBox.warning(self, '提示', '未选择数据库')
            return
        conn = sqlite3.connect(self.active_db)
        c = conn.cursor()
        c.execute("SELECT id, filename FROM images WHERE image_type='image'")
        rows = c.fetchall()
        conn.close()
        if not rows:
            QMessageBox.information(self, '提示', '数据库中没有图片')
            return
        # 显示图片选择对话框
        from PyQt5.QtWidgets import QDialog, QListWidget, QPushButton, QVBoxLayout, QLabel
        dlg = QDialog(self)
        dlg.setWindowTitle('选择图片')
        dlg.resize(300, 400)
        layout = QVBoxLayout(dlg)
        list_widget = QListWidget()
        layout.addWidget(list_widget)
        for row in rows:
            list_widget.addItem(row[1])
        btn_ok = QPushButton('确定')
        btn_ok.clicked.connect(dlg.accept)
        layout.addWidget(btn_ok)
        if dlg.exec_() == QDialog.Accepted and list_widget.currentItem():
            filename = list_widget.currentItem().text()
            img_path = os.path.join(APP_DIR, 'images', filename)
            if os.path.exists(img_path):
                self.inline_text_edit.insertHtml(
                    f'<img src="{img_path}" style="max-width:100%; margin:8px 0;">'
                )
            else:
                # 尝试从数据库解密
                if self.user_password:
                    try:
                        key = generate_key(self.user_password)
                        conn = sqlite3.connect(self.active_db)
                        c = conn.cursor()
                        c.execute("SELECT image_data FROM images WHERE filename=?", (filename,))
                        row = c.fetchone()
                        conn.close()
                        if row and row[0]:
                            raw_data = decrypt_data(row[0], key)
                            temp_path = os.path.join(tempfile.gettempdir(), filename)
                            with open(temp_path, 'wb') as f:
                                f.write(raw_data)
                            self.inline_text_edit.insertHtml(
                                f'<img src="{temp_path}" style="max-width:100%; margin:8px 0;">'
                            )
                    except Exception as e:
                        print(f"从数据库加载图片失败: {e}")

    def _inline_image_align(self, align):
        """调整图片对齐方式"""
        cursor = self.inline_text_edit.textCursor()
        if align == 'left':
            cursor.insertHtml('<p style="text-align:left; margin:8px 0;">')
        elif align == 'center':
            cursor.insertHtml('<p style="text-align:center; margin:8px 0;">')
        elif align == 'right':
            cursor.insertHtml('<p style="text-align:right; margin:8px 0;">')
        elif align == 'full':
            cursor.insertHtml('<p style="text-align:center; margin:8px 0;">')

    def _inline_image_resize(self, factor):
        """调整图片大小"""
        # QTextEdit中直接调整图片大小比较复杂，这里简单插入一个带有缩放的图片占位
        cursor = self.inline_text_edit.textCursor()
        cursor.insertHtml(
            f'<img src="" style="transform: scale({factor}); margin:8px 0;" '
            f'width="200" height="150" alt="图片">'
        )

    def _inline_insert_todo(self):
        """插入待办事项（Trilium Notes 风格）"""
        cursor = self.inline_text_edit.textCursor()
        cursor.insertHtml('<span style="color:#3f7bf7; font-family:Segoe UI Emoji;">☐</span> <span style="color:#c8c8cc;">待办事项</span><br>')

    def _inline_insert_note(self):
        """插入嵌入笔记引用（Trilium Notes 风格）"""
        from PyQt5.QtWidgets import QInputDialog
        note_name, ok = QInputDialog.getText(self, '嵌入笔记', '笔记名称:')
        if ok and note_name:
            cursor = self.inline_text_edit.textCursor()
            cursor.insertHtml(
                f'<div style="background:#2a2a30; border-left:3px solid #3f7bf7; '
                f'padding:8px 12px; margin:8px 0; border-radius:0 4px 4px 0;">'
                f'<span style="color:#3f7bf7; font-weight:bold;">📎 {note_name}</span>'
                f'<p style="color:#888; margin:4px 0 0 0; font-size:12px;">嵌入笔记内容将在此显示</p>'
                f'</div>'
            )

    def _show_inline_context_menu(self, pos):
        """显示内嵌编辑器右键菜单（支持剪贴板粘贴图片和表格）"""
        from PyQt5.QtWidgets import QMenu
        menu = QMenu(self)
        menu.setStyleSheet("QMenu { background-color: #2a2a30; color: #c8c8cc; border: 1px solid #444; padding: 4px; }"
                          "QMenu::item { padding: 6px 24px; border-radius: 4px; }"
                          "QMenu::item:selected { background-color: #3f7bf7; }"
                          "QMenu::separator { height: 1px; background: #444; margin: 4px 8px; }")
        act_undo = menu.addAction('撤销')
        act_redo = menu.addAction('重做')
        menu.addSeparator()
        act_cut = menu.addAction('剪切')
        act_copy = menu.addAction('复制')
        act_paste = menu.addAction('粘贴')
        act_paste_image = menu.addAction('粘贴图片')
        act_paste_table = menu.addAction('粘贴表格')
        menu.addSeparator()
        # 格式选项
        act_bold = menu.addAction('加粗')
        act_italic = menu.addAction('斜体')
        act_underline = menu.addAction('下划线')
        act_strikethrough = menu.addAction('删除线')
        menu.addSeparator()
        # 插入选项
        insert_menu = menu.addMenu('插入')
        insert_menu.setStyleSheet(menu.styleSheet())
        act_insert_image = insert_menu.addAction('插入图片')
        act_insert_table = insert_menu.addAction('插入表格')
        act_insert_link = insert_menu.addAction('插入链接')
        act_insert_hr = insert_menu.addAction('插入分隔线')
        act_insert_code = insert_menu.addAction('插入代码块')
        menu.addSeparator()
        act_select_all = menu.addAction('全选')

        action = menu.exec_(self.inline_text_edit.viewport().mapToGlobal(pos))
        if action == act_undo:
            self.inline_text_edit.undo()
        elif action == act_redo:
            self.inline_text_edit.redo()
        elif action == act_cut:
            self.inline_text_edit.cut()
        elif action == act_copy:
            self.inline_text_edit.copy()
        elif action == act_paste:
            self.inline_text_edit.paste()
        elif action == act_paste_image:
            self._inline_paste_image()
        elif action == act_paste_table:
            self._inline_paste_table()
        elif action == act_bold:
            self._inline_toggle_bold()
        elif action == act_italic:
            self._inline_toggle_italic()
        elif action == act_underline:
            self._inline_toggle_underline()
        elif action == act_strikethrough:
            self._inline_toggle_strikethrough()
        elif action == act_insert_image:
            self._inline_insert_image()
        elif action == act_insert_table:
            self._inline_insert_table()
        elif action == act_insert_link:
            self._inline_insert_link()
        elif action == act_insert_hr:
            self._inline_insert_hr()
        elif action == act_insert_code:
            self._inline_insert_code_block()
        elif action == act_select_all:
            self.inline_text_edit.selectAll()

    def _inline_paste_image(self):
        """从剪贴板粘贴图片到内嵌编辑器"""
        clipboard = QApplication.clipboard()
        mime_data = clipboard.mimeData()
        if mime_data.hasImage():
            pixmap = clipboard.pixmap()
            if not pixmap.isNull():
                import uuid
                img_dir = os.path.join(APP_DIR, 'images')
                os.makedirs(img_dir, exist_ok=True)
                filename = str(uuid.uuid4())[:8] + '.png'
                dest = os.path.join(img_dir, filename)
                pixmap.save(dest)
                self.inline_text_edit._compress_image_if_needed(dest)
                self.inline_text_edit.insertHtml(f'<img src="{dest}" style="max-width:100%; margin:8px 0;">')

    def _inline_paste_table(self):
        """从剪贴板粘贴表格到内嵌编辑器"""
        clipboard = QApplication.clipboard()
        mime_data = clipboard.mimeData()
        if mime_data.hasText():
            text = mime_data.text()
            rows = text.strip().split('\n')
            if len(rows) > 0:
                html = '<table border="1" style="border-collapse:collapse; margin:8px 0;">'
                for row in rows:
                    cells = row.split('\t')
                    html += '<tr>'
                    for cell in cells:
                        html += f'<td style="padding:4px 8px; border:1px solid #444;">{cell}</td>'
                    html += '</tr>'
                html += '</table>'
                self.inline_text_edit.insertHtml(html)

    # ── 编辑模式切换 ──
    def _enter_edit_mode(self, record_id=None):
        """进入编辑模式"""
        self._is_edit_mode = True
        self._editing_record_id = record_id

        if record_id:
            # 加载现有记录
            conn = sqlite3.connect(self.active_db)
            c = conn.cursor()
            if self._current_record_type == 'medical':
                c.execute("SELECT title, content FROM medical_records WHERE id=?", (record_id,))
            elif self._current_record_type == 'anatomy':
                c.execute("SELECT title, content FROM anatomy_records WHERE id=?", (record_id,))
            else:  # imaging
                c.execute("SELECT title, content FROM imaging_records WHERE id=?", (record_id,))
            row_data = c.fetchone()
            conn.close()
            if row_data:
                self.inline_title_input.setText(row_data[0] or '')
                self.inline_text_edit.setHtml(row_data[1] or '')
            else:
                self.inline_title_input.setText('')
                self.inline_text_edit.setHtml('')
        else:
            self.inline_title_input.setText('')
            self.inline_text_edit.setHtml('')

        self.inline_save_indicator.setText('编辑中...')
        self.inline_save_indicator.setStyleSheet("color: #f0a030; font-size: 11px; background: transparent;")
        self.record_stack.setCurrentIndex(1)
        self.inline_title_input.setFocus()
        # 禁用列表切换避免丢失编辑，但保持列表可见
        self.record_list.setEnabled(False)
        # 确保笔记列表面板可见且不被遮盖
        self.record_list_panel.setVisible(True)
        self.record_splitter.setSizes([160, max(400, self.record_splitter.width() - 160)])

    def _exit_edit_mode(self, save=True):
        """退出编辑模式"""
        if save and self._editing_record_id is not None:
            self._inline_save_current()
        self._is_edit_mode = False
        self._editing_record_id = None
        self._inline_save_timer.stop()
        self.record_stack.setCurrentIndex(0)
        self.record_list.setEnabled(True)
        self._image_edit_popup.hide()
        self._refresh_content()

    def _inline_save_current(self):
        """保存当前编辑内容到数据库"""
        if self._editing_record_id is None:
            return
        title = self.inline_title_input.text().strip() or '无标题'
        content = self.inline_text_edit.toHtml()
        try:
            conn = sqlite3.connect(self.active_db)
            c = conn.cursor()
            if self._current_record_type == 'medical':
                c.execute("UPDATE medical_records SET title=?, content=? WHERE id=?",
                          (title, content, self._editing_record_id))
            elif self._current_record_type == 'anatomy':
                c.execute("UPDATE anatomy_records SET title=?, content=? WHERE id=?",
                          (title, content, self._editing_record_id))
            else:  # imaging
                c.execute("UPDATE imaging_records SET title=?, content=? WHERE id=?",
                          (title, content, self._editing_record_id))
            conn.commit()
            conn.close()
            self.inline_save_indicator.setText('已保存')
            self.inline_save_indicator.setStyleSheet("color: #4ade80; font-size: 11px; background: transparent;")
        except Exception as e:
            self.inline_save_indicator.setText(f'保存失败: {e}')
            self.inline_save_indicator.setStyleSheet("color: #e64a3a; font-size: 11px; background: transparent;")

    def _inline_auto_save(self):
        """内嵌编辑器自动保存"""
        self._inline_save_current()

    def keyPressEvent(self, event):
        if self._is_edit_mode and event.key() == Qt.Key_S and event.modifiers() & Qt.ControlModifier:
            self._inline_save_timer.stop()
            self._inline_save_current()
            event.accept()
        else:
            super().keyPressEvent(event)

    # ── 标签切换 ──
    def switch_tab(self, idx):
        # 如果正在编辑，先保存
        if self._is_edit_mode:
            self._exit_edit_mode(save=True)
        self.current_tab = idx
        for i, btn in enumerate(self.tab_buttons):
            btn.setProperty('active', i == idx)
            btn.style().unpolish(btn)
            btn.style().polish(btn)
        # 临床与诊断、影像所见及标准报告参考使用 content_browser
        is_simple = (idx == 0 or idx == 1)
        # 影像所见、医学资料、影像解剖使用 record_splitter
        is_record = (idx == 2 or idx == 3 or idx == 4)
        self.content_browser.setVisible(is_simple)
        self.record_splitter.setVisible(is_record)
        self._refresh_content()

    def load_disease(self, disease_id, db_path):
        self.current_disease_id = disease_id
        self.active_db = db_path
        data = get_disease(db_path, disease_id)
        if data:
            self.disease_data = data
        else:
            self.disease_data = None
        self._refresh_content()

    def _refresh_content(self):
        if not self.disease_data:
            self.content_browser.setHtml('<p style="color:#666;">请从左侧选择一个疾病</p>')
            self.record_content_browser.setHtml('<p style="color:#666;">请从左侧选择一个疾病</p>')
            return
        d = self.disease_data
        if self.current_tab == 0:
            self._show_clinical_tab(d)
        elif self.current_tab == 1:
            self._show_report_tab(d)
        elif self.current_tab == 2:
            self._show_imaging_tab()
        elif self.current_tab == 3:
            self._show_medical_tab()
        elif self.current_tab == 4:
            self._show_anatomy_tab()

    def _html_heading(self, text, level=2):
        colors = {1: '#3f7bf7', 2: '#5a91ff', 3: '#7aadff'}
        color = colors.get(level, '#5a91ff')
        return f'<h{level} style="color:{color}; margin-top:16px; margin-bottom:8px;">{text}</h{level}>'

    def _html_body(self, text):
        return f'<div style="color:#c8c8cc; line-height:1.8; white-space:pre-wrap;">{text}</div>'

    def _show_clinical_tab(self, d):
        html = f'''
        <h1 style="color:#3f7bf7; margin-bottom:4px;">{d.get("name_cn","")}</h1>
        <p style="color:#888890; font-style:italic; margin-top:0;">{d.get("name_en","")}</p>
        <p style="color:#888890;">系统: {d.get("system","")} &nbsp;|&nbsp; 分类: {d.get("category","")}</p>
        <hr style="border:1px solid #2a2a30;">
        {self._html_heading("临床表现")}
        {self._html_body(d.get("clinical",""))}
        {self._html_heading("诊断要点")}
        {self._html_body(d.get("diagnosis",""))}
        {self._html_heading("鉴别诊断")}
        {self._html_body(d.get("differential_diagnosis",""))}
        {self._html_heading("治疗原则")}
        {self._html_body(d.get("treatment",""))}
        '''
        self.content_browser.setHtml(html)

    def _show_report_tab(self, d):
        html = ''
        # 影像所见字段分区显示
        imaging_findings = [
            ('X线所见', d.get('xray_finding', '')),
            ('CT所见', d.get('ct_finding', '')),
            ('MRI所见', d.get('mri_finding', '')),
            ('PET所见', d.get('pet_finding', '')),
        ]
        has_finding = False
        for label, text in imaging_findings:
            if text and text.strip():
                has_finding = True
                html += self._html_heading(label, 2)
                html += self._html_body(text)
        if has_finding:
            html += '<hr style="border:1px solid #2a2a30; margin:16px 0;">'
        # 标准报告模板
        html += self._html_heading("影像报告模板")
        html += self._html_body(d.get("report_template", ""))
        self.content_browser.setHtml(html)

    def _show_imaging_tab(self):
        """影像相关病例标签页 - 使用记录列表方式"""
        self._current_record_type = 'imaging'
        self.record_list.clear()
        if not self.current_disease_id:
            self.record_content_browser.setHtml('<p style="color:#666;">请先选择疾病</p>')
            return
        conn = sqlite3.connect(self.active_db)
        c = conn.cursor()
        c.execute("SELECT id, title, content, COALESCE(parent_id,0) FROM imaging_records WHERE disease_id=? ORDER BY id DESC",
                  (self.current_disease_id,))
        rows = c.fetchall()
        conn.close()
        self._current_records = [(r[0], r[1], r[2]) for r in rows]
        if not rows:
            self.record_content_browser.setHtml('<p style="color:#666;">暂无影像资料，点击笔记列表标题栏菜单中的"添加笔记"创建</p>')
            return
        self._populate_record_list(rows)

    def _show_medical_tab(self):
        self._current_record_type = 'medical'
        self.record_list.clear()
        if not self.current_disease_id:
            self.record_content_browser.setHtml('<p style="color:#666;">请先选择疾病</p>')
            return
        conn = sqlite3.connect(self.active_db)
        c = conn.cursor()
        c.execute("SELECT id, title, content, COALESCE(parent_id,0) FROM medical_records WHERE disease_id=? ORDER BY id DESC",
                  (self.current_disease_id,))
        rows = c.fetchall()
        conn.close()
        self._current_records = [(r[0], r[1], r[2]) for r in rows]
        if not rows:
            self.record_content_browser.setHtml('<p style="color:#666;">暂无医学资料，点击"添加图文"创建</p>')
            return
        self._populate_record_list(rows)

    def _show_anatomy_tab(self):
        self._current_record_type = 'anatomy'
        self.record_list.clear()
        conn = sqlite3.connect(self.active_db)
        c = conn.cursor()
        c.execute("SELECT id, title, content, COALESCE(parent_id,0) FROM anatomy_records ORDER BY id DESC")
        rows = c.fetchall()
        conn.close()
        self._current_records = [(r[0], r[1], r[2]) for r in rows]
        if not rows:
            self.record_content_browser.setHtml('<p style="color:#666;">暂无影像解剖图谱资料，点击"添加图文"创建</p>')
            return
        self._populate_record_list(rows)

    def _populate_record_list(self, rows):
        """填充笔记列表，支持层级显示和+按钮"""
        # 构建父子关系
        children_map = {}  # parent_id -> [(id, title)]
        root_items = []    # [(id, title)]
        for r in rows:
            rid, title, content, parent_id = r[0], r[1], r[2], r[3]
            if parent_id and parent_id != 0:
                children_map.setdefault(parent_id, []).append((rid, title))
            else:
                root_items.append((rid, title))
        
        def add_items(items, depth=0):
            for rid, title in items:
                item = QListWidgetItem()
                indent = '  ' * depth
                has_children = rid in children_map
                prefix = f'{indent}▶ ' if has_children else f'{indent}'
                item.setText(f'{prefix}{title or "无标题"}')
                item.setData(Qt.UserRole, rid)
                item.setData(Qt.UserRole + 1, depth)
                # 子笔记缩进后字体稍小
                if depth > 0:
                    f = item.font()
                    f.setPointSize(11)
                    item.setFont(f)
                self.record_list.addItem(item)
                if has_children:
                    add_items(children_map[rid], depth + 1)
        
        add_items(root_items)
        if self.record_list.count() > 0:
            self.record_list.setCurrentRow(0)

    def _on_record_list_clicked(self, row):
        """点击笔记列表项，在右侧显示内容"""
        if self._is_edit_mode:
            return  # 编辑模式下不切换
        if row < 0:
            return
        item = self.record_list.item(row)
        if not item:
            return
        rid = item.data(Qt.UserRole)
        conn = sqlite3.connect(self.active_db)
        c = conn.cursor()
        if self._current_record_type == 'medical':
            c.execute("SELECT title, content FROM medical_records WHERE id=?", (rid,))
        elif self._current_record_type == 'anatomy':
            c.execute("SELECT title, content FROM anatomy_records WHERE id=?", (rid,))
        else:  # imaging
            c.execute("SELECT title, content FROM imaging_records WHERE id=?", (rid,))
        row_data = c.fetchone()
        conn.close()
        if row_data:
            title, content = row_data
            html = f'<h2 style="color:#5a91ff; margin-bottom:12px;">{title or "无标题"}</h2>'
            html += f'<div style="color:#c8c8cc; line-height:1.8;">{content or ""}</div>'
            self.record_content_browser.setHtml(html)

    def _on_record_list_double_clicked(self, item):
        """双击笔记列表项，进入编辑模式"""
        if self._is_edit_mode:
            return
        if not item:
            return
        rid = item.data(Qt.UserRole)
        if rid:
            self._enter_edit_mode(record_id=rid)

    def _add_record(self, parent_id=0):
        """添加新记录并进入编辑模式"""
        if not self.current_disease_id:
            QMessageBox.warning(self, '提示', '请先选择一个疾病')
            return
        conn = sqlite3.connect(self.active_db)
        c = conn.cursor()
        if self._current_record_type == 'medical':
            c.execute("INSERT INTO medical_records (disease_id, parent_id, title, content) VALUES (?, ?, ?, ?)",
                      (self.current_disease_id, parent_id, '新建笔记', ''))
        elif self._current_record_type == 'anatomy':
            c.execute("INSERT INTO anatomy_records (parent_id, title, content) VALUES (?, ?, ?)",
                      (parent_id, '新建笔记', ''))
        else:  # imaging
            c.execute("INSERT INTO imaging_records (disease_id, title, content) VALUES (?, ?, ?)",
                      (self.current_disease_id, '新建笔记', ''))
        conn.commit()
        new_id = c.lastrowid
        conn.close()
        # 刷新列表
        self._refresh_content()
        # 在列表中选中新建的记录
        for i in range(self.record_list.count()):
            if self.record_list.item(i).data(Qt.UserRole) == new_id:
                self.record_list.setCurrentRow(i)
                break
        # 进入编辑模式
        self._enter_edit_mode(record_id=new_id)

    def _show_record_item_menu(self, pos):
        """显示笔记列表项右键菜单（支持添加子笔记）"""
        item = self.record_list.itemAt(pos)
        menu = QMenu(self)
        menu.setStyleSheet("""
            QMenu { background-color: #2a2a30; color: #c8c8cc; border: 1px solid #444; padding: 4px; }
            QMenu::item { padding: 6px 24px; border-radius: 4px; }
            QMenu::item:selected { background-color: #3f7bf7; }
        """)
        act_add_child = menu.addAction('+ 添加子笔记')
        menu.addSeparator()
        act_edit = menu.addAction('编辑')
        act_rename = menu.addAction('重命名')
        menu.addSeparator()
        act_delete = menu.addAction('删除')
        
        action = menu.exec_(self.record_list.mapToGlobal(pos))
        if action == act_add_child and item:
            parent_id = item.data(Qt.UserRole)
            self._add_record(parent_id=parent_id)
        elif action == act_edit and item:
            record_id = item.data(Qt.UserRole)
            self._enter_edit_mode(record_id=record_id)
        elif action == act_rename and item:
            self._rename_record_item(item)
        elif action == act_delete and item:
            self._delete_record_item(item)

    def _rename_record_item(self, item):
        """重命名笔记"""
        from PyQt5.QtWidgets import QInputDialog
        current_title = item.text().replace('▶ ', '').strip()
        new_title, ok = QInputDialog.getText(self, '重命名', '输入新标题:', text=current_title)
        if ok and new_title.strip():
            record_id = item.data(Qt.UserRole)
            conn = sqlite3.connect(self.active_db)
            c = conn.cursor()
            if self._current_record_type == 'medical':
                c.execute("UPDATE medical_records SET title=? WHERE id=?", (new_title.strip(), record_id))
            elif self._current_record_type == 'anatomy':
                c.execute("UPDATE anatomy_records SET title=? WHERE id=?", (new_title.strip(), record_id))
            else:
                c.execute("UPDATE imaging_records SET title=? WHERE id=?", (new_title.strip(), record_id))
            conn.commit()
            conn.close()
            self._refresh_content()

    def _delete_record_item(self, item):
        """删除笔记及其子笔记"""
        reply = QMessageBox.question(self, '确认删除', '确定要删除这条笔记及其所有子笔记吗？',
                                     QMessageBox.Yes | QMessageBox.No)
        if reply == QMessageBox.No:
            return
        record_id = item.data(Qt.UserRole)
        conn = sqlite3.connect(self.active_db)
        c = conn.cursor()
        if self._current_record_type == 'medical':
            c.execute("DELETE FROM medical_records WHERE id=? OR parent_id=?", (record_id, record_id))
        elif self._current_record_type == 'anatomy':
            c.execute("DELETE FROM anatomy_records WHERE id=? OR parent_id=?", (record_id, record_id))
        else:
            c.execute("DELETE FROM imaging_records WHERE id=?", (record_id,))
        conn.commit()
        conn.close()
        self._refresh_content()

    def _edit_record(self):
        """编辑选中的记录 - 进入内嵌编辑模式"""
        current_item = self.record_list.currentItem()
        if not current_item:
            QMessageBox.warning(self, '提示', '请先选择要编辑的笔记')
            return
        rid = current_item.data(Qt.UserRole)
        self._enter_edit_mode(record_id=rid)

    def _delete_record(self):
        """删除选中的记录"""
        current_item = self.record_list.currentItem()
        if not current_item:
            QMessageBox.warning(self, '提示', '请先选择要删除的笔记')
            return
        rid = current_item.data(Qt.UserRole)
        title = current_item.text()
        reply = QMessageBox.question(
            self, '确认删除', f'确定要删除笔记"{title}"吗？',
            QMessageBox.Yes | QMessageBox.No, QMessageBox.No
        )
        if reply == QMessageBox.Yes:
            conn = sqlite3.connect(self.active_db)
            c = conn.cursor()
            if self._current_record_type == 'medical':
                c.execute("DELETE FROM medical_records WHERE id=?", (rid,))
            elif self._current_record_type == 'anatomy':
                c.execute("DELETE FROM anatomy_records WHERE id=?", (rid,))
            else:  # imaging
                c.execute("DELETE FROM imaging_records WHERE id=?", (rid,))
            conn.commit()
            conn.close()
            self._refresh_content()

    def _detach_tab(self, tab_idx):
        """将指定标签页脱离为独立窗口"""
        if tab_idx in self._detached_windows and self._detached_windows[tab_idx] is not None:
            # 已脱离，激活窗口
            w = self._detached_windows[tab_idx]
            w.showNormal()
            w.activateWindow()
            w.raise_()
            return
        tab_names = ['临床与诊断', '影像所见及标准报告参考', '影像相关病例', '医学资料', '影像解剖图谱与资料']
        title = tab_names[tab_idx] if tab_idx < len(tab_names) else '详情'
        if self.disease_data:
            title += f' - {self.disease_data.get("name_cn", "")}'
        dialog = _DetachedTabWindow(self, tab_idx, title, self.current_disease_id,
                                     self.active_db, self.disease_data,
                                     self.user_info, self.user_password)
        dialog._on_close_callback = lambda idx=tab_idx: self._on_detached_closed(idx)
        self._detached_windows[tab_idx] = dialog
        # 更新图钉按钮状态
        self.tab_buttons[tab_idx].setEnabled(False)
        self.tab_pin_btns[tab_idx].set_pinned(True)
        dialog.show()
        # 如果当前显示的是被脱离的标签，切换到第一个可用标签
        if self.current_tab == tab_idx:
            for i in range(len(self.tab_buttons)):
                if self.tab_buttons[i].isEnabled():
                    self.switch_tab(i)
                    break

    def _on_detached_closed(self, tab_idx):
        """脱离窗口关闭后恢复标签按钮"""
        if tab_idx in self._detached_windows:
            self._detached_windows[tab_idx] = None
        if tab_idx < len(self.tab_buttons):
            self.tab_buttons[tab_idx].setEnabled(True)
            self.tab_pin_btns[tab_idx].set_pinned(False)


class _DetachedTabWindow(QDialog):
    """脱离主窗口的独立标签窗口"""
    def __init__(self, parent_panel, tab_idx, title, disease_id, db_path, disease_data,
                 user_info=None, user_password=None):
        super().__init__(parent_panel)
        self.parent_panel = parent_panel
        self.tab_idx = tab_idx
        self.disease_id = disease_id
        self.active_db = db_path
        self.disease_data = disease_data
        self.user_info = user_info
        self.user_password = user_password
        self._is_edit_mode = False
        self._editing_record_id = None
        self._current_record_type = 'medical'
        self._current_records = []

        self.setWindowTitle(title)
        self.resize(800, 600)
        self.setWindowFlags(Qt.Window | Qt.FramelessWindowHint)

        layout = QVBoxLayout(self)
        layout.setContentsMargins(0, 0, 0, 0)
        layout.setSpacing(0)

        # 自定义标题栏（带返回按钮）
        self.title_bar = _CustomTitleBar(self, title, show_return=True)
        self.title_bar.return_clicked.connect(self._return_to_main)
        layout.addWidget(self.title_bar)

        # 内容区域
        if tab_idx <= 1:
            # 临床与诊断 / 影像所见及标准报告参考：简单文本浏览
            self.content_browser = QTextBrowser()
            self.content_browser.setOpenExternalLinks(True)
            self.content_browser.setStyleSheet("border: none; padding: 16px;")
            layout.addWidget(self.content_browser, stretch=1)
            self._load_simple_content()
        else:
            # 影像所见 / 医学资料 / 影像解剖：列表+内容
            self._build_record_panel(layout)
            if tab_idx == 2:
                self._current_record_type = 'imaging'
            elif tab_idx == 3:
                self._current_record_type = 'medical'
            else:
                self._current_record_type = 'anatomy'
            self._load_records()

        # QSizeGrip - 右下角调整大小手柄
        self._size_grip = QSizeGrip(self)
        self._size_grip.setFixedSize(16, 16)
        self._size_grip.setStyleSheet("background: transparent;")

    def resizeEvent(self, event):
        super().resizeEvent(event)
        # 确保 QSizeGrip 在右下角
        if hasattr(self, '_size_grip'):
            self._size_grip.move(self.width() - 16, self.height() - 16)

    def _return_to_main(self):
        """返回主窗口"""
        # 保存编辑
        if self._is_edit_mode:
            self._exit_edit_mode(save=True)
        # 关闭窗口并通知主面板
        self.close()

    def _load_simple_content(self):
        """加载简单文本内容（临床与诊断、影像所见及标准报告参考）"""
        d = self.disease_data
        if not d:
            self.content_browser.setHtml('<p style="color:#666;">无数据</p>')
            return
        if self.tab_idx == 0:
            html = f'''
            <h1 style="color:#3f7bf7; margin-bottom:4px;">{d.get("name_cn","")}</h1>
            <p style="color:#888890; font-style:italic; margin-top:0;">{d.get("name_en","")}</p>
            <p style="color:#888890;">系统: {d.get("system","")} &nbsp;|&nbsp; 分类: {d.get("category","")}</p>
            <hr style="border:1px solid #2a2a30;">
            <h2 style="color:#5a91ff;">临床表现</h2>
            <div style="color:#c8c8cc; line-height:1.8; white-space:pre-wrap;">{d.get("clinical","")}</div>
            <h2 style="color:#5a91ff;">诊断要点</h2>
            <div style="color:#c8c8cc; line-height:1.8; white-space:pre-wrap;">{d.get("diagnosis","")}</div>
            <h2 style="color:#5a91ff;">鉴别诊断</h2>
            <div style="color:#c8c8cc; line-height:1.8; white-space:pre-wrap;">{d.get("differential_diagnosis","")}</div>
            <h2 style="color:#5a91ff;">治疗原则</h2>
            <div style="color:#c8c8cc; line-height:1.8; white-space:pre-wrap;">{d.get("treatment","")}</div>
            '''
        else:
            # 影像所见及标准报告参考：显示影像所见字段 + 报告模板
            html = ''
            imaging_findings = [
                ('X线所见', d.get('xray_finding', '')),
                ('CT所见', d.get('ct_finding', '')),
                ('MRI所见', d.get('mri_finding', '')),
                ('PET所见', d.get('pet_finding', '')),
            ]
            has_finding = False
            for label, text in imaging_findings:
                if text and text.strip():
                    has_finding = True
                    html += f'<h2 style="color:#5a91ff;">{label}</h2>'
                    html += f'<div style="color:#c8c8cc; line-height:1.8; white-space:pre-wrap;">{text}</div>'
            if has_finding:
                html += '<hr style="border:1px solid #2a2a30; margin:16px 0;">'
            html += '<h2 style="color:#5a91ff;">影像报告模板</h2>'
            html += f'<div style="color:#c8c8cc; line-height:1.8; white-space:pre-wrap;">{d.get("report_template","")}</div>'
        self.content_browser.setHtml(html)

    def _build_record_panel(self, parent_layout):
        """构建记录面板（列表+内容+编辑）- Trilium Notes 风格"""
        splitter = QSplitter(Qt.Horizontal)
        splitter.setStyleSheet("QSplitter::handle { background: #2a2a30; width: 1px; }")

        # 左侧列表（紧凑风格）
        list_panel = QWidget()
        list_panel.setStyleSheet("background: transparent;")
        list_layout = QVBoxLayout(list_panel)
        list_layout.setContentsMargins(0, 0, 0, 0)
        list_layout.setSpacing(0)
        # 标题栏 + 菜单
        header = QWidget()
        header.setFixedHeight(28)
        header.setStyleSheet("background: transparent; border-bottom: 1px solid rgba(255,255,255,15);")
        hl = QHBoxLayout(header)
        hl.setContentsMargins(8, 0, 4, 0)
        hl.setSpacing(0)
        hl_label = QLabel('笔记')
        hl_label.setFont(QFont("Microsoft YaHei", 9, QFont.Bold))
        hl_label.setStyleSheet("background: transparent; color: #888890;")
        hl.addWidget(hl_label)
        hl.addStretch()
        menu_btn = QPushButton('⋮')
        menu_btn.setFixedSize(24, 24)
        menu_btn.setFont(QFont("Microsoft YaHei", 12))
        menu_btn.setCursor(Qt.PointingHandCursor)
        menu_btn.setStyleSheet("""
            QPushButton { background: transparent; border: none; color: #666670; border-radius: 3px; }
            QPushButton:hover { background: rgba(255,255,255,15); color: #c8c8cc; }
        """)
        menu_btn.clicked.connect(self._show_record_menu)
        hl.addWidget(menu_btn)
        list_layout.addWidget(header)

        self.record_list = QListWidget()
        self.record_list.currentRowChanged.connect(self._on_record_clicked)
        self.record_list.setStyleSheet("""
            QListWidget { border: none; background: transparent; outline: none; }
            QListWidget::item { padding: 6px 10px; border-bottom: 1px solid rgba(255,255,255,8);
                                color: #a0a0a8; font-size: 12px; }
            QListWidget::item:selected { background-color: rgba(63,123,247,30); color: #e0e0e4; }
            QListWidget::item:hover { background-color: rgba(255,255,255,8); color: #c8c8cc; }
        """)
        list_layout.addWidget(self.record_list)
        list_panel.setFixedWidth(140)
        splitter.addWidget(list_panel)

        # 右侧内容
        self.record_stack = QStackedWidget()
        # 查看模式（优化排版）
        self.content_browser = QTextBrowser()
        self.content_browser.setOpenExternalLinks(True)
        self.content_browser.setStyleSheet("""
            QTextBrowser { border: none; background: transparent; padding: 20px 28px; color: #d0d0d4; }
            QTextBrowser h1 { color: #3f7bf7; margin-bottom: 4px; }
            QTextBrowser h2 { color: #5a91ff; margin-top: 20px; margin-bottom: 8px; }
            QTextBrowser a { color: #3f7bf7; }
        """)
        self.record_stack.addWidget(self.content_browser)
        # 编辑模式（工具固定在侧方）
        edit_widget = QWidget()
        edit_layout = QHBoxLayout(edit_widget)
        edit_layout.setContentsMargins(0, 0, 0, 0)
        edit_layout.setSpacing(0)
        # 左侧工具面板
        self._tool_panel = QWidget()
        self._tool_panel.setFixedWidth(40)
        self._tool_panel.setStyleSheet("background: #222226;")
        tool_panel_layout = QVBoxLayout(self._tool_panel)
        tool_panel_layout.setContentsMargins(4, 8, 4, 8)
        tool_panel_layout.setSpacing(4)
        format_tools = [
            ('bold', '加粗'), ('italic', '斜体'), ('underline', '下划线'), ('strikethrough', '删除线'),
            ('superscript', '上标'), ('subscript', '下标'), ('code', '行内代码'),
        ]
        for fmt_id, tip in format_tools:
            btn = _FormatBtn(fmt_id, fg_color='#c8c8cc')
            btn.setToolTip(tip)
            btn.clicked.connect(lambda checked, fid=fmt_id: self._on_format_action(fid))
            tool_panel_layout.addWidget(btn)
        tool_panel_layout.addStretch()
        align_tools = [('align_left', '左对齐'), ('align_center', '居中'), ('align_right', '右对齐')]
        for fmt_id, tip in align_tools:
            btn = _FormatBtn(fmt_id, fg_color='#c8c8cc')
            btn.setToolTip(tip)
            btn.clicked.connect(lambda checked, fid=fmt_id: self._on_format_action(fid))
            tool_panel_layout.addWidget(btn)
        tool_panel_layout.addStretch()
        list_tools = [('bullet_list', '无序列表'), ('ordered_list', '有序列表')]
        for fmt_id, tip in list_tools:
            btn = _FormatBtn(fmt_id, fg_color='#c8c8cc')
            btn.setToolTip(tip)
            btn.clicked.connect(lambda checked, fid=fmt_id: self._on_format_action(fid))
            tool_panel_layout.addWidget(btn)
        tool_panel_layout.addStretch()
        insert_tools = [('insert_image', '插入图片'), ('insert_table', '插入表格'), ('insert_hr', '分隔线'), ('insert_link', '插入链接')]
        for fmt_id, tip in insert_tools:
            btn = _FormatBtn(fmt_id, fg_color='#c8c8cc')
            btn.setToolTip(tip)
            btn.clicked.connect(lambda checked, fid=fmt_id: self._on_block_action(fid))
            tool_panel_layout.addWidget(btn)
        tool_panel_layout.addStretch()
        undo_redo_tools = [('undo', '撤销'), ('redo', '重做')]
        for fmt_id, tip in undo_redo_tools:
            btn = _FormatBtn(fmt_id, fg_color='#c8c8cc')
            btn.setToolTip(tip)
            btn.clicked.connect(lambda checked, fid=fmt_id: self._on_format_action(fid))
            tool_panel_layout.addWidget(btn)
        edit_layout.addWidget(self._tool_panel)
        # 分隔线
        tool_sep = QFrame()
        tool_sep.setFrameShape(QFrame.VLine)
        tool_sep.setStyleSheet("background-color: #2a2a30; width: 1px;")
        edit_layout.addWidget(tool_sep)
        # 右侧编辑区域
        edit_content = QWidget()
        edit_content_layout = QVBoxLayout(edit_content)
        edit_content_layout.setContentsMargins(0, 0, 0, 0)
        edit_content_layout.setSpacing(0)
        # 编辑标题栏（紧凑）
        edit_title_bar = QWidget()
        edit_title_bar.setFixedHeight(42)
        etl = QHBoxLayout(edit_title_bar)
        etl.setContentsMargins(12, 4, 12, 4)
        etl.setSpacing(8)
        self.inline_title_input = QLineEdit()
        self.inline_title_input.setPlaceholderText('输入标题...')
        self.inline_title_input.setFont(QFont("Microsoft YaHei", 13, QFont.Bold))
        self.inline_title_input.setStyleSheet("""
            QLineEdit { background: transparent; border: none; border-bottom: 1px solid #3a3a40;
                        color: #e0e0e4; padding: 2px 0; }
            QLineEdit:focus { border-bottom-color: #3f7bf7; }
        """)
        etl.addWidget(self.inline_title_input, stretch=1)
        self.save_indicator = QLabel('已保存')
        self.save_indicator.setStyleSheet("color: #4ade80; font-size: 11px; background: transparent;")
        etl.addWidget(self.save_indicator)
        btn_save = QPushButton('保存')
        btn_save.setFixedHeight(26)
        btn_save.setFixedWidth(52)
        btn_save.setObjectName('accentBtn')
        btn_save.clicked.connect(lambda checked: self._exit_edit_mode(save=True))
        etl.addWidget(btn_save)
        btn_cancel = QPushButton('取消')
        btn_cancel.setFixedHeight(26)
        btn_cancel.setFixedWidth(52)
        btn_cancel.clicked.connect(lambda checked: self._exit_edit_mode(save=False))
        etl.addWidget(btn_cancel)
        edit_content_layout.addWidget(edit_title_bar)
        # 分隔线
        title_sep = QFrame()
        title_sep.setFrameShape(QFrame.HLine)
        title_sep.setStyleSheet("background-color: #2a2a30; max-height: 1px;")
        edit_content_layout.addWidget(title_sep)
        # 编辑器（支持图片选中、缩放、拖动）
        self.inline_text_edit = _ImageEditableTextEdit()
        self.inline_text_edit.setAcceptRichText(True)
        self.inline_text_edit.setStyleSheet("""
            QTextEdit { border: none; background: transparent; padding: 16px 24px; color: #d0d0d4; }
        """)
        edit_content_layout.addWidget(self.inline_text_edit, stretch=1)
        edit_layout.addWidget(edit_content, stretch=1)
        self.record_stack.addWidget(edit_widget)

        splitter.addWidget(self.record_stack)
        splitter.setStretchFactor(0, 0)
        splitter.setStretchFactor(1, 1)
        parent_layout.addWidget(splitter)

    def _load_records(self):
        """加载记录列表"""
        self.record_list.clear()
        if not self.disease_id and self._current_record_type != 'anatomy':
            self.content_browser.setHtml('<p style="color:#666;">请先选择疾病</p>')
            return
        # 加载记录
        conn = sqlite3.connect(self.active_db)
        c = conn.cursor()
        if self._current_record_type == 'medical':
            c.execute("SELECT id, title, content FROM medical_records WHERE disease_id=? ORDER BY id DESC",
                      (self.disease_id,))
        elif self._current_record_type == 'anatomy':
            c.execute("SELECT id, title, content FROM anatomy_records ORDER BY id DESC")
        else:
            c.execute("SELECT id, title, content FROM imaging_records WHERE disease_id=? ORDER BY id DESC",
                      (self.disease_id,))
        rows = c.fetchall()
        conn.close()
        self._current_records = [(r[0], r[1], r[2]) for r in rows]
        for rid, title, content in rows:
            item = QListWidgetItem(title or '无标题')
            item.setData(Qt.UserRole, rid)
            self.record_list.addItem(item)
        if self.record_list.count() > 0:
            self.record_list.setCurrentRow(0)
        else:
            self.content_browser.setHtml('<p style="color:#666;">暂无记录</p>')

    def _on_record_clicked(self, row):
        """点击记录列表项"""
        if self._is_edit_mode:
            return
        if row < 0:
            return
        item = self.record_list.item(row)
        if not item:
            return
        rid = item.data(Qt.UserRole)
        conn = sqlite3.connect(self.active_db)
        c = conn.cursor()
        if self._current_record_type == 'medical':
            c.execute("SELECT title, content FROM medical_records WHERE id=?", (rid,))
        elif self._current_record_type == 'anatomy':
            c.execute("SELECT title, content FROM anatomy_records WHERE id=?", (rid,))
        else:
            c.execute("SELECT title, content FROM imaging_records WHERE id=?", (rid,))
        row_data = c.fetchone()
        conn.close()
        if row_data:
            title, content = row_data
            html = f'<h2 style="color:#5a91ff; margin-bottom:12px;">{title or "无标题"}</h2>'
            html += f'<div style="color:#c8c8cc; line-height:1.8;">{content or ""}</div>'
            self.content_browser.setHtml(html)

    def _show_record_menu(self, checked=False):
        """显示记录操作菜单"""
        menu = QMenu(self)
        menu.setStyleSheet("""
            QMenu { background-color: #2a2a30; color: #c8c8cc; border: 1px solid #444; padding: 4px; }
            QMenu::item { padding: 6px 24px; border-radius: 4px; }
            QMenu::item:selected { background-color: #3f7bf7; }
        """)
        act_add = menu.addAction('添加笔记')
        act_edit = menu.addAction('编辑选中')
        menu.addSeparator()
        act_delete = menu.addAction('删除选中')
        action = menu.exec_(self.mapToGlobal(QPoint(100, 50)))
        if action == act_add:
            self._add_record()
        elif action == act_edit:
            self._edit_record()
        elif action == act_delete:
            self._delete_record()

    def _add_record(self):
        if not self.disease_id and self._current_record_type != 'anatomy':
            QMessageBox.warning(self, '提示', '请先选择一个疾病')
            return
        conn = sqlite3.connect(self.active_db)
        c = conn.cursor()
        if self._current_record_type == 'medical':
            c.execute("INSERT INTO medical_records (disease_id, title, content) VALUES (?, ?, ?)",
                      (self.disease_id, '新建笔记', ''))
        elif self._current_record_type == 'anatomy':
            c.execute("INSERT INTO anatomy_records (title, content) VALUES (?, ?)",
                      ('新建笔记', ''))
        else:
            c.execute("INSERT INTO imaging_records (disease_id, title, content) VALUES (?, ?, ?)",
                      (self.disease_id, '新建笔记', ''))
        conn.commit()
        new_id = c.lastrowid
        conn.close()
        self._load_records()
        for i in range(self.record_list.count()):
            if self.record_list.item(i).data(Qt.UserRole) == new_id:
                self.record_list.setCurrentRow(i)
                break
        self._enter_edit_mode(record_id=new_id)

    def _edit_record(self):
        current_item = self.record_list.currentItem()
        if not current_item:
            QMessageBox.warning(self, '提示', '请先选择要编辑的笔记')
            return
        rid = current_item.data(Qt.UserRole)
        self._enter_edit_mode(record_id=rid)

    def _delete_record(self):
        current_item = self.record_list.currentItem()
        if not current_item:
            QMessageBox.warning(self, '提示', '请先选择要删除的笔记')
            return
        rid = current_item.data(Qt.UserRole)
        title = current_item.text()
        reply = QMessageBox.question(
            self, '确认删除', f'确定要删除笔记"{title}"吗？',
            QMessageBox.Yes | QMessageBox.No, QMessageBox.No
        )
        if reply == QMessageBox.Yes:
            conn = sqlite3.connect(self.active_db)
            c = conn.cursor()
            if self._current_record_type == 'medical':
                c.execute("DELETE FROM medical_records WHERE id=?", (rid,))
            elif self._current_record_type == 'anatomy':
                c.execute("DELETE FROM anatomy_records WHERE id=?", (rid,))
            else:
                c.execute("DELETE FROM imaging_records WHERE id=?", (rid,))
            conn.commit()
            conn.close()
            self._load_records()

    def _enter_edit_mode(self, record_id=None):
        self._is_edit_mode = True
        self._editing_record_id = record_id
        if record_id:
            conn = sqlite3.connect(self.active_db)
            c = conn.cursor()
            if self._current_record_type == 'medical':
                c.execute("SELECT title, content FROM medical_records WHERE id=?", (record_id,))
            elif self._current_record_type == 'anatomy':
                c.execute("SELECT title, content FROM anatomy_records WHERE id=?", (record_id,))
            else:
                c.execute("SELECT title, content FROM imaging_records WHERE id=?", (record_id,))
            row_data = c.fetchone()
            conn.close()
            if row_data:
                self.inline_title_input.setText(row_data[0] or '')
                self.inline_text_edit.setHtml(row_data[1] or '')
            else:
                self.inline_title_input.setText('')
                self.inline_text_edit.setHtml('')
        else:
            self.inline_title_input.setText('')
            self.inline_text_edit.setHtml('')
        self.save_indicator.setText('编辑中...')
        self.save_indicator.setStyleSheet("color: #f0a030; font-size: 11px; background: transparent;")
        self.record_stack.setCurrentIndex(1)
        self.inline_title_input.setFocus()
        self.record_list.setEnabled(False)

    def _exit_edit_mode(self, save=True):
        if save and self._editing_record_id is not None:
            self._save_current()
        self._is_edit_mode = False
        self._editing_record_id = None
        self.record_stack.setCurrentIndex(0)
        self.record_list.setEnabled(True)
        self._load_records()

    def _save_current(self):
        if self._editing_record_id is None:
            return
        title = self.inline_title_input.text().strip() or '无标题'
        content = self.inline_text_edit.toHtml()
        try:
            conn = sqlite3.connect(self.active_db)
            c = conn.cursor()
            if self._current_record_type == 'medical':
                c.execute("UPDATE medical_records SET title=?, content=? WHERE id=?",
                          (title, content, self._editing_record_id))
            elif self._current_record_type == 'anatomy':
                c.execute("UPDATE anatomy_records SET title=?, content=? WHERE id=?",
                          (title, content, self._editing_record_id))
            else:
                c.execute("UPDATE imaging_records SET title=?, content=? WHERE id=?",
                          (title, content, self._editing_record_id))
            conn.commit()
            conn.close()
            self.save_indicator.setText('已保存')
            self.save_indicator.setStyleSheet("color: #4ade80; font-size: 11px; background: transparent;")
        except Exception as e:
            self.save_indicator.setText(f'保存失败: {e}')
            self.save_indicator.setStyleSheet("color: #e64a3a; font-size: 11px; background: transparent;")

    def keyPressEvent(self, event):
        if self._is_edit_mode and event.key() == Qt.Key_S and event.modifiers() & Qt.ControlModifier:
            self._save_current()
            event.accept()
        else:
            super().keyPressEvent(event)

    def closeEvent(self, event):
        if self._is_edit_mode:
            self._exit_edit_mode(save=True)
        if hasattr(self, '_on_close_callback') and self._on_close_callback:
            self._on_close_callback()
        super().closeEvent(event)


# ── 主窗口 ────────────────────────────────────────────────
class MainWindow(QMainWindow):
    def __init__(self, user_info, active_db, user_password=None):
        super().__init__()
        self.user_info = user_info
        self.active_db = active_db
        self.user_password = user_password
        self.current_theme = '暖沙'
        self.setWindowTitle(f'RadAtlas {APP_VERSION} 影像图鉴助手')
        self.setMinimumSize(1000, 650)
        self.resize(1200, 750)
        # 使用无边框窗口 + 自定义标题栏
        self.setWindowFlags(Qt.Window | Qt.FramelessWindowHint)

        # 中央容器
        central = QWidget()
        self.setCentralWidget(central)
        root_layout = QVBoxLayout(central)
        root_layout.setContentsMargins(0, 0, 0, 0)
        root_layout.setSpacing(0)

        # 自定义标题栏
        self.title_bar = _CustomTitleBar(self, f'RadAtlas {APP_VERSION} 影像图鉴助手  |  楚雄州人民医院 医学影像中心 张兴文')
        root_layout.addWidget(self.title_bar)

        # 菜单栏（手动添加到布局中，确保在标题栏之下）
        menubar = QMenuBar(self)
        root_layout.addWidget(menubar)

        # 文件菜单
        file_menu = menubar.addMenu('文件')
        add_disease_action = QAction('添加疾病', self)
        add_disease_action.triggered.connect(self.show_add_disease)
        file_menu.addAction(add_disease_action)

        edit_disease_action = QAction('编辑当前疾病', self)
        edit_disease_action.triggered.connect(self.show_edit_disease)
        file_menu.addAction(edit_disease_action)

        del_disease_action = QAction('删除当前疾病', self)
        del_disease_action.triggered.connect(self.show_delete_disease)
        file_menu.addAction(del_disease_action)

        file_menu.addSeparator()

        add_img_action = QAction('添加图片到当前疾病', self)
        add_img_action.triggered.connect(self.show_add_image)
        file_menu.addAction(add_img_action)

        file_menu.addSeparator()

        user_action = QAction('用户管理', self)
        user_action.triggered.connect(self.show_user_manage)
        file_menu.addAction(user_action)

        file_menu.addSeparator()

        ai_config_action = QAction('AI 配置', self)
        ai_config_action.triggered.connect(self.show_ai_config)
        file_menu.addAction(ai_config_action)

        # AI查询菜单
        ai_menu = menubar.addMenu('AI查询')
        ai_diagnosis_action = QAction('影像诊断查询', self)
        ai_diagnosis_action.triggered.connect(self.show_ai_diagnosis)
        ai_menu.addAction(ai_diagnosis_action)

        ai_kb_config_action = QAction('知识库配置', self)
        ai_kb_config_action.triggered.connect(self.show_kb_config)
        ai_menu.addAction(ai_kb_config_action)

        file_menu.addSeparator()

        export_single_action = QAction('导出当前疾病', self)
        export_single_action.triggered.connect(self._export_current_disease)
        file_menu.addAction(export_single_action)

        export_multi_action = QAction('批量导出疾病', self)
        export_multi_action.triggered.connect(self._export_multi_disease)
        file_menu.addAction(export_multi_action)

        # 主题菜单
        theme_menu = menubar.addMenu('主题')
        for theme_name in THEMES:
            action = QAction(theme_name, self)
            action.triggered.connect(lambda checked, name=theme_name: self.switch_theme(name))
            theme_menu.addAction(action)

        # 中央部件
        # 顶部搜索栏
        top_bar = QHBoxLayout()
        top_bar.setContentsMargins(12, 8, 12, 8)
        top_bar.setSpacing(8)

        logo = QLabel('RadAtlas')
        logo.setFont(QFont("Microsoft YaHei", 14, QFont.Bold))
        logo.setStyleSheet("color: #3f7bf7; background: transparent;")
        logo.setFixedWidth(100)
        top_bar.addWidget(logo)

        # 搜索框 + 自动补全
        self.search_input = QLineEdit()
        self.search_input.setPlaceholderText('搜索疾病...')
        self.search_input.setFixedHeight(36)
        self.search_input.returnPressed.connect(self.do_search)
        self.search_input.textChanged.connect(self.on_search_text_changed)

        self.completer_model = QStringListModel()
        self.completer = QCompleter()
        self.completer.setModel(self.completer_model)
        self.completer.setCaseSensitivity(Qt.CaseInsensitive)
        self.completer.setFilterMode(Qt.MatchContains)
        self.completer.activated.connect(self.on_completer_activated)
        self.search_input.setCompleter(self.completer)

        top_bar.addWidget(self.search_input, stretch=1)

        search_btn = QPushButton('搜索')
        search_btn.setFixedSize(70, 36)
        search_btn.clicked.connect(self.do_search)
        top_bar.addWidget(search_btn)

        clear_btn = QPushButton('清除')
        clear_btn.setObjectName('mutedBtn')
        clear_btn.setFixedSize(70, 36)
        clear_btn.clicked.connect(self.clear_search)
        top_bar.addWidget(clear_btn)

        top_bar.addStretch()

        user_label = QLabel(f'{user_info.get("username", "admin")}')
        user_label.setStyleSheet("color: #888890; font-size: 12px; background: transparent;")
        top_bar.addWidget(user_label)

        logout_btn = QPushButton('退出')
        logout_btn.setObjectName('dangerBtn')
        logout_btn.setFixedSize(60, 36)
        logout_btn.clicked.connect(self.do_logout)
        top_bar.addWidget(logout_btn)

        root_layout.addLayout(top_bar)

        # 分隔线
        line = QFrame()
        line.setFrameShape(QFrame.HLine)
        line.setStyleSheet("background-color: #2a2a30; max-height: 1px;")
        root_layout.addWidget(line)

        # 主内容区
        splitter = QSplitter(Qt.Horizontal)

        # 左侧疾病列表
        left_panel = QWidget()
        left_layout = QVBoxLayout(left_panel)
        left_layout.setContentsMargins(8, 8, 4, 8)
        left_layout.setSpacing(6)

        list_header = QLabel('疾病列表')
        list_header.setStyleSheet("color: #777780; font-size: 12px; font-weight: bold; background: transparent;")
        left_layout.addWidget(list_header)

        self.disease_list = QListWidget()
        self.disease_list.currentItemChanged.connect(self.on_disease_selected)
        self.disease_list.itemDoubleClicked.connect(self.on_disease_double_clicked)
        self.disease_list.setStyleSheet("""
            QListWidget {
                border: none; background: transparent; outline: none;
                font-size: 13px;
            }
            QListWidget::item {
                padding: 6px 10px; border-bottom: 1px solid rgba(255,255,255,6);
                color: #c0c0c8; border-radius: 4px; margin: 1px 2px;
            }
            QListWidget::item:selected {
                background-color: rgba(63,123,247,40); color: #e0e0e4;
            }
            QListWidget::item:hover {
                background-color: rgba(255,255,255,8); color: #e0e0e4;
            }
        """)
        left_layout.addWidget(self.disease_list)

        splitter.addWidget(left_panel)

        # 右侧详情面板
        right_panel = QWidget()
        right_layout = QVBoxLayout(right_panel)
        right_layout.setContentsMargins(4, 8, 8, 8)

        self.detail = DetailPanel(user_info=self.user_info, user_password=self.user_password)
        right_layout.addWidget(self.detail)

        splitter.addWidget(right_panel)

        splitter.setStretchFactor(0, 35)
        splitter.setStretchFactor(1, 65)
        splitter.setSizes([350, 650])

        root_layout.addWidget(splitter, stretch=1)

        # 状态栏
        self.statusBar().showMessage(f'就绪  |  {APP_AUTHOR}  |  {APP_VERSION}')

        # 加载疾病列表
        self.load_disease_list()
        self._update_completer()
        # 初始化标题栏主题
        self.title_bar.update_theme(THEMES[self.current_theme])

        # 迁移现有图片文件到加密数据库存储
        if self.user_password and self.active_db:
            try:
                key = generate_key(self.user_password)
                migrated = migrate_images_to_db(self.active_db, key)
                if migrated > 0:
                    self.statusBar().showMessage(f'已迁移 {migrated} 张图片到加密数据库  |  {APP_AUTHOR}  |  {APP_VERSION}')
            except Exception as e:
                print(f"图片迁移失败: {e}")

    # ── 主题切换 ──
    def switch_theme(self, theme_name):
        if theme_name not in THEMES:
            return
        self.current_theme = theme_name
        t = THEMES[theme_name]
        # 淡入过渡动画
        if not hasattr(self, '_theme_anim'):
            self._theme_opacity = QGraphicsOpacityEffect(self)
            self.setGraphicsEffect(self._theme_opacity)
            self._theme_anim = QPropertyAnimation(self._theme_opacity, b"opacity")
        self._theme_anim.stop()
        self._theme_anim.setDuration(300)
        self._theme_anim.setStartValue(0.6)
        self._theme_anim.setEndValue(1.0)
        self._theme_anim.setEasingCurve(QEasingCurve.InOutCubic)
        QApplication.instance().setStyleSheet(build_stylesheet(t))
        # 更新调色板
        palette = QPalette()
        is_light = theme_name == '浅色经典'
        if is_light:
            palette.setColor(QPalette.Window, QColor(245, 245, 248))
            palette.setColor(QPalette.WindowText, QColor(34, 34, 48))
            palette.setColor(QPalette.Base, QColor(255, 255, 255))
            palette.setColor(QPalette.Text, QColor(34, 34, 48))
        else:
            palette.setColor(QPalette.Window, QColor(t['bg'].lstrip('#')))
            palette.setColor(QPalette.WindowText, QColor(t['fg'].lstrip('#')))
            palette.setColor(QPalette.Base, QColor(t['bg2'].lstrip('#')))
            palette.setColor(QPalette.Text, QColor(t['fg'].lstrip('#')))
        QApplication.instance().setPalette(palette)
        # 更新自定义标题栏
        if hasattr(self, 'title_bar'):
            self.title_bar.update_theme(t)
        self._theme_anim.start()
        self.statusBar().showMessage(f'已切换主题: {theme_name}')

    # ── 搜索自动补全 ──
    def _update_completer(self):
        conn = sqlite3.connect(self.active_db)
        c = conn.cursor()
        c.execute("SELECT name_cn, name_en FROM diseases")
        rows = c.fetchall()
        conn.close()
        names = []
        for cn, en in rows:
            if cn and en:
                names.append(f'{cn} ({en})')
            elif cn:
                names.append(cn)
            elif en:
                names.append(en)
        self.completer_model.setStringList(names)
        # 设置补全弹出列表样式：14px楷体灰蓝色斜体，背景透明
        self.completer.popup().setStyleSheet("""
            QListView {
                background-color: transparent;
                border: 1px solid #3a3e55;
                border-radius: 6px;
                padding: 4px;
                font-family: KaiTi;
                font-style: italic;
                color: #8899bb;
                font-size: 14px;
                outline: none;
            }
            QListView::item {
                padding: 5px 10px;
                border-radius: 4px;
                background-color: rgba(20, 22, 36, 180);
            }
            QListView::item:selected {
                background-color: #3f7bf7;
                color: #ffffff;
                font-style: italic;
            }
            QListView::item:hover {
                background-color: rgba(42, 46, 69, 200);
            }
        """)
        # 补全弹出淡入动画（延长到500ms）
        self._completer_anim = QPropertyAnimation(self.completer.popup(), b"windowOpacity")
        self._completer_anim.setDuration(500)
        self._completer_anim.setStartValue(0.0)
        self._completer_anim.setEndValue(1.0)
        self._completer_anim.setEasingCurve(QEasingCurve.InOutCubic)

    def on_search_text_changed(self, text):
        # 补全弹出时播放淡入动画
        if text and self.completer_model.stringList():
            try:
                self._completer_anim.start()
            except Exception:
                pass

    def on_completer_activated(self, text):
        # 从 "中文名(英文名)" 格式中提取中文名用于搜索
        search_text = text.split('(')[0].strip() if '(' in text else text
        self.search_input.setText(search_text)
        self.do_search()
        # 自动选中第一个匹配的疾病条目
        if self.disease_list.count() > 0:
            self.disease_list.setCurrentRow(0)

    # ── 疾病列表 ──
    def load_disease_list(self, search=None):
        self.disease_list.clear()
        conn = sqlite3.connect(self.active_db)
        c = conn.cursor()
        if search:
            c.execute("""SELECT id, name_cn, name_en, system FROM diseases
                         WHERE name_cn LIKE ? OR name_en LIKE ? OR clinical LIKE ?
                         OR diagnosis LIKE ? OR system LIKE ? OR category LIKE ?
                         ORDER BY system, name_cn""",
                      (f'%{search}%',) * 6)
        else:
            c.execute("SELECT id, name_cn, name_en, system FROM diseases ORDER BY system, name_cn")
        rows = c.fetchall()
        conn.close()

        if not rows:
            item = QListWidgetItem('暂无数据')
            item.setFlags(item.flags() & ~Qt.ItemIsSelectable & ~Qt.ItemIsEnabled)
            self.disease_list.addItem(item)
            return

        current_system = None
        for did, name_cn, name_en, system in rows:
            if system != current_system:
                current_system = system
                header = QListWidgetItem(f'── {system} ──')
                header.setFlags(header.flags() & ~Qt.ItemIsSelectable & ~Qt.ItemIsEnabled)
                header.setData(Qt.UserRole, None)
                header.setForeground(QColor('#666670'))
                header.setFont(QFont("Microsoft YaHei", 10))
                self.disease_list.addItem(header)

            display = name_cn or name_en or f'疾病 #{did}'
            if name_en and name_cn:
                display = f'{name_cn} ({name_en})'
            item = QListWidgetItem(display)
            item.setData(Qt.UserRole, did)

            # 搜索时高亮匹配项
            if search:
                search_lower = search.lower()
                display_lower = display.lower()
                if search_lower in display_lower:
                    # 高亮匹配项：使用醒目的前景色和背景色
                    t = THEMES.get(self.current_theme, THEMES['深蓝暗夜'])
                    item.setForeground(QColor(t.get('select_fg', '#5a91ff')))
                    item.setBackground(QColor(t.get('select_bg', '#2a3a5a')))
                    item.setFont(QFont("Microsoft YaHei", 13, QFont.Bold))
                else:
                    item.setForeground(QColor('#888890'))
            self.disease_list.addItem(item)

    def on_disease_selected(self, current, previous):
        if current is None:
            return
        disease_id = current.data(Qt.UserRole)
        if disease_id is None:
            return
        self.detail.load_disease(disease_id, self.active_db)
        self.statusBar().showMessage(f'已选择疾病 ID: {disease_id}')

    def on_disease_double_clicked(self, item):
        self.show_edit_disease()

    # ── 添加疾病 ──
    def show_add_disease(self):
        dlg = DiseaseDialog(self)
        if dlg.exec_() == QDialog.Accepted:
            data = dlg.get_data()
            if not data.get('name_cn'):
                QMessageBox.warning(self, '提示', '中文名为必填项')
                return
            try:
                add_disease(self.active_db, data)
                self.load_disease_list()
                self._update_completer()
                self.statusBar().showMessage(f'已添加疾病: {data["name_cn"]}')
            except Exception as e:
                QMessageBox.critical(self, '错误', f'添加失败:\n{e}')

    # ── 编辑疾病 ──
    def show_edit_disease(self):
        item = self.disease_list.currentItem()
        if not item:
            QMessageBox.warning(self, '提示', '请先选择一个疾病')
            return
        disease_id = item.data(Qt.UserRole)
        if disease_id is None:
            return
        data = get_disease(self.active_db, disease_id)
        if not data:
            return
        dlg = DiseaseDialog(self, disease_data=data)
        if dlg.exec_() == QDialog.Accepted:
            new_data = dlg.get_data()
            try:
                update_disease(self.active_db, disease_id, new_data)
                self.load_disease_list()
                self._update_completer()
                self.detail.load_disease(disease_id, self.active_db)
                self.statusBar().showMessage(f'已更新疾病: {new_data.get("name_cn", "")}')
            except Exception as e:
                QMessageBox.critical(self, '错误', f'更新失败:\n{e}')

    # ── 删除疾病 ──
    def show_delete_disease(self):
        item = self.disease_list.currentItem()
        if not item:
            QMessageBox.warning(self, '提示', '请先选择一个疾病')
            return
        disease_id = item.data(Qt.UserRole)
        if disease_id is None:
            return
        if QMessageBox.question(self, '确认', '确定删除该疾病及其所有关联数据？') == QMessageBox.Yes:
            try:
                delete_disease(self.active_db, disease_id)
                self.load_disease_list()
                self._update_completer()
                self.statusBar().showMessage('疾病已删除')
            except Exception as e:
                QMessageBox.critical(self, '错误', f'删除失败:\n{e}')

    # ── 添加图片 ──
    def show_add_image(self):
        item = self.disease_list.currentItem()
        if not item:
            QMessageBox.warning(self, '提示', '请先选择一个疾病')
            return
        disease_id = item.data(Qt.UserRole)
        if disease_id is None:
            return

        file_paths, _ = QFileDialog.getOpenFileNames(self, '选择图片', '', 'Images (*.png *.jpg *.jpeg *.bmp *.gif)')
        if not file_paths:
            return

        # 确保images目录存在
        img_dir = os.path.join(APP_DIR, 'images')
        os.makedirs(img_dir, exist_ok=True)

        for fp in file_paths:
            filename = os.path.basename(fp)
            dest = os.path.join(img_dir, filename)
            # 避免重名
            if os.path.exists(dest) and dest != fp:
                base, ext = os.path.splitext(filename)
                i = 1
                while os.path.exists(os.path.join(img_dir, f'{base}_{i}{ext}')):
                    i += 1
                filename = f'{base}_{i}{ext}'
                dest = os.path.join(img_dir, filename)
            if fp != dest:
                shutil.copy2(fp, dest)

            # 加密并存储到数据库
            image_data_blob = None
            if self.user_password:
                try:
                    key = generate_key(self.user_password)
                    with open(dest, 'rb') as f:
                        raw_data = f.read()
                    image_data_blob = encrypt_data(raw_data, key)
                except Exception as e:
                    print(f"加密图片失败: {e}")

            try:
                img_id = add_image(self.active_db, disease_id, filename, '', '', '', 'image',
                                   encrypted=bool(image_data_blob), image_data=image_data_blob)
            except Exception as e:
                QMessageBox.critical(self, '错误', f'添加图片失败:\n{e}')
                return

            # 如果成功存储到数据库，删除本地文件
            if image_data_blob and img_id:
                try:
                    os.remove(dest)
                except Exception as e:
                    print(f"删除本地图片文件失败: {e}")

        self.detail.load_disease(disease_id, self.active_db)
        self.statusBar().showMessage(f'已添加 {len(file_paths)} 张图片')

    # ── 用户管理 ──
    def show_user_manage(self):
        if self.user_info.get('role') != 'admin':
            QMessageBox.warning(self, '提示', '仅管理员可管理用户')
            return
        dlg = UserManageDialog(self)
        dlg.exec_()

    def show_ai_config(self):
        dlg = AIConfigDialog(self)
        dlg.exec_()

    def show_ai_diagnosis(self):
        """显示AI影像诊断查询对话框"""
        dlg = AIDiagnosisDialog(self)
        dlg.diagnosis_clicked.connect(self._on_diagnosis_selected)
        dlg.exec_()

    def _on_diagnosis_selected(self, data):
        """AI诊断条目被选中 - 在疾病列表中搜索匹配的疾病"""
        disease_name = data.get('disease_name', '')
        if disease_name:
            self.search_input.setText(disease_name)
            self.do_search()

    def show_kb_config(self):
        """显示知识库配置"""
        dlg = _KBConfigDialog(self)
        dlg.exec_()

    # ── 搜索 ──
    def do_search(self, checked=False):
        keyword = self.search_input.text().strip()
        self.load_disease_list(search=keyword if keyword else None)
        self.statusBar().showMessage(f'搜索: {keyword}' if keyword else '就绪')

    def clear_search(self, checked=False):
        self.search_input.clear()
        self.load_disease_list()
        self.statusBar().showMessage(f'就绪  |  {APP_AUTHOR}  |  {APP_VERSION}')

    def _get_disease_full_data(self, disease_id):
        """获取疾病完整数据（基本信息+影像+医学资料+解剖资料）"""
        conn = sqlite3.connect(self.active_db)
        c = conn.cursor()
        # 基本信息
        c.execute("""SELECT name_cn, name_en, system, category, clinical, diagnosis,
                     xray_finding, ct_finding, mri_finding, pet_finding,
                     report_template, differential_diagnosis, treatment
                     FROM diseases WHERE id=?""", (disease_id,))
        row = c.fetchone()
        if not row:
            conn.close()
            return None
        data = {
            'name_cn': row[0], 'name_en': row[1], 'system': row[2], 'category': row[3],
            'clinical': row[4], 'diagnosis': row[5],
            'xray_finding': row[6], 'ct_finding': row[7], 'mri_finding': row[8], 'pet_finding': row[9],
            'report_template': row[10], 'differential_diagnosis': row[11], 'treatment': row[12],
            'images': [], 'medical_records': [], 'anatomy_records': [], 'imaging_records': []
        }
        # 影像图片
        c.execute("SELECT id, filename, image_type, caption, media_type FROM images WHERE disease_id=?", (disease_id,))
        for img_row in c.fetchall():
            img_data = {'id': img_row[0], 'filename': img_row[1], 'image_type': img_row[2],
                       'caption': img_row[3], 'media_type': img_row[4], 'path': None}
            # 尝试从数据库加载加密图片
            if self.user_password:
                try:
                    key = generate_key(self.user_password)
                    raw = load_image_decrypted(self.active_db, img_row[0], key)
                    if raw:
                        tmp = tempfile.NamedTemporaryFile(suffix='.png', delete=False)
                        tmp.write(raw)
                        tmp.close()
                        img_data['path'] = tmp.name
                except Exception:
                    pass
            # 回退到本地文件
            if not img_data['path']:
                local_path = os.path.join(APP_DIR, 'images', img_row[1]) if img_row[1] else ''
                if local_path and os.path.exists(local_path):
                    img_data['path'] = local_path
                else:
                    # 尝试预览图
                    base = os.path.splitext(local_path)[0] if local_path else ''
                    preview = base + '.preview.png' if base else ''
                    if preview and os.path.exists(preview):
                        img_data['path'] = preview
            data['images'].append(img_data)
        # 医学资料
        c.execute("SELECT title, content FROM medical_records WHERE disease_id=?", (disease_id,))
        data['medical_records'] = [{'title': r[0], 'content': r[1]} for r in c.fetchall()]
        # 解剖资料
        c.execute("SELECT title, content FROM anatomy_records")
        data['anatomy_records'] = [{'title': r[0], 'content': r[1]} for r in c.fetchall()]
        # 影像资料
        c.execute("SELECT title, content FROM imaging_records WHERE disease_id=?", (disease_id,))
        data['imaging_records'] = [{'title': r[0], 'content': r[1]} for r in c.fetchall()]
        conn.close()
        return data

    def _export_current_disease(self):
        """导出当前选中的疾病"""
        item = self.disease_list.currentItem()
        if not item or not item.data(Qt.UserRole):
            QMessageBox.warning(self, '提示', '请先选择一个疾病条目')
            return
        disease_id = item.data(Qt.UserRole)
        self._do_export([disease_id])

    def _export_multi_disease(self):
        """批量导出疾病"""
        # 弹出多选对话框
        conn = sqlite3.connect(self.active_db)
        c = conn.cursor()
        c.execute("SELECT id, name_cn, name_en, system FROM diseases ORDER BY system, name_cn")
        rows = c.fetchall()
        conn.close()
        if not rows:
            QMessageBox.warning(self, '提示', '没有可导出的疾病条目')
            return
        dlg = QDialog(self)
        dlg.setWindowTitle('批量导出疾病')
        dlg.setMinimumSize(400, 500)
        dlg.setWindowFlags(Qt.Dialog | Qt.FramelessWindowHint)
        layout = QVBoxLayout(dlg)
        # 标题栏
        title_bar = _CustomTitleBar(dlg, '批量导出疾病')
        layout.addWidget(title_bar)
        # 提示
        tip = QLabel('选择要导出的疾病条目：')
        tip.setStyleSheet("color: #aaa; padding: 8px;")
        layout.addWidget(tip)
        # 列表
        list_widget = QListWidget()
        list_widget.setSelectionMode(QAbstractItemView.MultiSelection)
        for did, name_cn, name_en, system in rows:
            display = name_cn or f'疾病 #{did}'
            if name_en and name_cn:
                display = f'{name_cn} ({name_en})'
            item = QListWidgetItem(f'[{system}] {display}')
            item.setData(Qt.UserRole, did)
            list_widget.addItem(item)
        layout.addWidget(list_widget)
        # 按钮
        btn_layout = QHBoxLayout()
        btn_select_all = QPushButton('全选')
        btn_select_all.clicked.connect(list_widget.selectAll)
        btn_layout.addWidget(btn_select_all)
        btn_layout.addStretch()
        btn_cancel = QPushButton('取消')
        btn_cancel.clicked.connect(dlg.reject)
        btn_layout.addWidget(btn_cancel)
        btn_ok = QPushButton('导出')
        btn_ok.setObjectName('accentBtn')
        btn_ok.clicked.connect(dlg.accept)
        btn_layout.addWidget(btn_ok)
        layout.addLayout(btn_layout)
        # 应用主题样式
        dlg.setStyleSheet(self.styleSheet())
        if dlg.exec_() == QDialog.Accepted:
            selected_ids = [list_widget.item(i).data(Qt.UserRole)
                          for i in range(list_widget.count())
                          if list_widget.item(i).isSelected()]
            if not selected_ids:
                QMessageBox.warning(self, '提示', '请至少选择一个疾病条目')
                return
            self._do_export(selected_ids)

    def _do_export(self, disease_ids):
        """执行导出操作"""
        # 选择导出格式
        format_dlg = QDialog(self)
        format_dlg.setWindowTitle('选择导出格式')
        format_dlg.setFixedWidth(300)
        format_dlg.setWindowFlags(Qt.Dialog | Qt.FramelessWindowHint)
        flayout = QVBoxLayout(format_dlg)
        ftb = _CustomTitleBar(format_dlg, '选择导出格式')
        flayout.addWidget(ftb)
        fmt_label = QLabel('请选择导出格式：')
        fmt_label.setStyleSheet("color: #aaa; padding: 8px;")
        flayout.addWidget(fmt_label)
        fmt_combo = QComboBox()
        fmt_combo.addItems(['Word 文档 (.docx)', 'PDF 文档 (.pdf)', 'PowerPoint 演示文稿 (.pptx)'])
        fmt_combo.setFixedHeight(30)
        flayout.addWidget(fmt_combo)

        # 加密选项
        encrypt_check = QCheckBox('加密导出文件')
        encrypt_check.setStyleSheet("color: #c8c8cc; padding: 4px 8px;")
        flayout.addWidget(encrypt_check)

        pwd_widget = QWidget()
        pwd_layout = QHBoxLayout(pwd_widget)
        pwd_layout.setContentsMargins(8, 0, 8, 0)
        pwd_label = QLabel('密码：')
        pwd_label.setStyleSheet("color: #aaa;")
        pwd_layout.addWidget(pwd_label)
        pwd_input = _PasswordLineEdit()
        pwd_input.setPlaceholderText('输入加密密码')
        pwd_input.setFixedHeight(28)
        pwd_layout.addWidget(pwd_input)
        pwd_widget.setVisible(False)
        flayout.addWidget(pwd_widget)

        encrypt_check.toggled.connect(pwd_widget.setVisible)

        btn_row = QHBoxLayout()
        btn_row.addStretch()
        btn_cancel2 = QPushButton('取消')
        btn_cancel2.clicked.connect(format_dlg.reject)
        btn_row.addWidget(btn_cancel2)
        btn_ok2 = QPushButton('导出')
        btn_ok2.setObjectName('accentBtn')
        btn_ok2.clicked.connect(format_dlg.accept)
        btn_row.addWidget(btn_ok2)
        flayout.addLayout(btn_row)
        format_dlg.setStyleSheet(self.styleSheet())
        if format_dlg.exec_() != QDialog.Accepted:
            return
        fmt_map = {0: 'docx', 1: 'pdf', 2: 'pptx'}
        fmt = fmt_map.get(fmt_combo.currentIndex(), 'docx')
        # 选择保存路径
        username = self.user_info.get('username', '') if self.user_info else ''
        if len(disease_ids) == 1:
            # 获取疾病名称
            disease_name = ''
            try:
                conn = sqlite3.connect(self.active_db)
                c = conn.cursor()
                c.execute("SELECT name_cn, name_en FROM diseases WHERE id=?", (disease_ids[0],))
                row = c.fetchone()
                conn.close()
                if row:
                    disease_name = row[0] or row[1] or ''
            except Exception:
                pass
            name_part = f'{disease_name}_{username}' if disease_name else f'疾病导出_{username}'
            default_name = f'{name_part}_{datetime.now().strftime("%Y%m%d")}.{fmt}'
        else:
            default_name = f'批量导出_{username}_{datetime.now().strftime("%Y%m%d")}.{fmt}'
        ext_filter_map = {'docx': 'Word 文档 (*.docx)', 'pdf': 'PDF 文档 (*.pdf)', 'pptx': 'PowerPoint 演示文稿 (*.pptx)'}
        ext_filter = ext_filter_map.get(fmt, f'文件 (*.{fmt})')
        save_path, _ = QFileDialog.getSaveFileName(self, '保存导出文件', default_name, ext_filter)
        if not save_path:
            return
        # 收集数据
        all_data = []
        for did in disease_ids:
            d = self._get_disease_full_data(did)
            if d:
                all_data.append(d)
        if not all_data:
            QMessageBox.warning(self, '提示', '没有可导出的数据')
            return
        # 执行导出
        try:
            if fmt == 'docx':
                self._export_to_docx(all_data, save_path)
            elif fmt == 'pdf':
                self._export_to_pdf(all_data, save_path)
            elif fmt == 'pptx':
                self._export_to_pptx(all_data, save_path)

            # 加密导出文件（使用原生密码保护）
            if encrypt_check.isChecked() and pwd_input.text().strip():
                try:
                    pwd = pwd_input.text().strip()
                    if fmt == 'pdf':
                        import pikepdf
                        pdf_doc = pikepdf.open(save_path)
                        pdf_doc.save(save_path, encryption=pikepdf.Encryption(
                            owner=pwd, user=pwd, R=4
                        ))
                        pdf_doc.close()
                    else:
                        # docx 和 pptx 使用 msoffcrypto-python
                        import msoffcrypto
                        from io import BytesIO
                        temp_buf = BytesIO()
                        with open(save_path, 'rb') as f:
                            file_obj = msoffcrypto.OfficeFile(f)
                            file_obj.encrypt(password=pwd, outfile=temp_buf)
                        with open(save_path, 'wb') as f:
                            f.write(temp_buf.getvalue())
                except Exception as enc_err:
                    QMessageBox.warning(self, '加密失败', f'文件已导出但加密失败：{enc_err}')

            QMessageBox.information(self, '导出成功', f'已导出到：{save_path}')
            # 清理临时文件
            for d in all_data:
                for img in d['images']:
                    if img['path'] and img['path'].startswith(tempfile.gettempdir()):
                        try:
                            os.remove(img['path'])
                        except Exception:
                            pass
        except Exception as e:
            QMessageBox.critical(self, '导出失败', f'导出时出错：{e}')
            import traceback
            traceback.print_exc()

    def _export_to_docx(self, all_data, save_path):
        """导出为Word文档"""
        from docx import Document
        from docx.shared import Inches, Pt, Cm, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.enum.table import WD_TABLE_ALIGNMENT
        from docx.oxml.ns import qn

        doc = Document()
        # 设置默认字体
        style = doc.styles['Normal']
        font = style.font
        font.name = 'SimSun'
        font.size = Pt(11)
        style.element.rPr.rFonts.set(qn('w:eastAsia'), 'SimSun')

        for idx, data in enumerate(all_data):
            if idx > 0:
                doc.add_page_break()

            # 标题
            title = data['name_cn'] or '未命名疾病'
            if data['name_en']:
                title += f' ({data["name_en"]})'
            h = doc.add_heading(title, level=1)
            for run in h.runs:
                run.font.color.rgb = RGBColor(0x1F, 0x4E, 0xA3)

            # 基本信息表格
            doc.add_heading('基本信息', level=2)
            table = doc.add_table(rows=2, cols=4)
            table.style = 'Light Grid Accent 1'
            table.alignment = WD_TABLE_ALIGNMENT.CENTER
            cells = [
                ('系统', data.get('system', '')),
                ('分类', data.get('category', '')),
                ('中文名', data.get('name_cn', '')),
                ('英文名', data.get('name_en', '')),
            ]
            for i, (label, value) in enumerate(cells):
                row_idx = i // 2
                col_idx = (i % 2) * 2
                table.cell(row_idx, col_idx).text = label
                table.cell(row_idx, col_idx + 1).text = value or ''

            # 临床表现
            sections = [
                ('临床表现', 'clinical'),
                ('诊断要点', 'diagnosis'),
                ('X线表现', 'xray_finding'),
                ('CT表现', 'ct_finding'),
                ('MRI表现', 'mri_finding'),
                ('PET表现', 'pet_finding'),
                ('报告模板', 'report_template'),
                ('鉴别诊断', 'differential_diagnosis'),
                ('治疗', 'treatment'),
            ]
            for sec_title, sec_key in sections:
                content = data.get(sec_key, '')
                if content and content.strip():
                    doc.add_heading(sec_title, level=2)
                    # 处理HTML内容，提取纯文本
                    import re
                    text = re.sub(r'<[^>]+>', '', content)
                    text = text.replace('&nbsp;', ' ').replace('&lt;', '<').replace('&gt;', '>')
                    for para_text in text.split('\n'):
                        if para_text.strip():
                            doc.add_paragraph(para_text.strip())

            # 影像图片
            if data['images']:
                doc.add_heading('影像图片', level=2)
                for img in data['images']:
                    if img['path'] and os.path.exists(img['path']):
                        try:
                            doc.add_picture(img['path'], width=Inches(5))
                            last_paragraph = doc.paragraphs[-1]
                            last_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
                        except Exception:
                            doc.add_paragraph(f'[图片无法加载: {img.get("filename", "")}]')
                    caption_parts = []
                    if img.get('image_type'):
                        caption_parts.append(img['image_type'])
                    if img.get('caption'):
                        caption_parts.append(img['caption'])
                    if caption_parts:
                        cap = doc.add_paragraph(' - '.join(caption_parts))
                        cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
                        for run in cap.runs:
                            run.font.size = Pt(9)
                            run.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

            # 医学资料
            if data['medical_records']:
                doc.add_heading('医学资料', level=2)
                for rec in data['medical_records']:
                    doc.add_heading(rec.get('title', '无标题'), level=3)
                    import re
                    text = re.sub(r'<[^>]+>', '', rec.get('content', ''))
                    text = text.replace('&nbsp;', ' ').replace('&lt;', '<').replace('&gt;', '>')
                    for para_text in text.split('\n'):
                        if para_text.strip():
                            doc.add_paragraph(para_text.strip())

            # 解剖资料
            if data['anatomy_records']:
                doc.add_heading('影像解剖图谱与资料', level=2)
                for rec in data['anatomy_records']:
                    doc.add_heading(rec.get('title', '无标题'), level=3)
                    import re
                    text = re.sub(r'<[^>]+>', '', rec.get('content', ''))
                    text = text.replace('&nbsp;', ' ').replace('&lt;', '<').replace('&gt;', '>')
                    for para_text in text.split('\n'):
                        if para_text.strip():
                            doc.add_paragraph(para_text.strip())

            # 影像资料
            if data['imaging_records']:
                doc.add_heading('影像相关病例', level=2)
                for rec in data['imaging_records']:
                    doc.add_heading(rec.get('title', '无标题'), level=3)
                    import re
                    text = re.sub(r'<[^>]+>', '', rec.get('content', ''))
                    text = text.replace('&nbsp;', ' ').replace('&lt;', '<').replace('&gt;', '>')
                    for para_text in text.split('\n'):
                        if para_text.strip():
                            doc.add_paragraph(para_text.strip())

        doc.save(save_path)

    def _export_to_pdf(self, all_data, save_path):
        """导出为PDF文档"""
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.units import cm, mm
        from reportlab.lib.colors import HexColor
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Image, PageBreak
        from reportlab.lib.enums import TA_CENTER, TA_LEFT
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.ttfonts import TTFont
        import re

        # 注册中文字体
        font_paths = [
            ('SimSun', 'C:/Windows/Fonts/simsun.ttc'),
            ('SimHei', 'C:/Windows/Fonts/simhei.ttf'),
            ('Microsoft YaHei', 'C:/Windows/Fonts/msyh.ttc'),
        ]
        chinese_font = 'SimSun'
        for fname, fpath in font_paths:
            if os.path.exists(fpath):
                try:
                    pdfmetrics.registerFont(TTFont(fname, fpath))
                    if fname == 'Microsoft YaHei':
                        chinese_font = fname
                except Exception:
                    pass
                break

        doc = SimpleDocTemplate(save_path, pagesize=A4,
                               leftMargin=2*cm, rightMargin=2*cm,
                               topMargin=2*cm, bottomMargin=2*cm)

        styles = getSampleStyleSheet()
        # 自定义样式
        title_style = ParagraphStyle('Title_CN', parent=styles['Title'],
                                     fontName=chinese_font, fontSize=18,
                                     textColor=HexColor('#1F4EA3'),
                                     spaceAfter=12)
        h2_style = ParagraphStyle('H2_CN', parent=styles['Heading2'],
                                  fontName=chinese_font, fontSize=14,
                                  textColor=HexColor('#2E5E8E'),
                                  spaceBefore=16, spaceAfter=8)
        h3_style = ParagraphStyle('H3_CN', parent=styles['Heading3'],
                                  fontName=chinese_font, fontSize=12,
                                  textColor=HexColor('#3A7BBF'),
                                  spaceBefore=10, spaceAfter=6)
        body_style = ParagraphStyle('Body_CN', parent=styles['Normal'],
                                    fontName=chinese_font, fontSize=10,
                                    leading=16, spaceAfter=4)
        caption_style = ParagraphStyle('Caption_CN', parent=styles['Normal'],
                                       fontName=chinese_font, fontSize=8,
                                       textColor=HexColor('#666666'),
                                       alignment=TA_CENTER, spaceAfter=8)

        elements = []

        for idx, data in enumerate(all_data):
            if idx > 0:
                elements.append(PageBreak())

            # 标题
            title = data['name_cn'] or '未命名疾病'
            if data['name_en']:
                title += f' ({data["name_en"]})'
            elements.append(Paragraph(title, title_style))
            elements.append(Spacer(1, 8))

            # 基本信息表格
            elements.append(Paragraph('基本信息', h2_style))
            table_data = [
                ['系统', data.get('system', ''), '分类', data.get('category', '')],
                ['中文名', data.get('name_cn', ''), '英文名', data.get('name_en', '')],
            ]
            t = Table(table_data, colWidths=[2.5*cm, 5*cm, 2.5*cm, 5*cm])
            t.setStyle(TableStyle([
                ('FONTNAME', (0, 0), (-1, -1), chinese_font),
                ('FONTSIZE', (0, 0), (-1, -1), 10),
                ('BACKGROUND', (0, 0), (0, -1), HexColor('#E8EFF7')),
                ('BACKGROUND', (2, 0), (2, -1), HexColor('#E8EFF7')),
                ('GRID', (0, 0), (-1, -1), 0.5, HexColor('#CCCCCC')),
                ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
                ('TOPPADDING', (0, 0), (-1, -1), 4),
                ('BOTTOMPADDING', (0, 0), (-1, -1), 4),
            ]))
            elements.append(t)
            elements.append(Spacer(1, 8))

            # 文本内容区
            sections = [
                ('临床表现', 'clinical'),
                ('诊断要点', 'diagnosis'),
                ('X线表现', 'xray_finding'),
                ('CT表现', 'ct_finding'),
                ('MRI表现', 'mri_finding'),
                ('PET表现', 'pet_finding'),
                ('报告模板', 'report_template'),
                ('鉴别诊断', 'differential_diagnosis'),
                ('治疗', 'treatment'),
            ]
            for sec_title, sec_key in sections:
                content = data.get(sec_key, '')
                if content and content.strip():
                    elements.append(Paragraph(sec_title, h2_style))
                    text = re.sub(r'<[^>]+>', '', content)
                    text = text.replace('&nbsp;', ' ').replace('&lt;', '<').replace('&gt;', '>')
                    text = text.replace('&amp;', '&')
                    for line in text.split('\n'):
                        if line.strip():
                            # 转义XML特殊字符
                            line = line.strip().replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                            elements.append(Paragraph(line, body_style))

            # 影像图片
            if data['images']:
                elements.append(Paragraph('影像图片', h2_style))
                for img in data['images']:
                    if img['path'] and os.path.exists(img['path']):
                        try:
                            from PIL import Image as PILImage
                            pil_img = PILImage.open(img['path'])
                            w, h = pil_img.size
                            max_w = 14 * cm
                            if w > max_w:
                                ratio = max_w / w
                                w = max_w
                                h = h * ratio
                            img_elem = Image(img['path'], width=w, height=h)
                            elements.append(img_elem)
                        except Exception:
                            elements.append(Paragraph(f'[图片无法加载: {img.get("filename", "")}]', body_style))
                    caption_parts = []
                    if img.get('image_type'):
                        caption_parts.append(img['image_type'])
                    if img.get('caption'):
                        caption_parts.append(img['caption'])
                    if caption_parts:
                        elements.append(Paragraph(' - '.join(caption_parts), caption_style))

            # 医学资料
            if data['medical_records']:
                elements.append(Paragraph('医学资料', h2_style))
                for rec in data['medical_records']:
                    elements.append(Paragraph(rec.get('title', '无标题'), h3_style))
                    text = re.sub(r'<[^>]+>', '', rec.get('content', ''))
                    text = text.replace('&nbsp;', ' ').replace('&lt;', '<').replace('&gt;', '>')
                    text = text.replace('&amp;', '&')
                    for line in text.split('\n'):
                        if line.strip():
                            line = line.strip().replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                            elements.append(Paragraph(line, body_style))

            # 解剖资料
            if data['anatomy_records']:
                elements.append(Paragraph('影像解剖图谱与资料', h2_style))
                for rec in data['anatomy_records']:
                    elements.append(Paragraph(rec.get('title', '无标题'), h3_style))
                    text = re.sub(r'<[^>]+>', '', rec.get('content', ''))
                    text = text.replace('&nbsp;', ' ').replace('&lt;', '<').replace('&gt;', '>')
                    text = text.replace('&amp;', '&')
                    for line in text.split('\n'):
                        if line.strip():
                            line = line.strip().replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                            elements.append(Paragraph(line, body_style))

            # 影像资料
            if data['imaging_records']:
                elements.append(Paragraph('影像相关病例', h2_style))
                for rec in data['imaging_records']:
                    elements.append(Paragraph(rec.get('title', '无标题'), h3_style))
                    text = re.sub(r'<[^>]+>', '', rec.get('content', ''))
                    text = text.replace('&nbsp;', ' ').replace('&lt;', '<').replace('&gt;', '>')
                    text = text.replace('&amp;', '&')
                    for line in text.split('\n'):
                        if line.strip():
                            line = line.strip().replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                            elements.append(Paragraph(line, body_style))

        doc.build(elements)

    def _export_to_pptx(self, all_data, save_path):
        """导出为PowerPoint演示文稿"""
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        import re

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        DARK_BLUE = RGBColor(0x1F, 0x4E, 0xA3)
        LIGHT_BLUE = RGBColor(0xE8, 0xEF, 0xF7)
        BODY_COLOR = RGBColor(0x33, 0x33, 0x33)
        CAPTION_COLOR = RGBColor(0x66, 0x66, 0x66)

        def _add_text_box(slide, left, top, width, height, text, font_size=14,
                          bold=False, color=BODY_COLOR, alignment=PP_ALIGN.LEFT,
                          font_name='SimSun'):
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = text
            p.font.size = Pt(font_size)
            p.font.bold = bold
            p.font.color.rgb = color
            p.font.name = font_name
            p.alignment = alignment
            if font_name == 'Microsoft YaHei':
                from pptx.oxml.ns import qn
                rPr = p.runs[0]._r.get_or_add_rPr()
                rPr.set(qn('a:ea'), font_name)
            return tf

        def _add_content_slide(title_text, body_text):
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
            # 标题
            _add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.8),
                          title_text, font_size=28, bold=True, color=DARK_BLUE,
                          font_name='Microsoft YaHei')
            # 分隔线
            from pptx.util import Emu as _Emu
            line = slide.shapes.add_shape(
                1, Inches(0.8), Inches(1.3), Inches(11.7), Pt(2)
            )
            line.fill.solid()
            line.fill.fore_color.rgb = DARK_BLUE
            line.line.fill.background()
            # 正文
            if body_text and body_text.strip():
                clean = re.sub(r'<[^>]+>', '', body_text)
                clean = clean.replace('&nbsp;', ' ').replace('&lt;', '<').replace('&gt;', '>').replace('&amp;', '&')
                lines = [l.strip() for l in clean.split('\n') if l.strip()]
                txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.4))
                tf = txBox.text_frame
                tf.word_wrap = True
                for i, line in enumerate(lines):
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.text = line
                    p.font.size = Pt(14)
                    p.font.color.rgb = BODY_COLOR
                    p.font.name = 'SimSun'
                    p.space_after = Pt(4)
            return slide

        for idx, data in enumerate(all_data):
            # ── 标题幻灯片 ──
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
            # 背景色块
            bg_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(4.2))
            bg_shape.fill.solid()
            bg_shape.fill.fore_color.rgb = DARK_BLUE
            bg_shape.line.fill.background()
            # 疾病中文名
            name_cn = data.get('name_cn', '') or '未命名疾病'
            _add_text_box(slide, Inches(1.5), Inches(1.2), Inches(10.3), Inches(1.5),
                          name_cn, font_size=40, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
                          alignment=PP_ALIGN.CENTER, font_name='Microsoft YaHei')
            # 英文名
            name_en = data.get('name_en', '')
            if name_en:
                _add_text_box(slide, Inches(1.5), Inches(2.7), Inches(10.3), Inches(0.8),
                              name_en, font_size=20, color=RGBColor(0xCC, 0xDD, 0xFF),
                              alignment=PP_ALIGN.CENTER)
            # 副标题
            subtitle = f'{data.get("system", "")} | {data.get("category", "")}'
            _add_text_box(slide, Inches(1.5), Inches(4.8), Inches(10.3), Inches(0.6),
                          subtitle, font_size=16, color=CAPTION_COLOR,
                          alignment=PP_ALIGN.CENTER)

            # ── 基本信息幻灯片 ──
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            _add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.8),
                          '基本信息', font_size=28, bold=True, color=DARK_BLUE,
                          font_name='Microsoft YaHei')
            info_rows = [
                ['系统', data.get('system', ''), '分类', data.get('category', '')],
                ['中文名', data.get('name_cn', ''), '英文名', data.get('name_en', '')],
            ]
            table_shape = slide.shapes.add_table(2, 4, Inches(1.5), Inches(1.6), Inches(10.3), Inches(1.2))
            table = table_shape.table
            for r, row_data in enumerate(info_rows):
                for c, val in enumerate(row_data):
                    cell = table.cell(r, c)
                    cell.text = val or ''
                    for p in cell.text_frame.paragraphs:
                        p.font.size = Pt(14)
                        p.font.name = 'SimSun'
                        if c % 2 == 0:
                            p.font.bold = True
                            p.font.color.rgb = DARK_BLUE
                        else:
                            p.font.color.rgb = BODY_COLOR
                    if c % 2 == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = LIGHT_BLUE

            # ── 内容幻灯片 ──
            sections = [
                ('临床表现', 'clinical'),
                ('诊断要点', 'diagnosis'),
                ('影像所见 - X线', 'xray_finding'),
                ('影像所见 - CT', 'ct_finding'),
                ('影像所见 - MRI', 'mri_finding'),
                ('影像所见 - PET', 'pet_finding'),
                ('报告模板', 'report_template'),
                ('鉴别诊断', 'differential_diagnosis'),
                ('治疗', 'treatment'),
            ]
            for sec_title, sec_key in sections:
                content = data.get(sec_key, '')
                if content and content.strip():
                    _add_content_slide(sec_title, content)

            # ── 图片幻灯片 ──
            if data['images']:
                for img in data['images']:
                    if img['path'] and os.path.exists(img['path']):
                        slide = prs.slides.add_slide(prs.slide_layouts[6])
                        try:
                            from PIL import Image as PILImage
                            pil_img = PILImage.open(img['path'])
                            w, h = pil_img.size
                            max_w, max_h = 9.0, 5.0
                            ratio = min(max_w / (w / 96), max_h / (h / 96))
                            img_w = (w / 96) * ratio
                            img_h = (h / 96) * ratio
                            left = (13.333 - img_w) / 2
                            slide.shapes.add_picture(img['path'],
                                                     Inches(left), Inches(0.5),
                                                     Inches(img_w), Inches(img_h))
                        except Exception:
                            _add_text_box(slide, Inches(2), Inches(2.5), Inches(9), Inches(1),
                                          f'[图片无法加载: {img.get("filename", "")}]',
                                          font_size=14, color=CAPTION_COLOR)
                        caption_parts = []
                        if img.get('image_type'):
                            caption_parts.append(img['image_type'])
                        if img.get('caption'):
                            caption_parts.append(img['caption'])
                        if caption_parts:
                            _add_text_box(slide, Inches(1), Inches(6.2), Inches(11.3), Inches(0.5),
                                          ' - '.join(caption_parts), font_size=12,
                                          color=CAPTION_COLOR, alignment=PP_ALIGN.CENTER)

            # ── 医学资料幻灯片 ──
            if data['medical_records']:
                for rec in data['medical_records']:
                    slide = _add_content_slide(
                        f'医学资料 - {rec.get("title", "无标题")}',
                        rec.get('content', '')
                    )

            # ── 解剖资料幻灯片 ──
            if data['anatomy_records']:
                for rec in data['anatomy_records']:
                    slide = _add_content_slide(
                        f'影像解剖图谱 - {rec.get("title", "无标题")}',
                        rec.get('content', '')
                    )

            # ── 影像资料幻灯片 ──
            if data['imaging_records']:
                for rec in data['imaging_records']:
                    slide = _add_content_slide(
                        f'影像资料 - {rec.get("title", "无标题")}',
                        rec.get('content', '')
                    )

        prs.save(save_path)

    def do_logout(self, checked=False):
        self.close()
        QApplication.instance().quit()


# ── 应用入口 ──────────────────────────────────────────────
class RadAtlasApp:
    def __init__(self):
        self.app = QApplication(sys.argv)
        self.app.setStyleSheet(DARK_STYLE)
        self.app.setStyle('Fusion')

        palette = QPalette()
        palette.setColor(QPalette.Window, QColor(20, 20, 22))
        palette.setColor(QPalette.WindowText, QColor(217, 217, 220))
        palette.setColor(QPalette.Base, QColor(26, 26, 30))
        palette.setColor(QPalette.AlternateBase, QColor(30, 30, 34))
        palette.setColor(QPalette.ToolTipBase, QColor(26, 26, 30))
        palette.setColor(QPalette.ToolTipText, QColor(217, 217, 220))
        palette.setColor(QPalette.Text, QColor(217, 217, 220))
        palette.setColor(QPalette.Button, QColor(42, 42, 48))
        palette.setColor(QPalette.ButtonText, QColor(217, 217, 220))
        palette.setColor(QPalette.BrightText, QColor(255, 255, 255))
        palette.setColor(QPalette.Highlight, QColor(63, 123, 247))
        palette.setColor(QPalette.HighlightedText, QColor(255, 255, 255))
        self.app.setPalette(palette)

    def run(self):
        try:
            init_user_db()
        except Exception as e:
            QMessageBox.critical(None, '错误', f'初始化用户数据库失败:\n{e}')
            return

        login = LoginDialog()
        if login.exec_() != QDialog.Accepted:
            return

        user_info = login.user_info
        user_password = login.user_password
        active_db = user_info.get('database') or DATABASE
        try:
            init_db(active_db)
        except Exception as e:
            QMessageBox.critical(None, '错误', f'初始化数据库失败:\n{e}')
            return

        try:
            if active_db == DATABASE:
                load_data(active_db)
        except Exception as e:
            QMessageBox.critical(None, '错误', f'加载数据失败:\n{e}')
            return

        try:
            window = MainWindow(user_info, active_db, user_password)
            window.show()
        except Exception as e:
            QMessageBox.critical(None, '错误', f'创建主窗口失败:\n{e}')
            return

        sys.exit(self.app.exec_())


if __name__ == '__main__':
    RadAtlasApp().run()
